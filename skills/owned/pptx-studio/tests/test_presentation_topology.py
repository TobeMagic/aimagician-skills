"""Focused tests for the independent strict PresentationML helper."""

from __future__ import annotations

import base64
import copy
import sys
from pathlib import Path
from types import SimpleNamespace
import zipfile
import xml.etree.ElementTree as ET

import pytest


SKILL_ROOT = Path(__file__).resolve().parents[1]
SCRIPTS_ROOT = SKILL_ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS_ROOT))

from window_pptx import presentation_topology as topology_module  # noqa: E402
from window_pptx.independent_validation_security import (  # noqa: E402
    ZipResourceLimits,
)
from window_pptx.presentation_topology import (  # noqa: E402
    OFFICE_RELATIONSHIP_NS,
    OFFICE_DOCUMENT_RELATIONSHIP_TYPE,
    DRAWINGML_NS,
    PACKAGE_RELATIONSHIP_NS,
    PACKAGE_ROOT_RELS_PART,
    PRESENTATIONML_NS,
    PRESENTATION_PART,
    PRESENTATION_RELS_PART,
    RELATIONSHIP_ID_ATTR,
    RELATIONSHIP_TAG,
    SLIDE_ID_LIST_TAG,
    SLIDE_RELATIONSHIP_TYPE,
    inspect_presentation_topology,
)


def _make_deck(path: Path, slide_count: int = 1) -> None:
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation()
    for ordinal in range(1, slide_count + 1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(100_000, 100_000, 2_000_000, 500_000)
        textbox.text = f"Slide {ordinal}"
    presentation.save(path)


def _package_parts(path: Path) -> dict[str, bytes]:
    with zipfile.ZipFile(path, "r") as archive:
        return {name: archive.read(name) for name in archive.namelist()}


def _rewrite_package(
    path: Path,
    updates: dict[str, bytes],
    *,
    remove: set[str] | None = None,
) -> None:
    parts = _package_parts(path)
    for name in remove or set():
        parts.pop(name, None)
    parts.update(updates)
    replacement = path.with_suffix(".rewrite.pptx")
    with zipfile.ZipFile(replacement, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, payload in sorted(parts.items()):
            archive.writestr(name, payload)
    replacement.replace(path)


def _codes(result: object) -> set[str]:
    return {issue.code for issue in result.issues}


def _presentation_root(path: Path) -> ET.Element:
    return ET.fromstring(_package_parts(path)[PRESENTATION_PART])


def _relationships_root(path: Path) -> ET.Element:
    return ET.fromstring(_package_parts(path)[PRESENTATION_RELS_PART])


def _root_relationships_root(path: Path) -> ET.Element:
    return ET.fromstring(_package_parts(path)[PACKAGE_ROOT_RELS_PART])


def _slide_relationships(root: ET.Element) -> list[ET.Element]:
    return [
        node
        for node in list(root)
        if node.tag == RELATIONSHIP_TAG
        and node.attrib.get("Type") == SLIDE_RELATIONSHIP_TYPE
    ]


def test_valid_topology_is_ordered_reopened_and_counted(tmp_path: Path) -> None:
    output = tmp_path / "valid.pptx"
    _make_deck(output, slide_count=2)

    result = inspect_presentation_topology(output)

    assert result.status == "pass"
    assert result.issues == ()
    assert len(result.ordered_relationship_ids) == 2
    assert len(set(result.ordered_relationship_ids)) == 2
    assert result.ordered_slide_parts == (
        "ppt/slides/slide1.xml",
        "ppt/slides/slide2.xml",
    )
    assert result.statistics is not None
    assert result.statistics.slide_count == 2
    assert result.statistics.python_pptx_slide_count == 2
    assert result.statistics.native_editable_slide_count == 2
    assert result.statistics.native_editable_coverage == 1
    assert result.statistics.text_run_count == 2
    assert [slide.part_name for slide in result.statistics.slides] == list(
        result.ordered_slide_parts
    )


def test_resource_preflight_stops_before_xml_or_python_pptx_parse(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    output = tmp_path / "oversized-xml.pptx"
    _make_deck(output)
    limits = ZipResourceLimits(
        max_entries=10_000,
        max_entry_uncompressed_bytes=1_000_000,
        max_total_uncompressed_bytes=10_000_000,
        max_compression_ratio=1_000,
        max_xml_uncompressed_bytes=64,
        max_relationship_uncompressed_bytes=32,
    )

    def _unexpected_parse(*_args: object, **_kwargs: object) -> object:
        raise AssertionError("python-pptx must not run after resource rejection")

    monkeypatch.setattr(
        topology_module,
        "_python_pptx_statistics",
        _unexpected_parse,
    )
    monkeypatch.setattr(
        topology_module,
        "_parse_xml",
        _unexpected_parse,
    )

    result = inspect_presentation_topology(output, archive_limits=limits)

    assert result.status == "fail"
    assert "ZIP_RESOURCE_XML_SIZE_EXCEEDED" in _codes(result)
    assert result.statistics is None


def test_wrong_presentationml_namespace_is_rejected(tmp_path: Path) -> None:
    output = tmp_path / "wrong-namespace.pptx"
    _make_deck(output)
    parts = _package_parts(output)
    parts[PRESENTATION_PART] = parts[PRESENTATION_PART].replace(
        PRESENTATIONML_NS.encode(),
        b"urn:not-presentationml",
    )
    _rewrite_package(output, {PRESENTATION_PART: parts[PRESENTATION_PART]})

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PRESENTATION_ROOT_INVALID" in _codes(result)


def test_utf16_dtd_and_entity_are_rejected_without_expansion(tmp_path: Path) -> None:
    output = tmp_path / "utf16-entity.pptx"
    _make_deck(output)
    malicious = (
        '<?xml version="1.0" encoding="utf-16"?>'
        '<!DOCTYPE p:presentation [<!ENTITY payload "expanded">]>'
        f'<p:presentation xmlns:p="{PRESENTATIONML_NS}">'
        "&payload;</p:presentation>"
    ).encode("utf-16")
    _rewrite_package(output, {PRESENTATION_PART: malicious})

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PRESENTATION_XML_MALFORMED" in _codes(result)


def test_duplicate_slide_id_list_is_rejected(tmp_path: Path) -> None:
    output = tmp_path / "duplicate-list.pptx"
    _make_deck(output)
    root = _presentation_root(output)
    slide_id_list = root.find(SLIDE_ID_LIST_TAG)
    assert slide_id_list is not None
    root.append(copy.deepcopy(slide_id_list))
    _rewrite_package(
        output,
        {PRESENTATION_PART: ET.tostring(root, encoding="utf-8", xml_declaration=True)},
    )

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PRESENTATION_SLIDE_ID_LIST_INVALID" in _codes(result)


def test_duplicate_ordered_slide_relationship_id_is_rejected(tmp_path: Path) -> None:
    output = tmp_path / "duplicate-rid.pptx"
    _make_deck(output, slide_count=2)
    root = _presentation_root(output)
    slide_id_list = root.find(SLIDE_ID_LIST_TAG)
    assert slide_id_list is not None
    slide_ids = list(slide_id_list)
    slide_ids[1].set(RELATIONSHIP_ID_ATTR, slide_ids[0].get(RELATIONSHIP_ID_ATTR, ""))
    _rewrite_package(
        output,
        {PRESENTATION_PART: ET.tostring(root, encoding="utf-8", xml_declaration=True)},
    )

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PRESENTATION_SLIDE_RELATIONSHIP_ID_DUPLICATE" in _codes(result)


def test_presentation_relationship_namespace_is_strict(tmp_path: Path) -> None:
    output = tmp_path / "relationship-namespace.pptx"
    _make_deck(output)
    relationships = _package_parts(output)[PRESENTATION_RELS_PART].replace(
        PACKAGE_RELATIONSHIP_NS.encode(),
        b"urn:not-package-relationships",
    )
    _rewrite_package(output, {PRESENTATION_RELS_PART: relationships})

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PRESENTATION_RELATIONSHIPS_ROOT_INVALID" in _codes(result)


@pytest.mark.parametrize(
    ("mutation", "expected_code"),
    [
        ("missing", "PRESENTATION_REQUIRED_PART_MISSING"),
        ("duplicate", "PACKAGE_ROOT_OFFICE_DOCUMENT_RELATIONSHIP_INVALID"),
        ("external", "PACKAGE_ROOT_OFFICE_DOCUMENT_RELATIONSHIP_EXTERNAL"),
        ("wrong-type", "PACKAGE_ROOT_PRESENTATION_RELATIONSHIP_TYPE_INVALID"),
        ("wrong-target", "PACKAGE_ROOT_OFFICE_DOCUMENT_TARGET_INVALID"),
    ],
)
def test_package_root_has_one_internal_office_document_relationship(
    tmp_path: Path,
    mutation: str,
    expected_code: str,
) -> None:
    output = tmp_path / f"root-{mutation}.pptx"
    _make_deck(output)
    if mutation == "missing":
        _rewrite_package(output, {}, remove={PACKAGE_ROOT_RELS_PART})
    else:
        root = _root_relationships_root(output)
        office_documents = [
            node
            for node in list(root)
            if node.attrib.get("Type") == OFFICE_DOCUMENT_RELATIONSHIP_TYPE
        ]
        assert len(office_documents) == 1
        office_document = office_documents[0]
        if mutation == "duplicate":
            duplicate = copy.deepcopy(office_document)
            duplicate.set("Id", "rDuplicateOfficeDocument")
            root.append(duplicate)
        elif mutation == "external":
            office_document.set("TargetMode", "External")
        elif mutation == "wrong-type":
            office_document.set("Type", f"{OFFICE_RELATIONSHIP_NS}/image")
        elif mutation == "wrong-target":
            office_document.set("Target", "ppt/not-presentation.xml")
        _rewrite_package(
            output,
            {
                PACKAGE_ROOT_RELS_PART: ET.tostring(
                    root,
                    encoding="utf-8",
                    xml_declaration=True,
                )
            },
        )

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert expected_code in _codes(result)


@pytest.mark.parametrize(
    "target",
    [
        "/ppt/presentation.xml",
        "./ppt/presentation.xml",
        "ppt%2Fpresentation.xml",
    ],
)
def test_wrong_root_relationship_type_cannot_alias_presentation_part(
    tmp_path: Path,
    target: str,
) -> None:
    output = tmp_path / "root-wrong-type-alias.pptx"
    _make_deck(output)
    root = _root_relationships_root(output)
    office_document = next(
        node
        for node in list(root)
        if node.attrib.get("Type") == OFFICE_DOCUMENT_RELATIONSHIP_TYPE
    )
    office_document.set("Type", f"{OFFICE_RELATIONSHIP_NS}/image")
    office_document.set("Target", target)
    _rewrite_package(
        output,
        {
            PACKAGE_ROOT_RELS_PART: ET.tostring(
                root,
                encoding="utf-8",
                xml_declaration=True,
            )
        },
    )

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PACKAGE_ROOT_PRESENTATION_RELATIONSHIP_TYPE_INVALID" in _codes(result)


@pytest.mark.parametrize(
    ("relationships_part", "payload_location"),
    [
        (PRESENTATION_RELS_PART, "root-text"),
        (PRESENTATION_RELS_PART, "child-tail"),
        (PACKAGE_ROOT_RELS_PART, "root-text"),
        (PACKAGE_ROOT_RELS_PART, "child-tail"),
    ],
)
def test_relationship_mixed_text_is_rejected(
    tmp_path: Path,
    relationships_part: str,
    payload_location: str,
) -> None:
    output = tmp_path / f"relationship-{payload_location}.pptx"
    _make_deck(output)
    root = (
        _relationships_root(output)
        if relationships_part == PRESENTATION_RELS_PART
        else _root_relationships_root(output)
    )
    prefix = (
        "PRESENTATION"
        if relationships_part == PRESENTATION_RELS_PART
        else "PACKAGE_ROOT"
    )
    if payload_location == "root-text":
        root.text = "payload"
        expected_code = f"{prefix}_RELATIONSHIPS_ROOT_INVALID"
    else:
        list(root)[0].tail = "payload"
        expected_code = f"{prefix}_RELATIONSHIP_ELEMENT_INVALID"
    _rewrite_package(
        output,
        {
            relationships_part: ET.tostring(
                root,
                encoding="utf-8",
                xml_declaration=True,
            )
        },
    )

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert expected_code in _codes(result)


def test_duplicate_presentation_relationship_id_is_rejected(tmp_path: Path) -> None:
    output = tmp_path / "duplicate-relationship.pptx"
    _make_deck(output)
    root = _relationships_root(output)
    slide_relationship = _slide_relationships(root)[0]
    duplicate = copy.deepcopy(slide_relationship)
    duplicate.set("Target", "slides/slide2.xml")
    root.append(duplicate)
    _rewrite_package(
        output,
        {
            PRESENTATION_RELS_PART: ET.tostring(
                root,
                encoding="utf-8",
                xml_declaration=True,
            )
        },
    )

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PRESENTATION_RELATIONSHIP_ID_DUPLICATE" in _codes(result)


@pytest.mark.parametrize(
    ("entry_name", "expected_code"),
    [
        ("ppt/./presentation.xml", "PRESENTATION_PACKAGE_ENTRY_NONCANONICAL"),
        ("ppt/%70resentation.xml", "PRESENTATION_PACKAGE_ENTRY_NONCANONICAL"),
        ("ppt/query?.xml", "PRESENTATION_PACKAGE_ENTRY_NONCANONICAL"),
        ("ppt/fragment#.xml", "PRESENTATION_PACKAGE_ENTRY_NONCANONICAL"),
        ("ppt/cafe\u0301.xml", "PRESENTATION_PACKAGE_ENTRY_NONCANONICAL"),
        ("ppt/hidden\u007f.xml", "PRESENTATION_PACKAGE_ENTRY_NONCANONICAL"),
        ("ppt/hidden\u202e.xml", "PRESENTATION_PACKAGE_ENTRY_NONCANONICAL"),
        ("PPT/presentation.xml", "PRESENTATION_PACKAGE_ENTRY_CASE_COLLISION"),
    ],
)
def test_zip_entry_names_are_canonical_and_collision_free(
    tmp_path: Path,
    entry_name: str,
    expected_code: str,
) -> None:
    output = tmp_path / "noncanonical-entry.pptx"
    _make_deck(output)
    _rewrite_package(output, {entry_name: b"not a package part"})

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert expected_code in _codes(result)


def test_duplicate_zip_entry_name_is_rejected(tmp_path: Path) -> None:
    output = tmp_path / "duplicate-entry.pptx"
    _make_deck(output)
    presentation_payload = _package_parts(output)[PRESENTATION_PART]
    with pytest.warns(UserWarning, match="Duplicate name"):
        with zipfile.ZipFile(output, "a", zipfile.ZIP_DEFLATED) as archive:
            archive.writestr(PRESENTATION_PART, presentation_payload)

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PRESENTATION_PACKAGE_ENTRY_DUPLICATE" in _codes(result)


@pytest.mark.parametrize(
    "target",
    [
        "/ppt/slides/slide1.xml",
        "./slides/slide1.xml",
        "slides%2Fslide1.xml",
        "slides/%2e%2e/slides/slide1.xml",
        "../ppt/slides/slide1.xml",
    ],
)
def test_non_slide_relationship_cannot_target_a_slide_part(
    tmp_path: Path,
    target: str,
) -> None:
    output = tmp_path / "wrong-type-slide-target.pptx"
    _make_deck(output)
    root = _relationships_root(output)
    ET.SubElement(
        root,
        f"{{{PACKAGE_RELATIONSHIP_NS}}}Relationship",
        {
            "Id": "rWrongType",
            "Type": f"{OFFICE_RELATIONSHIP_NS}/image",
            "Target": target,
        },
    )
    _rewrite_package(
        output,
        {
            PRESENTATION_RELS_PART: ET.tostring(
                root,
                encoding="utf-8",
                xml_declaration=True,
            )
        },
    )

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PRESENTATION_RELATIONSHIP_TARGET_INVALID" in _codes(result)
    assert "PRESENTATION_SLIDE_RELATIONSHIP_TYPE_INVALID" in _codes(result)


@pytest.mark.parametrize(
    "target",
    [
        "/ppt/theme/theme1.xml",
        "./theme/theme1.xml",
        "theme%2Ftheme1.xml",
        "../theme/theme1.xml",
        "theme//theme1.xml",
        "theme\\theme1.xml",
        "theme/theme1.xml\t",
    ],
)
def test_all_internal_presentation_targets_must_be_canonical(
    tmp_path: Path,
    target: str,
) -> None:
    output = tmp_path / "noncanonical-internal-target.pptx"
    _make_deck(output)
    root = _relationships_root(output)
    non_slide_relationship = next(
        node
        for node in list(root)
        if node.attrib.get("Type") != SLIDE_RELATIONSHIP_TYPE
        and not node.attrib.get("TargetMode")
    )
    non_slide_relationship.set("Target", target)
    _rewrite_package(
        output,
        {
            PRESENTATION_RELS_PART: ET.tostring(
                root,
                encoding="utf-8",
                xml_declaration=True,
            )
        },
    )

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PRESENTATION_RELATIONSHIP_TARGET_INVALID" in _codes(result)


def test_slide_targets_must_follow_slide_list_order(tmp_path: Path) -> None:
    output = tmp_path / "swapped-targets.pptx"
    _make_deck(output, slide_count=2)
    root = _relationships_root(output)
    slide_relationships = _slide_relationships(root)
    slide_relationships[0].set("Target", "slides/slide2.xml")
    slide_relationships[1].set("Target", "slides/slide1.xml")
    _rewrite_package(
        output,
        {
            PRESENTATION_RELS_PART: ET.tostring(
                root,
                encoding="utf-8",
                xml_declaration=True,
            )
        },
    )

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PRESENTATION_SLIDE_TARGET_SEQUENCE_INVALID" in _codes(result)


@pytest.mark.parametrize("extra_kind", ["relationship", "part"])
def test_extra_slide_relationship_or_part_is_rejected(
    tmp_path: Path,
    extra_kind: str,
) -> None:
    output = tmp_path / f"extra-{extra_kind}.pptx"
    _make_deck(output)
    parts = _package_parts(output)
    updates: dict[str, bytes] = {
        "ppt/slides/slide2.xml": parts["ppt/slides/slide1.xml"],
    }
    if extra_kind == "relationship":
        root = _relationships_root(output)
        ET.SubElement(
            root,
            f"{{{PACKAGE_RELATIONSHIP_NS}}}Relationship",
            {
                "Id": "rExtraSlide",
                "Type": f"{OFFICE_RELATIONSHIP_NS}/slide",
                "Target": "slides/slide2.xml",
            },
        )
        updates[PRESENTATION_RELS_PART] = ET.tostring(
            root,
            encoding="utf-8",
            xml_declaration=True,
        )
    _rewrite_package(output, updates)

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    if extra_kind == "relationship":
        assert "PRESENTATION_EXTRA_SLIDE_RELATIONSHIP" in _codes(result)
    assert "PRESENTATION_EXTRA_SLIDE_PART" in _codes(result)


def test_slide_root_requires_exact_presentationml_namespace(tmp_path: Path) -> None:
    output = tmp_path / "wrong-slide-root.pptx"
    _make_deck(output)
    slide = _package_parts(output)["ppt/slides/slide1.xml"].replace(
        PRESENTATIONML_NS.encode(),
        b"urn:not-presentationml",
    )
    _rewrite_package(output, {"ppt/slides/slide1.xml": slide})

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PRESENTATION_SLIDE_ROOT_INVALID" in _codes(result)


def test_python_pptx_reopen_count_is_independently_compared(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    pptx = pytest.importorskip("pptx")
    output = tmp_path / "python-count.pptx"
    _make_deck(output)
    fake = SimpleNamespace(slides=[], slide_width=1, slide_height=1)
    monkeypatch.setattr(pptx, "Presentation", lambda _path: fake)

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PYTHON_PPTX_SLIDE_COUNT_MISMATCH" in _codes(result)
    assert result.statistics is not None
    assert result.statistics.slide_count == 1
    assert result.statistics.python_pptx_slide_count == 0


@pytest.mark.parametrize("shape_type", [7, 8, 10, 12, 16, 18, 26, -2, 999])
def test_unsafe_or_unknown_objects_do_not_count_as_native_editable(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    shape_type: int,
) -> None:
    pptx = pytest.importorskip("pptx")
    output = tmp_path / f"unsafe-shape-{shape_type}.pptx"
    _make_deck(output)
    fake_shape = SimpleNamespace(
        shape_type=shape_type,
        element=SimpleNamespace(tag=f"{{{PRESENTATIONML_NS}}}graphicFrame"),
        left=0,
        top=0,
        width=100,
        height=100,
        has_text_frame=False,
    )
    fake_slide = SimpleNamespace(shapes=[fake_shape])
    fake = SimpleNamespace(slides=[fake_slide], slide_width=100, slide_height=100)
    monkeypatch.setattr(pptx, "Presentation", lambda _path: fake)

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "EDITABILITY_UNSUPPORTED_OBJECT" in _codes(result)
    assert "EDITABILITY_EMPTY_SLIDE" in _codes(result)
    assert result.statistics is not None
    assert result.statistics.native_object_count == 0


def test_shape_type_not_implemented_fails_closed_as_unsupported(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    pptx = pytest.importorskip("pptx")
    output = tmp_path / "shape-type-not-implemented.pptx"
    _make_deck(output)

    class UnsupportedShape:
        element = SimpleNamespace(tag=f"{{{PRESENTATIONML_NS}}}sp")
        left = 0
        top = 0
        width = 100
        height = 100
        has_text_frame = False

        @property
        def shape_type(self) -> int:
            raise NotImplementedError("shape type is not implemented")

    fake_slide = SimpleNamespace(shapes=[UnsupportedShape()])
    fake = SimpleNamespace(slides=[fake_slide], slide_width=100, slide_height=100)
    monkeypatch.setattr(pptx, "Presentation", lambda _path: fake)

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "EDITABILITY_UNSUPPORTED_OBJECT" in _codes(result)
    assert result.statistics is not None
    assert result.statistics.native_object_count == 0


def test_invalid_python_pptx_slide_shapes_fail_closed(tmp_path: Path) -> None:
    output = tmp_path / "missing-shape-tree.pptx"
    _make_deck(output)
    malformed_slide = (
        f'<p:sld xmlns:p="{PRESENTATIONML_NS}"><p:cSld/></p:sld>'
    ).encode()
    _rewrite_package(output, {"ppt/slides/slide1.xml": malformed_slide})

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "PYTHON_PPTX_SLIDE_SHAPES_INVALID" in _codes(result)
    assert result.statistics is not None
    assert result.statistics.native_object_count == 0


def test_empty_slide_is_not_accepted_as_editable(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    output = tmp_path / "empty-slide.pptx"
    presentation = pptx.Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(output)

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "EDITABILITY_EMPTY_SLIDE" in _codes(result)
    assert result.statistics is not None
    assert result.statistics.native_editable_slide_count == 0


def test_full_slide_picture_exposes_raster_coverage_statistics(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_SHAPE

    image_path = tmp_path / "pixel.png"
    image_path.write_bytes(
        base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
            "/w8AAusB9Y9Z4SIAAAAASUVORK5CYII="
        )
    )
    output = tmp_path / "raster.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        presentation.slide_width,
        presentation.slide_height,
    )
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 1_000, 1_000, 10_000, 10_000)
    presentation.save(output)

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "EDITABILITY_RASTER_DOMINANT_SLIDE" in _codes(result)
    assert result.statistics is not None
    assert result.statistics.picture_count == 1
    assert result.statistics.native_object_count == 1
    assert result.statistics.full_slide_raster_count == 1
    assert result.statistics.raster_dominant_slide_count == 1
    assert result.statistics.slides[0].picture_coverage == 1
    assert result.statistics.slides[0].native_coverage < 0.01


def test_full_bleed_background_with_material_native_editorial_system_is_not_raster_dominant(
    tmp_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_SHAPE

    image_path = tmp_path / "editorial-background.png"
    image_path.write_bytes(
        base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
            "/w8AAusB9Y9Z4SIAAAAASUVORK5CYII="
        )
    )
    output = tmp_path / "editorial-background.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path), 0, 0, presentation.slide_width, presentation.slide_height,
    )
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width, presentation.slide_height // 2,
    )
    heading = slide.shapes.add_textbox(914400, 914400, 3657600, 457200)
    heading.text = "章节一"
    presentation.save(output)

    result = inspect_presentation_topology(output)

    assert result.status == "pass"
    assert result.statistics is not None
    assert result.statistics.full_slide_raster_count == 1
    assert result.statistics.raster_dominant_slide_count == 0


def test_scaled_group_picture_uses_composed_group_geometry(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    image_path = tmp_path / "group-pixel.png"
    image_path.write_bytes(
        base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
            "/w8AAusB9Y9Z4SIAAAAASUVORK5CYII="
        )
    )
    output = tmp_path / "scaled-group-raster.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    group.shapes.add_picture(str(image_path), 0, 0, 100, 100)
    group.left = 0
    group.top = 0
    group.width = presentation.slide_width
    group.height = presentation.slide_height
    presentation.save(output)

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "EDITABILITY_RASTER_DOMINANT_SLIDE" in _codes(result)
    assert result.statistics is not None
    assert result.statistics.slides[0].picture_coverage == 1
    assert result.statistics.slides[0].full_slide_raster is True


def test_nested_group_affines_are_composed_once(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_SHAPE

    image_path = tmp_path / "nested-group-pixel.png"
    image_path.write_bytes(
        base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
            "/w8AAusB9Y9Z4SIAAAAASUVORK5CYII="
        )
    )
    output = tmp_path / "nested-group-affine.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    outer = slide.shapes.add_group_shape()
    inner = outer.shapes.add_group_shape()
    inner.shapes.add_picture(str(image_path), 0, 0, 100, 100)
    inner.left = 25
    inner.top = 25
    inner.width = 50
    inner.height = 50
    outer.left = 0
    outer.top = 0
    outer.width = presentation.slide_width
    outer.height = presentation.slide_height
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 1_000, 1_000)
    presentation.save(output)

    result = inspect_presentation_topology(output)

    assert result.status == "pass"
    assert result.statistics is not None
    assert result.statistics.slides[0].picture_coverage == 0.25
    assert result.statistics.slides[0].full_slide_raster is False


def test_overlapping_picture_coverage_uses_union_not_area_sum(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_SHAPE

    image_path = tmp_path / "overlap-pixel.png"
    image_path.write_bytes(
        base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
            "/w8AAusB9Y9Z4SIAAAAASUVORK5CYII="
        )
    )
    output = tmp_path / "overlap-union.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    picture_width = round(presentation.slide_width * 0.60)
    for _ in range(2):
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            picture_width,
            presentation.slide_height,
        )
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 1_000, 1_000)
    presentation.save(output)

    result = inspect_presentation_topology(output)

    assert result.status == "pass"
    assert result.statistics is not None
    assert result.statistics.slides[0].picture_coverage == 0.6
    assert result.statistics.slides[0].raster_dominant is False


def test_nonpositive_group_child_extent_does_not_count_as_editable(
    tmp_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_SHAPE

    output = tmp_path / "invalid-group-child-extent.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    group.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 100, 100)
    group.left = 0
    group.top = 0
    group.width = 1_000_000
    group.height = 1_000_000
    presentation.save(output)

    parts = _package_parts(output)
    slide_root = ET.fromstring(parts["ppt/slides/slide1.xml"])
    child_extent = slide_root.find(
        ".//p:grpSp/p:sp/p:spPr/a:xfrm/a:ext",
        {"p": PRESENTATIONML_NS, "a": DRAWINGML_NS},
    )
    assert child_extent is not None
    child_extent.set("cx", "0")
    _rewrite_package(
        output,
        {
            "ppt/slides/slide1.xml": ET.tostring(
                slide_root,
                encoding="utf-8",
                xml_declaration=True,
            )
        },
    )

    result = inspect_presentation_topology(output)

    assert result.status == "fail"
    assert "EDITABILITY_EMPTY_SLIDE" in _codes(result)
    assert result.statistics is not None
    assert result.statistics.native_object_count == 0


def test_axis_aligned_connectors_count_as_native_without_area(
    tmp_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_CONNECTOR

    output = tmp_path / "axis-aligned-connectors.pptx"
    presentation = pptx.Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[6])
    first.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        100_000,
        100_000,
        900_000,
        100_000,
    )
    first.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        1_000_000,
        100_000,
        1_000_000,
        900_000,
    )
    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    second.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        100_000,
        100_000,
        900_000,
        100_000,
    )
    textbox = second.shapes.add_textbox(
        100_000,
        300_000,
        1_000_000,
        300_000,
    )
    textbox.text = "Editable"
    presentation.save(output)

    result = inspect_presentation_topology(output)

    assert result.status == "pass"
    assert result.issues == ()
    assert result.statistics is not None
    assert result.statistics.shape_count == 4
    assert result.statistics.native_object_count == 4
    assert result.statistics.native_editable_slide_count == 2
    assert [
        slide.native_object_count for slide in result.statistics.slides
    ] == [2, 2]
    assert result.statistics.slides[0].native_coverage == 0
    assert result.statistics.slides[1].native_coverage > 0
