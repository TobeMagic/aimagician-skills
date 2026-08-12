"""Focused regression tests for the v6.1 cross-package OPC importer."""

from __future__ import annotations

import hashlib
import io
import json
import sys
import zipfile
import xml.etree.ElementTree as ET
from dataclasses import replace
from pathlib import Path

import pytest

SKILL_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_ROOT / "scripts"))

from window_pptx.page_template_library import (
    DEFAULT_SCORING,
    LibraryIndex,
    PageTemplate,
    _compile_governed_content_inventory,
    _discover_slots,
    _slot_graph_for,
    query_page_template_candidates,
)
from window_pptx.cli import parse_args as parse_window_pptx_args
from window_pptx.physical_assembly import (
    AssetBindingSpec,
    AssemblyImportContext,
    AssemblyPlan,
    AssemblyTargetSlide,
    AuthorityLock,
    FragmentGroupContract,
    PhysicalAssemblyError,
    SelectionEvidence,
    StyleCloneSpec,
    TextBindingSpec,
    _SourcePackageContext,
    _build_source_graph,
    _adapt_slide_text,
    _auto_authorize_governed_value,
    _apply_governed_style_clones,
    _cover_crop_values,
    _discover_picture_slots,
    _inspect_all_relationships,
    _parse_content_types,
    _prepare_governed_content_replacements,
    _resolve_output_governed_target,
    _sanitize_layout_master_fields,
    _semantic_character_count,
    _style_clone_scope_sha256,
    _style_clone_target_guard_sha256,
    _slide_structure_signature,
    _validate_fragment_group_bindings,
    _validate_query_selection_evidence,
    _verify_all_relationships,
    assemble_physical_deck,
    load_assembly_plan,
    resolve_project_file,
    verify_physical_assembly,
)
from render_window_pptx_assembly import parse_args as parse_assembly_renderer_args
from window_pptx.weak_model import (
    WeakModelValidationError,
    load_fact_store,
    validate_fact_store,
)


REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
OFFICE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def test_semantic_character_count_ignores_template_layout_whitespace() -> None:
    assert _semantic_character_count("深入" + (" " * 23) + "临床科室") == 6
    assert _semantic_character_count("A\nB\tC") == 3


def test_unbound_short_integer_template_marker_is_cleared_not_authorized() -> None:
    replacement, fact_refs, connective_ref, mode = _auto_authorize_governed_value(
        "1", facts={}, connective_copy={},
    )
    assert replacement == ""
    assert fact_refs == ()
    assert connective_ref == ""
    assert mode == "source-decoration-numeric"


def _fit_policy_slide_xml(*, fit_node: str = "spAutoFit") -> bytes:
    return (
        f'<p:sld xmlns:p="{PML_NS}" xmlns:a="{DML_NS}"><p:cSld><p:spTree>'
        '<p:sp><p:nvSpPr><p:cNvPr id="2" name="target"/></p:nvSpPr>'
        '<p:spPr><a:xfrm><a:off x="101" y="202"/><a:ext cx="303" cy="404"/>'
        '</a:xfrm></p:spPr><p:txBody><a:bodyPr wrap="none">'
        f'<a:{fit_node}/></a:bodyPr><a:p><a:r><a:rPr sz="1800"/>'
        '<a:t>Old</a:t></a:r></a:p></p:txBody></p:sp>'
        '<p:sp><p:nvSpPr><p:cNvPr id="3" name="peer"/></p:nvSpPr>'
        '<p:spPr/><p:txBody><a:bodyPr><a:spAutoFit/></a:bodyPr>'
        '<a:p><a:r><a:t>Peer</a:t></a:r></a:p></p:txBody></p:sp>'
        '</p:spTree></p:cSld></p:sld>'
    ).encode()


def test_no_autofit_mutates_only_target_body_policy_and_text() -> None:
    source = _fit_policy_slide_xml()

    actual = _adapt_slide_text(
        source,
        {"shape_2": "New"},
        allowed_slots=("shape_2", "shape_3"),
        fit_policies={"shape_2": "no-autofit"},
    )

    expected = source.replace(b"<a:spAutoFit/>", b"<a:noAutofit/>", 1).replace(
        b">Old<", b">New<", 1
    )
    assert actual == expected
    assert b'<a:off x="101" y="202"/>' in actual
    assert b'<a:rPr sz="1800"/>' in actual
    assert actual.count(b"<a:spAutoFit/>") == 1


def test_shrink_to_fit_mutates_only_target_body_policy_and_text() -> None:
    source = _fit_policy_slide_xml()

    actual = _adapt_slide_text(
        source,
        {"shape_2": "New value"},
        allowed_slots=("shape_2", "shape_3"),
        fit_policies={"shape_2": "shrink-to-fit"},
    )

    expected = source.replace(
        b"<a:spAutoFit/>",
        b'<a:normAutofit fontScale="40000" lnSpcReduction="20000"/>',
        1,
    ).replace(b'<a:rPr sz="1800"/>', b'<a:rPr sz="800"/>', 1).replace(
        b">Old<", b">New value<", 1
    )
    assert actual == expected
    assert b'<a:off x="101" y="202"/>' in actual
    assert b'<a:rPr sz="800"/>' in actual
    assert actual.count(b"<a:spAutoFit/>") == 1


def test_text_adaptation_preserves_existing_run_styles() -> None:
    source = (
        f'<p:sld xmlns:p="{PML_NS}" xmlns:a="{DML_NS}"><p:cSld><p:spTree>'
        '<p:sp><p:nvSpPr><p:cNvPr id="2" name="percent"/></p:nvSpPr>'
        '<p:txBody><a:bodyPr><a:spAutoFit/></a:bodyPr><a:p>'
        '<a:r><a:rPr sz="4000"/><a:t>44.6</a:t></a:r>'
        '<a:r><a:rPr sz="2400"/><a:t>%</a:t></a:r>'
        '</a:p></p:txBody></p:sp>'
        '<p:sp><p:nvSpPr><p:cNvPr id="3" name="title"/></p:nvSpPr>'
        '<p:txBody><a:bodyPr><a:spAutoFit/></a:bodyPr><a:p>'
        '<a:r><a:rPr sz="3200"/><a:t>202x</a:t></a:r>'
        '<a:r><a:rPr sz="3200"/><a:t>年财务决算：</a:t></a:r>'
        '<a:r><a:rPr sz="3200" b="1"/><a:t>财政项目及政府债支出</a:t></a:r>'
        '</a:p></p:txBody></p:sp>'
        '</p:spTree></p:cSld></p:sld>'
    ).encode()

    actual = _adapt_slide_text(
        source,
        {
            "shape_2": "49.1%",
            "shape_3": "2025年财政项目及政府债支出",
        },
        allowed_slots=("shape_2", "shape_3"),
    )
    root = ET.fromstring(actual)
    runs_by_shape: dict[int, list[str]] = {}
    for shape in root.iter(f"{{{PML_NS}}}sp"):
        marker = shape.find(f".//{{{PML_NS}}}cNvPr")
        if marker is None:
            continue
        runs_by_shape[int(marker.attrib["id"])] = [
            node.text or "" for node in shape.iter(f"{{{DML_NS}}}t")
        ]

    assert runs_by_shape[2] == ["49.1", "%"]
    assert runs_by_shape[3] == ["2025", "年", "财政项目及政府债支出"]
    assert b'<a:rPr sz="2400"/>' in actual
    assert b'<a:rPr sz="3200" b="1"/>' in actual


def test_structure_signature_treats_autofit_variants_as_one_governed_node() -> None:
    signatures = {
        _slide_structure_signature(_fit_policy_slide_xml(fit_node=fit_node))
        for fit_node in ("spAutoFit", "normAutofit", "noAutofit")
    }

    assert len(signatures) == 1


def test_structure_signature_treats_autofit_run_size_as_governed() -> None:
    source = _fit_policy_slide_xml()
    target = source.replace(b'<a:rPr sz="1800"/>', b'<a:rPr sz="990"/>', 1)

    assert _slide_structure_signature(source) == _slide_structure_signature(target)


def _style_clone_slide_xml() -> bytes:
    return (
        f'<p:sld xmlns:p="{PML_NS}" xmlns:a="{DML_NS}" '
        f'xmlns:r="{OFFICE_REL}"><p:cSld><p:spTree>'
        '<p:sp><p:nvSpPr><p:cNvPr id="10" name="green panel"/></p:nvSpPr>'
        '<p:spPr><a:xfrm><a:off x="10" y="20"/><a:ext cx="30" cy="40"/>'
        '</a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        '<a:solidFill><a:schemeClr val="accent1"/></a:solidFill></p:spPr>'
        '<p:txBody><a:bodyPr/><a:p><a:r><a:rPr><a:solidFill>'
        '<a:schemeClr val="accent1"/></a:solidFill></a:rPr><a:t>Source</a:t>'
        '</a:r></a:p></p:txBody></p:sp>'
        '<p:sp><p:nvSpPr><p:cNvPr id="11" name="brown panel"/></p:nvSpPr>'
        '<p:spPr><a:xfrm><a:off x="110" y="120"/><a:ext cx="130" cy="140"/>'
        '</a:xfrm><a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>'
        '<a:solidFill><a:schemeClr val="accent2"/></a:solidFill></p:spPr>'
        '<p:txBody><a:bodyPr/><a:p><a:r><a:rPr><a:solidFill>'
        '<a:schemeClr val="accent2"/></a:solidFill></a:rPr><a:t>Target</a:t>'
        '</a:r></a:p></p:txBody></p:sp>'
        '<p:sp><p:nvSpPr><p:cNvPr id="12" name="green text"/></p:nvSpPr>'
        '<p:spPr/><p:txBody><a:bodyPr/><a:p><a:r><a:rPr sz="2800">'
        '<a:solidFill><a:schemeClr val="accent1"/></a:solidFill></a:rPr>'
        '<a:t>39.83%</a:t></a:r></a:p></p:txBody></p:sp>'
        '<p:sp><p:nvSpPr><p:cNvPr id="13" name="brown text"/></p:nvSpPr>'
        '<p:spPr/><p:txBody><a:bodyPr/><a:p><a:r><a:rPr sz="2400">'
        '<a:solidFill><a:schemeClr val="accent2"/></a:solidFill></a:rPr>'
        '<a:t>27.28%</a:t></a:r></a:p></p:txBody></p:sp>'
        '<p:pic><p:nvPicPr><p:cNvPr id="14" name="green check"/>'
        '<p:cNvPicPr/><p:nvPr/></p:nvPicPr><p:blipFill><a:blip r:embed="rId7">'
        '<a:duotone><a:prstClr val="black"/><a:schemeClr val="accent1"/>'
        '</a:duotone></a:blip><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        '<p:spPr><a:xfrm><a:off x="1" y="2"/><a:ext cx="3" cy="4"/>'
        '</a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'
        '<p:pic><p:nvPicPr><p:cNvPr id="15" name="brown check"/>'
        '<p:cNvPicPr/><p:nvPr/></p:nvPicPr><p:blipFill><a:blip r:embed="rId8">'
        '<a:duotone><a:prstClr val="black"/><a:schemeClr val="accent2"/>'
        '</a:duotone></a:blip><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        '<p:spPr><a:xfrm><a:off x="11" y="12"/><a:ext cx="13" cy="14"/>'
        '</a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'
        '</p:spTree></p:cSld></p:sld>'
    ).encode()


def test_governed_style_clones_are_paint_only_and_hash_guarded() -> None:
    source = _style_clone_slide_xml()
    specs = tuple(
        StyleCloneSpec(
            source_shape_id=source_id,
            target_shape_id=target_id,
            scope=scope,
            source_style_sha256=_style_clone_scope_sha256(source, source_id, scope),
            target_guard_sha256=_style_clone_target_guard_sha256(
                source, target_id, scope
            ),
        )
        for source_id, target_id, scope in (
            (10, 11, "shape-fill"),
            (12, 13, "text-color"),
            (14, 15, "picture-color-effects"),
        )
    )

    actual = _apply_governed_style_clones(source, specs)

    for spec in specs:
        assert _style_clone_scope_sha256(
            actual, spec.source_shape_id, spec.scope
        ) == _style_clone_scope_sha256(actual, spec.target_shape_id, spec.scope)
        assert _style_clone_target_guard_sha256(
            actual, spec.target_shape_id, spec.scope
        ) == spec.target_guard_sha256
    assert b'>Target<' in actual
    assert b'x="110" y="120"' in actual
    assert b'<a:prstGeom prst="roundRect">' in actual
    assert b'r:embed="rId8"' in actual
    assert b'<a:rPr sz="2400">' in actual


def test_governed_style_clone_rejects_source_anchor_drift() -> None:
    source = _style_clone_slide_xml()
    spec = StyleCloneSpec(10, 11, "shape-fill", "0" * 64, "1" * 64)

    with pytest.raises(PhysicalAssemblyError, match="STYLE_CLONE_SOURCE_STYLE_DRIFT"):
        _apply_governed_style_clones(source, (spec,))


def _content_types(*overrides: tuple[str, str]) -> bytes:
    rows = [
        '<?xml version="1.0" encoding="UTF-8"?>',
        f'<Types xmlns="{CT_NS}">',
        '<Default ContentType="application/xml" Extension="xml"/>',
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>',
        '<Default Extension="xlsx" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"/>',
        '<Default Extension="png" ContentType="image/png"/>',
    ]
    rows.extend(
        f'<Override ContentType="{content_type}" PartName="/{part_name}"/>'
        for part_name, content_type in overrides
    )
    rows.append("</Types>")
    return "".join(rows).encode()


def _rels(*entries: tuple[str, str, str, str]) -> bytes:
    body = []
    for rel_id, rel_type, target, target_mode in entries:
        mode = f' TargetMode="{target_mode}"' if target_mode else ""
        body.append(
            f'<Relationship Id="{rel_id}" Type="{rel_type}" Target="{target}"{mode}/>'
        )
    return (
        f'<?xml version="1.0" encoding="UTF-8"?><Relationships xmlns="{REL_NS}">'
        + "".join(body)
        + "</Relationships>"
    ).encode()


def _write_zip(path: Path, parts: dict[str, bytes]) -> None:
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, data in parts.items():
            archive.writestr(name, data)


def _chart_source(path: Path, *, external_target: str | None = None) -> None:
    chart_type = "application/vnd.openxmlformats-officedocument.drawingml.chart+xml"
    style_type = "application/vnd.ms-office.chartstyle+xml"
    color_type = "application/vnd.ms-office.chartcolorstyle+xml"
    slide_type = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
    slide_entries = [
        ("rId1", f"{OFFICE_REL}/chart", "../charts/chart1.xml", ""),
    ]
    if external_target:
        slide_entries.append(
            ("rId2", f"{OFFICE_REL}/hyperlink", external_target, "External")
        )
    parts = {
        "[Content_Types].xml": _content_types(
            ("ppt/slides/slide1.xml", slide_type),
            ("ppt/charts/chart1.xml", chart_type),
            ("ppt/charts/style1.xml", style_type),
            ("ppt/charts/colors1.xml", color_type),
        ),
        "ppt/slides/slide1.xml": b"<p:sld xmlns:p='p'/>",
        "ppt/slides/_rels/slide1.xml.rels": _rels(*slide_entries),
        "ppt/charts/chart1.xml": b"<c:chartSpace xmlns:c='c'/>",
        "ppt/charts/_rels/chart1.xml.rels": _rels(
            ("rId1", f"{OFFICE_REL}/chartStyle", "style1.xml", ""),
            ("rId2", f"{OFFICE_REL}/chartColorStyle", "colors1.xml", ""),
            ("rId3", f"{OFFICE_REL}/package", "../embeddings/book.xlsx", ""),
        ),
        "ppt/charts/style1.xml": b"<cs:chartStyle xmlns:cs='cs'/>",
        "ppt/charts/colors1.xml": b"<cs:colorStyle xmlns:cs='cs'/>",
        "ppt/embeddings/book.xlsx": b"workbook-bytes",
    }
    _write_zip(path, parts)


def _governed_chart_table_source(path: Path, *, include_tag: bool = False) -> None:
    """Create one editable slide with chart caches, XLSX peers, and a table."""

    pptx = pytest.importorskip("pptx")
    chart_data_module = pytest.importorskip("pptx.chart.data")
    chart_enum_module = pytest.importorskip("pptx.enum.chart")
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    heading = slide.shapes.add_textbox(914400, 457200, 4572000, 685800)
    heading.text = "Template heading"
    chart_data = chart_data_module.ChartData()
    chart_data.categories = ("A", "B")
    chart_data.add_series("Revenue", (1, 2))
    slide.shapes.add_chart(
        chart_enum_module.XL_CHART_TYPE.COLUMN_CLUSTERED,
        914400,
        1371600,
        4572000,
        2743200,
        chart_data,
    )
    table = slide.shapes.add_table(
        2,
        2,
        5943600,
        1371600,
        2743200,
        1371600,
    ).table
    for row, values in enumerate((("Metric", "Value"), ("Revenue", "2"))):
        for column, value in enumerate(values):
            table.cell(row, column).text = value
    presentation.save(path)
    if not include_tag:
        return

    with zipfile.ZipFile(path, "r") as archive:
        parts = {name: archive.read(name) for name in archive.namelist()}
    rels_name = "ppt/slides/_rels/slide1.xml.rels"
    relationship = (
        f'<Relationship Id="rId999" Type="{OFFICE_REL}/tags" '
        'Target="../tags/tag1.xml"/>'
    ).encode()
    assert b"</Relationships>" in parts[rels_name]
    parts[rels_name] = parts[rels_name].replace(
        b"</Relationships>", relationship + b"</Relationships>", 1
    )
    parts["ppt/tags/tag1.xml"] = (
        b'<p14:tagLst xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"/>'
    )
    _write_zip(path, parts)


def _governed_replacement_specs(
    template: PageTemplate,
    replacements: dict[str, str],
    fact_refs: dict[str, str],
) -> dict[str, TextBindingSpec]:
    specs: dict[str, TextBindingSpec] = {}
    for record in template.governed_content_inventory["slots"]:
        replacement = replacements[record["source_text"]]
        binding_id = record.get("peer_group_id") or record["slot_id"]
        spec = TextBindingSpec(replacement, (fact_refs[replacement],), "auto")
        if binding_id in specs:
            assert specs[binding_id] == spec
        else:
            specs[binding_id] = spec
    return specs


def _governed_authority_files(
    tmp_path: Path,
    replacements: tuple[str, ...],
) -> tuple[Path, str, Path, str, Path, str, dict[str, str]]:
    unique_replacements = tuple(dict.fromkeys(replacements))
    fact_refs = {
        replacement: f"governed-fact-{ordinal:02d}"
        for ordinal, replacement in enumerate(unique_replacements, 1)
    }
    fact_store = tmp_path / "governed-fact-store.json"
    fact_store.write_text(
        json.dumps(
            {
                "schema_version": "1.0",
                "project": {"title": "Governed test", "language": "en-US"},
                "sources": [
                    {"id": "client-request", "kind": "request", "locator": "REQUEST.md"}
                ],
                "facts": [
                    {
                        "id": fact_refs[replacement],
                        "kind": "claim",
                        "text": replacement,
                        "language": "en-US",
                        "source_id": "client-request",
                        "locator": f"REQUEST.md#fact-{ordinal:02d}",
                        "required": False,
                    }
                    for ordinal, replacement in enumerate(unique_replacements, 1)
                ],
            },
            sort_keys=True,
        ),
        encoding="utf-8",
    )
    asset_manifest = tmp_path / "governed-asset-manifest.json"
    asset_manifest.write_text(
        json.dumps({"schema_version": "1.0", "bindings": {}}, sort_keys=True),
        encoding="utf-8",
    )
    connective_copy = tmp_path / "governed-connective-copy.json"
    connective_copy.write_text(
        json.dumps({"schema_version": "1.0", "entries": []}, sort_keys=True),
        encoding="utf-8",
    )
    return (
        fact_store,
        hashlib.sha256(fact_store.read_bytes()).hexdigest(),
        asset_manifest,
        hashlib.sha256(asset_manifest.read_bytes()).hexdigest(),
        connective_copy,
        hashlib.sha256(connective_copy.read_bytes()).hexdigest(),
        fact_refs,
    )


def _effective_worksheet_values(workbook_bytes: bytes) -> set[str]:
    with zipfile.ZipFile(io.BytesIO(workbook_bytes), "r") as workbook:
        shared_strings: list[str] = []
        if "xl/sharedStrings.xml" in workbook.namelist():
            shared_root = ET.fromstring(workbook.read("xl/sharedStrings.xml"))
            shared_strings = [
                "".join(
                    node.text or ""
                    for node in item.iter()
                    if node.tag.rsplit("}", 1)[-1] == "t"
                )
                for item in shared_root.iter()
                if item.tag.rsplit("}", 1)[-1] == "si"
            ]
        values: set[str] = set()
        for name in workbook.namelist():
            if not name.startswith("xl/worksheets/") or not name.endswith(".xml"):
                continue
            root = ET.fromstring(workbook.read(name))
            for cell in (
                node for node in root.iter() if node.tag.rsplit("}", 1)[-1] == "c"
            ):
                cell_type = cell.attrib.get("t", "")
                if cell_type == "inlineStr":
                    value = "".join(
                        node.text or ""
                        for node in cell.iter()
                        if node.tag.rsplit("}", 1)[-1] == "t"
                    )
                else:
                    value_node = next(
                        (
                            node
                            for node in list(cell)
                            if node.tag.rsplit("}", 1)[-1] == "v"
                        ),
                        None,
                    )
                    value = (value_node.text or "") if value_node is not None else ""
                    if cell_type == "s" and value.isdigit():
                        value = shared_strings[int(value)]
                if value:
                    values.add(value)
        return values


def _chart_values(chart_xml: bytes) -> set[str]:
    root = ET.fromstring(chart_xml)
    return {
        (node.text or "").strip()
        for node in root.iter()
        if node.tag.rsplit("}", 1)[-1] == "v" and (node.text or "").strip()
    }


def test_content_types_parser_preserves_source_types_independent_of_attribute_order() -> None:
    defaults, overrides = _parse_content_types(
        _content_types(("ppt/charts/style1.xml", "application/vnd.example.style+xml"))
    )
    assert ("xml", "application/xml") in defaults
    assert (
        "/ppt/charts/style1.xml",
        "application/vnd.example.style+xml",
    ) in overrides


def test_cover_crop_is_deterministic_and_preserves_aspect_ratio() -> None:
    assert _cover_crop_values(
        frame_width=16,
        frame_height=9,
        image_width=100,
        image_height=100,
    ) == (0, 21875, 0, 21875)
    assert _cover_crop_values(
        frame_width=1,
        frame_height=1,
        image_width=16,
        image_height=9,
    ) == (21875, 0, 21875, 0)


def test_source_graph_traverses_relationships_relative_to_each_owner(tmp_path: Path) -> None:
    source = tmp_path / "chart-source.pptx"
    _chart_source(source)

    graph = _build_source_graph(source, slide_number=1)

    assert set(graph.extra_parts) == {
        "ppt/charts/chart1.xml",
        "ppt/charts/style1.xml",
        "ppt/charts/colors1.xml",
        "ppt/embeddings/book.xlsx",
    }
    assert "ppt/charts/_rels/chart1.xml.rels" in graph.rels
    assert graph.content_types["ppt/charts/style1.xml"] == "application/vnd.ms-office.chartstyle+xml"


def test_source_graph_preserves_https_and_rejects_other_external_targets(tmp_path: Path) -> None:
    safe = tmp_path / "safe.pptx"
    unsafe = tmp_path / "unsafe.pptx"
    _chart_source(safe, external_target="https://example.com/evidence")
    _chart_source(unsafe, external_target="file:///tmp/secret.txt")

    assert _build_source_graph(safe).extra_parts
    with pytest.raises(PhysicalAssemblyError, match="unsafe external"):
        _build_source_graph(unsafe)


def test_output_context_shares_same_source_and_deduplicates_cross_package_media(tmp_path: Path) -> None:
    payload = b"identical-image-bytes"
    paths = [tmp_path / "a.pptx", tmp_path / "b.pptx"]
    for ordinal, path in enumerate(paths, start=1):
        _write_zip(
            path,
            {
                "[Content_Types].xml": _content_types(),
                "ppt/media/image1.png": payload,
                f"marker-{ordinal}.xml": str(ordinal).encode(),
            },
        )

    context = AssemblyImportContext()
    try:
        sources = [
            context.open_source(path, hashlib.sha256(path.read_bytes()).hexdigest())
            for path in paths
        ]
        first = context.allocate_dependency(
            sources[0], "ppt/media/image1.png", payload, "image/png", relationship_free=True
        )
        repeated = context.allocate_dependency(
            sources[0], "ppt/media/image1.png", payload, "image/png", relationship_free=True
        )
        cross_package = context.allocate_dependency(
            sources[1], "ppt/media/image1.png", payload, "image/png", relationship_free=True
        )
        scoped_a = context.allocate_dependency(
            sources[0], "ppt/notesSlides/notesSlide1.xml", b"notes", "application/xml", mutation_scope="slide_001"
        )
        scoped_b = context.allocate_dependency(
            sources[0], "ppt/notesSlides/notesSlide1.xml", b"notes", "application/xml", mutation_scope="slide_002"
        )
    finally:
        context.close()

    assert first == repeated == cross_package
    assert scoped_a != scoped_b
    assert context.same_source_reuse_count == 1
    assert context.same_source_reuse_bytes == len(payload)
    assert context.cross_source_safe_dedup_count == 1
    assert context.cross_source_safe_dedup_bytes == len(payload)
    assert context.deduplicated_part_count == 2
    assert context.deduplicated_bytes == 2 * len(payload)


def test_recursive_verifier_catches_nested_unresolved_relationship(tmp_path: Path) -> None:
    output = tmp_path / "broken.pptx"
    slide_type = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
    chart_type = "application/vnd.openxmlformats-officedocument.drawingml.chart+xml"
    _write_zip(
        output,
        {
            "[Content_Types].xml": _content_types(
                ("ppt/slides/slide1.xml", slide_type),
                ("ppt/charts/chart1.xml", chart_type),
            ),
            "ppt/slides/slide1.xml": b"<p:sld xmlns:p='p'/>",
            "ppt/slides/_rels/slide1.xml.rels": _rels(
                ("rId1", f"{OFFICE_REL}/chart", "../charts/chart1.xml", "")
            ),
            "ppt/charts/chart1.xml": b"<c:chartSpace xmlns:c='c'/>",
            "ppt/charts/_rels/chart1.xml.rels": _rels(
                ("rId1", f"{OFFICE_REL}/chartStyle", "missing-style.xml", "")
            ),
        },
    )

    ok, details, unresolved = _verify_all_relationships(output)

    assert ok is False
    assert unresolved == 1
    assert "missing-style.xml" in details


def test_relationship_audit_keeps_unsafe_external_separate_from_unresolved(
    tmp_path: Path,
) -> None:
    output = tmp_path / "unsafe.pptx"
    slide_type = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
    _write_zip(
        output,
        {
            "[Content_Types].xml": _content_types(
                ("ppt/slides/slide1.xml", slide_type),
            ),
            "ppt/slides/slide1.xml": b"<p:sld xmlns:p='p'/>",
            "ppt/slides/_rels/slide1.xml.rels": _rels(
                ("rId1", f"{OFFICE_REL}/hyperlink", "file:///tmp/secret", "External")
            ),
        },
    )

    audit = _inspect_all_relationships(output)

    assert audit.status == "fail"
    assert audit.total_relationship_count == 1
    assert audit.internal_relationship_count == 0
    assert audit.external_relationship_count == 1
    assert audit.unresolved_internal_relationships == ()
    assert len(audit.unsafe_relationships) == 1
    assert audit.unsafe_relationships[0]["reason"] == "unsafe-external-target"


def _template_for_slide(source: Path, slide_number: int) -> PageTemplate:
    package_sha = hashlib.sha256(source.read_bytes()).hexdigest()
    with zipfile.ZipFile(source) as archive:
        slide_xml = archive.read(f"ppt/slides/slide{slide_number}.xml")
        governed_content_inventory = _compile_governed_content_inventory(
            archive,
            slide_number,
        )
    slots = _discover_slots(slide_xml.decode())
    assert slots
    return PageTemplate(
        schema_version="1.0",
        page_id=f"{package_sha}:{slide_number:03d}",
        package_sha256=package_sha,
        slide_number=slide_number,
        source_path=str(source),
        source_sha256=package_sha,
        source_slide_sha256=hashlib.sha256(slide_xml).hexdigest(),
        page_role="content-blocks",
        category_names=("test",),
        style_cluster_id="test-style",
        deck_family_id="test-family",
        theme_palette=("#000000", "#FFFFFF", "#336699"),
        capacity={
            "text_slot_count": len(slots),
            "max_text_chars": sum(slot.max_chars for slot in slots),
        },
        editability="native_editable",
        certification="certified",
        visual_quality=1.0,
        structure={
            "slide_count": 1,
            "shape_count": len(slots),
            "layout_count": 1,
            "master_count": 1,
            "theme_count": 1,
            "media_count": 0,
            "page_shape_count": len(slots),
            "slide_relationship_count": 0,
            "linked_style_part_count": 0,
            "page_image_count": 0,
            "page_media_count": 0,
            "page_chart_count": 0,
            "page_table_count": 0,
            "page_native_object_count": len(slots),
        },
        slot_graph=_slot_graph_for(slots),
        requires_customer_asset=False,
        media_retention_policy="no-page-media",
        governed_content_inventory=governed_content_inventory,
    )


def test_locked_governed_chart_workbook_table_mutation_and_tag_stripping(
    tmp_path: Path,
) -> None:
    source = tmp_path / "governed-source.pptx"
    _governed_chart_table_source(source, include_tag=True)
    template = _template_for_slide(source, 1)
    inventory = template.governed_content_inventory
    assert inventory["complete"] is True
    assert inventory["closure_metadata"]["tag_part_count"] == 1
    assert inventory["closure_metadata"]["chart_part_count"] == 1
    assert inventory["closure_metadata"]["workbook_part_count"] == 1
    assert inventory["closure_metadata"]["table_count"] == 1

    governed_replacements = {
        "A": "North",
        "B": "South",
        "Revenue": "Net income",
        "1": "101",
        "2": "202",
        "Metric": "Measure",
        "Value": "Result",
    }
    heading_replacement = "Approved heading"
    (
        fact_path,
        fact_sha,
        asset_path,
        asset_sha,
        connective_path,
        connective_sha,
        fact_refs,
    ) = _governed_authority_files(
        tmp_path,
        (*governed_replacements.values(), heading_replacement),
    )
    heading_slot = template.slot_graph["text_slot_ids"][0]
    plan = AssemblyPlan(
        schema_version="1.0",
        plan_id="governed-content-plan",
        scenario_id="governed-content",
        dominant_style_cluster_id=template.style_cluster_id,
        created_at="2026-08-08T00:00:00Z",
        target_slide_count=1,
        target_slides=(
            AssemblyTargetSlide(
                1,
                template,
                {heading_slot: heading_replacement},
                "content",
                "Governed content",
                "Governed content",
                text_binding_specs={
                    heading_slot: TextBindingSpec(
                        heading_replacement,
                        (fact_refs[heading_replacement],),
                        "auto",
                    )
                },
                governed_content_binding_specs=_governed_replacement_specs(
                    template,
                    governed_replacements,
                    fact_refs,
                ),
            ),
        ),
        library_index_sha256="c" * 64,
        authority=AuthorityLock(
            fact_path.name,
            fact_sha,
            asset_path.name,
            asset_sha,
            connective_path.name,
            connective_sha,
        ),
    )
    output = tmp_path / "governed-output.pptx"

    report = assemble_physical_deck(
        plan,
        output,
        library_index_sha256="c" * 64,
        **_locked_call_kwargs(
            plan,
            tmp_path,
            fact_path,
            fact_sha,
            asset_path,
            asset_sha,
            production_checks=True,
        ),
    )

    assert report.status == "pass"
    assert report.source_residue.status == "pass"
    assert report.source_residue.governed_content_binding_count == inventory[
        "content_slot_count"
    ]
    assert report.source_residue.verified_governed_content_count == inventory[
        "content_slot_count"
    ]
    assert report.source_residue.governed_content_mismatch_count == 0
    assert report.source_residue.peer_group_mismatch_count == 0
    assert len(report.source_residue.mutation_manifest_sha256) == 64
    assert report.source_residue.tag_part_count == 0
    assert report.source_residue.tag_relationship_count == 0
    assert report.source_residue.layout_master_cached_field_count == 0
    assert report.source_residue.orphan_media_count == 0
    with zipfile.ZipFile(output, "r") as archive:
        names = archive.namelist()
        assert not any("/tags/" in name.lower() for name in names)
        assert all(
            b"/tags" not in archive.read(name).lower()
            for name in names
            if name.endswith(".rels")
        )
        chart_name = next(
            name for name in names if "/charts/" in name and name.endswith(".xml")
        )
        workbook_name = next(
            name for name in names if "/embeddings/" in name and name.endswith(".xlsx")
        )
        assert {"North", "South", "Net income", "101", "202"}.issubset(
            _chart_values(archive.read(chart_name))
        )
        assert {"North", "South", "Net income", "101", "202"}.issubset(
            _effective_worksheet_values(archive.read(workbook_name))
        )
        slide_root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
        slide_text = {
            (node.text or "").strip()
            for node in slide_root.iter()
            if node.tag.rsplit("}", 1)[-1] == "t" and (node.text or "").strip()
        }
        assert {
            heading_replacement,
            "Measure",
            "Result",
            "Net income",
            "202",
        }.issubset(slide_text)
        inventory_by_slot = {
            record["slot_id"]: record
            for record in inventory["slots"]
        }
        assert len(report.source_residue.governed_mutations) == inventory[
            "content_slot_count"
        ]
        for mutation in report.source_residue.governed_mutations:
            actual_lineage = _resolve_output_governed_target(
                archive,
                ordinal=mutation["ordinal"],
                record=inventory_by_slot[mutation["slot_id"]],
            )
            for field in (
                "slide_part",
                "shape_id",
                "slide_relationship_id",
                "chart_part",
                "chart_relationship_id",
                "target_part",
                "target_part_sha256",
            ):
                assert mutation[field] == actual_lineage[field]


def test_governed_mutations_are_isolated_when_one_source_page_is_reused(
    tmp_path: Path,
) -> None:
    source = tmp_path / "reused-governed-source.pptx"
    _governed_chart_table_source(source)
    template = _template_for_slide(source, 1)
    first_replacements = {
        "A": "North",
        "B": "South",
        "Revenue": "Net income",
        "1": "101",
        "2": "202",
        "Metric": "Measure",
        "Value": "Result",
    }
    second_replacements = {
        "A": "East",
        "B": "West",
        "Revenue": "Operating margin",
        "1": "303",
        "2": "404",
        "Metric": "KPI",
        "Value": "Actual",
    }
    fact_path, _, _, _, _, _, fact_refs = _governed_authority_files(
        tmp_path,
        (*first_replacements.values(), *second_replacements.values()),
    )
    fact_store = load_fact_store(fact_path)
    graph = _build_source_graph(source)
    first_slide = AssemblyTargetSlide(
        1,
        template,
        {},
        "content",
        "First",
        "First",
        governed_content_binding_specs=_governed_replacement_specs(
            template,
            first_replacements,
            fact_refs,
        ),
    )
    second_slide = replace(
        first_slide,
        ordinal=2,
        title="Second",
        headline="Second",
        governed_content_binding_specs=_governed_replacement_specs(
            template,
            second_replacements,
            fact_refs,
        ),
    )

    first_mutation, _ = _prepare_governed_content_replacements(
        first_slide,
        graph,
        fact_store,
        {},
    )
    second_mutation, _ = _prepare_governed_content_replacements(
        second_slide,
        graph,
        fact_store,
        {},
    )
    first_repeat, _ = _prepare_governed_content_replacements(
        first_slide,
        graph,
        fact_store,
        {},
    )

    chart_part = next(name for name in first_mutation if "/charts/" in name)
    workbook_part = next(name for name in first_mutation if name.endswith(".xlsx"))
    assert first_mutation == first_repeat
    assert first_mutation[chart_part] != second_mutation[chart_part]
    assert {"North", "South", "Net income", "101", "202"}.issubset(
        _chart_values(first_mutation[chart_part])
    )
    assert {"East", "West", "Operating margin", "303", "404"}.issubset(
        _chart_values(second_mutation[chart_part])
    )
    assert {"North", "South", "Net income", "101", "202"}.issubset(
        _effective_worksheet_values(first_mutation[workbook_part])
    )
    assert {"East", "West", "Operating margin", "303", "404"}.issubset(
        _effective_worksheet_values(second_mutation[workbook_part])
    )


def test_legacy_assembly_still_sanitizes_every_reachable_xlsx(tmp_path: Path) -> None:
    source = tmp_path / "legacy-chart-source.pptx"
    _governed_chart_table_source(source)
    template = _template_for_slide(source, 1)
    bindings = {
        slot_id: f"Approved {ordinal}"
        for ordinal, slot_id in enumerate(
            template.slot_graph["text_slot_ids"],
            1,
        )
    }
    plan = AssemblyPlan(
        schema_version="1.0",
        plan_id="legacy-workbook-sanitization",
        scenario_id="security-test",
        dominant_style_cluster_id=template.style_cluster_id,
        created_at="2026-08-08T00:00:00Z",
        target_slide_count=1,
        target_slides=(
            AssemblyTargetSlide(
                1,
                template,
                bindings,
                "content",
                "Security test",
                "Security test",
            ),
        ),
        library_index_sha256="a" * 64,
    )
    output = tmp_path / "legacy-chart-output.pptx"

    report = assemble_physical_deck(
        plan,
        output,
        library_index_sha256="a" * 64,
    )

    assert report.status == "pass"
    with zipfile.ZipFile(output, "r") as archive:
        workbook_name = next(
            name for name in archive.namelist() if name.lower().endswith(".xlsx")
        )
        with zipfile.ZipFile(io.BytesIO(archive.read(workbook_name)), "r") as workbook:
            nested_names = set(workbook.namelist())
    assert not any(name.startswith("docProps/") for name in nested_names)
    assert not any(name.startswith("xl/tables/") for name in nested_names)


def test_layout_master_fields_fail_closed_on_unknown_type_and_cache_drift() -> None:
    source_part = "ppt/slideLayouts/slideLayout1.xml"
    locator = (
        "/sldLayout[1]/cSld[1]/spTree[1]/sp[1]/txBody[1]/p[1]/fld[1]"
    )

    def layout_xml(field_type: str) -> bytes:
        return (
            f'<p:sldLayout xmlns:p="{PML_NS}" xmlns:a="{DML_NS}">'
            "<p:cSld><p:spTree><p:sp><p:txBody><a:p>"
            f'<a:fld type="{field_type}"><a:t>4</a:t></a:fld>'
            "</a:p></p:txBody></p:sp></p:spTree></p:cSld></p:sldLayout>"
        ).encode()

    with pytest.raises(
        PhysicalAssemblyError,
        match="LAYOUT_MASTER_FIELD_TYPE_UNSUPPORTED",
    ):
        _sanitize_layout_master_fields(
            source_part,
            layout_xml("customField"),
            (
                {
                    "field_id": "unused",
                    "source_part": source_part,
                    "locator": locator,
                    "field_type": "customField",
                    "source_text_sha256": hashlib.sha256(b"4").hexdigest(),
                },
            ),
        )

    field_id = "field_" + hashlib.sha256(
        f"{source_part}\0{locator}".encode()
    ).hexdigest()[:24]
    with pytest.raises(
        PhysicalAssemblyError,
        match="LAYOUT_MASTER_FIELD_CACHE_DRIFT",
    ):
        _sanitize_layout_master_fields(
            source_part,
            layout_xml("slidenum"),
            (
                {
                    "field_id": field_id,
                    "source_part": source_part,
                    "locator": locator,
                    "field_type": "slidenum",
                    "source_text_sha256": "0" * 64,
                },
            ),
        )


def test_assembly_reuses_static_dependencies_and_is_byte_deterministic(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    pptx = pytest.importorskip("pptx")
    source = tmp_path / "source.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(914400, 914400, 4572000, 914400)
    textbox.text = "Original"
    slide.notes_slide.notes_text_frame.text = "Speaker note"
    presentation.save(source)
    template = _template_for_slide(source, 1)
    slot_ids = template.slot_graph["text_slot_ids"]
    bindings = {slot_id: f"Replacement {index}" for index, slot_id in enumerate(slot_ids)}
    plan = AssemblyPlan(
        schema_version="1.0",
        plan_id="test-plan",
        scenario_id="test",
        dominant_style_cluster_id="test-style",
        created_at="2026-08-08T00:00:00Z",
        target_slide_count=2,
        target_slides=(
            AssemblyTargetSlide(1, template, bindings, "content", "One", "One"),
            AssemblyTargetSlide(2, template, bindings, "content", "Two", "Two"),
        ),
        library_index_sha256="f" * 64,
    )
    first = tmp_path / "first.pptx"
    second = tmp_path / "second.pptx"
    source_open_calls: list[Path] = []
    original_open = _SourcePackageContext.open.__func__

    def counting_open(cls: type[_SourcePackageContext], path: Path, package_sha: str):
        source_open_calls.append(path)
        return original_open(cls, path, package_sha)

    monkeypatch.setattr(_SourcePackageContext, "open", classmethod(counting_open))

    first_report = assemble_physical_deck(plan, first, library_index_sha256="f" * 64)
    second_report = assemble_physical_deck(plan, second, library_index_sha256="f" * 64)

    assert first_report.status == "pass"
    assert first_report.opc_integrity.unresolved_internal_relationship_count == 0
    assert first_report.assembly_metrics.deduplicated_part_count > 0
    assert first_report.assembly_metrics.same_source_reuse_count > 0
    assert first_report.assembly_metrics.cross_source_safe_dedup_count == 0
    assert first_report.assembly_metrics.imported_part_count > 0
    assert first_report.assembly_metrics.static_duplicate_bytes == 0
    assert first_report.libreoffice.status == "not_run"
    assert first_report.size_check.status == "not_run"
    assert first_report.distinct_page_id_count == 1
    assert len(first_report.duplicate_page_records) == 1
    size_failure = verify_physical_assembly(
        first,
        plan=plan,
        max_output_size_bytes=1,
    )
    assert size_failure.status == "fail"
    assert size_failure.size_check.status == "fail"
    assert first.read_bytes() == second.read_bytes()
    assert source_open_calls == [source.resolve(), source.resolve()]
    with zipfile.ZipFile(first) as archive:
        assert sum("/slideMasters/" in name and name.endswith(".xml") for name in archive.namelist()) == 1
        assert sum("/notesMasters/" in name and name.endswith(".xml") for name in archive.namelist()) == 1
        output_theme_count = sum(
            "/theme/" in name and name.endswith(".xml") for name in archive.namelist()
        )
        presentation_rels = archive.read("ppt/_rels/presentation.xml.rels")
        assert b"/notesMaster" in presentation_rels
    with zipfile.ZipFile(source) as archive:
        source_theme_count = sum(
            name.startswith("ppt/theme/") and name.endswith(".xml")
            for name in archive.namelist()
        )
    assert output_theme_count == source_theme_count


def test_verifier_rejects_full_slide_raster_as_native_editable(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    image_module = pytest.importorskip("PIL.Image")
    image_path = tmp_path / "full-slide.png"
    image_module.new("RGB", (1600, 900), (24, 48, 72)).save(image_path)
    output = tmp_path / "raster-only.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        presentation.slide_width,
        presentation.slide_height,
    )
    presentation.save(output)
    package_sha = hashlib.sha256(output.read_bytes()).hexdigest()
    template = PageTemplate(
        schema_version="1.0",
        page_id=f"{package_sha}:001",
        package_sha256=package_sha,
        slide_number=1,
        source_path=str(output),
        source_sha256=package_sha,
        source_slide_sha256=hashlib.sha256(
            zipfile.ZipFile(output).read("ppt/slides/slide1.xml")
        ).hexdigest(),
        page_role="content",
        category_names=("test",),
        style_cluster_id="test-style",
        deck_family_id="test-family",
        theme_palette=("#000000", "#FFFFFF", "#336699"),
        capacity={"max_text_chars": 0},
        editability="image_only",
        certification="certified",
        visual_quality=1.0,
        structure={},
        slot_graph={
            "text_slot_ids": [],
            "text_slot_count": 0,
            "reading_order": [],
            "fragment_groups": [],
            "slots": [],
        },
        requires_customer_asset=False,
        media_retention_policy="no-page-media",
    )
    plan = AssemblyPlan(
        schema_version="1.0",
        plan_id="raster-plan",
        scenario_id="raster",
        dominant_style_cluster_id="test-style",
        created_at="2026-08-08T00:00:00Z",
        target_slide_count=1,
        target_slides=(
            AssemblyTargetSlide(1, template, {}, "content", "Raster", ""),
        ),
        library_index_sha256="f" * 64,
    )

    report = verify_physical_assembly(output, plan=plan)

    assert report.status == "fail"
    assert report.editability.native_editable is False
    assert report.editability.native_editable_coverage == 0
    assert report.editability.full_slide_raster_count == 1


def test_verifier_rejects_raster_dominant_slide_with_decorative_shape(
    tmp_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    image_module = pytest.importorskip("PIL.Image")
    from pptx.enum.shapes import MSO_SHAPE

    image_path = tmp_path / "raster-dominant.png"
    image_module.new("RGB", (1600, 900), (24, 48, 72)).save(image_path)
    output = tmp_path / "raster-plus-decoration.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        presentation.slide_width,
        presentation.slide_height,
    )
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 1000, 1000, 10000, 10000)
    presentation.save(output)
    package_sha = hashlib.sha256(output.read_bytes()).hexdigest()
    with zipfile.ZipFile(output) as archive:
        source_slide_sha = hashlib.sha256(
            archive.read("ppt/slides/slide1.xml")
        ).hexdigest()
    template = PageTemplate(
        schema_version="1.0",
        page_id=f"{package_sha}:001",
        package_sha256=package_sha,
        slide_number=1,
        source_path=str(output),
        source_sha256=package_sha,
        source_slide_sha256=source_slide_sha,
        page_role="content-blocks",
        category_names=("test",),
        style_cluster_id="test-style",
        deck_family_id="test-family",
        theme_palette=("#000000", "#FFFFFF", "#336699"),
        capacity={"max_text_chars": 0},
        editability="image_only",
        certification="certified",
        visual_quality=1.0,
        structure={},
        slot_graph={
            "text_slot_ids": [],
            "text_slot_count": 0,
            "reading_order": [],
            "fragment_groups": [],
            "slots": [],
        },
        requires_customer_asset=False,
        media_retention_policy="certified-decorative-retain",
    )
    plan = AssemblyPlan(
        schema_version="1.0",
        plan_id="raster-decoration-plan",
        scenario_id="raster",
        dominant_style_cluster_id="test-style",
        created_at="2026-08-08T00:00:00Z",
        target_slide_count=1,
        target_slides=(
            AssemblyTargetSlide(1, template, {}, "content", "Raster", ""),
        ),
        library_index_sha256="f" * 64,
    )

    report = verify_physical_assembly(output, plan=plan)

    assert report.status == "fail"
    assert report.editability.raster_dominant_slide_count == 1
    assert report.editability.native_editable is False


def test_verifier_rejects_full_screenshot_with_tiny_native_decoys(
    tmp_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    image_module = pytest.importorskip("PIL.Image")
    from pptx.enum.shapes import MSO_SHAPE

    image_path = tmp_path / "screenshot.png"
    image_module.new("RGB", (1600, 900), (50, 60, 70)).save(image_path)
    output = tmp_path / "screenshot-decoys.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path), 0, 0, presentation.slide_width, presentation.slide_height
    )
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 1000, 1000, 10000, 10000)
    slide.shapes.add_shape(MSO_SHAPE.OVAL, 20000, 20000, 10000, 10000)
    textbox = slide.shapes.add_textbox(30000, 30000, 20000, 10000)
    textbox.text = "x"
    presentation.save(output)
    template = _template_for_slide(output, 1)
    plan = AssemblyPlan(
        schema_version="1.0",
        plan_id="screenshot-decoys-plan",
        scenario_id="raster",
        dominant_style_cluster_id=template.style_cluster_id,
        created_at="2026-08-08T00:00:00Z",
        target_slide_count=1,
        target_slides=(
            AssemblyTargetSlide(
                1,
                template,
                {slot_id: "x" for slot_id in template.slot_graph["text_slot_ids"]},
                "content",
                "Raster",
                "",
            ),
        ),
        library_index_sha256="f" * 64,
    )

    report = verify_physical_assembly(output, plan=plan)

    assert report.editability.raster_dominant_slide_count == 1
    assert report.editability.native_editable is False


def test_verifier_rejects_two_half_slide_screenshots_as_one_raster_surface(
    tmp_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    image_module = pytest.importorskip("PIL.Image")

    left_image = tmp_path / "left.png"
    right_image = tmp_path / "right.png"
    image_module.new("RGB", (800, 900), (20, 30, 40)).save(left_image)
    image_module.new("RGB", (800, 900), (40, 30, 20)).save(right_image)
    output = tmp_path / "split-raster.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    half_width = presentation.slide_width // 2
    slide.shapes.add_picture(
        str(left_image), 0, 0, half_width, presentation.slide_height
    )
    slide.shapes.add_picture(
        str(right_image), half_width, 0, half_width, presentation.slide_height
    )
    textbox = slide.shapes.add_textbox(1000, 1000, 20000, 10000)
    textbox.text = "x"
    presentation.save(output)
    template = _template_for_slide(output, 1)
    plan = AssemblyPlan(
        schema_version="1.0",
        plan_id="split-raster-plan",
        scenario_id="raster",
        dominant_style_cluster_id=template.style_cluster_id,
        created_at="2026-08-08T00:00:00Z",
        target_slide_count=1,
        target_slides=(
            AssemblyTargetSlide(
                1,
                template,
                {slot_id: "x" for slot_id in template.slot_graph["text_slot_ids"]},
                "content",
                "Raster",
                "",
            ),
        ),
        library_index_sha256="f" * 64,
    )

    report = verify_physical_assembly(output, plan=plan)

    assert report.editability.full_slide_raster_count == 0
    assert report.editability.raster_dominant_slide_count == 1
    assert report.editability.native_editable is False


def test_cross_package_assembly_deduplicates_identical_images(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    image_module = pytest.importorskip("PIL.Image")
    image_path = tmp_path / "shared.png"
    image_module.new("RGB", (64, 64), (12, 120, 90)).save(image_path)
    sources: list[Path] = []
    templates: list[PageTemplate] = []
    for ordinal, label in enumerate(("Alpha", "Beta"), start=1):
        source = tmp_path / f"source-{ordinal}.pptx"
        presentation = pptx.Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(914400, 914400, 4572000, 914400)
        textbox.text = label
        slide.shapes.add_picture(str(image_path), 914400, 1828800, 914400, 914400)
        presentation.save(source)
        sources.append(source)
        templates.append(_template_for_slide(source, 1))

    target_slides = []
    for ordinal, template in enumerate(templates, start=1):
        bindings = {
            slot_id: f"Replacement {ordinal}-{index}"
            for index, slot_id in enumerate(template.slot_graph["text_slot_ids"])
        }
        target_slides.append(
            AssemblyTargetSlide(
                ordinal,
                template,
                bindings,
                "content",
                f"Slide {ordinal}",
                f"Slide {ordinal}",
            )
        )
    plan = AssemblyPlan(
        schema_version="1.0",
        plan_id="cross-package-plan",
        scenario_id="test",
        dominant_style_cluster_id="test-style",
        created_at="2026-08-08T00:00:00Z",
        target_slide_count=2,
        target_slides=tuple(target_slides),
        library_index_sha256="e" * 64,
    )
    output = tmp_path / "cross-package.pptx"

    report = assemble_physical_deck(plan, output, library_index_sha256="e" * 64)

    assert report.status == "pass"
    assert report.assembly_metrics.unique_source_package_count == 2
    assert report.assembly_metrics.deduplicated_part_count > 0
    assert report.assembly_metrics.cross_source_safe_dedup_count > 0
    assert report.assembly_metrics.static_duplicate_bytes == 0
    with zipfile.ZipFile(output) as archive:
        assert sum("/media/" in name and name.endswith(".png") for name in archive.namelist()) == 1


def _locked_authority_files(
    tmp_path: Path,
    *,
    fact_text: str,
    asset_path: Path,
) -> tuple[Path, str, Path, str, Path, str]:
    fact_store_path = tmp_path / "fact-store.json"
    fact_store_path.write_text(
        json.dumps(
            {
                "schema_version": "1.0",
                "project": {"title": "Locked test", "language": "en-US"},
                "sources": [
                    {"id": "client-request", "kind": "request", "locator": "REQUEST.md"}
                ],
                "facts": [
                    {
                        "id": "approved-copy",
                        "kind": "claim",
                        "text": fact_text,
                        "language": "en-US",
                        "source_id": "client-request",
                        "locator": "REQUEST.md#copy",
                        "required": True,
                    }
                ],
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
        encoding="utf-8",
    )
    asset_sha = hashlib.sha256(asset_path.read_bytes()).hexdigest()
    asset_manifest_path = tmp_path / "asset-manifest.json"
    asset_manifest_path.write_text(
        json.dumps(
            {
                "schema_version": "1.0",
                "bindings": {
                    "approved-image": {
                        "path": asset_path.name,
                        "sha256": asset_sha,
                        "record": {
                            "id": "approved-image",
                            "kind": "photo",
                            "quality": 1.0,
                            "source": "client",
                            "license": "client-owned",
                            "retrieved_at": "2026-08-08",
                            "width_px": 64,
                            "height_px": 64,
                        },
                    }
                },
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
        encoding="utf-8",
    )
    connective_copy_path = tmp_path / "connective-copy.json"
    connective_copy_path.write_text(
        json.dumps(
            {
                "schema_version": "1.0",
                "entries": [
                    {"id": "connective-section", "text": "Approved section"},
                    {"id": "connective-clear", "text": ""},
                    {"id": "connective-source", "text": "Original"},
                ],
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
        encoding="utf-8",
    )
    return (
        fact_store_path,
        hashlib.sha256(fact_store_path.read_bytes()).hexdigest(),
        asset_manifest_path,
        hashlib.sha256(asset_manifest_path.read_bytes()).hexdigest(),
        connective_copy_path,
        hashlib.sha256(connective_copy_path.read_bytes()).hexdigest(),
    )


def _locked_plan_fixture(
    tmp_path: Path,
    *,
    replacement: str = "Approved replacement",
) -> tuple[AssemblyPlan, Path, str, Path, str, Path]:
    pptx = pytest.importorskip("pptx")
    image_module = pytest.importorskip("PIL.Image")
    original_image = tmp_path / "original.png"
    approved_image = tmp_path / "approved.png"
    image_module.new("RGB", (64, 64), (180, 40, 40)).save(original_image)
    image_module.new("RGB", (64, 64), (20, 160, 80)).save(approved_image)
    source = tmp_path / "locked-source.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(914400, 914400, 4572000, 914400)
    textbox.text = "Original"
    slide.shapes.add_picture(str(original_image), 914400, 1828800, 914400, 914400)
    presentation.save(source)
    template = _template_for_slide(source, 1)
    text_slot = template.slot_graph["text_slot_ids"][0]
    with zipfile.ZipFile(source) as archive:
        picture_slots = _discover_picture_slots(archive.read("ppt/slides/slide1.xml"))
    picture_slot = next(iter(picture_slots))
    fact_path, fact_sha, asset_path, asset_sha, connective_path, connective_sha = _locked_authority_files(
        tmp_path,
        fact_text="Approved replacement",
        asset_path=approved_image,
    )
    plan = AssemblyPlan(
        schema_version="1.0",
        plan_id="locked-plan",
        scenario_id="locked",
        dominant_style_cluster_id="test-style",
        created_at="2026-08-08T00:00:00Z",
        target_slide_count=1,
        target_slides=(
            AssemblyTargetSlide(
                1,
                template,
                {text_slot: replacement},
                "content",
                "Locked",
                "Locked",
                text_binding_specs={
                    text_slot: TextBindingSpec(
                        replacement=replacement,
                        fact_refs=("approved-copy",),
                        mode="auto",
                    )
                },
                asset_binding_specs={
                    picture_slot: AssetBindingSpec("approved-image", "cover")
                },
            ),
        ),
        library_index_sha256="d" * 64,
        authority=AuthorityLock(
            fact_path.name,
            fact_sha,
            asset_path.name,
            asset_sha,
            connective_path.name,
            connective_sha,
        ),
    )
    return plan, fact_path, fact_sha, asset_path, asset_sha, approved_image


def _locked_call_kwargs(
    plan: AssemblyPlan,
    tmp_path: Path,
    fact_path: Path,
    fact_sha: str,
    asset_path: Path,
    asset_sha: str,
    *,
    production_checks: bool = False,
) -> dict[str, object]:
    assert plan.authority is not None
    return {
        "fact_store_path": fact_path,
        "fact_store_sha256": fact_sha,
        "asset_manifest_path": asset_path,
        "asset_manifest_sha256": asset_sha,
        "connective_copy_path": tmp_path / plan.authority.connective_copy_path,
        "connective_copy_sha256": plan.authority.connective_copy_sha256,
        "project_root": tmp_path,
        "require_locked_authority": True,
        "require_libreoffice": production_checks,
        "max_output_size_bytes": 33_941_179 if production_checks else None,
    }


def test_locked_authority_binds_fact_and_asset_with_shape_evidence(tmp_path: Path) -> None:
    jsonschema = pytest.importorskip("jsonschema")
    plan, fact_path, fact_sha, asset_path, asset_sha, approved_image = _locked_plan_fixture(
        tmp_path
    )
    output = tmp_path / "locked-output.pptx"

    report = assemble_physical_deck(
        plan,
        output,
        library_index_sha256="d" * 64,
        **_locked_call_kwargs(
            plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha,
            production_checks=True,
        ),
    )

    assert report.status == "pass"
    assert report.authority.mode == "locked"
    assert report.authority.status == "pass"
    assert report.libreoffice.status == "pass"
    assert report.size_check.status == "pass"
    assert report.distinct_page_id_count == 1
    assert report.duplicate_page_records == ()
    assert report.opc_integrity.total_relationship_count == (
        report.opc_integrity.internal_relationship_count
        + report.opc_integrity.external_relationship_count
    )
    assert report.opc_integrity.unresolved_internal_relationships == ()
    assert report.opc_integrity.unsafe_relationships == ()
    assert report.assembly_metrics.imported_part_count == len(
        report.assembly_metrics.imported_parts
    )
    assert report.assembly_metrics.imported_parts == tuple(
        sorted(report.assembly_metrics.imported_parts)
    )
    assert all(
        "source_text" not in evidence
        for evidence in report.to_dict()["binding_evidence"]
    )
    assert {item.binding_kind for item in report.binding_evidence} == {"text", "asset"}
    text_evidence = next(item for item in report.binding_evidence if item.binding_kind == "text")
    asset_evidence = next(item for item in report.binding_evidence if item.binding_kind == "asset")
    assert text_evidence.fact_refs == ("approved-copy",)
    assert text_evidence.char_used <= text_evidence.char_limit
    assert asset_evidence.asset_refs == ("approved-image",)
    assert asset_evidence.replacement_sha256 == hashlib.sha256(
        approved_image.read_bytes()
    ).hexdigest()
    assert asset_evidence.image_used == asset_evidence.image_limit == 1
    assert report.source_residue.replacement_asset_count == 1
    assert report.source_residue.replacement_asset_hash_mismatch_count == 0
    assert report.source_residue.asset_slot_mismatch_count == 0
    assert report.source_residue.orphan_media_count == 0
    assert report.source_residue.tag_relationship_count == 0
    with zipfile.ZipFile(plan.target_slides[0].page_template.source_path) as source_zip:
        source_media_hashes = {
            hashlib.sha256(source_zip.read(name)).hexdigest()
            for name in source_zip.namelist()
            if "/media/" in name
        }
    with zipfile.ZipFile(output) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml")
        assert b'<a:srcRect l="0" t="0" r="0" b="0"/>' in slide_xml
        output_media_hashes = {
            hashlib.sha256(archive.read(name)).hexdigest()
            for name in archive.namelist()
            if "/media/" in name
        }
        assert source_media_hashes.isdisjoint(output_media_hashes)
        assert any(
            name.endswith(f"authority_{asset_evidence.replacement_sha256[:20]}.png")
            for name in archive.namelist()
        )
    assembly_schema = json.loads(
        (SKILL_ROOT / "schemas" / "assembly-plan.v1.schema.json").read_text()
    )
    report_schema = json.loads(
        (SKILL_ROOT / "schemas" / "physical-assembly-report.v1.schema.json").read_text()
    )
    jsonschema.validate(plan.to_dict(), assembly_schema)
    jsonschema.validate(report.to_dict(), report_schema)
    plan_path = tmp_path / "assembly-plan.json"
    plan_path.write_text(json.dumps(plan.to_dict(), ensure_ascii=False), encoding="utf-8")
    loaded = load_assembly_plan(
        plan_path,
        {plan.target_slides[0].page_template.page_id: plan.target_slides[0].page_template},
        project_root=tmp_path,
    )
    assert loaded.authority == plan.authority
    loaded_spec = next(iter(loaded.target_slides[0].text_binding_specs.values()))
    source_spec = next(iter(plan.target_slides[0].text_binding_specs.values()))
    assert loaded_spec.replacement == source_spec.replacement
    assert loaded_spec.fact_refs == source_spec.fact_refs
    assert loaded_spec.mode == "auto"
    assert loaded_spec.fit_policy == "preserve"


def test_plan_loader_rejects_unknown_text_fit_policy(tmp_path: Path) -> None:
    plan, *_ = _locked_plan_fixture(tmp_path)
    payload = plan.to_dict()
    text_slot = next(iter(plan.target_slides[0].text_binding_specs))
    payload["target_slides"][0]["bindings"][text_slot]["fit_policy"] = "shrink"
    plan_path = tmp_path / "invalid-fit-policy-plan.json"
    plan_path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")

    with pytest.raises(PhysicalAssemblyError, match="fit_policy|schema"):
        load_assembly_plan(
            plan_path,
            {
                plan.target_slides[0].page_template.page_id:
                    plan.target_slides[0].page_template
            },
            project_root=tmp_path,
        )


@pytest.mark.parametrize(
    ("replacement", "expected_ref"),
    [
        ("Approved section", "connective-section"),
        ("", "connective-clear"),
        ("Original", "connective-source"),
    ],
)
def test_locked_connective_copy_authorizes_exact_text_and_safe_clear(
    tmp_path: Path,
    replacement: str,
    expected_ref: str,
) -> None:
    plan, fact_path, _, asset_path, asset_sha, _ = _locked_plan_fixture(tmp_path)
    fact_payload = json.loads(fact_path.read_text(encoding="utf-8"))
    fact_payload["facts"][0]["required"] = False
    fact_path.write_text(
        json.dumps(fact_payload, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    fact_sha = hashlib.sha256(fact_path.read_bytes()).hexdigest()
    slide = plan.target_slides[0]
    slot_id = next(iter(slide.bindings))
    slide = replace(
        slide,
        bindings={slot_id: replacement},
        text_binding_specs={
            slot_id: TextBindingSpec(replacement=replacement, fact_refs=(), mode="auto")
        },
    )
    plan = replace(
        plan,
        target_slides=(slide,),
        authority=replace(plan.authority, fact_store_sha256=fact_sha),
    )

    report = assemble_physical_deck(
        plan,
        tmp_path / f"connective-{expected_ref}.pptx",
        library_index_sha256="d" * 64,
        **_locked_call_kwargs(
            plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha,
            production_checks=True,
        ),
    )

    evidence = next(item for item in report.binding_evidence if item.binding_kind == "text")
    assert evidence.mode == "connective"
    assert evidence.connective_ref == expected_ref
    assert evidence.fact_refs == ()


def test_locked_connective_copy_rejects_unregistered_literal(tmp_path: Path) -> None:
    plan, fact_path, _, asset_path, asset_sha, _ = _locked_plan_fixture(tmp_path)
    fact_payload = json.loads(fact_path.read_text(encoding="utf-8"))
    fact_payload["facts"][0]["required"] = False
    fact_path.write_text(
        json.dumps(fact_payload, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    fact_sha = hashlib.sha256(fact_path.read_bytes()).hexdigest()
    slide = plan.target_slides[0]
    slot_id = next(iter(slide.bindings))
    slide = replace(
        slide,
        bindings={slot_id: "Invented connective"},
        text_binding_specs={
            slot_id: TextBindingSpec(
                replacement="Invented connective",
                fact_refs=(),
                mode="auto",
            )
        },
    )
    plan = replace(
        plan,
        target_slides=(slide,),
        authority=replace(plan.authority, fact_store_sha256=fact_sha),
    )

    with pytest.raises(PhysicalAssemblyError, match="not registered connective"):
        assemble_physical_deck(
            plan,
            tmp_path / "unregistered-connective.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(
                plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha
            ),
        )


def test_locked_source_residue_requires_exact_connective_registration(
    tmp_path: Path,
) -> None:
    plan, fact_path, _, asset_path, asset_sha, _ = _locked_plan_fixture(tmp_path)
    fact_payload = json.loads(fact_path.read_text(encoding="utf-8"))
    fact_payload["facts"][0]["required"] = False
    fact_path.write_text(
        json.dumps(fact_payload, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    fact_sha = hashlib.sha256(fact_path.read_bytes()).hexdigest()
    assert plan.authority is not None
    connective_path = tmp_path / plan.authority.connective_copy_path
    connective_payload = json.loads(connective_path.read_text(encoding="utf-8"))
    connective_payload["entries"] = [
        entry
        for entry in connective_payload["entries"]
        if entry["text"] != "Original"
    ]
    connective_path.write_text(
        json.dumps(connective_payload, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    connective_sha = hashlib.sha256(connective_path.read_bytes()).hexdigest()
    slide = plan.target_slides[0]
    slot_id = next(iter(slide.bindings))
    slide = replace(
        slide,
        bindings={slot_id: "Original"},
        text_binding_specs={
            slot_id: TextBindingSpec(replacement="Original", fact_refs=(), mode="auto")
        },
    )
    plan = replace(
        plan,
        target_slides=(slide,),
        authority=replace(
            plan.authority,
            fact_store_sha256=fact_sha,
            connective_copy_sha256=connective_sha,
        ),
    )

    with pytest.raises(PhysicalAssemblyError, match="not registered connective"):
        assemble_physical_deck(
            plan,
            tmp_path / "unregistered-source-residue.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(
                plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha
            ),
        )


def test_locked_authority_accepts_whitespace_variant_of_registered_rendering(
    tmp_path: Path,
) -> None:
    replacement = "Approved   short form"
    plan, fact_path, _, asset_path, asset_sha, _ = _locked_plan_fixture(
        tmp_path,
        replacement=replacement,
    )
    fact_payload = json.loads(fact_path.read_text(encoding="utf-8"))
    fact_payload["facts"][0]["allowed_renderings"] = ["Approved short form"]
    fact_path.write_text(
        json.dumps(fact_payload, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    fact_sha = hashlib.sha256(fact_path.read_bytes()).hexdigest()
    plan = replace(
        plan,
        authority=replace(plan.authority, fact_store_sha256=fact_sha),
    )

    report = assemble_physical_deck(
        plan,
        tmp_path / "registered-whitespace.pptx",
        library_index_sha256="d" * 64,
        **_locked_call_kwargs(
            plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha,
            production_checks=True,
        ),
    )

    evidence = next(item for item in report.binding_evidence if item.binding_kind == "text")
    assert evidence.mode == "whitespace"


def test_locked_authority_rejects_invented_replacement(tmp_path: Path) -> None:
    plan, fact_path, fact_sha, asset_path, asset_sha, _ = _locked_plan_fixture(
        tmp_path,
        replacement="Invented replacement",
    )

    with pytest.raises(PhysicalAssemblyError, match="not an allowed rendering"):
        assemble_physical_deck(
            plan,
            tmp_path / "must-not-exist.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha),
        )


def test_locked_authority_requires_all_required_facts_to_be_bound(tmp_path: Path) -> None:
    plan, fact_path, fact_sha, asset_path, asset_sha, _ = _locked_plan_fixture(tmp_path)
    slide = plan.target_slides[0]
    slot_id = next(iter(slide.bindings))
    cleared_slide = replace(
        slide,
        bindings={slot_id: ""},
        text_binding_specs={
            slot_id: TextBindingSpec(
                replacement="",
                fact_refs=(),
                mode="auto",
            )
        },
    )
    cleared_plan = replace(plan, target_slides=(cleared_slide,))

    with pytest.raises(PhysicalAssemblyError, match="REQUIRED_FACTS_NOT_BOUND"):
        assemble_physical_deck(
            cleared_plan,
            tmp_path / "missing-required-fact.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(
                cleared_plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha
            ),
        )


def test_locked_authority_accepts_registered_short_fact_rendering(tmp_path: Path) -> None:
    plan, fact_path, fact_sha, asset_path, asset_sha, _ = _locked_plan_fixture(
        tmp_path,
        replacement="A",
    )
    fact_payload = json.loads(fact_path.read_text(encoding="utf-8"))
    fact_payload["facts"][0]["allowed_renderings"] = ["A"]
    jsonschema = pytest.importorskip("jsonschema")
    fact_schema = json.loads(
        (SKILL_ROOT / "schemas" / "fact-store.v1.schema.json").read_text()
    )
    jsonschema.validate(fact_payload, fact_schema)
    fact_path.write_text(
        json.dumps(fact_payload, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    fact_sha = hashlib.sha256(fact_path.read_bytes()).hexdigest()
    plan = replace(
        plan,
        authority=replace(plan.authority, fact_store_sha256=fact_sha),
    )
    output = tmp_path / "slice-output.pptx"

    report = assemble_physical_deck(
        plan,
        output,
        library_index_sha256="d" * 64,
        **_locked_call_kwargs(
            plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha,
            production_checks=True,
        ),
    )

    text_evidence = next(item for item in report.binding_evidence if item.binding_kind == "text")
    assert text_evidence.mode == "exact"
    assert text_evidence.fact_refs == ("approved-copy",)
    assert load_fact_store(fact_path).fact("approved-copy").allowed_renderings == ("A",)


def test_locked_authority_rejects_character_outside_locked_fragment_group(
    tmp_path: Path,
) -> None:
    plan, fact_path, fact_sha, asset_path, asset_sha, _ = _locked_plan_fixture(
        tmp_path,
        replacement="A",
    )

    with pytest.raises(PhysicalAssemblyError, match="not an allowed rendering"):
        assemble_physical_deck(
            plan,
            tmp_path / "character-fragment.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(
                plan,
                tmp_path,
                fact_path,
                fact_sha,
                asset_path,
                asset_sha,
                production_checks=True,
            ),
        )


def test_locked_fragment_group_reconstructs_one_fact_before_mutation(
    tmp_path: Path,
) -> None:
    plan, fact_path, _, _, _, _ = _locked_plan_fixture(tmp_path)
    payload = json.loads(fact_path.read_text(encoding="utf-8"))
    payload["facts"][0]["text"] = "AB"
    fact_path.write_text(
        json.dumps(payload, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    facts = {fact.id: fact for fact in load_fact_store(fact_path).active_facts()}
    source_slide = plan.target_slides[0]
    slot_ids = ("shape_101", "shape_102", "shape_103")
    slide = replace(
        source_slide,
        bindings={slot_ids[0]: "A", slot_ids[1]: "B", slot_ids[2]: ""},
        text_binding_specs={
            slot_ids[0]: TextBindingSpec("A", ("approved-copy",), "auto"),
            slot_ids[1]: TextBindingSpec("B", ("approved-copy",), "auto"),
            slot_ids[2]: TextBindingSpec("", (), "auto"),
        },
    )
    contract = FragmentGroupContract(
        ordinal=1,
        page_id=slide.page_template.page_id,
        group_id="fragment_01",
        ordered_slot_ids=slot_ids,
    )

    authorized = _validate_fragment_group_bindings(
        slide,
        contracts=(contract,),
        facts=facts,
        connective_copy={"": "connective-clear"},
    )

    assert authorized == {
        slot_ids[0]: ("character", ""),
        slot_ids[1]: ("character", ""),
        slot_ids[2]: ("connective", "connective-clear"),
    }


def test_locked_fragment_group_rejects_mixed_fact_refs(tmp_path: Path) -> None:
    plan, fact_path, _, _, _, _ = _locked_plan_fixture(tmp_path)
    payload = json.loads(fact_path.read_text(encoding="utf-8"))
    payload["facts"][0]["text"] = "AB"
    payload["facts"].append(
        {
            **payload["facts"][0],
            "id": "other-copy",
        }
    )
    fact_path.write_text(
        json.dumps(payload, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    facts = {fact.id: fact for fact in load_fact_store(fact_path).active_facts()}
    source_slide = plan.target_slides[0]
    slot_ids = ("shape_101", "shape_102")
    slide = replace(
        source_slide,
        bindings={slot_ids[0]: "A", slot_ids[1]: "B"},
        text_binding_specs={
            slot_ids[0]: TextBindingSpec("A", ("approved-copy",), "auto"),
            slot_ids[1]: TextBindingSpec("B", ("other-copy",), "auto"),
        },
    )

    with pytest.raises(PhysicalAssemblyError, match="FACT_REF_DRIFT"):
        _validate_fragment_group_bindings(
            slide,
            contracts=(
                FragmentGroupContract(
                    1,
                    slide.page_template.page_id,
                    "fragment_01",
                    slot_ids,
                ),
            ),
            facts=facts,
            connective_copy={"": "connective-clear"},
        )


def test_locked_authority_rejects_unregistered_fact_substring(tmp_path: Path) -> None:
    plan, fact_path, fact_sha, asset_path, asset_sha, _ = _locked_plan_fixture(
        tmp_path,
        replacement="Approved",
    )

    with pytest.raises(PhysicalAssemblyError, match="not an allowed rendering"):
        assemble_physical_deck(
            plan,
            tmp_path / "unregistered-substring.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha),
        )


def test_locked_authority_rejects_slot_capacity_overflow(tmp_path: Path) -> None:
    plan, fact_path, fact_sha, asset_path, asset_sha, _ = _locked_plan_fixture(tmp_path)
    slide = plan.target_slides[0]
    slot_id = next(iter(slide.bindings))
    slots = [dict(slot) for slot in slide.page_template.slot_graph["slots"]]
    slots[0]["max_chars"] = 4
    constrained_template = replace(
        slide.page_template,
        slot_graph={**slide.page_template.slot_graph, "slots": slots},
        capacity={**slide.page_template.capacity, "max_text_chars": 4},
    )
    constrained_slide = replace(slide, page_template=constrained_template)
    constrained_plan = replace(plan, target_slides=(constrained_slide,))

    with pytest.raises(PhysicalAssemblyError, match="max_chars exceeded"):
        assemble_physical_deck(
            constrained_plan,
            tmp_path / "capacity-fail.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(
                constrained_plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha
            ),
        )


def test_locked_authority_rejects_item_capacity_overflow(tmp_path: Path) -> None:
    plan, fact_path, _, asset_path, asset_sha, _ = _locked_plan_fixture(tmp_path)
    fact_payload = json.loads(fact_path.read_text(encoding="utf-8"))
    fact_payload["facts"][0]["text"] = "Line one\nLine two"
    fact_path.write_text(
        json.dumps(fact_payload, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    fact_sha = hashlib.sha256(fact_path.read_bytes()).hexdigest()
    slide = plan.target_slides[0]
    slot_id = next(iter(slide.bindings))
    replacement = "Line one\nLine two"
    updated_slide = replace(
        slide,
        bindings={slot_id: replacement},
        text_binding_specs={
            slot_id: TextBindingSpec(
                replacement=replacement,
                fact_refs=("approved-copy",),
                mode="exact",
            )
        },
    )
    updated_plan = replace(
        plan,
        target_slides=(updated_slide,),
        authority=replace(plan.authority, fact_store_sha256=fact_sha),
    )

    with pytest.raises(PhysicalAssemblyError, match="item capacity exceeded"):
        assemble_physical_deck(
            updated_plan,
            tmp_path / "item-capacity-fail.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(
                updated_plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha
            ),
        )


def test_locked_authority_rejects_manifest_digest_drift(tmp_path: Path) -> None:
    plan, fact_path, fact_sha, asset_path, asset_sha, _ = _locked_plan_fixture(tmp_path)
    asset_path.write_text(asset_path.read_text() + "\n", encoding="utf-8")

    with pytest.raises(PhysicalAssemblyError, match="ASSET_MANIFEST_FINGERPRINT_MISMATCH"):
        assemble_physical_deck(
            plan,
            tmp_path / "manifest-fail.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha),
        )


@pytest.mark.parametrize("raw_asset_path", ["../outside.png", "ABSOLUTE"])
def test_locked_asset_path_must_be_project_relative(
    tmp_path: Path,
    raw_asset_path: str,
) -> None:
    plan, fact_path, fact_sha, manifest_path, _, approved_image = _locked_plan_fixture(
        tmp_path
    )
    payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    payload["bindings"]["approved-image"]["path"] = (
        str(approved_image) if raw_asset_path == "ABSOLUTE" else raw_asset_path
    )
    manifest_path.write_text(
        json.dumps(payload, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    manifest_sha = hashlib.sha256(manifest_path.read_bytes()).hexdigest()
    assert plan.authority is not None
    plan = replace(
        plan,
        authority=replace(plan.authority, asset_manifest_sha256=manifest_sha),
    )

    with pytest.raises(PhysicalAssemblyError, match="PATH_NOT_PROJECT_RELATIVE"):
        assemble_physical_deck(
            plan,
            tmp_path / "asset-path-escape.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(
                plan,
                tmp_path,
                fact_path,
                fact_sha,
                manifest_path,
                manifest_sha,
            ),
        )


def test_locked_asset_path_rejects_symlink_component(tmp_path: Path) -> None:
    plan, fact_path, fact_sha, manifest_path, _, approved_image = _locked_plan_fixture(
        tmp_path
    )
    linked_dir = tmp_path / "linked-assets"
    real_dir = tmp_path / "real-assets"
    real_dir.mkdir()
    copied = real_dir / "approved.png"
    copied.write_bytes(approved_image.read_bytes())
    linked_dir.symlink_to(real_dir, target_is_directory=True)
    payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    payload["bindings"]["approved-image"]["path"] = "linked-assets/approved.png"
    payload["bindings"]["approved-image"]["sha256"] = hashlib.sha256(
        copied.read_bytes()
    ).hexdigest()
    manifest_path.write_text(
        json.dumps(payload, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    manifest_sha = hashlib.sha256(manifest_path.read_bytes()).hexdigest()
    assert plan.authority is not None
    plan = replace(
        plan,
        authority=replace(plan.authority, asset_manifest_sha256=manifest_sha),
    )

    with pytest.raises(PhysicalAssemblyError, match="SYMLINK_REJECTED"):
        assemble_physical_deck(
            plan,
            tmp_path / "asset-symlink.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(
                plan,
                tmp_path,
                fact_path,
                fact_sha,
                manifest_path,
                manifest_sha,
            ),
        )


def test_locked_authority_rejects_duplicate_page_ids(tmp_path: Path) -> None:
    plan, fact_path, fact_sha, asset_path, asset_sha, _ = _locked_plan_fixture(tmp_path)
    duplicate_slide = replace(plan.target_slides[0], ordinal=2)
    duplicate_plan = replace(
        plan,
        target_slide_count=2,
        target_slides=(plan.target_slides[0], duplicate_slide),
    )

    with pytest.raises(PhysicalAssemblyError, match="DUPLICATE_PAGE_ID"):
        assemble_physical_deck(
            duplicate_plan,
            tmp_path / "duplicate-page.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(
                duplicate_plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha
            ),
        )


def test_locked_authority_rejects_non_direct_page_at_assembly_boundary(
    tmp_path: Path,
) -> None:
    plan, fact_path, fact_sha, asset_path, asset_sha, _ = _locked_plan_fixture(tmp_path)
    slide = plan.target_slides[0]
    blocked_template = replace(
        slide.page_template,
        direct_use=False,
        pool="reference-only",
        decision="reference-only",
    )
    blocked_plan = replace(
        plan,
        target_slides=(replace(slide, page_template=blocked_template),),
    )

    with pytest.raises(PhysicalAssemblyError, match="PAGE_NOT_DIRECT_USE"):
        assemble_physical_deck(
            blocked_plan,
            tmp_path / "non-direct.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(
                blocked_plan,
                tmp_path,
                fact_path,
                fact_sha,
                asset_path,
                asset_sha,
            ),
        )


def test_phase49_profile_mechanically_requires_exactly_fifteen_slides(
    tmp_path: Path,
) -> None:
    plan, fact_path, fact_sha, asset_path, asset_sha, _ = _locked_plan_fixture(tmp_path)

    with pytest.raises(
        PhysicalAssemblyError,
        match="PHASE49_SEQUENCE_REQUIRES_15_SLIDES",
    ):
        assemble_physical_deck(
            plan,
            tmp_path / "not-fifteen.pptx",
            library_index_sha256="d" * 64,
            acceptance_profile="phase49-work-report-15",
            **_locked_call_kwargs(
                plan,
                tmp_path,
                fact_path,
                fact_sha,
                asset_path,
                asset_sha,
            ),
        )


def test_failed_candidate_is_not_promoted_to_requested_output(tmp_path: Path) -> None:
    plan, fact_path, fact_sha, asset_path, asset_sha, _ = _locked_plan_fixture(tmp_path)
    output = tmp_path / "rejected-candidate.pptx"

    report = assemble_physical_deck(
        plan,
        output,
        library_index_sha256="d" * 64,
        **_locked_call_kwargs(
            plan,
            tmp_path,
            fact_path,
            fact_sha,
            asset_path,
            asset_sha,
            production_checks=False,
        ),
    )

    assert report.status == "fail"
    assert output.exists() is False
    assert not tuple(tmp_path.glob(".*.candidate.pptx"))


def test_assembler_rejects_preexisting_output_before_building_candidate(
    tmp_path: Path,
) -> None:
    plan, fact_path, fact_sha, asset_path, asset_sha, _ = _locked_plan_fixture(tmp_path)
    output = tmp_path / "stale-output.pptx"
    stale = b"stale-user-owned-output"
    output.write_bytes(stale)

    with pytest.raises(PhysicalAssemblyError, match="OUTPUT_ALREADY_EXISTS"):
        assemble_physical_deck(
            plan,
            output,
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(
                plan,
                tmp_path,
                fact_path,
                fact_sha,
                asset_path,
                asset_sha,
            ),
        )

    assert output.read_bytes() == stale


def test_query_bundle_selection_is_hash_bound_to_candidate_rank(
    tmp_path: Path,
) -> None:
    plan, _, _, _, _, _ = _locked_plan_fixture(tmp_path)
    slide = plan.target_slides[0]
    library_index = LibraryIndex(
        schema_version="4.0",
        library_id="test-library",
        compiled_at="2026-08-08T00:00:00Z",
        source_core_schema="test",
        private_root_sha256="b" * 64,
        source_package_count=1,
        source_package_index={
            slide.page_template.package_sha256: {
                "page_count": 1,
                "source_sha256": slide.page_template.package_sha256,
                "source_size_bytes": Path(slide.page_template.source_path).stat().st_size,
            }
        },
        page_template_count=1,
        role_index={slide.page_template.page_role: 1},
        style_cluster_index={slide.page_template.style_cluster_id: 1},
        deck_family_index={slide.page_template.deck_family_id: 1},
        category_index={"test": 1},
        scoring=dict(DEFAULT_SCORING),
        dominant_style_cluster_id=slide.page_template.style_cluster_id,
        compatible_style_cluster_ids=(slide.page_template.style_cluster_id,),
        page_templates=(slide.page_template,),
    )
    candidates = query_page_template_candidates(
        library_index,
        role=slide.page_template.page_role,
        capacity_budget=0,
        semantic_categories=(),
        style_cluster=slide.page_template.style_cluster_id,
        asset_requirements=(),
        customer_assets_available=False,
        limit=6,
        allow_fallback=False,
        direct_use_only=True,
        include_ineligible=False,
    )
    assert len(candidates) == 1
    score_total = candidates[0].scores.total
    selection = SelectionEvidence(
        query_id="slide-01",
        candidate_rank=1,
        score_total=score_total,
        selection_reason="Best certified content match",
        fallback_reason=None,
    )
    selected_plan = replace(
        plan,
        target_slides=(replace(slide, selection_evidence=selection),),
    )
    bundle = {
        "schema_version": "page-template-query-bundle.v1",
        "request_sha256": "a" * 64,
        "library_index_sha256": plan.library_index_sha256,
        "library_resolution_source": "explicit-private-root",
        "query_count": 1,
        "queries": [
            {
                "target_ordinal": 1,
                "query_id": "slide-01",
                "result": {
                    "schema_version": "page-template-query-result.v1",
                    "library_index_sha256": plan.library_index_sha256,
                    "required_source_ordinal": 1,
                    "role": slide.page_template.page_role,
                    "capacity_budget": 0,
                    "semantic_categories": [],
                    "style_cluster": slide.page_template.style_cluster_id,
                    "asset_requirements": [],
                    "customer_assets_available": False,
                    "limit": 6,
                    "allow_fallback": False,
                    "direct_use_only": True,
                    "include_ineligible": False,
                    "weights": dict(DEFAULT_SCORING),
                    "count": 1,
                    "eligible_count": 1,
                    "candidates": [candidates[0].to_dict()],
                },
            }
        ],
    }
    bundle_path = tmp_path / "template-query-results.v1.json"
    bundle_path.write_text(
        json.dumps(bundle, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    selected_plan = replace(
        selected_plan,
        query_bundle_path=bundle_path.name,
        query_bundle_sha256=hashlib.sha256(bundle_path.read_bytes()).hexdigest(),
    )

    evidence = _validate_query_selection_evidence(
        selected_plan,
        project_root=tmp_path,
        library_index=library_index,
    )

    assert evidence.status == "pass"
    assert evidence.query_bundle_sha256 == selected_plan.query_bundle_sha256
    assert evidence.query_count == evidence.selected_count == 1
    phase49_evidence = _validate_query_selection_evidence(
        selected_plan,
        project_root=tmp_path,
        library_index=library_index,
        require_phase49_ordinals=True,
    )
    assert phase49_evidence.status == "pass"

    ordinal_tamper = json.loads(json.dumps(bundle))
    ordinal_tamper["queries"][0]["result"]["required_source_ordinal"] = 2
    bundle_path.write_text(
        json.dumps(ordinal_tamper, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    ordinal_tampered_plan = replace(
        selected_plan,
        query_bundle_sha256=hashlib.sha256(bundle_path.read_bytes()).hexdigest(),
    )
    with pytest.raises(PhysicalAssemblyError, match="ORDINAL_NO_MATCH"):
        _validate_query_selection_evidence(
            ordinal_tampered_plan,
            project_root=tmp_path,
            library_index=library_index,
        )

    missing_phase49_ordinal = json.loads(json.dumps(bundle))
    del missing_phase49_ordinal["queries"][0]["result"][
        "required_source_ordinal"
    ]
    bundle_path.write_text(
        json.dumps(missing_phase49_ordinal, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    missing_ordinal_plan = replace(
        selected_plan,
        query_bundle_sha256=hashlib.sha256(bundle_path.read_bytes()).hexdigest(),
    )
    with pytest.raises(
        PhysicalAssemblyError,
        match="PHASE49_QUERY_SOURCE_ORDINAL_MISMATCH",
    ):
        _validate_query_selection_evidence(
            missing_ordinal_plan,
            project_root=tmp_path,
            library_index=library_index,
            require_phase49_ordinals=True,
        )

    bundle_path.write_text(
        json.dumps(bundle, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    bundle["queries"][0]["result"]["candidates"][0]["scores"]["total"] = 0.1
    bundle_path.write_text(
        json.dumps(bundle, ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    tampered_plan = replace(
        selected_plan,
        query_bundle_sha256=hashlib.sha256(bundle_path.read_bytes()).hexdigest(),
    )
    with pytest.raises(PhysicalAssemblyError, match="RECOMPUTE_MISMATCH"):
        _validate_query_selection_evidence(
            tampered_plan,
            project_root=tmp_path,
            library_index=library_index,
        )


def test_production_assembly_cli_requires_external_authority_locks() -> None:
    base = [
        "--project-dir",
        "/tmp/client",
        "--render-assembly-plan",
        "--assembly-plan",
        "assembly-plan.json",
        "--output",
        "output.pptx",
    ]
    with pytest.raises(SystemExit):
        parse_window_pptx_args(base)

    args = parse_window_pptx_args(
        base
        + [
            "--fact-store",
            "fact-store.json",
            "--fact-store-sha256",
            "a" * 64,
            "--asset-manifest",
            "asset-manifest.json",
            "--asset-manifest-sha256",
            "b" * 64,
            "--connective-copy",
            "connective-copy.json",
            "--connective-copy-sha256",
            "c" * 64,
        ]
    )
    assert args.fact_store_sha256 == "a" * 64
    assert args.asset_manifest_sha256 == "b" * 64
    assert args.connective_copy_sha256 == "c" * 64


def test_standalone_renderer_requires_project_root_and_connective_lock() -> None:
    with pytest.raises(SystemExit):
        parse_assembly_renderer_args(
            [
                "--assembly-plan", "assembly-plan.json",
                "--fact-store", "fact-store.json",
                "--fact-store-sha256", "a" * 64,
                "--asset-manifest", "asset-manifest.json",
                "--asset-manifest-sha256", "b" * 64,
                "--output", "output/final.pptx",
            ]
        )
    args = parse_assembly_renderer_args(
        [
            "--project-root", "/tmp/client",
            "--assembly-plan", "assembly-plan.json",
            "--fact-store", "fact-store.json",
            "--fact-store-sha256", "a" * 64,
            "--asset-manifest", "asset-manifest.json",
            "--asset-manifest-sha256", "b" * 64,
            "--connective-copy", "connective-copy.json",
            "--connective-copy-sha256", "c" * 64,
            "--output", "output/final.pptx",
        ]
    )
    assert args.max_output_size_bytes == 33_941_179
    assert args.connective_copy_sha256 == "c" * 64


def test_project_file_resolver_rejects_escape_and_symlink(tmp_path: Path) -> None:
    project = tmp_path / "project"
    project.mkdir()
    authority = project / "fact-store.json"
    authority.write_text("{}", encoding="utf-8")
    assert resolve_project_file(
        "fact-store.json",
        project,
        label="FACT_STORE",
    ) == authority
    with pytest.raises(PhysicalAssemblyError, match="NOT_PROJECT_RELATIVE"):
        resolve_project_file("../fact-store.json", project, label="FACT_STORE")
    with pytest.raises(PhysicalAssemblyError, match="NOT_PROJECT_RELATIVE"):
        resolve_project_file(str(authority), project, label="FACT_STORE")
    linked = project / "linked.json"
    linked.symlink_to(authority)
    with pytest.raises(PhysicalAssemblyError, match="SYMLINK_REJECTED"):
        resolve_project_file("linked.json", project, label="FACT_STORE")


@pytest.mark.parametrize(
    "invalid_renderings",
    [
        ["duplicate", "duplicate"],
        [""],
        ["x"] * 65,
        ["x" * 4001],
    ],
)
def test_fact_store_rejects_invalid_allowed_renderings(
    tmp_path: Path,
    invalid_renderings: list[str],
) -> None:
    image_module = pytest.importorskip("PIL.Image")
    asset = tmp_path / "asset.png"
    image_module.new("RGB", (2, 2), (1, 2, 3)).save(asset)
    fact_path, _, _, _, _, _ = _locked_authority_files(
        tmp_path,
        fact_text="Authoritative fact",
        asset_path=asset,
    )
    payload = json.loads(fact_path.read_text(encoding="utf-8"))
    payload["facts"][0]["allowed_renderings"] = invalid_renderings

    with pytest.raises(WeakModelValidationError, match="allowed_renderings"):
        validate_fact_store(payload)


def test_locked_authority_rejects_connective_digest_drift(tmp_path: Path) -> None:
    plan, fact_path, fact_sha, asset_path, asset_sha, _ = _locked_plan_fixture(tmp_path)
    assert plan.authority is not None
    connective_path = tmp_path / plan.authority.connective_copy_path
    connective_path.write_text(connective_path.read_text() + "\n", encoding="utf-8")

    with pytest.raises(PhysicalAssemblyError, match="CONNECTIVE_COPY_FINGERPRINT_MISMATCH"):
        assemble_physical_deck(
            plan,
            tmp_path / "connective-drift.pptx",
            library_index_sha256="d" * 64,
            **_locked_call_kwargs(plan, tmp_path, fact_path, fact_sha, asset_path, asset_sha),
        )
