"""Plan-aware release QA for PPTX Studio physical assemblies.

The only automatic visual intervention is the compiler-owned
``shrink-to-fit`` text policy applied before assembly.  This module is
intentionally fail-closed for every other visual defect: changing source
geometry, palette, or arbitrary shapes after a review would destroy template
lineage.  The caller must select another certified page or shorten approved
copy and reassemble instead.
"""

from __future__ import annotations

import hashlib
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Mapping, Sequence

from window_pptx.physical_assembly import AssemblyPlan, PhysicalAssemblyReport
from window_pptx.physical_rule_qa import RuleFinding, run_physical_rule_qa


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _area_overlap(first: Any, second: Any) -> float:
    left = max(first.left, second.left)
    top = max(first.top, second.top)
    right = min(first.left + first.width, second.left + second.width)
    bottom = min(first.top + first.height, second.top + second.height)
    if right <= left or bottom <= top:
        return 0.0
    overlap = (right - left) * (bottom - top)
    smaller = min(max(1, first.width * first.height), max(1, second.width * second.height))
    return overlap / smaller


@dataclass(frozen=True)
class StudioQAReport:
    schema_version: str
    status: str
    output_sha256: str
    physical_assembly_status: str
    rule_qa_status: str
    repairs: tuple[Mapping[str, Any], ...]
    blockers: tuple[Mapping[str, Any], ...]
    warnings: tuple[Mapping[str, Any], ...]
    checks: tuple[str, ...]

    def to_dict(self) -> dict[str, Any]:
        return {
            "schema_version": self.schema_version,
            "status": self.status,
            "output_sha256": self.output_sha256,
            "physical_assembly_status": self.physical_assembly_status,
            "rule_qa_status": self.rule_qa_status,
            "repairs": [dict(item) for item in self.repairs],
            "blockers": [dict(item) for item in self.blockers],
            "warnings": [dict(item) for item in self.warnings],
            "checks": list(self.checks),
        }


def _finding(finding: RuleFinding) -> dict[str, Any]:
    return finding.to_dict()


def _intentional_fragment_overlap_ids(
    lineage: Mapping[str, Any], *, ordinal: int,
) -> set[frozenset[str]]:
    """Return only lineage-proven title-lockup overlap pairs for one slide.

    A certified editorial title may deliberately overlap its one-character
    boxes.  That is not a post-assembly collision if (and only if) the exact
    shapes were produced by one ``replace_fragment_text`` operation.  This is
    deliberately narrower than a geometry whitelist: ordinary text boxes,
    separate fragment components, and unbound source text still fail closed.
    """

    slides = lineage.get("slides")
    if not isinstance(slides, list):
        return set()
    for slide in slides:
        if not isinstance(slide, Mapping) or slide.get("ordinal") != ordinal:
            continue
        allowed: set[frozenset[str]] = set()
        for binding in slide.get("fragment_title_bindings", []):
            if not isinstance(binding, Mapping):
                continue
            shape_ids = binding.get("shape_ids")
            if not isinstance(shape_ids, list):
                continue
            normalized = [item for item in shape_ids if isinstance(item, str)]
            for index, first in enumerate(normalized):
                for second in normalized[index + 1:]:
                    allowed.add(frozenset((first, second)))
        return allowed
    return set()


def _visual_checks(output: Path, *, lineage: Mapping[str, Any]) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    blockers: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    try:
        from pptx import Presentation
    except ImportError:
        return blockers, [{"rule": "python-pptx", "severity": "warning", "slide": None, "message": "python-pptx unavailable; visual geometry check skipped"}]
    try:
        deck = Presentation(str(output))
    except Exception as exc:
        return [{"rule": "open", "severity": "blocker", "slide": None, "message": str(exc)}], warnings
    for number, slide in enumerate(deck.slides, 1):
        intentional_fragment_pairs = _intentional_fragment_overlap_ids(
            lineage, ordinal=number,
        )
        populated = [shape for shape in slide.shapes if bool(getattr(shape, "has_text_frame", False)) and shape.text.strip()]
        character_count = sum(len(shape.text.strip()) for shape in populated)
        if character_count > 700:
            blockers.append({"rule": "density", "severity": "blocker", "slide": number, "message": f"visible text is too dense ({character_count} characters)"})
        elif character_count < 12:
            warnings.append({"rule": "density", "severity": "warning", "slide": number, "message": "very low text density; confirm deliberate visual page"})
        for index, first in enumerate(populated):
            for second in populated[index + 1:]:
                if _area_overlap(first, second) > 0.55:
                    pair = frozenset((f"shape_{first.shape_id}", f"shape_{second.shape_id}"))
                    if pair in intentional_fragment_pairs:
                        continue
                    blockers.append({"rule": "text-overlap", "severity": "blocker", "slide": number, "message": f"populated text shapes {first.shape_id} and {second.shape_id} materially overlap"})
    lineage_slides = lineage.get("slides")
    if not isinstance(lineage_slides, list):
        blockers.append({"rule": "lineage", "severity": "blocker", "slide": None, "message": "lineage slides are missing"})
    else:
        for previous, current in zip(lineage_slides, lineage_slides[1:]):
            if isinstance(previous, Mapping) and isinstance(current, Mapping) and previous.get("catalog_page_id") == current.get("catalog_page_id"):
                warnings.append({"rule": "repetition", "severity": "warning", "slide": current.get("ordinal"), "message": "adjacent output slides use the same certified catalog page"})
    return blockers, warnings


def run_studio_qa(
    output_path: Path | str,
    *,
    plan: AssemblyPlan,
    physical_report: PhysicalAssemblyReport,
    lineage: Mapping[str, Any],
) -> StudioQAReport:
    """Run release QA and return only evidence-bound repair decisions.

    A caller can promote only a ``pass`` report.  ``repairs`` documents the
    pre-assembly safe text-fit policy; a blocker requires re-planning rather
    than an uncontrolled post-hoc edit.
    """

    output = Path(output_path).expanduser().resolve(strict=False)
    blockers: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    repairs: list[dict[str, Any]] = []
    if not output.is_file():
        blockers.append({"rule": "output", "severity": "blocker", "slide": None, "message": "assembly output is missing"})
        output_sha = ""
    else:
        output_sha = _sha256_file(output)
    if physical_report.status != "pass":
        blockers.append({"rule": "physical-assembly", "severity": "blocker", "slide": None, "message": "portable OPC/import verification did not pass"})
    if output_sha and physical_report.output_sha256 != output_sha:
        blockers.append({"rule": "output-lineage", "severity": "blocker", "slide": None, "message": "physical report and actual output fingerprints differ"})
    rule_report = run_physical_rule_qa(output, plan=plan)
    blockers.extend(_finding(item) for item in rule_report.blocking_findings)
    warnings.extend(_finding(item) for item in rule_report.warnings)
    visual_blockers, visual_warnings = _visual_checks(output, lineage=lineage) if output_sha else ([], [])
    blockers.extend(visual_blockers)
    warnings.extend(visual_warnings)
    for slide in plan.target_slides:
        for slot_id, spec in slide.text_binding_specs.items():
            if spec.fit_policy == "shrink-to-fit" and spec.replacement != "":
                repairs.append({"kind": "shrink-to-fit", "ordinal": slide.ordinal, "slot_id": slot_id, "status": "applied-before-assembly"})
    checks = (
        "physical-opc-lineage", "pptx-open-editability", "text-overflow-fit-policy",
        "text-bounds", "text-overlap", "placeholder-and-source-residue",
        "typography", "image-cover-crop", "density", "adjacent-repetition", "style-coherence",
    )
    return StudioQAReport(
        "1.0", "pass" if not blockers else "fail", output_sha,
        physical_report.status, rule_report.status, tuple(repairs),
        tuple(blockers), tuple(warnings), checks,
    )


__all__ = ["StudioQAReport", "run_studio_qa"]
