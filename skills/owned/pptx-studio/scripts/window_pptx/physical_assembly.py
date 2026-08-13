"""Physical cross-package OPC assembly for v6.1.

Each target slide is physically copied from a different certified source
package, including its slide layout, slide master, theme, media, charts,
chart styles/colors, embedded workbooks, diagrams, notes, notes masters,
comments, and any other required custom XML.

The OPC graph is namespaced once per source package.  Slide-scoped notes and
comments receive copy-on-write names, while immutable media and relationship-
free theme/chart-style parts may be deduplicated across packages.

After dependency closure, declared text bindings are applied to each slide
through the existing ``adapt_template_pack`` OOXML patcher (the single-page
TemplatePack). Finally the target is committed and a
``physical-assembly-report.v1`` is emitted.
"""

from __future__ import annotations

import hashlib
import html
import io
import json
import math
import os
import posixpath
import re
import shutil
import subprocess
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from copy import deepcopy
from collections import Counter, deque
from dataclasses import dataclass, field, replace
from decimal import Decimal, InvalidOperation
from datetime import datetime, timezone
from itertools import product
from pathlib import Path
from typing import Any, Iterable, Mapping, Sequence
from urllib.parse import urlsplit

from .page_template_library import (
    DEFAULT_SCORING,
    DEFAULT_DOMINANT_STYLE_CLUSTER,
    LibraryIndex,
    PageTemplate,
    _CNVPR_RE,
    _LAYOUT_RE,
    _MASTER_RE,
    _SLIDE_RE,
    _TEXT_RE,
    _discover_slots as _library_discover_slots,
    query_page_template_candidates,
)
from .template_pack import _replace_shape_text
from .weak_model import Fact, FactStore, WeakModelValidationError, load_fact_store
from .workbook_security import (
    WorkbookSecurityError,
    mutate_governed_xlsx,
    read_governed_xlsx_slot,
)


class PhysicalAssemblyError(ValueError):
    """Cross-package physical assembly has failed."""


DEFAULT_MAX_OUTPUT_SIZE_BYTES = 33_941_179

PHASE49_TEMPLATE_ROLE_SEQUENCE = (
    "cover",
    "contents",
    "section",
    "section",
    "data",
    "data",
    "table",
    "case-study",
    "kpi",
    "section",
    "people",
    "content-blocks",
    "section",
    "process",
    "closing",
)

PHASE49_NARRATIVE_ROLE_SEQUENCE = (
    "cover",
    "contents",
    "section-governance",
    "policy-evidence",
    "revenue-composition",
    "medical-revenue-comparison",
    "expenditure-table",
    "projects-debt",
    "kpi-dashboard",
    "section-innovation",
    "team",
    "efficiency-comparison",
    "section-roadmap",
    "roadmap",
    "closing",
)


@dataclass(frozen=True)
class AuthorityLock:
    fact_store_path: str
    fact_store_sha256: str
    asset_manifest_path: str
    asset_manifest_sha256: str
    connective_copy_path: str
    connective_copy_sha256: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "fact_store": {
                "path": self.fact_store_path,
                "sha256": self.fact_store_sha256,
            },
            "asset_manifest": {
                "path": self.asset_manifest_path,
                "sha256": self.asset_manifest_sha256,
            },
            "connective_copy": {
                "path": self.connective_copy_path,
                "sha256": self.connective_copy_sha256,
            },
        }


@dataclass(frozen=True)
class BindingProfileAuthorityLock:
    profile_id: str
    profile_sha256: str

    def to_dict(self) -> dict[str, str]:
        return {
            "profile_id": self.profile_id,
            "profile_sha256": self.profile_sha256,
        }


@dataclass(frozen=True)
class StyleCloneSpec:
    source_shape_id: int
    target_shape_id: int
    scope: str
    source_style_sha256: str
    target_guard_sha256: str


@dataclass(frozen=True)
class BindingProfileAuthorityEvidence:
    profile_id: str
    profile_sha256: str
    acceptance_profile: str
    style_clone_count: int
    status: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "profile_id": self.profile_id,
            "profile_sha256": self.profile_sha256,
            "acceptance_profile": self.acceptance_profile,
            "style_clone_count": self.style_clone_count,
            "status": self.status,
        }


@dataclass(frozen=True)
class StyleCloneEvidence:
    ordinal: int
    page_id: str
    source_shape_id: int
    target_shape_id: int
    scope: str
    expected_style_sha256: str
    actual_source_style_sha256: str
    actual_target_style_sha256: str
    actual_target_guard_sha256: str
    status: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "ordinal": self.ordinal,
            "page_id": self.page_id,
            "source_shape_id": self.source_shape_id,
            "target_shape_id": self.target_shape_id,
            "scope": self.scope,
            "expected_style_sha256": self.expected_style_sha256,
            "actual_source_style_sha256": self.actual_source_style_sha256,
            "actual_target_style_sha256": self.actual_target_style_sha256,
            "actual_target_guard_sha256": self.actual_target_guard_sha256,
            "status": self.status,
        }


@dataclass(frozen=True)
class TextBindingSpec:
    replacement: str
    fact_refs: tuple[str, ...]
    mode: str
    fit_policy: str = "preserve"
    field: str = "text"
    separator: str = ""
    slice_start: int | None = None
    slice_end: int | None = None


@dataclass(frozen=True)
class AssetBindingSpec:
    asset_ref: str
    fit: str = "cover"


@dataclass(frozen=True)
class LockedAsset:
    asset_ref: str
    path: Path
    sha256: str
    kind: str
    width_px: int | None
    height_px: int | None


@dataclass(frozen=True)
class BindingEvidence:
    ordinal: int
    page_id: str
    slot_id: str
    shape_id: int
    binding_kind: str
    mode: str
    source_text: str
    source_sha256: str
    replacement_sha256: str
    fact_refs: tuple[str, ...]
    asset_refs: tuple[str, ...]
    connective_ref: str
    char_used: int
    char_limit: int
    item_used: int
    item_limit: int
    image_used: int
    image_limit: int
    status: str
    fit_policy: str = "preserve"
    relationship_id: str = ""
    target_part: str = ""

    def to_dict(self) -> dict[str, Any]:
        return {
            "ordinal": self.ordinal,
            "page_id": self.page_id,
            "slot_id": self.slot_id,
            "shape_id": self.shape_id,
            "binding_kind": self.binding_kind,
            "mode": self.mode,
            "fit_policy": self.fit_policy,
            "source_sha256": self.source_sha256,
            "replacement_sha256": self.replacement_sha256,
            "fact_refs": list(self.fact_refs),
            "asset_refs": list(self.asset_refs),
            "connective_ref": self.connective_ref,
            "relationship_id": self.relationship_id,
            "target_part": self.target_part,
            "capacity": {
                "chars": {"used": self.char_used, "limit": self.char_limit},
                "items": {"used": self.item_used, "limit": self.item_limit},
                "images": {"used": self.image_used, "limit": self.image_limit},
            },
            "status": self.status,
        }


@dataclass(frozen=True)
class AuthorityEvidence:
    mode: str
    fact_store_path: str
    fact_store_sha256: str
    asset_manifest_path: str
    asset_manifest_sha256: str
    connective_copy_path: str
    connective_copy_sha256: str
    status: str

    def to_dict(self) -> dict[str, str]:
        return {
            "mode": self.mode,
            "fact_store_path": self.fact_store_path,
            "fact_store_sha256": self.fact_store_sha256,
            "asset_manifest_path": self.asset_manifest_path,
            "asset_manifest_sha256": self.asset_manifest_sha256,
            "connective_copy_path": self.connective_copy_path,
            "connective_copy_sha256": self.connective_copy_sha256,
            "status": self.status,
        }


@dataclass(frozen=True)
class FragmentGroupContract:
    """One immutable character-fragment group from the locked query bundle."""

    ordinal: int
    page_id: str
    group_id: str
    ordered_slot_ids: tuple[str, ...]


@dataclass(frozen=True)
class SelectionAuthorityEvidence:
    mode: str
    query_bundle_path: str
    query_bundle_sha256: str
    library_index_sha256: str
    query_count: int
    selected_count: int
    distinct_query_id_count: int
    distinct_page_id_count: int
    status: str
    fragment_group_contracts: tuple[FragmentGroupContract, ...] = field(
        default=(),
        repr=False,
        compare=False,
    )

    def to_dict(self) -> dict[str, Any]:
        return {
            "mode": self.mode,
            "query_bundle_path": self.query_bundle_path,
            "query_bundle_sha256": self.query_bundle_sha256,
            "library_index_sha256": self.library_index_sha256,
            "query_count": self.query_count,
            "selected_count": self.selected_count,
            "distinct_query_id_count": self.distinct_query_id_count,
            "distinct_page_id_count": self.distinct_page_id_count,
            "status": self.status,
        }


@dataclass(frozen=True)
class SlideBinding:
    ordinal: int
    narrative_role: str
    title: str
    headline: str
    bindings: Mapping[str, str]
    page_id: str
    package_sha256: str
    slide_number: int


@dataclass(frozen=True)
class SelectionEvidence:
    query_id: str
    candidate_rank: int
    score_total: float
    selection_reason: str
    fallback_reason: str | None = None

    def to_dict(self) -> dict[str, Any]:
        return {
            "query_id": self.query_id,
            "candidate_rank": self.candidate_rank,
            "score_total": self.score_total,
            "selection_reason": self.selection_reason,
            "fallback_reason": self.fallback_reason,
        }


@dataclass(frozen=True)
class AssemblyTargetSlide:
    ordinal: int
    page_template: PageTemplate
    bindings: Mapping[str, str]
    narrative_role: str
    title: str
    headline: str
    text_binding_specs: Mapping[str, TextBindingSpec] = field(default_factory=dict)
    governed_content_binding_specs: Mapping[str, TextBindingSpec] = field(
        default_factory=dict
    )
    asset_binding_specs: Mapping[str, AssetBindingSpec] = field(default_factory=dict)
    selection_evidence: SelectionEvidence | None = None

    def to_dict(self) -> dict[str, Any]:
        payload = {
            "ordinal": self.ordinal,
            "page_id": self.page_template.page_id,
            "package_sha256": self.page_template.package_sha256,
            "slide_number": self.page_template.slide_number,
            "narrative_role": self.narrative_role,
            "title": self.title,
            "headline": self.headline,
            "bindings": (
                ({
                    slot_id: {
                        "text": spec.replacement,
                        "fact_refs": list(spec.fact_refs),
                        "asset_refs": [],
                        "fit_policy": spec.fit_policy,
                    }
                    for slot_id, spec in self.text_binding_specs.items()
                } | {
                    slot_id: {
                        "text": spec.replacement,
                        "fact_refs": list(spec.fact_refs),
                        "asset_refs": [],
                        "fit_policy": spec.fit_policy,
                    }
                    for slot_id, spec in self.governed_content_binding_specs.items()
                })
                | {
                    slot_id: {
                        "text": "",
                        "fact_refs": [],
                        "asset_refs": [spec.asset_ref],
                    }
                    for slot_id, spec in self.asset_binding_specs.items()
                }
                if (
                    self.text_binding_specs
                    or self.governed_content_binding_specs
                    or self.asset_binding_specs
                )
                else dict(self.bindings)
            ),
        }
        if self.selection_evidence is not None:
            payload["selection"] = self.selection_evidence.to_dict()
        return payload


@dataclass(frozen=True)
class AssemblyPlan:
    schema_version: str
    plan_id: str
    scenario_id: str
    dominant_style_cluster_id: str
    created_at: str
    target_slide_count: int
    target_slides: tuple[AssemblyTargetSlide, ...]
    library_index_sha256: str
    binding_profile_authority: BindingProfileAuthorityLock | None = None
    query_bundle_path: str | None = None
    query_bundle_sha256: str | None = None
    authority: AuthorityLock | None = None

    def to_dict(self) -> dict[str, Any]:
        return {
            "schema_version": self.schema_version,
            "plan_id": self.plan_id,
            "scenario_id": self.scenario_id,
            "dominant_style_cluster_id": self.dominant_style_cluster_id,
            "created_at": self.created_at,
            "target_slide_count": self.target_slide_count,
            "target_slides": [item.to_dict() for item in self.target_slides],
            "library_index_sha256": self.library_index_sha256,
            **(
                {
                    "binding_profile_authority": (
                        self.binding_profile_authority.to_dict()
                    )
                }
                if self.binding_profile_authority is not None
                else {}
            ),
            **(
                {
                    "query_bundle": {
                        "path": self.query_bundle_path,
                        "sha256": self.query_bundle_sha256,
                    }
                }
                if self.query_bundle_path and self.query_bundle_sha256
                else {}
            ),
            **(
                {"authority": self.authority.to_dict()}
                if self.authority is not None
                else {}
            ),
        }


@dataclass(frozen=True)
class LineageRecord:
    ordinal: int
    page_id: str
    package_sha256: str
    slide_number: int
    source_sha256: str
    source_slide_sha256: str
    target_slide_sha256: str
    source_package_verified: bool
    source_slide_verified: bool
    structure_signature_source: str
    structure_signature_target: str
    structure_match: bool
    imported_part_map_sha256: str
    imported_part_count: int
    narrative_role: str
    title: str
    status: str
    binding_count: int
    byte_match_score: float

    def to_dict(self) -> dict[str, Any]:
        return {
            "ordinal": self.ordinal,
            "page_id": self.page_id,
            "package_sha256": self.package_sha256,
            "slide_number": self.slide_number,
            "source_sha256": self.source_sha256,
            "source_slide_sha256": self.source_slide_sha256,
            "target_slide_sha256": self.target_slide_sha256,
            "source_package_verified": self.source_package_verified,
            "source_slide_verified": self.source_slide_verified,
            "structure_signature_source": self.structure_signature_source,
            "structure_signature_target": self.structure_signature_target,
            "structure_match": self.structure_match,
            "imported_part_map_sha256": self.imported_part_map_sha256,
            "imported_part_count": self.imported_part_count,
            "narrative_role": self.narrative_role,
            "title": self.title,
            "status": self.status,
            "binding_count": self.binding_count,
            "byte_match_score": self.byte_match_score,
        }


@dataclass(frozen=True)
class OPCIntegrity:
    zip_open: bool
    content_types_parsed: bool
    slide_rels_resolved: bool
    package_entry_count: int
    media_count: int
    total_relationship_count: int
    internal_relationship_count: int
    external_relationship_count: int
    unresolved_internal_relationship_count: int
    unresolved_internal_relationships: tuple[Mapping[str, str], ...]
    unsafe_relationship_count: int
    unsafe_relationships: tuple[Mapping[str, str], ...]
    status: str
    details: str = ""

    def to_dict(self) -> dict[str, Any]:
        return {
            "zip_open": self.zip_open,
            "content_types_parsed": self.content_types_parsed,
            "slide_rels_resolved": self.slide_rels_resolved,
            "package_entry_count": self.package_entry_count,
            "media_count": self.media_count,
            "total_relationship_count": self.total_relationship_count,
            "internal_relationship_count": self.internal_relationship_count,
            "external_relationship_count": self.external_relationship_count,
            "unresolved_internal_relationship_count": self.unresolved_internal_relationship_count,
            "unresolved_internal_relationships": [
                dict(record) for record in self.unresolved_internal_relationships
            ],
            "unsafe_relationship_count": self.unsafe_relationship_count,
            "unsafe_relationships": [dict(record) for record in self.unsafe_relationships],
            "status": self.status,
            "details": self.details,
        }


@dataclass(frozen=True)
class Editability:
    native_editable: bool
    python_pptx_open: bool
    slide_count: int
    text_run_count: int
    shape_count: int
    native_object_count: int
    picture_count: int
    native_editable_slide_count: int
    full_slide_raster_count: int
    raster_dominant_slide_count: int
    native_editable_coverage: float
    status: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "native_editable": self.native_editable,
            "python_pptx_open": self.python_pptx_open,
            "slide_count": self.slide_count,
            "text_run_count": self.text_run_count,
            "shape_count": self.shape_count,
            "native_object_count": self.native_object_count,
            "picture_count": self.picture_count,
            "native_editable_slide_count": self.native_editable_slide_count,
            "full_slide_raster_count": self.full_slide_raster_count,
            "raster_dominant_slide_count": self.raster_dominant_slide_count,
            "native_editable_coverage": self.native_editable_coverage,
            "status": self.status,
        }


@dataclass(frozen=True)
class StyleClusterAdherence:
    dominant_style_cluster_id: str
    matches: int
    total: int
    status: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "dominant_style_cluster_id": self.dominant_style_cluster_id,
            "matches": self.matches,
            "total": self.total,
            "status": self.status,
        }


@dataclass(frozen=True)
class AssemblyMetrics:
    """Deterministic package-size and dependency-deduplication evidence."""

    output_size_bytes: int
    source_size_bytes: int
    unique_source_package_count: int
    imported_part_count: int
    imported_parts: tuple[str, ...]
    unique_dependency_part_count: int
    same_source_reuse_count: int
    same_source_reuse_bytes: int
    cross_source_safe_dedup_count: int
    cross_source_safe_dedup_bytes: int
    deduplicated_part_count: int
    deduplicated_bytes: int
    static_duplicate_bytes: int
    unresolved_internal_relationship_count: int
    amplification_ratio: float
    parts_by_kind: Mapping[str, int]

    def to_dict(self) -> dict[str, Any]:
        return {
            "output_size_bytes": self.output_size_bytes,
            "source_size_bytes": self.source_size_bytes,
            "unique_source_package_count": self.unique_source_package_count,
            "imported_part_count": self.imported_part_count,
            "imported_parts": list(self.imported_parts),
            "unique_dependency_part_count": self.unique_dependency_part_count,
            "same_source_reuse_count": self.same_source_reuse_count,
            "same_source_reuse_bytes": self.same_source_reuse_bytes,
            "cross_source_safe_dedup_count": self.cross_source_safe_dedup_count,
            "cross_source_safe_dedup_bytes": self.cross_source_safe_dedup_bytes,
            "deduplicated_part_count": self.deduplicated_part_count,
            "deduplicated_bytes": self.deduplicated_bytes,
            "static_duplicate_bytes": self.static_duplicate_bytes,
            "unresolved_internal_relationship_count": self.unresolved_internal_relationship_count,
            "amplification_ratio": self.amplification_ratio,
            "parts_by_kind": dict(self.parts_by_kind),
        }


@dataclass(frozen=True)
class LibreOfficeEvidence:
    available: bool
    executable: str
    open_result: str
    render_result: str
    status: str
    details: str = ""

    def to_dict(self) -> dict[str, Any]:
        return {
            "available": self.available,
            "executable": self.executable,
            "open_result": self.open_result,
            "render_result": self.render_result,
            "status": self.status,
            "details": self.details,
        }


@dataclass(frozen=True)
class SizeCheck:
    output_size_bytes: int
    max_output_size_bytes: int | None
    status: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "output_size_bytes": self.output_size_bytes,
            "max_output_size_bytes": self.max_output_size_bytes,
            "status": self.status,
        }


@dataclass(frozen=True)
class SourceResidueEvidence:
    governed_content_slot_count: int
    governed_content_binding_count: int
    verified_governed_content_count: int
    governed_content_mismatch_count: int
    peer_group_mismatch_count: int
    mutation_manifest_sha256: str
    governed_mutations: tuple[Mapping[str, Any], ...]
    unauthorized_content_count: int
    tag_part_count: int
    tag_relationship_count: int
    layout_master_cached_field_count: int
    certified_media_count: int
    media_hash_mismatch_count: int
    replacement_asset_count: int
    replacement_asset_hash_mismatch_count: int
    asset_slot_mismatch_count: int
    orphan_media_count: int
    status: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "governed_content_slot_count": self.governed_content_slot_count,
            "governed_content_binding_count": self.governed_content_binding_count,
            "verified_governed_content_count": self.verified_governed_content_count,
            "governed_content_mismatch_count": self.governed_content_mismatch_count,
            "peer_group_mismatch_count": self.peer_group_mismatch_count,
            "mutation_manifest_sha256": self.mutation_manifest_sha256,
            "governed_mutations": [dict(item) for item in self.governed_mutations],
            "unauthorized_content_count": self.unauthorized_content_count,
            "tag_part_count": self.tag_part_count,
            "tag_relationship_count": self.tag_relationship_count,
            "layout_master_cached_field_count": self.layout_master_cached_field_count,
            "certified_media_count": self.certified_media_count,
            "media_hash_mismatch_count": self.media_hash_mismatch_count,
            "replacement_asset_count": self.replacement_asset_count,
            "replacement_asset_hash_mismatch_count": self.replacement_asset_hash_mismatch_count,
            "asset_slot_mismatch_count": self.asset_slot_mismatch_count,
            "orphan_media_count": self.orphan_media_count,
            "status": self.status,
        }


@dataclass(frozen=True)
class PhysicalAssemblyReport:
    schema_version: str
    report_id: str
    plan_id: str
    binding_profile_authority: BindingProfileAuthorityEvidence | None
    output_path: str
    output_sha256: str
    acceptance_profile: str
    expected_slide_count: int | None
    status: str
    target_slide_count: int
    distinct_page_id_count: int
    duplicate_page_records: tuple[Mapping[str, Any], ...]
    lineage_records: tuple[LineageRecord, ...]
    opc_integrity: OPCIntegrity
    editability: Editability
    style_cluster_adherence: StyleClusterAdherence
    assembly_metrics: AssemblyMetrics
    authority: AuthorityEvidence
    selection_authority: SelectionAuthorityEvidence
    binding_evidence: tuple[BindingEvidence, ...]
    style_clone_evidence: tuple[StyleCloneEvidence, ...]
    source_residue: SourceResidueEvidence
    libreoffice: LibreOfficeEvidence
    size_check: SizeCheck

    def to_dict(self) -> dict[str, Any]:
        return {
            "schema_version": self.schema_version,
            "report_id": self.report_id,
            "plan_id": self.plan_id,
            **(
                {
                    "binding_profile_authority": (
                        self.binding_profile_authority.to_dict()
                    )
                }
                if self.binding_profile_authority is not None
                else {}
            ),
            "output_path": self.output_path,
            "output_sha256": self.output_sha256,
            "acceptance_profile": self.acceptance_profile,
            "expected_slide_count": self.expected_slide_count,
            "status": self.status,
            "target_slide_count": self.target_slide_count,
            "distinct_page_id_count": self.distinct_page_id_count,
            "duplicate_page_records": [
                dict(record) for record in self.duplicate_page_records
            ],
            "lineage_records": [record.to_dict() for record in self.lineage_records],
            "opc_integrity": self.opc_integrity.to_dict(),
            "editability": self.editability.to_dict(),
            "style_cluster_adherence": self.style_cluster_adherence.to_dict(),
            "assembly_metrics": self.assembly_metrics.to_dict(),
            "authority": self.authority.to_dict(),
            "selection_authority": self.selection_authority.to_dict(),
            "binding_evidence": [item.to_dict() for item in self.binding_evidence],
            "style_clone_evidence": [
                item.to_dict() for item in self.style_clone_evidence
            ],
            "source_residue": self.source_residue.to_dict(),
            "libreoffice": self.libreoffice.to_dict(),
            "size_check": self.size_check.to_dict(),
        }


def _sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def _validate_schema_payload(
    payload: Mapping[str, Any],
    schema_name: str,
    *,
    label: str,
) -> None:
    """Validate against the local schema graph without network resolution."""

    try:
        import jsonschema
        from referencing import Registry, Resource
    except ImportError as exc:  # pragma: no cover - installation contract
        raise PhysicalAssemblyError(
            f"{label}_SCHEMA_RUNTIME_UNAVAILABLE: {exc}"
        ) from exc
    schemas_root = Path(__file__).resolve().parents[2] / "schemas"
    schema_path = schemas_root / schema_name
    try:
        schemas = [
            json.loads(path.read_text(encoding="utf-8"))
            for path in sorted(schemas_root.glob("*.schema.json"))
        ]
        schema = next(item for item in schemas if item.get("$id") == json.loads(
            schema_path.read_text(encoding="utf-8")
        ).get("$id"))
        registry = Registry().with_resources(
            (
                item["$id"],
                Resource.from_contents(item),
            )
            for item in schemas
            if isinstance(item.get("$id"), str)
        )
        jsonschema.Draft202012Validator(
            schema,
            registry=registry,
            format_checker=jsonschema.FormatChecker(),
        ).validate(dict(payload))
    except (OSError, StopIteration, json.JSONDecodeError) as exc:
        raise PhysicalAssemblyError(f"{label}_SCHEMA_LOAD_FAILED: {exc}") from exc
    except jsonschema.ValidationError as exc:
        location = "/".join(str(item) for item in exc.absolute_path)
        raise PhysicalAssemblyError(
            f"{label}_SCHEMA_INVALID:{location}:{exc.message}"
        ) from exc


def _now_utc() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


PACKAGE_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PACKAGE_RELATIONSHIPS_TAG = f"{{{PACKAGE_REL_NS}}}Relationships"
PACKAGE_RELATIONSHIP_TAG = f"{{{PACKAGE_REL_NS}}}Relationship"
PACKAGE_RELATIONSHIP_ATTRIBUTES = frozenset(
    {"Id", "Type", "Target", "TargetMode"}
)
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
OD_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

for _prefix, _namespace in (
    ("p", PML_NS),
    ("a", DML_NS),
    ("c", CHART_NS),
    ("r", OD_REL_NS),
):
    ET.register_namespace(_prefix, _namespace)
ET.register_namespace("", PACKAGE_REL_NS)

STYLE_CLONE_SCOPES = frozenset(
    {"shape-fill", "text-color", "picture-color-effects"}
)
_STYLE_FILL_TAGS = frozenset(
    f"{{{DML_NS}}}{name}"
    for name in ("noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill")
)
_PICTURE_COLOR_EFFECT_TAGS = frozenset(
    f"{{{DML_NS}}}{name}"
    for name in (
        "alphaBiLevel", "alphaCeiling", "alphaFloor", "alphaInv",
        "alphaMod", "alphaModFix", "alphaRepl", "biLevel", "blur",
        "clrChange", "clrRepl", "duotone", "fillOverlay", "grayscl",
        "hsl", "lum", "tint",
    )
)
_TEXT_STYLE_CARRIER_TAGS = frozenset(
    f"{{{DML_NS}}}{name}" for name in ("rPr", "defRPr", "endParaRPr")
)
_SUPPORTED_STYLE_SHAPE_TAGS = frozenset(
    f"{{{PML_NS}}}{name}" for name in ("sp", "cxnSp", "pic")
)
_PHASE49_PROFILE_BY_ACCEPTANCE = {
    "phase49-work-report-15": "phase49-work-report-15",
}
_TRUSTED_BINDING_PROFILE_FILES = {
    "phase49-work-report-15": (
        "phase49-work-report-15.binding-profile.v1.json"
    ),
}


def _canonical_style_node(node: ET.Element, *, guard: bool = False) -> Any:
    local = node.tag.rsplit("}", 1)[-1]
    if guard and local in {"spAutoFit", "normAutofit", "noAutofit"}:
        return ["__GOVERNED_FIT_POLICY__"]
    attributes = []
    for key, value in sorted(node.attrib.items()):
        if (
            guard
            and local in {"rPr", "defRPr", "endParaRPr"}
            and key.rsplit("}", 1)[-1] == "sz"
        ):
            value = "__GOVERNED_TEXT_SIZE__"
        attributes.append((key, value))
    text = "__GOVERNED_TEXT__" if guard and local == "t" else (node.text or "").strip()
    return [
        node.tag,
        attributes,
        text,
        [_canonical_style_node(child, guard=guard) for child in list(node)],
    ]


def _style_shape_index(root: ET.Element) -> dict[int, ET.Element]:
    result: dict[int, ET.Element] = {}
    for shape in root.iter():
        if shape.tag not in _SUPPORTED_STYLE_SHAPE_TAGS:
            continue
        markers = [
            node
            for node in shape.iter(f"{{{PML_NS}}}cNvPr")
        ]
        if len(markers) != 1:
            raise PhysicalAssemblyError("STYLE_CLONE_SHAPE_MARKER_INVALID")
        try:
            shape_id = int(markers[0].attrib["id"])
        except (KeyError, ValueError) as exc:
            raise PhysicalAssemblyError("STYLE_CLONE_SHAPE_MARKER_INVALID") from exc
        if shape_id in result:
            raise PhysicalAssemblyError(
                f"STYLE_CLONE_SHAPE_ID_DUPLICATE: {shape_id}"
            )
        result[shape_id] = shape
    return result


def _style_scope_nodes(shape: ET.Element, scope: str) -> list[ET.Element]:
    if scope not in STYLE_CLONE_SCOPES:
        raise PhysicalAssemblyError(f"STYLE_CLONE_SCOPE_INVALID: {scope}")
    if scope == "shape-fill":
        properties = shape.find(f"{{{PML_NS}}}spPr")
        nodes = [
            child
            for child in (list(properties) if properties is not None else ())
            if child.tag in _STYLE_FILL_TAGS
        ]
        if properties is None or len(nodes) != 1:
            raise PhysicalAssemblyError("STYLE_CLONE_SHAPE_FILL_INVALID")
        return nodes
    if scope == "text-color":
        carriers = [
            node for node in shape.iter() if node.tag in _TEXT_STYLE_CARRIER_TAGS
        ]
        if not carriers:
            raise PhysicalAssemblyError("STYLE_CLONE_TEXT_COLOR_INVALID")
        nodes: list[ET.Element] = []
        for carrier in carriers:
            colors = [child for child in list(carrier) if child.tag in _STYLE_FILL_TAGS]
            if len(colors) != 1:
                raise PhysicalAssemblyError("STYLE_CLONE_TEXT_COLOR_INVALID")
            nodes.extend(colors)
        return nodes
    blip = shape.find(f"{{{PML_NS}}}blipFill/{{{DML_NS}}}blip")
    nodes = [
        child
        for child in (list(blip) if blip is not None else ())
        if child.tag in _PICTURE_COLOR_EFFECT_TAGS
    ]
    if blip is None or not nodes:
        raise PhysicalAssemblyError("STYLE_CLONE_PICTURE_EFFECT_INVALID")
    return nodes


def _style_scope_sha(shape: ET.Element, scope: str) -> str:
    payload = [
        scope,
        [_canonical_style_node(node) for node in _style_scope_nodes(shape, scope)],
    ]
    return _sha256_bytes(
        json.dumps(
            payload,
            ensure_ascii=True,
            sort_keys=False,
            separators=(",", ":"),
        ).encode("utf-8")
    )


def _remove_style_scope(shape: ET.Element, scope: str) -> None:
    if scope == "shape-fill":
        parent = shape.find(f"{{{PML_NS}}}spPr")
        if parent is None:
            raise PhysicalAssemblyError("STYLE_CLONE_SHAPE_FILL_INVALID")
        for node in list(parent):
            if node.tag in _STYLE_FILL_TAGS:
                parent.remove(node)
        return
    if scope == "text-color":
        for carrier in shape.iter():
            if carrier.tag in _TEXT_STYLE_CARRIER_TAGS:
                for node in list(carrier):
                    if node.tag in _STYLE_FILL_TAGS:
                        carrier.remove(node)
        return
    blip = shape.find(f"{{{PML_NS}}}blipFill/{{{DML_NS}}}blip")
    if blip is None:
        raise PhysicalAssemblyError("STYLE_CLONE_PICTURE_EFFECT_INVALID")
    for node in list(blip):
        if node.tag in _PICTURE_COLOR_EFFECT_TAGS:
            blip.remove(node)


def _style_guard_sha(shape: ET.Element, scope: str) -> str:
    guarded = deepcopy(shape)
    _remove_style_scope(guarded, scope)
    return _sha256_bytes(
        json.dumps(
            _canonical_style_node(guarded, guard=True),
            ensure_ascii=True,
            sort_keys=False,
            separators=(",", ":"),
        ).encode("utf-8")
    )


def _style_clone_scope_sha256(
    slide_xml: bytes, shape_id: int, scope: str
) -> str:
    try:
        root = ET.fromstring(slide_xml)
    except ET.ParseError as exc:
        raise PhysicalAssemblyError(f"STYLE_CLONE_SLIDE_XML_INVALID: {exc}") from exc
    shape = _style_shape_index(root).get(shape_id)
    if shape is None:
        raise PhysicalAssemblyError(f"STYLE_CLONE_SHAPE_NOT_FOUND: {shape_id}")
    return _style_scope_sha(shape, scope)


def _style_clone_target_guard_sha256(
    slide_xml: bytes, shape_id: int, scope: str
) -> str:
    try:
        root = ET.fromstring(slide_xml)
    except ET.ParseError as exc:
        raise PhysicalAssemblyError(f"STYLE_CLONE_SLIDE_XML_INVALID: {exc}") from exc
    shape = _style_shape_index(root).get(shape_id)
    if shape is None:
        raise PhysicalAssemblyError(f"STYLE_CLONE_SHAPE_NOT_FOUND: {shape_id}")
    return _style_guard_sha(shape, scope)


def _replace_style_scope(
    target: ET.Element,
    scope: str,
    source_nodes: Sequence[ET.Element],
) -> None:
    if scope == "shape-fill":
        parent = target.find(f"{{{PML_NS}}}spPr")
        if parent is None:
            raise PhysicalAssemblyError("STYLE_CLONE_SHAPE_FILL_INVALID")
        existing = [node for node in list(parent) if node.tag in _STYLE_FILL_TAGS]
        if len(existing) != 1 or len(source_nodes) != 1:
            raise PhysicalAssemblyError("STYLE_CLONE_SCOPE_TOPOLOGY_MISMATCH")
        index = list(parent).index(existing[0])
        parent.remove(existing[0])
        parent.insert(index, deepcopy(source_nodes[0]))
        return
    if scope == "text-color":
        carriers = [
            node for node in target.iter() if node.tag in _TEXT_STYLE_CARRIER_TAGS
        ]
        if len(carriers) != len(source_nodes):
            raise PhysicalAssemblyError("STYLE_CLONE_SCOPE_TOPOLOGY_MISMATCH")
        for carrier, source_node in zip(carriers, source_nodes):
            existing = [node for node in list(carrier) if node.tag in _STYLE_FILL_TAGS]
            if len(existing) != 1:
                raise PhysicalAssemblyError("STYLE_CLONE_SCOPE_TOPOLOGY_MISMATCH")
            index = list(carrier).index(existing[0])
            carrier.remove(existing[0])
            carrier.insert(index, deepcopy(source_node))
        return
    blip = target.find(f"{{{PML_NS}}}blipFill/{{{DML_NS}}}blip")
    if blip is None:
        raise PhysicalAssemblyError("STYLE_CLONE_PICTURE_EFFECT_INVALID")
    for node in list(blip):
        if node.tag in _PICTURE_COLOR_EFFECT_TAGS:
            blip.remove(node)
    children = list(blip)
    insert_at = next(
        (
            index
            for index, child in enumerate(children)
            if child.tag.rsplit("}", 1)[-1] == "extLst"
        ),
        len(children),
    )
    for offset, node in enumerate(source_nodes):
        blip.insert(insert_at + offset, deepcopy(node))


def _apply_governed_style_clones(
    slide_xml: bytes,
    specs: Sequence[StyleCloneSpec],
) -> bytes:
    if not specs:
        return slide_xml
    try:
        root = ET.fromstring(slide_xml)
    except ET.ParseError as exc:
        raise PhysicalAssemblyError(f"STYLE_CLONE_SLIDE_XML_INVALID: {exc}") from exc
    shapes = _style_shape_index(root)
    sources = {spec.source_shape_id for spec in specs}
    targets = {spec.target_shape_id for spec in specs}
    if sources & targets:
        raise PhysicalAssemblyError("STYLE_CLONE_CHAIN_FORBIDDEN")
    target_scopes: set[tuple[int, str]] = set()
    snapshots: list[tuple[StyleCloneSpec, tuple[ET.Element, ...]]] = []
    for spec in specs:
        if spec.source_shape_id == spec.target_shape_id:
            raise PhysicalAssemblyError("STYLE_CLONE_SELF_TARGET")
        key = (spec.target_shape_id, spec.scope)
        if key in target_scopes:
            raise PhysicalAssemblyError("STYLE_CLONE_TARGET_DUPLICATE")
        target_scopes.add(key)
        source = shapes.get(spec.source_shape_id)
        target = shapes.get(spec.target_shape_id)
        if source is None or target is None:
            raise PhysicalAssemblyError(
                "STYLE_CLONE_SHAPE_NOT_FOUND: "
                f"{spec.source_shape_id}->{spec.target_shape_id}"
            )
        if source.tag != target.tag:
            raise PhysicalAssemblyError("STYLE_CLONE_SHAPE_KIND_MISMATCH")
        source_nodes = tuple(deepcopy(node) for node in _style_scope_nodes(source, spec.scope))
        if _style_scope_sha(source, spec.scope) != spec.source_style_sha256:
            raise PhysicalAssemblyError("STYLE_CLONE_SOURCE_STYLE_DRIFT")
        if _style_guard_sha(target, spec.scope) != spec.target_guard_sha256:
            raise PhysicalAssemblyError("STYLE_CLONE_TARGET_GUARD_DRIFT")
        snapshots.append((spec, source_nodes))
    for spec, source_nodes in snapshots:
        _replace_style_scope(shapes[spec.target_shape_id], spec.scope, source_nodes)
        if _style_guard_sha(
            shapes[spec.target_shape_id], spec.scope
        ) != spec.target_guard_sha256:
            raise PhysicalAssemblyError("STYLE_CLONE_TARGET_GUARD_MUTATED")
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _resolve_binding_profile_style_authority(
    plan: AssemblyPlan,
    *,
    acceptance_profile: str,
) -> tuple[
    BindingProfileAuthorityEvidence | None,
    dict[int, tuple[StyleCloneSpec, ...]],
]:
    """Resolve paint operations only from a fixed, installed Skill profile.

    The assembly plan carries a profile ID and byte fingerprint, never shape
    IDs or OOXML operations.  This resolver deliberately has no path argument:
    an agent-authored plan therefore cannot redirect the renderer to a client
    file or an arbitrary local profile.
    """

    lock = plan.binding_profile_authority
    required_profile_id = _PHASE49_PROFILE_BY_ACCEPTANCE.get(acceptance_profile)
    if lock is None:
        if required_profile_id is not None:
            raise PhysicalAssemblyError(
                "BINDING_PROFILE_AUTHORITY_REQUIRED: " + acceptance_profile
            )
        return None, {}
    if required_profile_id is not None and lock.profile_id != required_profile_id:
        raise PhysicalAssemblyError(
            "BINDING_PROFILE_AUTHORITY_MISMATCH: " + lock.profile_id
        )
    filename = _TRUSTED_BINDING_PROFILE_FILES.get(lock.profile_id)
    if filename is None:
        raise PhysicalAssemblyError(
            "BINDING_PROFILE_AUTHORITY_UNTRUSTED: " + lock.profile_id
        )
    registry_root = (
        Path(__file__).resolve().parents[2]
        / "registries"
        / "v61-binding-profiles"
    )
    literal_path = registry_root / filename
    if literal_path.is_symlink():
        raise PhysicalAssemblyError("BINDING_PROFILE_AUTHORITY_SYMLINK_FORBIDDEN")
    try:
        resolved_root = registry_root.resolve(strict=True)
        profile_path = literal_path.resolve(strict=True)
    except OSError as exc:
        raise PhysicalAssemblyError(
            f"BINDING_PROFILE_AUTHORITY_MISSING: {lock.profile_id}"
        ) from exc
    if profile_path.parent != resolved_root or not profile_path.is_file():
        raise PhysicalAssemblyError("BINDING_PROFILE_AUTHORITY_PATH_ESCAPE")
    raw = profile_path.read_bytes()
    actual_sha256 = _sha256_bytes(raw)
    if actual_sha256 != lock.profile_sha256:
        raise PhysicalAssemblyError("BINDING_PROFILE_AUTHORITY_FINGERPRINT_MISMATCH")
    try:
        profile = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise PhysicalAssemblyError(
            f"BINDING_PROFILE_AUTHORITY_INVALID: {exc}"
        ) from exc
    if not isinstance(profile, Mapping):
        raise PhysicalAssemblyError("BINDING_PROFILE_AUTHORITY_INVALID")
    _validate_schema_payload(
        profile,
        "binding-profile.v1.schema.json",
        label="BINDING_PROFILE_AUTHORITY",
    )
    if (
        profile.get("profile_id") != lock.profile_id
        or profile.get("acceptance_profile") != acceptance_profile
        or profile.get("scenario_id") != plan.scenario_id
        or profile.get("library_index_sha256") != plan.library_index_sha256
        or profile.get("dominant_style_cluster_id")
        != plan.dominant_style_cluster_id
    ):
        raise PhysicalAssemblyError("BINDING_PROFILE_AUTHORITY_APPLICABILITY_MISMATCH")

    profile_slides = profile.get("slides")
    if not isinstance(profile_slides, list):
        raise PhysicalAssemblyError("BINDING_PROFILE_AUTHORITY_INVALID")
    by_ordinal = {
        int(item["ordinal"]): item
        for item in profile_slides
        if isinstance(item, Mapping)
    }
    plan_ordinals = {slide.ordinal for slide in plan.target_slides}
    if set(by_ordinal) != plan_ordinals or len(by_ordinal) != len(profile_slides):
        raise PhysicalAssemblyError("BINDING_PROFILE_AUTHORITY_SLIDE_SET_MISMATCH")

    specs_by_ordinal: dict[int, tuple[StyleCloneSpec, ...]] = {}
    total = 0
    for slide in plan.target_slides:
        profile_slide = by_ordinal[slide.ordinal]
        if (
            profile_slide.get("page_id") != slide.page_template.page_id
            or profile_slide.get("narrative_role") != slide.narrative_role
        ):
            raise PhysicalAssemblyError(
                "BINDING_PROFILE_AUTHORITY_SLIDE_MISMATCH: "
                + str(slide.ordinal)
            )
        raw_specs = profile_slide.get("style_clones", ())
        if not isinstance(raw_specs, (list, tuple)):
            raise PhysicalAssemblyError(
                "BINDING_PROFILE_AUTHORITY_STYLE_CLONES_INVALID: "
                + str(slide.ordinal)
            )
        specs = tuple(
            StyleCloneSpec(
                source_shape_id=int(item["source_shape_id"]),
                target_shape_id=int(item["target_shape_id"]),
                scope=str(item["scope"]),
                source_style_sha256=str(item["source_style_sha256"]),
                target_guard_sha256=str(item["target_guard_sha256"]),
            )
            for item in raw_specs
        )
        if specs:
            sources = {item.source_shape_id for item in specs}
            targets = {item.target_shape_id for item in specs}
            target_scopes = {
                (item.target_shape_id, item.scope) for item in specs
            }
            if (
                sources & targets
                or len(target_scopes) != len(specs)
                or any(
                    item.source_shape_id == item.target_shape_id
                    for item in specs
                )
            ):
                raise PhysicalAssemblyError(
                    "BINDING_PROFILE_AUTHORITY_STYLE_CLONES_AMBIGUOUS: "
                    + str(slide.ordinal)
                )
            specs_by_ordinal[slide.ordinal] = specs
            total += len(specs)
    if required_profile_id is not None and total != 4:
        raise PhysicalAssemblyError(
            f"BINDING_PROFILE_AUTHORITY_STYLE_CLONE_COUNT: {total}"
        )
    return (
        BindingProfileAuthorityEvidence(
            profile_id=lock.profile_id,
            profile_sha256=actual_sha256,
            acceptance_profile=acceptance_profile,
            style_clone_count=total,
            status="pass",
        ),
        specs_by_ordinal,
    )


def _verify_style_clone_evidence(
    archive: zipfile.ZipFile,
    *,
    plan: AssemblyPlan,
    specs_by_ordinal: Mapping[int, Sequence[StyleCloneSpec]],
) -> tuple[StyleCloneEvidence, ...]:
    """Recompute the source anchor, cloned target paint, and target guard."""

    evidence: list[StyleCloneEvidence] = []
    slides = {slide.ordinal: slide for slide in plan.target_slides}
    for ordinal in sorted(specs_by_ordinal):
        slide = slides.get(ordinal)
        if slide is None:
            raise PhysicalAssemblyError(
                f"STYLE_CLONE_EVIDENCE_SLIDE_NOT_FOUND: {ordinal}"
            )
        try:
            slide_xml = archive.read(f"ppt/slides/slide{ordinal}.xml")
        except KeyError as exc:
            raise PhysicalAssemblyError(
                f"STYLE_CLONE_EVIDENCE_SLIDE_NOT_FOUND: {ordinal}"
            ) from exc
        for spec in specs_by_ordinal[ordinal]:
            source_style = _style_clone_scope_sha256(
                slide_xml, spec.source_shape_id, spec.scope
            )
            target_style = _style_clone_scope_sha256(
                slide_xml, spec.target_shape_id, spec.scope
            )
            target_guard = _style_clone_target_guard_sha256(
                slide_xml, spec.target_shape_id, spec.scope
            )
            status = (
                "pass"
                if source_style
                == target_style
                == spec.source_style_sha256
                and target_guard == spec.target_guard_sha256
                else "fail"
            )
            evidence.append(
                StyleCloneEvidence(
                    ordinal=ordinal,
                    page_id=slide.page_template.page_id,
                    source_shape_id=spec.source_shape_id,
                    target_shape_id=spec.target_shape_id,
                    scope=spec.scope,
                    expected_style_sha256=spec.source_style_sha256,
                    actual_source_style_sha256=source_style,
                    actual_target_style_sha256=target_style,
                    actual_target_guard_sha256=target_guard,
                    status=status,
                )
            )
    return tuple(evidence)
# The following relationship types are forbidden and must be rejected before
# any part is committed to the output.
FORBIDDEN_REL_TARGET_PATTERNS = (
    re.compile(r"\.exe$", re.IGNORECASE),
    re.compile(r"\.bat$", re.IGNORECASE),
    re.compile(r"\.cmd$", re.IGNORECASE),
    re.compile(r"\.scr$", re.IGNORECASE),
    re.compile(r"\.vbs$", re.IGNORECASE),
    re.compile(r"\.js$", re.IGNORECASE),
    re.compile(r"\.ps1$", re.IGNORECASE),
    re.compile(r"\.sh$", re.IGNORECASE),
    re.compile(r"\.py$", re.IGNORECASE),
    re.compile(r"\.hta$", re.IGNORECASE),
    re.compile(r"\.wsf$", re.IGNORECASE),
    re.compile(r"\.com$", re.IGNORECASE),
    re.compile(r"\.dll$", re.IGNORECASE),
    re.compile(r"\.msi$", re.IGNORECASE),
    re.compile(r"\.bin$", re.IGNORECASE),
    re.compile(r"\.(?:pptm|potm|ppsm|xlsm|xltm|xlam|docm|dotm)$", re.IGNORECASE),
)

FORBIDDEN_REL_TYPE_FRAGMENTS = (
    "oleobject",
    "vbaproject",
    "activex",
    "attachedtemplate",
    "externalworkbookpath",
    "javascript",
    "vbscript",
    "macro",
)

# Mandatory types we need to register after import.
CONTENT_TYPES = {
    "ppt": "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
    "slide": "application/vnd.openxmlformats-officedocument.presentationml.slide+xml",
    "slideLayout": "application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml",
    "slideMaster": "application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml",
    "notesSlide": "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml",
    "notesMaster": "application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml",
    "theme": "application/vnd.openxmlformats-officedocument.theme+xml",
    "chart": "application/vnd.openxmlformats-officedocument.drawingml.chart+xml",
    "image": "application/vnd.openxmlformats-officedocument.drawingml.picture+xml",
}


@dataclass
class _SourceGraph:
    """Owner-relative OPC closure rooted at one physical source slide."""

    root_slide_name: str
    slide_xml: bytes
    slide_sha: str
    rels: dict[str, bytes] = field(default_factory=dict)
    extra_parts: dict[str, bytes] = field(default_factory=dict)
    content_types: dict[str, str] = field(default_factory=dict)
    layout_paths: list[str] = field(default_factory=list)
    master_paths: list[str] = field(default_factory=list)
    theme_paths: list[str] = field(default_factory=list)


@dataclass
class _SourcePackageContext:
    """One open source archive, its content types, and cached slide closures."""

    package_path: Path
    package_sha256: str
    archive: zipfile.ZipFile
    names: frozenset[str]
    content_type_defaults: Mapping[str, str]
    content_type_overrides: Mapping[str, str]
    closures: dict[int, _SourceGraph] = field(default_factory=dict)

    @classmethod
    def open(cls, package_path: Path, package_sha256: str) -> "_SourcePackageContext":
        actual_sha256 = _sha256_file(package_path)
        if actual_sha256 != package_sha256:
            raise PhysicalAssemblyError(
                f"source package fingerprint mismatch: {package_path}"
            )
        archive = zipfile.ZipFile(package_path, "r")
        try:
            content_types = archive.read("[Content_Types].xml")
            defaults, overrides = _parse_content_types(content_types)
        except Exception:
            archive.close()
            raise
        return cls(
            package_path=package_path,
            package_sha256=package_sha256,
            archive=archive,
            names=frozenset(_normalise_zip_name(name) for name in archive.namelist()),
            content_type_defaults=dict(defaults),
            content_type_overrides={name.lstrip("/"): value for name, value in overrides},
        )

    def content_type_for(self, part_name: str) -> str:
        normalised = _normalise_zip_name(part_name).lstrip("/")
        override = self.content_type_overrides.get(normalised)
        if override:
            return override
        extension = normalised.rsplit(".", 1)[-1].lower() if "." in normalised else ""
        default = self.content_type_defaults.get(extension)
        if default:
            return default
        raise PhysicalAssemblyError(
            f"source content type missing for {normalised}: {self.package_path}"
        )

    def close(self) -> None:
        self.archive.close()


def _rels_path_for_part(part_name: str) -> str:
    part = _normalise_zip_name(part_name).lstrip("/")
    folder, filename = posixpath.split(part)
    return posixpath.join(folder, "_rels", f"{filename}.rels")


def _owner_part_from_rels_path(rels_path: str) -> str | None:
    rels = _normalise_zip_name(rels_path).lstrip("/")
    if rels == "_rels/.rels":
        return None
    parts = rels.split("/")
    if len(parts) < 2 or parts[-2] != "_rels" or not parts[-1].endswith(".rels"):
        return None
    return "/".join(parts[:-2] + [parts[-1][:-5]])


def _normalise_zip_name(name: str) -> str:
    return name.replace("\\", "/")


def _resolve_rel_target(rels_xml: bytes, source_rel_path: str, target: str) -> str | None:
    """Resolve a relationship target relative to the source rels file.

    ``source_rel_path`` looks like ``ppt/slides/_rels/slide1.xml.rels``. The
    relationships it carries describe files in ``ppt/slides/`` (one directory
    up from the rels file). Targets that begin with ``../`` therefore resolve
    to ``ppt/...`` siblings, not ``ppt/slides/...``.
    """

    del rels_xml  # Kept in the stable helper signature for existing callers.
    if not target or _looks_like_external_relationship_target(target):
        return None
    owner = _owner_part_from_rels_path(source_rel_path)
    base = posixpath.dirname(owner) if owner else ""
    if target.startswith("/"):
        return _normalise_zip_name(target.lstrip("/"))
    combined = _normalise_zip_name(posixpath.normpath(posixpath.join(base, target))).lstrip("/")
    if combined == ".." or combined.startswith("../"):
        return None
    return combined


def _parse_relationships(rels_xml: bytes) -> list[dict[str, str]]:
    """Strictly parse an OPC Relationships part.

    Relationship parts sit on a security boundary: accepting fragments from a
    malformed document can hide a second, unsafe relationship from dependency
    closure or output audit.  Parse the complete XML document, require the OPC
    package namespace (whether expressed as a default or prefixed namespace),
    and reject unknown structure or attributes.
    """

    if b"<!DOCTYPE" in rels_xml.upper():
        raise PhysicalAssemblyError("relationship XML must not contain a DOCTYPE")
    try:
        root = ET.fromstring(rels_xml)
    except (ET.ParseError, ValueError) as exc:
        raise PhysicalAssemblyError(f"relationship XML is invalid: {exc}") from exc
    if root.tag != PACKAGE_RELATIONSHIPS_TAG:
        raise PhysicalAssemblyError(
            "relationship XML root must be the OPC Relationships element"
        )
    if root.attrib:
        raise PhysicalAssemblyError("relationship XML root attributes are forbidden")
    if root.text and root.text.strip():
        raise PhysicalAssemblyError("relationship XML root text is forbidden")

    entries: list[dict[str, str]] = []
    seen_ids: set[str] = set()
    for node in list(root):
        if node.tag != PACKAGE_RELATIONSHIP_TAG:
            raise PhysicalAssemblyError(
                "relationship XML contains a non-Relationship child"
            )
        if list(node) or (node.text and node.text.strip()):
            raise PhysicalAssemblyError(
                "relationship XML Relationship elements must be empty"
            )
        if node.tail and node.tail.strip():
            raise PhysicalAssemblyError(
                "relationship XML contains text outside Relationship elements"
            )
        unknown_attributes = set(node.attrib) - PACKAGE_RELATIONSHIP_ATTRIBUTES
        if unknown_attributes:
            raise PhysicalAssemblyError(
                "relationship XML contains forbidden attributes: "
                + ",".join(sorted(unknown_attributes))
            )
        missing = [
            name
            for name in ("Id", "Type", "Target")
            if not node.attrib.get(name, "").strip()
        ]
        if missing:
            raise PhysicalAssemblyError(
                "relationship XML is missing required attributes: "
                + ",".join(missing)
            )
        relationship_id = node.attrib["Id"]
        if relationship_id in seen_ids:
            raise PhysicalAssemblyError(
                f"relationship XML contains duplicate Id: {relationship_id}"
            )
        seen_ids.add(relationship_id)
        target_mode = node.attrib.get("TargetMode", "")
        if target_mode not in {"", "Internal", "External"}:
            raise PhysicalAssemblyError(
                f"relationship XML has an unsafe TargetMode: {target_mode}"
            )
        entries.append(
            {
                "Id": relationship_id,
                "Type": node.attrib["Type"],
                "Target": node.attrib["Target"],
                "TargetMode": target_mode,
            }
        )
    return entries


def _discover_picture_slots(slide_xml: bytes) -> dict[str, tuple[int, str]]:
    """Map editable picture shape IDs to their embedded relationship IDs."""

    try:
        root = ET.fromstring(slide_xml)
    except ET.ParseError as exc:
        raise PhysicalAssemblyError(f"slide XML is invalid: {exc}") from exc
    result: dict[str, tuple[int, str]] = {}
    for picture in root.iter(f"{{{PML_NS}}}pic"):
        non_visual = picture.find(f".//{{{PML_NS}}}cNvPr")
        blip = picture.find(f".//{{{DML_NS}}}blip")
        if non_visual is None or blip is None:
            continue
        raw_shape_id = non_visual.attrib.get("id")
        relationship_id = blip.attrib.get(f"{{{OD_REL_NS}}}embed")
        if raw_shape_id and raw_shape_id.isdigit() and relationship_id:
            shape_id = int(raw_shape_id)
            result[f"shape_{shape_id}"] = (shape_id, relationship_id)
    return result


def _discover_picture_frames(slide_xml: bytes) -> dict[str, tuple[int, int]]:
    try:
        root = ET.fromstring(slide_xml)
    except ET.ParseError as exc:
        raise PhysicalAssemblyError(f"slide XML is invalid: {exc}") from exc
    result: dict[str, tuple[int, int]] = {}
    for picture in root.iter(f"{{{PML_NS}}}pic"):
        non_visual = picture.find(f".//{{{PML_NS}}}cNvPr")
        extent = picture.find(f".//{{{PML_NS}}}spPr/{{{DML_NS}}}xfrm/{{{DML_NS}}}ext")
        if non_visual is None or extent is None:
            continue
        raw_shape_id = non_visual.attrib.get("id")
        raw_cx = extent.attrib.get("cx")
        raw_cy = extent.attrib.get("cy")
        if (
            raw_shape_id
            and raw_shape_id.isdigit()
            and raw_cx
            and raw_cx.isdigit()
            and raw_cy
            and raw_cy.isdigit()
            and int(raw_cx) > 0
            and int(raw_cy) > 0
        ):
            result[f"shape_{int(raw_shape_id)}"] = (int(raw_cx), int(raw_cy))
    return result


def _cover_crop_values(
    *,
    frame_width: int,
    frame_height: int,
    image_width: int,
    image_height: int,
) -> tuple[int, int, int, int]:
    frame_ratio = frame_width / frame_height
    image_ratio = image_width / image_height
    left = top = right = bottom = 0
    if image_ratio > frame_ratio:
        crop = round((1.0 - (frame_ratio / image_ratio)) * 50000)
        left = right = max(0, min(50000, crop))
    elif image_ratio < frame_ratio:
        crop = round((1.0 - (image_ratio / frame_ratio)) * 50000)
        top = bottom = max(0, min(50000, crop))
    return left, top, right, bottom


def _apply_picture_cover_crops(
    slide_xml: bytes,
    crops: Mapping[str, tuple[int, int, int, int]],
) -> bytes:
    text = slide_xml.decode("utf-8", errors="strict")
    for slot_id, (left, top, right, bottom) in crops.items():
        shape_id = int(slot_id.removeprefix("shape_"))
        picture_pattern = re.compile(
            rf'(<p:pic\b.*?<p:cNvPr\b[^>]*\bid="{shape_id}"[^>]*>.*?</p:pic>)',
            re.DOTALL,
        )
        match = picture_pattern.search(text)
        if match is None:
            raise PhysicalAssemblyError(f"picture shape XML not found: {slot_id}")
        picture_xml = match.group(1)
        source_rect = (
            f'<a:srcRect l="{left}" t="{top}" r="{right}" b="{bottom}"/>'
        )
        if re.search(r"<a:srcRect\b[^>]*/>", picture_xml):
            updated_picture = re.sub(
                r"<a:srcRect\b[^>]*/>",
                source_rect,
                picture_xml,
                count=1,
            )
        else:
            blip_match = re.search(
                r"<a:blip\b[^>]*(?:/>|>.*?</a:blip>)",
                picture_xml,
                re.DOTALL,
            )
            if blip_match is None:
                raise PhysicalAssemblyError(f"picture blip XML not found: {slot_id}")
            updated_picture = (
                picture_xml[: blip_match.end()]
                + source_rect
                + picture_xml[blip_match.end() :]
            )
        text = text[: match.start()] + updated_picture + text[match.end() :]
    return text.encode("utf-8")


def _is_safe_https_target(target: str) -> bool:
    if any(ord(character) < 32 for character in target):
        return False
    try:
        parsed = urlsplit(target)
    except ValueError:
        return False
    return parsed.scheme.lower() == "https" and bool(parsed.netloc)


def _looks_like_external_relationship_target(target: str) -> bool:
    """Return whether ``target`` uses a URI/scheme instead of an OPC part path."""

    return bool(
        target.startswith("//")
        or re.match(r"^[A-Za-z][A-Za-z0-9+.-]*:", target)
    )


def _relationship_is_external(entry: Mapping[str, str]) -> bool:
    return entry.get("TargetMode", "") == "External"


def _relationship_security_issue(entry: Mapping[str, str]) -> str | None:
    """Return a stable security finding code, or ``None`` for a safe entry."""

    target = entry.get("Target", "")
    rel_type = entry.get("Type", "").lower()
    target_mode = entry.get("TargetMode", "")
    if target_mode not in {"", "Internal", "External"}:
        return "unsafe-target-mode"
    if any(fragment in rel_type for fragment in FORBIDDEN_REL_TYPE_FRAGMENTS):
        return "forbidden-relationship-type"
    if rel_type.rstrip("/").endswith("/script"):
        return "forbidden-relationship-type"
    target_path = target.split("?", 1)[0].split("#", 1)[0]
    if any(pattern.search(target_path) for pattern in FORBIDDEN_REL_TARGET_PATTERNS):
        return "forbidden-relationship-target"
    target_looks_external = _looks_like_external_relationship_target(target)
    if target_mode == "External":
        if not _is_safe_https_target(target):
            return "unsafe-external-target"
        return None
    if target_looks_external:
        return "external-target-mode-mismatch"
    return None


def _validate_relationship_security(
    entry: Mapping[str, str],
    *,
    context: str,
) -> bool:
    """Fail closed on unsafe semantics and return whether the target is external."""

    target = entry.get("Target", "")
    if not target:
        raise PhysicalAssemblyError(f"empty relationship target in {context}")
    issue = _relationship_security_issue(entry)
    if issue == "forbidden-relationship-type":
        raise PhysicalAssemblyError(
            f"forbidden relationship type in {context}: {entry.get('Type', '')}"
        )
    if issue == "forbidden-relationship-target":
        raise PhysicalAssemblyError(
            f"forbidden relationship target in {context}: {target}"
        )
    if issue == "unsafe-external-target":
        raise PhysicalAssemblyError(
            f"unsafe external relationship in {context}: {target}"
        )
    if issue in {"unsafe-target-mode", "external-target-mode-mismatch"}:
        raise PhysicalAssemblyError(
            f"unsafe relationship target mode in {context}: "
            f"{entry.get('TargetMode', '') or '<missing>'} for {target}"
        )
    return _relationship_is_external(entry)


def _relationship_is_discarded_source_metadata(entry: Mapping[str, str]) -> bool:
    """Return metadata relationships that must never cross the trust boundary.

    PowerPoint authoring add-ins commonly attach ``ppt/tags`` records to
    shapes.  They are neither visual dependencies nor editable customer
    content, and retaining them leaks source-specific authoring metadata.
    """

    rel_type = entry.get("Type", "").lower().rstrip("/")
    target = _normalise_zip_name(entry.get("Target", "")).lower()
    return rel_type.endswith("/tags") or "/tags/" in target


def _validate_relationship_entry(
    entry: Mapping[str, str],
    *,
    package_path: Path,
    rels_path: str,
) -> bool:
    """Validate one source relationship and return whether it is external."""

    return _validate_relationship_security(
        entry,
        context=f"{package_path}: {rels_path}",
    )


def _build_source_graph_from_context(
    source: _SourcePackageContext,
    slide_number: int,
) -> _SourceGraph:
    cached = source.closures.get(slide_number)
    if cached is not None:
        return cached
    slide_name = f"ppt/slides/slide{slide_number}.xml"
    if slide_name not in source.names:
        raise PhysicalAssemblyError(
            f"source missing {slide_name}: {source.package_path}"
        )
    slide_xml = source.archive.read(slide_name)
    rels_map: dict[str, bytes] = {}
    extras: dict[str, bytes] = {}
    content_types: dict[str, str] = {slide_name: source.content_type_for(slide_name)}
    layouts: list[str] = []
    masters: list[str] = []
    themes: list[str] = []
    seen_parts: set[str] = set()
    queue: deque[str] = deque([slide_name])

    while queue:
        owner = queue.popleft()
        if owner in seen_parts:
            continue
        seen_parts.add(owner)
        if owner != slide_name:
            if owner.startswith("ppt/slides/"):
                raise PhysicalAssemblyError(
                    f"cross-slide dependency is not importable: {owner}"
                )
            if owner not in source.names:
                raise PhysicalAssemblyError(
                    f"unresolved internal relationship in {source.package_path}: {owner}"
                )
            data = source.archive.read(owner)
            extras[owner] = data
            content_types[owner] = source.content_type_for(owner)
            if _LAYOUT_RE.match(owner):
                layouts.append(owner)
            elif _MASTER_RE.match(owner):
                masters.append(owner)
            elif owner.startswith("ppt/theme/"):
                themes.append(owner)

        rels_path = _rels_path_for_part(owner)
        if rels_path not in source.names:
            continue
        rels_xml = source.archive.read(rels_path)
        rels_map[rels_path] = rels_xml
        for entry in _parse_relationships(rels_xml):
            is_external = _validate_relationship_entry(
                entry,
                package_path=source.package_path,
                rels_path=rels_path,
            )
            if _relationship_is_discarded_source_metadata(entry):
                if is_external:
                    raise PhysicalAssemblyError(
                        f"external source metadata relationship is forbidden: {rels_path}"
                    )
                continue
            if is_external:
                continue
            resolved = _resolve_rel_target(rels_xml, rels_path, entry["Target"])
            if resolved is None or resolved not in source.names:
                raise PhysicalAssemblyError(
                    f"unresolved internal relationship in {source.package_path}: "
                    f"{rels_path} -> {entry['Target']}"
                )
            if resolved.startswith("ppt/slides/") and resolved != slide_name:
                raise PhysicalAssemblyError(
                    f"cross-slide dependency is not importable: {rels_path} -> {resolved}"
                )
            queue.append(resolved)

    graph = _SourceGraph(
        root_slide_name=slide_name,
        slide_xml=slide_xml,
        slide_sha=_sha256_bytes(slide_xml),
        rels=rels_map,
        extra_parts=extras,
        content_types=content_types,
        layout_paths=sorted(set(layouts)),
        master_paths=sorted(set(masters)),
        theme_paths=sorted(set(themes)),
    )
    source.closures[slide_number] = graph
    return graph


def _build_source_graph(package_path: Path, slide_number: int = 1) -> _SourceGraph:
    """Build a single-slide dependency graph from one source package.

    ``slide_number`` is deliberately explicit so a certified multi-page
    reference deck can contribute its original pages without first being
    redrawn into separate packages.
    """

    package_path = Path(package_path).expanduser().resolve(strict=True)
    package_sha = _sha256_file(package_path)
    source = _SourcePackageContext.open(package_path, package_sha)
    try:
        return _build_source_graph_from_context(source, slide_number)
    finally:
        source.close()


def _namespace_part_name(name: str, source_hash: str, ordinal: int = 0) -> str:
    """Return a stable package namespace, shared by all slides from a source."""

    normalised = _normalise_zip_name(name).lstrip("/")
    prefix = f"v61_{source_hash[:12]}"
    if normalised.startswith("ppt/"):
        return f"ppt/{prefix}/{normalised[len('ppt/'):]}"
    return f"ppt/{prefix}/package/{normalised}"


def _namespace_rels_name(name: str) -> str:
    return name  # We rely on part-name rewriting above to namespace rels parents.


def _namespace_relationship_path(
    rels_path: str,
    target_map: Mapping[str, str],
    *,
    source_slide_rels_name: str,
    ordinal: int,
) -> str:
    """Return the relationship part path for a namespaced owner part."""

    if rels_path.endswith(source_slide_rels_name):
        return f"ppt/slides/_rels/slide{ordinal}.xml.rels"
    parts = rels_path.split("/")
    if len(parts) >= 2 and parts[-2] == "_rels":
        owner = "/".join(parts[:-2] + [parts[-1][:-5]])
        mapped = target_map.get(owner)
        if mapped:
            return f"{os.path.dirname(mapped)}/_rels/{os.path.basename(mapped)}.rels"
    return rels_path


def _rewrite_relationship_targets(
    rels_xml: bytes,
    rels_path: str,
    target_map: Mapping[str, str],
    *,
    output_rels_path: str | None = None,
    relationship_overrides: Mapping[str, str] | None = None,
) -> bytes:
    """Rewrite ``Target`` attributes through ``target_map``.

    The supplied ``target_map`` keys may already be normalised (the caller
    stores resolved paths from the source OPC graph). We try both the
    normalised lookup and the un-normalised lookup, falling back to a path
    computed by relative-path math if neither hits.
    """

    entries = _parse_relationships(rels_xml)
    output_path = output_rels_path or rels_path
    output_owner = _owner_part_from_rels_path(output_path)
    new_base = posixpath.dirname(output_owner) if output_owner else ""
    overrides = dict(relationship_overrides or {})
    consumed_overrides: set[str] = set()
    rewritten_entries: list[dict[str, str]] = []

    for entry in entries:
        relationship_id = entry["Id"]
        raw_target = entry["Target"]
        is_external = _validate_relationship_security(
            entry,
            context=f"relationship rewrite {rels_path}",
        )
        if _relationship_is_discarded_source_metadata(entry):
            if is_external:
                raise PhysicalAssemblyError(
                    f"external source metadata relationship is forbidden: {rels_path}"
                )
            continue
        rewritten = dict(entry)
        if is_external:
            # Safe HTTPS relationships are copied semantically unchanged. XML
            # syntax is canonicalized below, so source quoting/prefixes cannot
            # influence the emitted package.
            rewritten_entries.append(rewritten)
            continue

        resolved = _resolve_rel_target(rels_xml, rels_path, raw_target)
        resolved_norm = _normalise_zip_name(resolved) if resolved else None
        new_target = overrides.get(relationship_id)
        if new_target is not None:
            consumed_overrides.add(relationship_id)
        elif resolved_norm and resolved_norm in target_map:
            new_target = target_map[resolved_norm]
        elif raw_target in target_map:
            new_target = target_map[raw_target]
        else:
            raise PhysicalAssemblyError(
                f"unmapped internal relationship: {rels_path} -> {raw_target}"
            )
        new_target = _normalise_zip_name(new_target)
        if _looks_like_external_relationship_target(new_target):
            raise PhysicalAssemblyError(
                f"relationship override must name an internal package part: {new_target}"
            )
        normalized_target = posixpath.normpath(new_target.lstrip("/"))
        if normalized_target in {"", ".", ".."} or normalized_target.startswith("../"):
            raise PhysicalAssemblyError(
                f"relationship override escapes the package: {new_target}"
            )
        rewritten["Target"] = posixpath.relpath(normalized_target, new_base or ".")
        rewritten_entries.append(rewritten)

    unused_overrides = set(overrides) - consumed_overrides
    if unused_overrides:
        raise PhysicalAssemblyError(
            "relationship overrides reference missing relationship IDs: "
            + ",".join(sorted(unused_overrides))
        )

    root = ET.Element(PACKAGE_RELATIONSHIPS_TAG)
    for entry in sorted(rewritten_entries, key=lambda item: item["Id"]):
        attributes = {
            "Id": entry["Id"],
            "Type": entry["Type"],
            "Target": entry["Target"],
        }
        if entry.get("TargetMode"):
            attributes["TargetMode"] = entry["TargetMode"]
        ET.SubElement(root, PACKAGE_RELATIONSHIP_TAG, attributes)
    return ET.tostring(
        root,
        encoding="utf-8",
        xml_declaration=True,
        short_empty_elements=True,
    )


def _rewrite_slide_references(slide_xml: bytes, target_map: Mapping[str, str]) -> bytes:
    """Rewrite XML references in slide body (r:embed, r:link, layout references)."""

    text = slide_xml.decode("utf-8", errors="replace")
    base = "ppt/slides"

    def relative(new_target: str) -> str:
        if new_target.startswith("/"):
            return new_target.lstrip("/")
        rel = os.path.relpath(new_target, base).replace("\\", "/")
        return rel

    for old, new in target_map.items():
        old_name = old.split("/")[-1]
        new_name = new.split("/")[-1]
        text = re.sub(
            r'(r:embed=")' + re.escape(old_name) + r'(")',
            r"\1" + new_name + r"\2",
            text,
        )
        text = re.sub(
            r'(r:link=")' + re.escape(old_name) + r'(")',
            r"\1" + new_name + r"\2",
            text,
        )
    return text.encode("utf-8")


_BODY_PR_RE = re.compile(
    r"<a:bodyPr\b[^>]*(?:/>|>.*?</a:bodyPr>)",
    re.DOTALL,
)
_AUTOFIT_PAIRED_RE = re.compile(
    r"<a:(?P<fit>spAutoFit|normAutofit|noAutofit)\b[^>]*>.*?</a:(?P=fit)>",
    re.DOTALL,
)
_AUTOFIT_EMPTY_RE = re.compile(
    r"<a:(?:spAutoFit|normAutofit|noAutofit)\b[^>]*/>",
    re.DOTALL,
)
_TEXT_SIZE_ATTRIBUTE_RE = re.compile(
    r'(<a:(?:rPr|defRPr|endParaRPr)\b[^>]*\bsz=")(\d+)(")',
    re.DOTALL,
)


def _scale_shape_text_runs(
    segment: str,
    *,
    slot_id: str,
    font_scale: int,
) -> str:
    """Scale only declared text-run sizes for a governed shrink-to-fit slot."""

    changed = 0

    def replace_size(match: re.Match[str]) -> str:
        nonlocal changed
        changed += 1
        source_size = int(match.group(2))
        scaled_size = max(800, round(source_size * font_scale / 100_000))
        return match.group(1) + str(scaled_size) + match.group(3)

    updated = _TEXT_SIZE_ATTRIBUTE_RE.sub(replace_size, segment)
    if changed == 0:
        # Some certified native text placeholders inherit their size from the
        # slide layout and therefore have no run-level ``sz`` attribute. The
        # governed caller still installs ``normAutofit`` on this exact text
        # body, which lets PowerPoint/LibreOffice fit the inherited style
        # without changing any geometry or inventing a font size. Treating
        # that safe case as a hard error makes an otherwise editable template
        # unusable for ordinary client copy.
        return segment
    return updated


def _normalise_shape_fit_policy(
    segment: str,
    *,
    slot_id: str,
    fit_policy: str,
    font_scale: int | None = None,
) -> str:
    """Set one shape bodyPr fit node without touching other shape styling."""

    fit_tag = {
        "no-autofit": "noAutofit",
        "shrink-to-fit": "normAutofit",
    }.get(fit_policy)
    if fit_tag is None:
        raise PhysicalAssemblyError(f"binding fit_policy is invalid: {slot_id}")
    if fit_policy == "shrink-to-fit":
        if type(font_scale) is not int or not 10_000 <= font_scale <= 100_000:
            raise PhysicalAssemblyError(
                f"shrink-to-fit scale is invalid: {slot_id}"
            )
        fit_xml = (
            f'<a:normAutofit fontScale="{font_scale}" '
            'lnSpcReduction="20000"/>'
        )
    else:
        fit_xml = "<a:noAutofit/>"

    matches = list(_BODY_PR_RE.finditer(segment))
    if len(matches) != 1:
        raise PhysicalAssemblyError(
            f"{fit_policy} binding requires one target bodyPr: {slot_id}"
        )
    match = matches[0]
    body_pr = match.group(0)
    if body_pr.endswith("/>"):
        replacement = body_pr[:-2] + f">{fit_xml}</a:bodyPr>"
    else:
        opening_end = body_pr.find(">") + 1
        opening = body_pr[:opening_end]
        inner = body_pr[opening_end:-len("</a:bodyPr>")]
        autofit_matches = sorted(
            [*_AUTOFIT_PAIRED_RE.finditer(inner), *_AUTOFIT_EMPTY_RE.finditer(inner)],
            key=lambda item: item.start(),
        )
        if autofit_matches:
            pieces: list[str] = []
            cursor = 0
            for index, autofit in enumerate(autofit_matches):
                pieces.append(inner[cursor:autofit.start()])
                if index == 0:
                    pieces.append(fit_xml)
                cursor = autofit.end()
            pieces.append(inner[cursor:])
            inner = "".join(pieces)
        else:
            inner = fit_xml + inner
        replacement = opening + inner + "</a:bodyPr>"
    return segment[:match.start()] + replacement + segment[match.end():]


def _adapt_slide_text(
    slide_xml: bytes,
    bindings: Mapping[str, str],
    *,
    allowed_slots: Iterable[str] | None = None,
    allowed_clear_alias_slots: Iterable[str] | None = None,
    fit_policies: Mapping[str, str] | None = None,
) -> bytes:
    """Apply declared bindings and fail closed on stale or invented slots."""

    text = slide_xml.decode("utf-8", errors="replace")
    allowed = set(allowed_slots or ())
    allowed_clear_aliases = set(allowed_clear_alias_slots or ())
    policies = dict(fit_policies or {})
    extra_policy_slots = set(policies) - set(bindings)
    if extra_policy_slots:
        raise PhysicalAssemblyError(
            "fit_policy targets an unbound slot: "
            + ",".join(sorted(extra_policy_slots))
        )
    for slot_id, replacement in bindings.items():
        if not slot_id.startswith("shape_"):
            raise PhysicalAssemblyError(f"invalid text slot id: {slot_id}")
        if allowed and slot_id not in allowed and slot_id not in allowed_clear_aliases:
            raise PhysicalAssemblyError(
                f"binding targets a slot outside the certified slot graph: {slot_id}"
            )
        try:
            shape_id = int(slot_id[len("shape_"):])
        except ValueError:
            raise PhysicalAssemblyError(f"invalid text slot id: {slot_id}") from None
        marker = re.compile(
            rf'<p:cNvPr\b[^>]*\bid="{shape_id}"[^>]*>'
        )
        m = marker.search(text)
        if m is None:
            raise PhysicalAssemblyError(f"declared text slot not found: {slot_id}")
        start = m.end()
        nxt = _CNVPR_RE.search(text, start)
        end = nxt.start() if nxt else len(text)
        segment = text[start:end]
        if not _TEXT_RE.search(segment):
            raise PhysicalAssemblyError(
                f"declared text slot has no editable text nodes: {slot_id}"
            )
        source_text = "".join(
            html.unescape(match.group(1)) for match in _TEXT_RE.finditer(segment)
        )
        new_segment = _replace_shape_text(segment, replacement)
        fit_policy = policies.get(slot_id, "preserve")
        if fit_policy not in TEXT_FIT_POLICIES:
            raise PhysicalAssemblyError(f"binding fit_policy is invalid: {slot_id}")
        if fit_policy in {"no-autofit", "shrink-to-fit"}:
            source_chars = len("".join(source_text.split()))
            replacement_chars = len("".join(replacement.split()))
            font_scale = (
                max(
                    40_000,
                    min(
                        100_000,
                        round(source_chars / replacement_chars * 100_000),
                    ),
                )
                if fit_policy == "shrink-to-fit" and replacement_chars > 0
                else None
            )
            if fit_policy == "shrink-to-fit" and font_scale is not None:
                new_segment = _scale_shape_text_runs(
                    new_segment,
                    slot_id=slot_id,
                    font_scale=font_scale,
                )
            new_segment = _normalise_shape_fit_policy(
                new_segment,
                slot_id=slot_id,
                fit_policy=fit_policy,
                font_scale=font_scale,
            )
        text = text[:start] + new_segment + text[end:]
    return text.encode("utf-8")


def _parse_content_types(xml_bytes: bytes) -> tuple[list[tuple[str, str]], list[tuple[str, str]]]:
    try:
        root = ET.fromstring(xml_bytes)
    except ET.ParseError as exc:
        raise PhysicalAssemblyError(f"invalid [Content_Types].xml: {exc}") from exc
    defaults: list[tuple[str, str]] = []
    overrides: list[tuple[str, str]] = []
    for child in root:
        local = child.tag.rsplit("}", 1)[-1]
        if local == "Default":
            extension = child.attrib.get("Extension", "").lower()
            content_type = child.attrib.get("ContentType", "")
            if extension and content_type:
                defaults.append((extension, content_type))
        elif local == "Override":
            part_name = child.attrib.get("PartName", "")
            content_type = child.attrib.get("ContentType", "")
            if part_name and content_type:
                overrides.append((part_name, content_type))
    return defaults, overrides


def _serialize_content_types(
    defaults: Sequence[tuple[str, str]],
    overrides: Sequence[tuple[str, str]],
) -> bytes:
    """Render [Content_Types].xml with stable field order."""

    seen_defaults: set[str] = set()
    seen_overrides: set[str] = set()
    parts: list[str] = [
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">',
    ]
    for ext, ct in defaults:
        ext = ext.lower()
        if ext in seen_defaults:
            continue
        seen_defaults.add(ext)
        parts.append(f'<Default Extension="{ext}" ContentType="{ct}"/>')
    for part_name, ct in overrides:
        if part_name in seen_overrides:
            continue
        seen_overrides.add(part_name)
        parts.append(f'<Override PartName="{part_name}" ContentType="{ct}"/>')
    parts.append("</Types>")
    return "".join(parts).encode("utf-8")


def _default_content_types() -> tuple[list[tuple[str, str]], list[tuple[str, str]]]:
    return [
        ("rels", "application/vnd.openxmlformats-package.relationships+xml"),
        ("xml", "application/xml"),
        ("png", "image/png"),
        ("jpeg", "image/jpeg"),
        ("jpg", "image/jpeg"),
        ("gif", "image/gif"),
        ("svg", "image/svg+xml"),
        ("wdp", "image/vnd.ms-photo"),
        ("emf", "image/x-emf"),
        ("wmf", "image/x-wmf"),
        ("bmp", "image/bmp"),
        ("tif", "image/tiff"),
        ("tiff", "image/tiff"),
    ], []


def _part_extension(part_name: str) -> str:
    filename = posixpath.basename(part_name)
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def _part_kind(part_name: str) -> str:
    name = _normalise_zip_name(part_name)
    if name.startswith("ppt/slides/") and name.endswith(".xml"):
        return "slide"
    if "slideLayouts/" in name:
        return "layout"
    if "slideMasters/" in name:
        return "master"
    if "/theme/" in name:
        return "theme"
    if "/media/" in name:
        return "media"
    if "/charts/" in name:
        return "chart"
    if "/embeddings/" in name:
        return "embedding"
    if "/diagrams/" in name:
        return "diagram"
    if "/notesSlides/" in name or "/notesMasters/" in name:
        return "notes"
    if "/comments/" in name or "/commentAuthors" in name:
        return "comments"
    if name.endswith(".rels"):
        return "relationships"
    return "other"


def _is_cross_package_dedup_candidate(
    source_part: str,
    content_type: str,
    *,
    relationship_free: bool,
) -> bool:
    """Allow only immutable binaries and relationship-free style/theme XML."""

    name = _normalise_zip_name(source_part).lower()
    ctype = content_type.lower()
    if (
        "/media/" in name
        or ctype.startswith("image/")
        or ctype.startswith("audio/")
        or ctype.startswith("video/")
    ):
        return True
    if not relationship_free:
        return False
    return (
        "/theme/" in name
        or "chartstyle" in name
        or "chartcolorstyle" in name
        or "chartstyle" in ctype
        or "chartcolorstyle" in ctype
    )


def _cross_package_target_name(source_part: str, data: bytes) -> str:
    digest = _sha256_bytes(data)[:20]
    extension = _part_extension(source_part) or "bin"
    lower = source_part.lower()
    if "/media/" in lower:
        return f"ppt/media/v61_dedup_{digest}.{extension}"
    if "/theme/" in lower:
        return f"ppt/theme/v61_dedup_{digest}.{extension}"
    if "chartcolorstyle" in lower:
        return f"ppt/charts/v61_color_{digest}.{extension}"
    return f"ppt/charts/v61_style_{digest}.{extension}"


def _is_slide_scoped_part(source_part: str) -> bool:
    name = _normalise_zip_name(source_part)
    return any(
        marker in name
        for marker in (
            "/notesSlides/",
            "/comments/",
            "/tags/",
        )
    )


@dataclass
class AssemblyImportContext:
    """Output-level importer state shared by every selected source slide."""

    parts: dict[str, bytes] = field(default_factory=dict)
    content_type_defaults: dict[str, str] = field(default_factory=dict)
    content_type_overrides: dict[str, str] = field(default_factory=dict)
    source_packages: dict[tuple[str, str], _SourcePackageContext] = field(default_factory=dict)
    closures: dict[tuple[str, int], _SourceGraph] = field(default_factory=dict)
    source_part_map: dict[tuple[str, str, str], str] = field(default_factory=dict)
    binary_hash_map: dict[tuple[str, str, str], str] = field(default_factory=dict)
    binary_hash_origin: dict[tuple[str, str, str], str] = field(default_factory=dict)
    target_parts: set[str] = field(default_factory=set)
    imported_parts: set[str] = field(default_factory=set)
    same_source_reuse_count: int = 0
    same_source_reuse_bytes: int = 0
    cross_source_safe_dedup_count: int = 0
    cross_source_safe_dedup_bytes: int = 0
    deduplicated_part_count: int = 0
    deduplicated_bytes: int = 0
    source_sizes: dict[str, int] = field(default_factory=dict)
    slide_target_maps: dict[int, dict[str, str]] = field(default_factory=dict)
    replaced_source_parts_by_slide: dict[int, set[str]] = field(default_factory=dict)
    pruned_parts: set[str] = field(default_factory=set)

    def __post_init__(self) -> None:
        if not self.content_type_defaults:
            defaults, _ = _default_content_types()
            self.content_type_defaults.update(defaults)

    def open_source(
        self,
        package_path: Path,
        package_sha256: str,
    ) -> _SourcePackageContext:
        resolved = package_path.expanduser().resolve(strict=True)
        key = (package_sha256, str(resolved))
        source = self.source_packages.get(key)
        if source is None:
            source = _SourcePackageContext.open(resolved, package_sha256)
            self.source_packages[key] = source
            self.source_sizes.setdefault(package_sha256, resolved.stat().st_size)
        return source

    def graph_for(
        self,
        package_path: Path,
        package_sha256: str,
        slide_number: int,
    ) -> tuple[_SourcePackageContext, _SourceGraph]:
        source = self.open_source(package_path, package_sha256)
        key = (package_sha256, slide_number)
        graph = self.closures.get(key)
        if graph is None:
            graph = _build_source_graph_from_context(source, slide_number)
            self.closures[key] = graph
        return source, graph

    def add_part(self, target_name: str, data: bytes, content_type: str | None = None) -> None:
        target = _normalise_zip_name(target_name).lstrip("/")
        existing = self.parts.get(target)
        if existing is not None and existing != data:
            raise PhysicalAssemblyError(f"target part collision: {target}")
        self.parts[target] = data
        self.target_parts.add(target)
        if content_type:
            self.content_type_overrides["/" + target] = content_type

    def allocate_dependency(
        self,
        source: _SourcePackageContext,
        source_part: str,
        data: bytes,
        content_type: str,
        *,
        mutation_scope: str = "shared",
        relationship_free: bool = False,
    ) -> str:
        map_key = (source.package_sha256, source_part, mutation_scope)
        mapped = self.source_part_map.get(map_key)
        if mapped is not None:
            self.same_source_reuse_count += 1
            self.same_source_reuse_bytes += len(data)
            self.deduplicated_part_count += 1
            self.deduplicated_bytes += len(data)
            return mapped

        cross_key: tuple[str, str, str] | None = None
        if mutation_scope == "shared" and _is_cross_package_dedup_candidate(
            source_part,
            content_type,
            relationship_free=relationship_free,
        ):
            cross_key = (content_type, _part_extension(source_part), _sha256_bytes(data))
            mapped = self.binary_hash_map.get(cross_key)
            if mapped is not None:
                self.source_part_map[map_key] = mapped
                if self.binary_hash_origin.get(cross_key) == source.package_sha256:
                    self.same_source_reuse_count += 1
                    self.same_source_reuse_bytes += len(data)
                else:
                    self.cross_source_safe_dedup_count += 1
                    self.cross_source_safe_dedup_bytes += len(data)
                self.deduplicated_part_count += 1
                self.deduplicated_bytes += len(data)
                return mapped

        if cross_key is not None:
            target = _cross_package_target_name(source_part, data)
        else:
            target = _namespace_part_name(source_part, source.package_sha256)
            if mutation_scope != "shared":
                prefix, dot, extension = target.rpartition(".")
                suffix = re.sub(r"[^A-Za-z0-9_-]+", "_", mutation_scope)
                target = f"{prefix}_{suffix}.{extension}" if dot else f"{target}_{suffix}"

        if target in self.parts and self.parts[target] != data:
            stem, dot, extension = target.rpartition(".")
            suffix = _sha256_bytes(data)[:12]
            target = f"{stem}_{suffix}.{extension}" if dot else f"{target}_{suffix}"
        self.add_part(target, data, content_type)
        self.imported_parts.add(target)
        self.source_part_map[map_key] = target
        if cross_key is not None:
            self.binary_hash_map[cross_key] = target
            self.binary_hash_origin[cross_key] = source.package_sha256
        return target

    def allocate_locked_asset(self, asset: LockedAsset) -> str:
        data = asset.path.read_bytes()
        content_type = _locked_asset_content_type(asset.path)
        extension = asset.path.suffix.lower().lstrip("/").lstrip(".")
        key = (content_type, extension, asset.sha256)
        mapped = self.binary_hash_map.get(key)
        if mapped is not None:
            if self.binary_hash_origin.get(key) == f"authority:{asset.sha256}":
                self.same_source_reuse_count += 1
                self.same_source_reuse_bytes += len(data)
            else:
                self.cross_source_safe_dedup_count += 1
                self.cross_source_safe_dedup_bytes += len(data)
            self.deduplicated_part_count += 1
            self.deduplicated_bytes += len(data)
            return mapped
        target = f"ppt/media/authority_{asset.sha256[:20]}.{extension}"
        self.add_part(target, data, content_type)
        self.imported_parts.add(target)
        self.binary_hash_map[key] = target
        self.binary_hash_origin[key] = f"authority:{asset.sha256}"
        return target

    def close(self) -> None:
        for source in self.source_packages.values():
            source.close()

    def metrics(self, output_path: Path, unresolved_count: int) -> AssemblyMetrics:
        output_size = output_path.stat().st_size if output_path.is_file() else 0
        source_size = sum(self.source_sizes.values())
        return AssemblyMetrics(
            output_size_bytes=output_size,
            source_size_bytes=source_size,
            unique_source_package_count=len(self.source_sizes),
            imported_part_count=len(self.imported_parts),
            imported_parts=tuple(sorted(self.imported_parts)),
            unique_dependency_part_count=len(
                set(self.source_part_map.values()) - self.pruned_parts
            ),
            same_source_reuse_count=self.same_source_reuse_count,
            same_source_reuse_bytes=self.same_source_reuse_bytes,
            cross_source_safe_dedup_count=self.cross_source_safe_dedup_count,
            cross_source_safe_dedup_bytes=self.cross_source_safe_dedup_bytes,
            deduplicated_part_count=self.deduplicated_part_count,
            deduplicated_bytes=self.deduplicated_bytes,
            static_duplicate_bytes=_static_duplicate_bytes(output_path),
            unresolved_internal_relationship_count=unresolved_count,
            amplification_ratio=round(output_size / source_size, 6) if source_size else 0.0,
            parts_by_kind=dict(sorted(Counter(_part_kind(name) for name in self.parts).items())),
        )


def _prune_unreachable_parts(context: AssemblyImportContext) -> None:
    """Remove imported OPC parts that are no longer reachable from package root.

    A picture relationship override intentionally disconnects the source image.
    Copying that now-orphaned media into the deliverable both inflates the file
    and leaks source material.  Reachability is therefore computed after every
    relationship has been rewritten and after the generated presentation/root
    relationships have been installed, but before content types are serialized.
    """

    root_rels_path = "_rels/.rels"
    if root_rels_path not in context.parts:
        raise PhysicalAssemblyError("OUTPUT_ROOT_RELATIONSHIPS_MISSING")

    reachable: set[str] = {root_rels_path}
    visited_owners: set[str] = set()
    queue: deque[str] = deque()

    def enqueue_relationship_targets(rels_path: str) -> None:
        rels_xml = context.parts.get(rels_path)
        if rels_xml is None:
            return
        reachable.add(rels_path)
        for entry in _parse_relationships(rels_xml):
            is_external = _validate_relationship_security(
                entry,
                context=f"assembled relationship graph {rels_path}",
            )
            if _relationship_is_discarded_source_metadata(entry):
                raise PhysicalAssemblyError(
                    f"OUTPUT_DISCARDED_METADATA_RELATIONSHIP: {rels_path}"
                )
            if is_external:
                continue
            target = _resolve_rel_target(
                rels_xml,
                rels_path,
                entry.get("Target", ""),
            )
            if target is None or target not in context.parts:
                raise PhysicalAssemblyError(
                    "OUTPUT_UNRESOLVED_RELATIONSHIP_BEFORE_PRUNE: "
                    f"{rels_path}->{entry.get('Target', '')}"
                )
            queue.append(target)

    enqueue_relationship_targets(root_rels_path)
    while queue:
        owner = queue.popleft()
        if owner in visited_owners:
            continue
        visited_owners.add(owner)
        reachable.add(owner)
        enqueue_relationship_targets(_rels_path_for_part(owner))

    removable = set(context.parts) - reachable
    # [Content_Types].xml is intentionally created after this pass.  Treating
    # an early copy as ordinary content would hide an ordering bug.
    if "[Content_Types].xml" in removable:
        raise PhysicalAssemblyError("CONTENT_TYPES_SERIALIZED_BEFORE_PRUNE")
    for part_name in sorted(removable):
        context.parts.pop(part_name, None)
        context.target_parts.discard(part_name)
        context.imported_parts.discard(part_name)
        context.content_type_overrides.pop("/" + part_name, None)
    context.pruned_parts.update(removable)


def _default_pres_rels() -> bytes:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>'
        '<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>'
        '<Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>'
        "</Relationships>"
    ).encode("utf-8")


def _root_rels() -> bytes:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>'
        "</Relationships>"
    ).encode("utf-8")


def _default_pres_xml(
    slide_paths: Sequence[str],
    master_paths: Sequence[str],
    notes_master_paths: Sequence[str] = (),
) -> bytes:
    """Build a presentation.xml pointing at the slides in stable order."""

    sld_ids = []
    slide_start = len(master_paths) + len(notes_master_paths) + 5
    for ordinal, path in enumerate(slide_paths, start=1):
        relationship_id = slide_start + ordinal - 1
        sld_ids.append(
            f'<p:sldId id="{255 + ordinal}" r:id="rId{relationship_id}"/>'
        )
    sld_id_list = "".join(sld_ids)
    master_ids = "".join(
        f'<p:sldMasterId id="{2147483648 + idx}" r:id="rId{idx + 1}"/>'
        for idx, _ in enumerate(master_paths)
    )
    notes_master_ids = "".join(
        f'<p:notesMasterId r:id="rId{len(master_paths) + idx}"/>'
        for idx, _ in enumerate(notes_master_paths, start=1)
    )
    notes_master_list = (
        f"<p:notesMasterIdLst>{notes_master_ids}</p:notesMasterIdLst>"
        if notes_master_ids
        else ""
    )
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        f"<p:sldMasterIdLst>{master_ids}</p:sldMasterIdLst>"
        f"{notes_master_list}"
        f"<p:sldIdLst>{sld_id_list}</p:sldIdLst>"
        '<p:sldSz cx="12192000" cy="6858000" type="screen16x9"/>'
        '<p:notesSz cx="6858000" cy="9144000"/>'
        "</p:presentation>"
    ).encode("utf-8")


def _pres_xml_rels(
    slide_rels: Sequence[str],
    master_paths: Sequence[str],
    notes_master_paths: Sequence[str] = (),
) -> bytes:
    entries: list[str] = [
        f'<Relationship Id="rId{idx}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="{target}"/>'
        for idx, target in enumerate(master_paths, start=1)
    ]
    entries.extend(
        f'<Relationship Id="rId{len(master_paths) + idx}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="{target}"/>'
        for idx, target in enumerate(notes_master_paths, start=1)
    )
    support_start = len(master_paths) + len(notes_master_paths) + 1
    entries.extend([
        f'<Relationship Id="rId{support_start}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/presProps" Target="presProps.xml"/>',
        f'<Relationship Id="rId{support_start + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/viewProps" Target="viewProps.xml"/>',
        f'<Relationship Id="rId{support_start + 2}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/tableStyles" Target="tableStyles.xml"/>',
    ])
    for idx, target in enumerate(
        slide_rels,
        start=len(master_paths) + len(notes_master_paths) + 5,
    ):
        entries.append(
            f'<Relationship Id="rId{idx}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="{target}"/>'
        )
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        + "".join(entries)
        + "</Relationships>"
    ).encode("utf-8")


def _empty_doc_props() -> tuple[bytes, bytes]:
    core = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" '
        b'xmlns:dc="http://purl.org/dc/elements/1.1/" '
        b'xmlns:dcterms="http://purl.org/dc/terms/" '
        b'xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">'
        b"<dc:title>Window-PPTX v6.1 Deck</dc:title>"
        b"<dc:creator>window-pptx</dc:creator>"
        b"</cp:coreProperties>"
    )
    app = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" '
        b'xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">'
        b"<Application>window-pptx</Application>"
        b"</Properties>"
    )
    return core, app


def _default_props_xml() -> bytes:
    return (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<p:presentationPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        b'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
    )


def _default_view_props_xml() -> bytes:
    return (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<p:viewPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
    )


def _default_table_styles_xml() -> bytes:
    return (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<a:tblStyleLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
    )


TEXT_RENDER_MODES = frozenset(
    {"exact", "whitespace", "slice", "join", "clear", "source"}
)
TEXT_RENDER_FIELDS = frozenset({"text", "value", "value_unit"})
TEXT_RENDER_SEPARATORS = frozenset({"", " ", "\n", " / ", " · ", "：", ": "})
TEXT_FIT_POLICIES = frozenset({"preserve", "no-autofit", "shrink-to-fit"})
ASSET_FIT_MODES = frozenset({"cover"})


def _parse_text_binding_spec(slot_id: str, value: Any) -> TextBindingSpec:
    if not isinstance(value, Mapping):
        raise PhysicalAssemblyError(
            f"production binding must be an evidence object: {slot_id}"
        )
    allowed = {"replacement", "fact_refs", "render", "fit_policy"}
    unknown = set(value) - allowed
    if unknown:
        raise PhysicalAssemblyError(
            f"unknown production binding fields for {slot_id}: {','.join(sorted(unknown))}"
        )
    replacement = value.get("replacement")
    refs = value.get("fact_refs")
    render = value.get("render")
    fit_policy = value.get("fit_policy", "preserve")
    if not isinstance(replacement, str):
        raise PhysicalAssemblyError(f"binding replacement must be text: {slot_id}")
    if not isinstance(refs, list) or not all(
        isinstance(ref, str) and ref for ref in refs
    ):
        raise PhysicalAssemblyError(f"binding fact_refs must be a string array: {slot_id}")
    if not isinstance(render, Mapping):
        raise PhysicalAssemblyError(f"binding render must be an object: {slot_id}")
    if fit_policy not in TEXT_FIT_POLICIES:
        raise PhysicalAssemblyError(f"binding fit_policy is invalid: {slot_id}")
    render_allowed = {"mode", "field", "separator", "slice_start", "slice_end"}
    render_unknown = set(render) - render_allowed
    if render_unknown:
        raise PhysicalAssemblyError(
            f"unknown render fields for {slot_id}: {','.join(sorted(render_unknown))}"
        )
    mode = render.get("mode")
    field_name = render.get("field", "text")
    separator = render.get("separator", "")
    if mode not in TEXT_RENDER_MODES:
        raise PhysicalAssemblyError(f"binding render mode is invalid: {slot_id}")
    if field_name not in TEXT_RENDER_FIELDS:
        raise PhysicalAssemblyError(f"binding render field is invalid: {slot_id}")
    if separator not in TEXT_RENDER_SEPARATORS:
        raise PhysicalAssemblyError(f"binding separator is not governed: {slot_id}")
    slice_start = render.get("slice_start")
    slice_end = render.get("slice_end")
    if slice_start is not None and (
        not isinstance(slice_start, int) or slice_start < 0
    ):
        raise PhysicalAssemblyError(f"slice_start is invalid: {slot_id}")
    if slice_end is not None and (
        not isinstance(slice_end, int) or slice_end < 0
    ):
        raise PhysicalAssemblyError(f"slice_end is invalid: {slot_id}")
    return TextBindingSpec(
        replacement=replacement,
        fact_refs=tuple(refs),
        mode=str(mode),
        fit_policy=str(fit_policy),
        field=str(field_name),
        separator=str(separator),
        slice_start=slice_start,
        slice_end=slice_end,
    )


def _parse_asset_binding_spec(slot_id: str, value: Any) -> AssetBindingSpec:
    if not isinstance(value, Mapping) or set(value) - {"asset_ref", "fit"}:
        raise PhysicalAssemblyError(f"asset binding is invalid: {slot_id}")
    asset_ref = value.get("asset_ref")
    fit = value.get("fit", "cover")
    if not isinstance(asset_ref, str) or not asset_ref:
        raise PhysicalAssemblyError(f"asset_ref is invalid: {slot_id}")
    if fit not in ASSET_FIT_MODES:
        raise PhysicalAssemblyError(f"asset fit mode is invalid: {slot_id}")
    return AssetBindingSpec(asset_ref=asset_ref, fit=str(fit))


def _resolve_authority_path(value: str, base_dir: Path | None) -> str:
    """Legacy resolver retained for in-process/unit callers."""

    path = Path(value).expanduser()
    if not path.is_absolute() and base_dir is not None:
        path = base_dir / path
    return str(path.resolve(strict=False))


def resolve_project_file(
    value: str | os.PathLike[str],
    project_root: str | os.PathLike[str],
    *,
    label: str,
    require_file: bool = True,
) -> Path:
    """Resolve one trusted project-relative file without following symlinks.

    This is intentionally stricter than ``Path.resolve``: production authority
    files must be ordinary files rooted directly in the clean client project.
    A symlink anywhere from the project root to the leaf is rejected before any
    bytes are read.
    """

    raw_root = Path(project_root).expanduser()
    if raw_root.is_symlink():
        raise PhysicalAssemblyError(f"PROJECT_ROOT_SYMLINK_REJECTED: {raw_root}")
    root = raw_root.resolve(strict=True)
    if not root.is_dir():
        raise PhysicalAssemblyError(f"PROJECT_ROOT_NOT_DIRECTORY: {root}")
    raw_value = os.fspath(value)
    relative = Path(raw_value)
    if (
        not raw_value
        or relative.is_absolute()
        or raw_value.startswith("~")
        or "\\" in raw_value
        or re.match(r"^[A-Za-z]:", raw_value)
        or any(part in {"", ".."} for part in relative.parts)
    ):
        raise PhysicalAssemblyError(f"{label}_PATH_NOT_PROJECT_RELATIVE: {raw_value}")
    candidate = root
    for part in relative.parts:
        if part == ".":
            continue
        candidate = candidate / part
        if candidate.is_symlink():
            raise PhysicalAssemblyError(f"{label}_SYMLINK_REJECTED: {raw_value}")
    lexical = Path(os.path.abspath(candidate))
    try:
        lexical.relative_to(root)
    except ValueError as exc:
        raise PhysicalAssemblyError(f"{label}_PATH_ESCAPE: {raw_value}") from exc
    if require_file and not lexical.is_file():
        raise PhysicalAssemblyError(f"{label}_MISSING: {raw_value}")
    return lexical


def _plan_from_payload(
    payload: Mapping[str, Any],
    library_lookup: Mapping[str, PageTemplate],
    *,
    base_dir: Path | None = None,
    project_root: Path | None = None,
) -> AssemblyPlan:
    """Build an AssemblyPlan from a payload dict."""

    slides: list[AssemblyTargetSlide] = []
    for entry in payload["target_slides"]:
        page_id = entry["page_id"]
        template = library_lookup.get(page_id)
        if template is None:
            raise PhysicalAssemblyError(f"library missing page_id={page_id}")
        if template.package_sha256 != entry["package_sha256"]:
            raise PhysicalAssemblyError(
                f"package_sha256 mismatch for page_id={page_id}"
            )
        raw_bindings = entry.get("bindings", {})
        if not isinstance(raw_bindings, Mapping):
            raise PhysicalAssemblyError("slide bindings must be an object")
        legacy_bindings: dict[str, str] = {}
        text_binding_specs: dict[str, TextBindingSpec] = {}
        governed_content_binding_specs: dict[str, TextBindingSpec] = {}
        asset_binding_specs: dict[str, AssetBindingSpec] = {}
        expected_text_slots = set(template.slot_graph.get("text_slot_ids", ()))
        inventory = template.governed_content_inventory
        raw_content_slots = inventory.get("slots", ()) if isinstance(inventory, Mapping) else ()
        expected_content_slots = {
            str(record.get("slot_id"))
            for record in raw_content_slots
            if isinstance(record, Mapping) and isinstance(record.get("slot_id"), str)
        }
        expected_content_groups = {
            str(record.get("peer_group_id"))
            for record in raw_content_slots
            if isinstance(record, Mapping)
            and isinstance(record.get("peer_group_id"), str)
            and record.get("peer_group_id")
        }
        for slot_id, value in raw_bindings.items():
            if not isinstance(slot_id, str):
                raise PhysicalAssemblyError("binding slot IDs must be strings")
            if isinstance(value, str):
                legacy_bindings[slot_id] = value
            elif isinstance(value, Mapping) and {
                "text",
                "fact_refs",
                "asset_refs",
            }.issubset(value):
                unknown_binding_fields = set(value) - {
                    "text",
                    "fact_refs",
                    "asset_refs",
                    "fit_policy",
                }
                if unknown_binding_fields:
                    raise PhysicalAssemblyError(
                        f"unknown binding fields for {slot_id}: "
                        + ",".join(sorted(unknown_binding_fields))
                    )
                text = value.get("text")
                fact_refs = value.get("fact_refs")
                asset_refs = value.get("asset_refs")
                fit_policy = value.get("fit_policy", "preserve")
                if not isinstance(text, str):
                    raise PhysicalAssemblyError(f"binding text is invalid: {slot_id}")
                if not isinstance(fact_refs, list) or not all(
                    isinstance(ref, str) and ref for ref in fact_refs
                ):
                    raise PhysicalAssemblyError(
                        f"binding fact_refs are invalid: {slot_id}"
                    )
                if not isinstance(asset_refs, list) or not all(
                    isinstance(ref, str) and ref for ref in asset_refs
                ):
                    raise PhysicalAssemblyError(
                        f"binding asset_refs are invalid: {slot_id}"
                    )
                if len(fact_refs) != len(set(fact_refs)) or len(asset_refs) != len(
                    set(asset_refs)
                ):
                    raise PhysicalAssemblyError(f"binding refs must be unique: {slot_id}")
                if fit_policy not in TEXT_FIT_POLICIES:
                    raise PhysicalAssemblyError(
                        f"binding fit_policy is invalid: {slot_id}"
                    )
                if slot_id in expected_text_slots:
                    if asset_refs:
                        raise PhysicalAssemblyError(
                            f"text slot cannot replace a picture asset: {slot_id}"
                        )
                    legacy_bindings[slot_id] = text
                    text_binding_specs[slot_id] = TextBindingSpec(
                        replacement=text,
                        fact_refs=tuple(fact_refs),
                        mode="auto",
                        fit_policy=str(fit_policy),
                    )
                elif slot_id in expected_content_slots or slot_id in expected_content_groups:
                    if asset_refs:
                        raise PhysicalAssemblyError(
                            f"governed content slot cannot replace an asset: {slot_id}"
                        )
                    governed_content_binding_specs[slot_id] = TextBindingSpec(
                        replacement=text,
                        fact_refs=tuple(fact_refs),
                        mode="auto",
                        fit_policy=str(fit_policy),
                    )
                else:
                    if text or fact_refs or len(asset_refs) != 1:
                        raise PhysicalAssemblyError(
                            f"picture slot requires empty text/fact_refs and one asset_ref: {slot_id}"
                        )
                    if fit_policy != "preserve":
                        raise PhysicalAssemblyError(
                            f"picture slot cannot declare text fit_policy: {slot_id}"
                        )
                    asset_binding_specs[slot_id] = AssetBindingSpec(asset_refs[0], "cover")
            else:
                spec = _parse_text_binding_spec(slot_id, value)
                if slot_id in expected_text_slots:
                    legacy_bindings[slot_id] = spec.replacement
                    text_binding_specs[slot_id] = spec
                elif slot_id in expected_content_slots or slot_id in expected_content_groups:
                    governed_content_binding_specs[slot_id] = spec
                else:
                    raise PhysicalAssemblyError(
                        f"binding targets an unknown slot: {slot_id}"
                    )
        raw_asset_bindings = entry.get("asset_bindings", {})
        if not isinstance(raw_asset_bindings, Mapping):
            raise PhysicalAssemblyError("slide asset_bindings must be an object")
        asset_binding_specs.update(
            {
                str(slot_id): _parse_asset_binding_spec(str(slot_id), value)
                for slot_id, value in raw_asset_bindings.items()
            }
        )
        selection_payload = entry.get("selection")
        selection_evidence: SelectionEvidence | None = None
        if selection_payload is not None:
            if not isinstance(selection_payload, Mapping) or set(selection_payload) != {
                "query_id",
                "candidate_rank",
                "score_total",
                "selection_reason",
                "fallback_reason",
            }:
                raise PhysicalAssemblyError("slide selection evidence is invalid")
            if (
                not isinstance(selection_payload.get("query_id"), str)
                or not selection_payload["query_id"]
                or not isinstance(selection_payload.get("candidate_rank"), int)
                or selection_payload["candidate_rank"] < 1
                or not isinstance(selection_payload.get("score_total"), (int, float))
                or isinstance(selection_payload.get("score_total"), bool)
                or not isinstance(selection_payload.get("selection_reason"), str)
                or not selection_payload["selection_reason"]
                or (
                    selection_payload.get("fallback_reason") is not None
                    and not isinstance(selection_payload.get("fallback_reason"), str)
                )
            ):
                raise PhysicalAssemblyError("slide selection evidence fields are invalid")
            selection_evidence = SelectionEvidence(
                query_id=selection_payload["query_id"],
                candidate_rank=selection_payload["candidate_rank"],
                score_total=float(selection_payload["score_total"]),
                selection_reason=selection_payload["selection_reason"],
                fallback_reason=selection_payload.get("fallback_reason"),
            )
        slides.append(
            AssemblyTargetSlide(
                ordinal=int(entry["ordinal"]),
                page_template=template,
                bindings=legacy_bindings,
                narrative_role=str(entry["narrative_role"]),
                title=str(entry["title"]),
                headline=str(entry.get("headline", "")),
                text_binding_specs=text_binding_specs,
                governed_content_binding_specs=governed_content_binding_specs,
                asset_binding_specs=asset_binding_specs,
                selection_evidence=selection_evidence,
            )
        )
    authority_payload = payload.get("authority")
    authority: AuthorityLock | None = None
    if authority_payload is not None:
        authority_names = {"fact_store", "asset_manifest", "connective_copy"}
        if not isinstance(authority_payload, Mapping) or set(authority_payload) != authority_names:
            raise PhysicalAssemblyError("assembly authority block is invalid")
        fact_lock = authority_payload.get("fact_store")
        asset_lock = authority_payload.get("asset_manifest")
        connective_lock = authority_payload.get("connective_copy")
        if not isinstance(fact_lock, Mapping) or set(fact_lock) != {"path", "sha256"}:
            raise PhysicalAssemblyError("authority.fact_store is invalid")
        if not isinstance(asset_lock, Mapping) or set(asset_lock) != {"path", "sha256"}:
            raise PhysicalAssemblyError("authority.asset_manifest is invalid")
        if not isinstance(connective_lock, Mapping) or set(connective_lock) != {"path", "sha256"}:
            raise PhysicalAssemblyError("authority.connective_copy is invalid")
        for label, lock in (
            ("fact_store", fact_lock),
            ("asset_manifest", asset_lock),
            ("connective_copy", connective_lock),
        ):
            if not isinstance(lock.get("path"), str) or not lock.get("path"):
                raise PhysicalAssemblyError(f"authority.{label}.path is invalid")
            if not isinstance(lock.get("sha256"), str) or not re.fullmatch(
                r"[0-9a-f]{64}", lock["sha256"]
            ):
                raise PhysicalAssemblyError(f"authority.{label}.sha256 is invalid")
            if project_root is not None:
                resolve_project_file(
                    lock["path"],
                    project_root,
                    label=f"AUTHORITY_{label.upper()}",
                )
        authority = AuthorityLock(
            fact_store_path=(
                fact_lock["path"]
                if project_root is not None
                else _resolve_authority_path(fact_lock["path"], base_dir)
            ),
            fact_store_sha256=fact_lock["sha256"],
            asset_manifest_path=(
                asset_lock["path"]
                if project_root is not None
                else _resolve_authority_path(asset_lock["path"], base_dir)
            ),
            asset_manifest_sha256=asset_lock["sha256"],
            connective_copy_path=(
                connective_lock["path"]
                if project_root is not None
                else _resolve_authority_path(connective_lock["path"], base_dir)
            ),
            connective_copy_sha256=connective_lock["sha256"],
        )
    query_bundle_payload = payload.get("query_bundle")
    query_bundle_path: str | None = None
    query_bundle_sha256: str | None = None
    if query_bundle_payload is not None:
        if not isinstance(query_bundle_payload, Mapping) or set(query_bundle_payload) != {
            "path", "sha256"
        }:
            raise PhysicalAssemblyError("assembly query_bundle block is invalid")
        if not isinstance(query_bundle_payload.get("path"), str) or not query_bundle_payload["path"]:
            raise PhysicalAssemblyError("assembly query_bundle.path is invalid")
        if not isinstance(query_bundle_payload.get("sha256"), str) or not re.fullmatch(
            r"[0-9a-f]{64}", query_bundle_payload["sha256"]
        ):
            raise PhysicalAssemblyError("assembly query_bundle.sha256 is invalid")
        if project_root is not None:
            resolve_project_file(
                query_bundle_payload["path"],
                project_root,
                label="QUERY_BUNDLE",
            )
            query_bundle_path = query_bundle_payload["path"]
        else:
            query_bundle_path = _resolve_authority_path(
                query_bundle_payload["path"], base_dir
            )
        query_bundle_sha256 = query_bundle_payload["sha256"]
    profile_authority_payload = payload.get("binding_profile_authority")
    binding_profile_authority: BindingProfileAuthorityLock | None = None
    if profile_authority_payload is not None:
        if (
            not isinstance(profile_authority_payload, Mapping)
            or set(profile_authority_payload) != {"profile_id", "profile_sha256"}
            or not isinstance(profile_authority_payload.get("profile_id"), str)
            or re.fullmatch(
                r"[a-z0-9_-]{1,120}",
                str(profile_authority_payload.get("profile_id", "")),
            )
            is None
            or not isinstance(profile_authority_payload.get("profile_sha256"), str)
            or re.fullmatch(
                r"[0-9a-f]{64}",
                str(profile_authority_payload.get("profile_sha256", "")),
            )
            is None
        ):
            raise PhysicalAssemblyError(
                "assembly binding_profile_authority is invalid"
            )
        binding_profile_authority = BindingProfileAuthorityLock(
            profile_id=str(profile_authority_payload["profile_id"]),
            profile_sha256=str(profile_authority_payload["profile_sha256"]),
        )
    return AssemblyPlan(
        schema_version=payload.get("schema_version", "1.0"),
        plan_id=str(payload["plan_id"]),
        scenario_id=str(payload["scenario_id"]),
        dominant_style_cluster_id=str(
            payload.get(
                "dominant_style_cluster_id",
                DEFAULT_DOMINANT_STYLE_CLUSTER,
            )
        ),
        created_at=str(payload.get("created_at", _now_utc())),
        target_slide_count=int(payload["target_slide_count"]),
        target_slides=tuple(slides),
        library_index_sha256=str(payload["library_index_sha256"]),
        binding_profile_authority=binding_profile_authority,
        query_bundle_path=query_bundle_path,
        query_bundle_sha256=query_bundle_sha256,
        authority=authority,
    )


def _locked_asset_content_type(path: Path) -> str:
    content_type = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
        ".tif": "image/tiff",
        ".tiff": "image/tiff",
        ".emf": "image/x-emf",
        ".wmf": "image/x-wmf",
    }.get(path.suffix.lower())
    if content_type is None:
        raise PhysicalAssemblyError(
            f"locked asset format is unsupported for physical assembly: {path.suffix}"
        )
    return content_type


def _load_locked_asset_manifest(
    path: Path,
    expected_sha256: str,
    *,
    project_root: str | os.PathLike[str] | None = None,
) -> dict[str, LockedAsset]:
    actual_manifest_sha = _sha256_file(path)
    if actual_manifest_sha != expected_sha256:
        raise PhysicalAssemblyError("ASSET_MANIFEST_FINGERPRINT_MISMATCH")
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise PhysicalAssemblyError(f"cannot load locked asset manifest: {exc}") from exc
    if not isinstance(payload, Mapping) or set(payload) != {"schema_version", "bindings"}:
        raise PhysicalAssemblyError("locked asset manifest root is invalid")
    if payload.get("schema_version") != "1.0" or not isinstance(
        payload.get("bindings"), Mapping
    ):
        raise PhysicalAssemblyError("locked asset manifest schema is invalid")
    result: dict[str, LockedAsset] = {}
    for asset_ref, raw in payload["bindings"].items():
        if not isinstance(asset_ref, str) or not asset_ref:
            raise PhysicalAssemblyError("locked asset reference is invalid")
        if not isinstance(raw, Mapping) or set(raw) != {"path", "sha256", "record"}:
            raise PhysicalAssemblyError(
                f"locked asset binding must contain path, sha256, record: {asset_ref}"
            )
        raw_path = raw.get("path")
        expected_asset_sha = raw.get("sha256")
        record = raw.get("record")
        if not isinstance(raw_path, str) or not raw_path:
            raise PhysicalAssemblyError(f"locked asset path is invalid: {asset_ref}")
        if not isinstance(expected_asset_sha, str) or not re.fullmatch(
            r"[0-9a-f]{64}", expected_asset_sha
        ):
            raise PhysicalAssemblyError(f"locked asset sha256 is invalid: {asset_ref}")
        if not isinstance(record, Mapping):
            raise PhysicalAssemblyError(f"locked asset record is invalid: {asset_ref}")
        required_record = {"id", "kind", "quality", "source", "license", "retrieved_at"}
        if not required_record.issubset(record):
            raise PhysicalAssemblyError(f"locked asset provenance is incomplete: {asset_ref}")
        if record.get("id") != asset_ref:
            raise PhysicalAssemblyError(f"locked asset record id mismatch: {asset_ref}")
        if not isinstance(record.get("kind"), str) or not record.get("kind"):
            raise PhysicalAssemblyError(f"locked asset kind is invalid: {asset_ref}")
        quality = record.get("quality")
        if (
            not isinstance(quality, (int, float))
            or isinstance(quality, bool)
            or not 0 <= float(quality) <= 1
        ):
            raise PhysicalAssemblyError(f"locked asset quality is invalid: {asset_ref}")
        for provenance_field in ("source", "license", "retrieved_at"):
            if not isinstance(record.get(provenance_field), str) or not record.get(
                provenance_field
            ):
                raise PhysicalAssemblyError(
                    f"locked asset provenance field is invalid: {asset_ref}.{provenance_field}"
                )
        if project_root is not None:
            asset_path = resolve_project_file(
                raw_path,
                project_root,
                label=f"LOCKED_ASSET_{asset_ref}",
            )
        else:
            # Legacy in-process callers retain the pre-v6.1 resolver. The
            # production renderer always supplies project_root and therefore
            # always receives the strict path/symlink boundary above.
            asset_path = Path(raw_path).expanduser()
            if not asset_path.is_absolute():
                asset_path = path.parent / asset_path
            asset_path = asset_path.resolve(strict=False)
            if not asset_path.is_file():
                raise PhysicalAssemblyError(
                    f"locked asset file is missing: {asset_ref}"
                )
        if _sha256_file(asset_path) != expected_asset_sha:
            raise PhysicalAssemblyError(f"locked asset fingerprint mismatch: {asset_ref}")
        _locked_asset_content_type(asset_path)
        width_px = record.get("width_px")
        height_px = record.get("height_px")
        if width_px is not None and (not isinstance(width_px, int) or width_px <= 0):
            raise PhysicalAssemblyError(f"locked asset width is invalid: {asset_ref}")
        if height_px is not None and (not isinstance(height_px, int) or height_px <= 0):
            raise PhysicalAssemblyError(f"locked asset height is invalid: {asset_ref}")
        if width_px is None or height_px is None:
            raise PhysicalAssemblyError(
                f"locked raster asset dimensions are required: {asset_ref}"
            )
        result[asset_ref] = LockedAsset(
            asset_ref=asset_ref,
            path=asset_path,
            sha256=expected_asset_sha,
            kind=str(record["kind"]),
            width_px=width_px,
            height_px=height_px,
        )
    return result


def _load_locked_connective_copy(
    path: Path,
    expected_sha256: str,
) -> dict[str, str]:
    """Load the exact non-factual copy allowlist as ``text -> id``."""

    if _sha256_file(path) != expected_sha256:
        raise PhysicalAssemblyError("CONNECTIVE_COPY_FINGERPRINT_MISMATCH")
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise PhysicalAssemblyError(f"cannot load locked connective copy: {exc}") from exc
    if not isinstance(payload, Mapping) or set(payload) != {"schema_version", "entries"}:
        raise PhysicalAssemblyError("locked connective copy root is invalid")
    if payload.get("schema_version") != "1.0" or not isinstance(
        payload.get("entries"), list
    ):
        raise PhysicalAssemblyError("locked connective copy schema is invalid")
    by_text: dict[str, str] = {}
    ids: set[str] = set()
    for entry in payload["entries"]:
        if not isinstance(entry, Mapping) or set(entry) != {"id", "text"}:
            raise PhysicalAssemblyError("locked connective copy entry is invalid")
        connective_id = entry.get("id")
        text = entry.get("text")
        if not isinstance(connective_id, str) or not connective_id:
            raise PhysicalAssemblyError("locked connective copy id is invalid")
        if not isinstance(text, str):
            raise PhysicalAssemblyError(
                f"locked connective copy text is invalid: {connective_id}"
            )
        if connective_id in ids:
            raise PhysicalAssemblyError(
                f"locked connective copy id is duplicated: {connective_id}"
            )
        if text in by_text:
            raise PhysicalAssemblyError(
                f"locked connective copy text is ambiguous: {connective_id}"
            )
        ids.add(connective_id)
        by_text[text] = connective_id
    return by_text


def _fact_render_value(fact: Fact, field_name: str) -> str:
    if field_name == "text":
        return fact.text
    if fact.value is None:
        raise PhysicalAssemblyError(f"fact has no renderable value: {fact.id}")
    value = str(fact.value)
    if field_name == "value_unit":
        return value + (fact.unit or "")
    return value


def _without_layout_whitespace(value: str) -> str:
    return "".join(value.split())


def _fact_registered_renderings(fact: Fact) -> tuple[str, ...]:
    """Return only renderings governed by immutable FactStore fields."""

    values: list[str] = [fact.text, *fact.allowed_renderings]
    if fact.value is not None:
        scalar = str(fact.value)
        values.append(scalar)
        values.append(scalar + (fact.unit or ""))
    return tuple(dict.fromkeys(values))


def _allowed_fact_rendering(replacement: str, fact: Fact) -> str | None:
    """Classify a complete rendering of one locked fact.

    Character fragments are deliberately excluded here.  They are authorized
    only by the locked, whole-group contract in
    ``_validate_fragment_group_bindings``.
    """

    replacement_compact = _without_layout_whitespace(replacement)
    registered_renderings = _fact_registered_renderings(fact)
    for registered in registered_renderings:
        if replacement == registered:
            return "exact"
        if replacement_compact == _without_layout_whitespace(registered):
            return "whitespace"
    return None


def _validate_text_binding(
    spec: TextBindingSpec,
    *,
    source_text: str,
    facts: Mapping[str, Fact],
    connective_copy: Mapping[str, str],
    slot_id: str,
) -> tuple[str, str]:
    if spec.mode == "auto" and not spec.fact_refs:
        connective_ref = connective_copy.get(spec.replacement)
        if connective_ref is not None:
            return "connective", connective_ref
        raise PhysicalAssemblyError(
            f"unreferenced text is not registered connective copy: {slot_id}"
        )
    if spec.mode in {"clear", "source"}:
        if spec.fact_refs:
            raise PhysicalAssemblyError(
                f"{spec.mode} binding cannot claim fact authority: {slot_id}"
            )
        expected = "" if spec.mode == "clear" else source_text
        if spec.replacement != expected:
            raise PhysicalAssemblyError(
                f"{spec.mode} binding diverges from its governed value: {slot_id}"
            )
        connective_ref = connective_copy.get(spec.replacement)
        if connective_ref is None:
            raise PhysicalAssemblyError(
                f"unreferenced text is not registered connective copy: {slot_id}"
            )
        return "connective", connective_ref
    if not spec.fact_refs:
        raise PhysicalAssemblyError(f"text replacement has no fact refs: {slot_id}")
    missing = [ref for ref in spec.fact_refs if ref not in facts]
    if missing:
        raise PhysicalAssemblyError(
            f"text replacement references unknown/inactive facts: {slot_id}:{','.join(missing)}"
        )
    if spec.mode == "auto":
        if len(spec.fact_refs) > 8:
            raise PhysicalAssemblyError(
                f"automatic fact rendering accepts at most 8 refs: {slot_id}"
            )
        if len(spec.fact_refs) == 1:
            allowed_mode = _allowed_fact_rendering(
                spec.replacement,
                facts[spec.fact_refs[0]],
            )
            if allowed_mode is not None:
                return allowed_mode, ""
            raise PhysicalAssemblyError(
                f"replacement is not an allowed rendering of locked facts: {slot_id}"
            )
        choices = [
            _fact_registered_renderings(facts[ref])
            for ref in spec.fact_refs
        ]
        combinations = list(product(*choices))
        replacement_compact = _without_layout_whitespace(spec.replacement)
        for values in combinations:
            for separator in TEXT_RENDER_SEPARATORS:
                canonical = separator.join(values)
                if spec.replacement == canonical:
                    return "join", ""
                if replacement_compact == _without_layout_whitespace(canonical):
                    return "join", ""
        raise PhysicalAssemblyError(
            f"replacement is not an allowed rendering of locked facts: {slot_id}"
        )
    rendered = [_fact_render_value(facts[ref], spec.field) for ref in spec.fact_refs]
    if spec.mode in {"exact", "whitespace", "slice"} and len(rendered) != 1:
        raise PhysicalAssemblyError(
            f"{spec.mode} binding requires exactly one fact: {slot_id}"
        )
    canonical = rendered[0] if len(rendered) == 1 else spec.separator.join(rendered)
    if spec.mode == "slice":
        if spec.slice_start is None or spec.slice_end is None:
            raise PhysicalAssemblyError(f"slice binding requires bounds: {slot_id}")
        if spec.slice_end <= spec.slice_start or spec.slice_end > len(canonical):
            raise PhysicalAssemblyError(f"slice binding bounds are invalid: {slot_id}")
        canonical = canonical[spec.slice_start:spec.slice_end]
    elif spec.mode == "join":
        canonical = spec.separator.join(rendered)
    elif spec.slice_start is not None or spec.slice_end is not None:
        raise PhysicalAssemblyError(f"slice bounds require slice mode: {slot_id}")
    if spec.mode in {"whitespace", "slice"}:
        matches = _without_layout_whitespace(spec.replacement) == _without_layout_whitespace(
            canonical
        )
    else:
        matches = spec.replacement == canonical
    if not matches:
        raise PhysicalAssemblyError(
            f"replacement is not an allowed rendering of locked facts: {slot_id}"
        )
    return spec.mode, ""


_NUMERIC_LITERAL_RE = re.compile(
    r"(?<![A-Za-z0-9])[-+]?\d[\d,]*(?:\.\d+)?(?:[Ee][-+]?\d+)?(?:[%％])?"
)
_TABLE_CONTENT_LOCATOR_RE = re.compile(
    r"^graphicFrame\[id=(\d+)\]/table\[(\d+)\]/row\[(\d+)\]/cell\[(\d+)\]$"
)
_WORKBOOK_CONTENT_LOCATOR_RE = re.compile(
    r"^chartFrame\[id=(\d+)\]/(xl/worksheets/[^!]+\.xml)!([^!]+)$"
)
_FRAME_PREFIX_RE = re.compile(r"^(?:graphicFrame|chartFrame)\[id=(\d+)\]")

# These are semantic aliases for generic chart/table sample labels, not
# customer facts.  The destination literal still has to exist in the locked
# connective-copy authority (or, for the year token, in the locked fact store).
_GOVERNED_CONNECTIVE_ALIASES = {
    "销售额": "占比",
    "明年预算": "年初预算",
    "今年决算": "本年决算",
    "去年决算": "上年决算",
    "去年": "上年",
    "列1": "项目",
    "列2": "金额",
}


def _normalise_embedded_text(value: str) -> str:
    return _without_layout_whitespace(value).replace("％", "%")


def _decimal_literal(value: str) -> tuple[Decimal, bool] | None:
    candidate = _normalise_embedded_text(value).replace(",", "")
    is_percent = candidate.endswith("%")
    if is_percent:
        candidate = candidate[:-1]
    if not re.fullmatch(r"[-+]?\d+(?:\.\d+)?(?:[Ee][-+]?\d+)?", candidate):
        return None
    try:
        return Decimal(candidate), is_percent
    except InvalidOperation:
        return None


def _decimal_close(left: Decimal, right: Decimal) -> bool:
    return abs(left - right) <= Decimal("0.000000000001")


def _numeric_rendering_matches(
    source: str,
    candidate: str,
    *,
    allow_percent_scale: bool = True,
) -> bool:
    source_number = _decimal_literal(source)
    candidate_number = _decimal_literal(candidate)
    if source_number is None or candidate_number is None:
        return False
    source_value, source_percent = source_number
    candidate_value, candidate_percent = candidate_number
    if source_percent != candidate_percent:
        if not allow_percent_scale:
            return False
        if source_percent:
            candidate_value *= Decimal(100)
        else:
            source_value *= Decimal(100)
    if _decimal_close(source_value, candidate_value):
        return True
    # Cached chart data is often rounded to the precision displayed by the
    # source token (for example 45063 for an authoritative 45063.1).  Only a
    # half-unit rounding window at that explicit precision is accepted.
    compact_source = _normalise_embedded_text(source).replace(",", "").rstrip("%")
    decimal_places = len(compact_source.split(".", 1)[1]) if "." in compact_source else 0
    quantum = Decimal(1).scaleb(-decimal_places)
    try:
        return source_value == candidate_value.quantize(quantum)
    except InvalidOperation:
        return False


def _fact_numeric_candidates(fact: Fact, *, scalar_only: bool) -> tuple[str, ...]:
    values: list[str] = []
    if fact.value is not None and not isinstance(fact.value, bool):
        scalar = str(fact.value)
        values.extend((scalar, scalar + (fact.unit or "")))
    if not scalar_only:
        for rendering in _fact_registered_renderings(fact):
            values.extend(match.group(0) for match in _NUMERIC_LITERAL_RE.finditer(rendering))
    return tuple(dict.fromkeys(values))


def _auto_authorize_governed_value(
    source_text: str,
    *,
    facts: Mapping[str, Fact],
    connective_copy: Mapping[str, str],
) -> tuple[str, tuple[str, ...], str, str]:
    """Resolve one retained source literal to immutable authority.

    Returns ``(replacement, fact_refs, connective_ref, mode)``.  Ambiguous
    matches deliberately fail; source content is never accepted merely because
    it came from a certified visual template.
    """

    connective_ref = connective_copy.get(source_text)
    if connective_ref is not None:
        return source_text, (), connective_ref, "source-connective"

    exact_matches = {
        fact.id
        for fact in facts.values()
        if any(
            _normalise_embedded_text(source_text)
            == _normalise_embedded_text(rendering)
            for rendering in _fact_registered_renderings(fact)
        )
    }
    if len(exact_matches) == 1:
        return source_text, tuple(exact_matches), "", "source-fact"
    if len(exact_matches) > 1:
        raise PhysicalAssemblyError(
            f"GOVERNED_CONTENT_AUTHORITY_AMBIGUOUS: {source_text}"
        )

    # A one- or two-digit isolated literal in an embedded governed-content
    # surface is a certified-template decoration (commonly a page marker), not
    # client data. Handle it before numeric matching: otherwise every genuine
    # customer value containing that digit makes the marker ambiguous. It must
    # not survive into a client deliverable. Longer, decimal, percent, signed
    # and all fact-matched numeric source content continues through the
    # fail-closed authority checks below.
    if re.fullmatch(r"[0-9]{1,2}", source_text):
        return "", (), "", "source-decoration-numeric"

    if _decimal_literal(source_text) is not None:
        for scalar_only, mode in (
            (True, "source-numeric-scalar"),
            (False, "source-numeric-rendering"),
        ):
            numeric_matches = {
                fact.id
                for fact in facts.values()
                if any(
                    _numeric_rendering_matches(
                        source_text,
                        candidate,
                        allow_percent_scale=(
                            fact.unit == "%"
                            or any(
                                "%" in rendering or "％" in rendering
                                for rendering in _fact_registered_renderings(fact)
                            )
                        ),
                    )
                    for candidate in _fact_numeric_candidates(
                        fact,
                        scalar_only=scalar_only,
                    )
                )
            }
            if len(numeric_matches) == 1:
                return source_text, tuple(numeric_matches), "", mode
            if len(numeric_matches) > 1:
                raise PhysicalAssemblyError(
                    f"GOVERNED_CONTENT_NUMERIC_AUTHORITY_AMBIGUOUS: {source_text}"
                )


    alias = _GOVERNED_CONNECTIVE_ALIASES.get(source_text)
    if alias is not None:
        alias_ref = connective_copy.get(alias)
        if alias_ref is None:
            raise PhysicalAssemblyError(
                f"GOVERNED_CONTENT_ALIAS_NOT_LOCKED: {source_text}->{alias}"
            )
        return alias, (), alias_ref, "normalized-connective"

    if source_text in {"202x", "202X", "202x年", "202X年"}:
        year_matches = [
            fact
            for fact in facts.values()
            if fact.id == "report-year" and fact.value is not None
        ]
        if len(year_matches) != 1:
            raise PhysicalAssemblyError("GOVERNED_CONTENT_REPORT_YEAR_MISSING")
        year = year_matches[0]
        suffix = "年" if source_text.lower().endswith("x年") else ""
        return f"{year.value}{suffix}", (year.id,), "", "normalized-year"

    raise PhysicalAssemblyError(
        "GOVERNED_CONTENT_UNBOUND_SOURCE_LITERAL: "
        + _sha256_bytes(source_text.encode("utf-8"))
    )


def _native_embedded_rendering(source_text: str, replacement: str) -> str:
    """Render a governed value in the source chart/workbook numeric domain."""

    source_number = _decimal_literal(source_text)
    if source_number is None:
        return replacement
    replacement_literals = [
        match.group(0) for match in _NUMERIC_LITERAL_RE.finditer(replacement)
    ]
    if len(replacement_literals) != 1:
        raise PhysicalAssemblyError(
            f"GOVERNED_CONTENT_NUMERIC_REPLACEMENT_INVALID: {replacement}"
        )
    replacement_number = _decimal_literal(replacement_literals[0])
    if replacement_number is None:
        raise PhysicalAssemblyError(
            f"GOVERNED_CONTENT_NUMERIC_REPLACEMENT_INVALID: {replacement}"
        )
    _, source_percent = source_number
    value, replacement_percent = replacement_number
    if replacement_percent and not source_percent:
        value /= Decimal(100)
    elif source_percent and not replacement_percent:
        value *= Decimal(100)
    rendered = format(value, "f")
    if "." in rendered:
        rendered = rendered.rstrip("0").rstrip(".")
    return rendered + ("%" if source_percent else "")


def _element_for_local_locator(root: ET.Element, locator: str) -> ET.Element:
    segments = re.findall(r"([^/\[]+)\[(\d+)\]", locator)
    if not segments:
        raise PhysicalAssemblyError(f"GOVERNED_CONTENT_LOCATOR_INVALID: {locator}")
    root_name, root_ordinal = segments[0]
    if root_ordinal != "1" or root.tag.rsplit("}", 1)[-1] != root_name:
        raise PhysicalAssemblyError(f"GOVERNED_CONTENT_LOCATOR_ROOT_MISMATCH: {locator}")
    current = root
    for local_name, raw_ordinal in segments[1:]:
        children = [
            child
            for child in list(current)
            if child.tag.rsplit("}", 1)[-1] == local_name
        ]
        ordinal = int(raw_ordinal)
        if ordinal < 1 or ordinal > len(children):
            raise PhysicalAssemblyError(
                f"GOVERNED_CONTENT_LOCATOR_NOT_FOUND: {locator}"
            )
        current = children[ordinal - 1]
    return current


def _table_cell_for_locator(root: ET.Element, locator: str) -> ET.Element:
    match = _TABLE_CONTENT_LOCATOR_RE.fullmatch(locator)
    if match is None:
        raise PhysicalAssemblyError(f"GOVERNED_TABLE_LOCATOR_INVALID: {locator}")
    shape_id, table_ordinal, row_ordinal, column_ordinal = map(int, match.groups())
    frames = []
    for frame in root.iter(f"{{{PML_NS}}}graphicFrame"):
        marker = frame.find(f".//{{{PML_NS}}}cNvPr")
        if marker is not None and marker.attrib.get("id") == str(shape_id):
            frames.append(frame)
    if len(frames) != 1:
        raise PhysicalAssemblyError(f"GOVERNED_TABLE_FRAME_NOT_FOUND: {locator}")
    tables = list(frames[0].iter(f"{{{DML_NS}}}tbl"))
    if table_ordinal > len(tables):
        raise PhysicalAssemblyError(f"GOVERNED_TABLE_NOT_FOUND: {locator}")
    rows = [node for node in list(tables[table_ordinal - 1]) if node.tag.rsplit("}", 1)[-1] == "tr"]
    if row_ordinal > len(rows):
        raise PhysicalAssemblyError(f"GOVERNED_TABLE_ROW_NOT_FOUND: {locator}")
    cells = [node for node in list(rows[row_ordinal - 1]) if node.tag.rsplit("}", 1)[-1] == "tc"]
    if column_ordinal > len(cells):
        raise PhysicalAssemblyError(f"GOVERNED_TABLE_CELL_NOT_FOUND: {locator}")
    return cells[column_ordinal - 1]


def _xml_slot_text_and_set(
    root: ET.Element,
    record: Mapping[str, Any],
    replacement: str | None,
) -> str:
    kind = str(record["kind"])
    locator = str(record["locator"])
    if kind == "table-cell":
        owner = _table_cell_for_locator(root, locator)
        text_nodes = [node for node in owner.iter() if node.tag.rsplit("}", 1)[-1] == "t"]
        source = "".join(node.text or "" for node in text_nodes).strip()
        if replacement is not None:
            if not text_nodes:
                raise PhysicalAssemblyError(f"GOVERNED_TABLE_TEXT_NODE_MISSING: {locator}")
            text_nodes[0].text = replacement
            for node in text_nodes[1:]:
                node.text = ""
        return source
    generic_locator = _FRAME_PREFIX_RE.sub("", locator, count=1)
    owner = _element_for_local_locator(root, generic_locator)
    source = (owner.text or "").strip()
    if replacement is not None:
        owner.text = replacement
    return source


def _xlsx_shared_strings(parts: Mapping[str, bytes]) -> tuple[str, ...]:
    raw = parts.get("xl/sharedStrings.xml")
    if raw is None:
        return ()
    try:
        root = ET.fromstring(raw)
    except ET.ParseError as exc:
        raise PhysicalAssemblyError("GOVERNED_WORKBOOK_SHARED_STRINGS_INVALID") from exc
    return tuple(
        "".join(
            child.text or ""
            for child in item.iter()
            if child.tag.rsplit("}", 1)[-1] == "t"
        ).strip()
        for item in root.iter()
        if item.tag.rsplit("}", 1)[-1] == "si"
    )


def _xlsx_cell_text(
    cell: ET.Element,
    shared_strings: Sequence[str],
) -> str:
    cell_type = cell.attrib.get("t", "")
    if cell_type == "inlineStr":
        return "".join(
            child.text or ""
            for child in cell.iter()
            if child.tag.rsplit("}", 1)[-1] == "t"
        ).strip()
    value = next(
        (child for child in list(cell) if child.tag.rsplit("}", 1)[-1] == "v"),
        None,
    )
    raw = (value.text or "").strip() if value is not None else ""
    if cell_type == "s":
        if not raw.isdigit() or int(raw) >= len(shared_strings):
            raise PhysicalAssemblyError("GOVERNED_WORKBOOK_SHARED_STRING_INDEX_INVALID")
        return shared_strings[int(raw)]
    return raw


def _set_xlsx_cell_text(cell: ET.Element, replacement: str) -> None:
    for child in list(cell):
        cell.remove(child)
    numeric = _decimal_literal(replacement)
    namespace = cell.tag.split("}", 1)[0] + "}" if "}" in cell.tag else ""
    if numeric is not None and not numeric[1]:
        cell.attrib.pop("t", None)
        value = ET.SubElement(cell, f"{namespace}v")
        value.text = replacement
    else:
        cell.attrib["t"] = "inlineStr"
        inline = ET.SubElement(cell, f"{namespace}is")
        text = ET.SubElement(inline, f"{namespace}t")
        text.text = replacement


def _mutate_governed_workbook(
    source_bytes: bytes,
    replacements: Sequence[tuple[Mapping[str, Any], str]],
) -> bytes:
    try:
        return mutate_governed_xlsx(source_bytes, replacements)
    except WorkbookSecurityError as exc:
        raise PhysicalAssemblyError(f"GOVERNED_WORKBOOK_SECURITY: {exc}") from exc


def _sanitize_layout_master_fields(
    source_part: str,
    source_bytes: bytes,
    expected_fields: Sequence[Mapping[str, Any]],
) -> bytes:
    try:
        root = ET.fromstring(source_bytes)
    except ET.ParseError as exc:
        raise PhysicalAssemblyError("LAYOUT_MASTER_XML_INVALID") from exc
    actual_fields = [
        node for node in root.iter() if node.tag.rsplit("}", 1)[-1] == "fld"
    ]
    if len(actual_fields) != len(expected_fields):
        raise PhysicalAssemblyError(
            f"LAYOUT_MASTER_FIELD_COUNT_DRIFT: {source_part}"
        )
    changed = False
    seen_locators: set[str] = set()
    for record in expected_fields:
        if record.get("source_part") != source_part:
            raise PhysicalAssemblyError(
                f"LAYOUT_MASTER_FIELD_PART_DRIFT: {source_part}"
            )
        locator = record.get("locator")
        if not isinstance(locator, str) or locator in seen_locators:
            raise PhysicalAssemblyError(
                f"LAYOUT_MASTER_FIELD_LOCATOR_INVALID: {source_part}"
            )
        seen_locators.add(locator)
        field_node = _element_for_local_locator(root, locator)
        if field_node.tag.rsplit("}", 1)[-1] != "fld":
            raise PhysicalAssemblyError(
                f"LAYOUT_MASTER_FIELD_LOCATOR_DRIFT: {source_part}"
            )
        field_type = field_node.attrib.get("type", "unknown")
        if field_type not in {"datetimeFigureOut", "slidenum"}:
            raise PhysicalAssemblyError(
                f"LAYOUT_MASTER_FIELD_TYPE_UNSUPPORTED: {field_type}"
            )
        if field_type != record.get("field_type"):
            raise PhysicalAssemblyError(
                f"LAYOUT_MASTER_FIELD_TYPE_DRIFT: {source_part}"
            )
        expected_id = "field_" + _sha256_bytes(
            f"{source_part}\0{locator}".encode("utf-8")
        )[:24]
        if record.get("field_id") != expected_id:
            raise PhysicalAssemblyError(
                f"LAYOUT_MASTER_FIELD_ID_DRIFT: {source_part}"
            )
        cached = "".join(
            node.text or ""
            for node in field_node.iter()
            if node.tag.rsplit("}", 1)[-1] == "t"
        ).strip()
        if _sha256_bytes(cached.encode("utf-8")) != record.get("source_text_sha256"):
            raise PhysicalAssemblyError(
                f"LAYOUT_MASTER_FIELD_CACHE_DRIFT: {source_part}"
            )
        for node in field_node.iter():
            if node.tag.rsplit("}", 1)[-1] == "t" and (node.text or ""):
                node.text = ""
                changed = True
    if not changed:
        return source_bytes
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _prepare_governed_content_replacements(
    slide: AssemblyTargetSlide,
    graph: _SourceGraph,
    fact_store: FactStore | None,
    connective_copy: Mapping[str, str],
) -> tuple[dict[str, bytes], list[BindingEvidence]]:
    inventory = slide.page_template.governed_content_inventory
    if not isinstance(inventory, Mapping) or not bool(inventory.get("complete")):
        raise PhysicalAssemblyError(
            f"GOVERNED_CONTENT_INVENTORY_INCOMPLETE: {slide.page_template.page_id}"
        )
    records = tuple(
        record
        for record in inventory.get("slots", ())
        if isinstance(record, Mapping)
    )
    source_parts = {graph.root_slide_name: graph.slide_xml, **graph.extra_parts}
    reachable_workbooks = tuple(
        source_part
        for source_part in sorted(source_parts)
        if source_part.lower().endswith((".xlsx", ".xlsm"))
    )
    forbidden_workbooks = tuple(
        source_part
        for source_part in reachable_workbooks
        if source_part.lower().endswith(".xlsm")
    )
    if forbidden_workbooks:
        raise PhysicalAssemblyError(
            "GOVERNED_WORKBOOK_XLSM_FORBIDDEN: " + ",".join(forbidden_workbooks)
        )
    if records and fact_store is None:
        raise PhysicalAssemblyError("GOVERNED_CONTENT_LOCKED_AUTHORITY_REQUIRED")
    facts = (
        {fact.id: fact for fact in fact_store.active_facts()}
        if fact_store is not None
        else {}
    )
    replacements_by_part: dict[str, list[tuple[Mapping[str, Any], str]]] = {}
    evidence: list[BindingEvidence] = []
    peer_replacements: dict[str, tuple[str, tuple[str, ...], str]] = {}
    for record in records:
        source_text = str(record.get("source_text", ""))
        if _sha256_bytes(source_text.encode("utf-8")) != record.get("source_text_sha256"):
            raise PhysicalAssemblyError(
                f"GOVERNED_CONTENT_MANIFEST_DRIFT: {record.get('slot_id', '')}"
            )
        slot_id = str(record.get("slot_id", ""))
        peer_group_id = record.get("peer_group_id")
        spec = slide.governed_content_binding_specs.get(slot_id)
        if spec is None and isinstance(peer_group_id, str):
            spec = slide.governed_content_binding_specs.get(peer_group_id)
        if spec is not None:
            mode, connective_ref = _validate_text_binding(
                spec,
                source_text=source_text,
                facts=facts,
                connective_copy=connective_copy,
                slot_id=f"{slide.ordinal}:{slot_id}",
            )
            replacement = _native_embedded_rendering(source_text, spec.replacement)
            fact_refs = spec.fact_refs
            mode = f"explicit-{mode}"
        else:
            replacement, fact_refs, connective_ref, mode = (
                _auto_authorize_governed_value(
                    source_text,
                    facts=facts,
                    connective_copy=connective_copy,
                )
            )
        if isinstance(peer_group_id, str):
            canonical = (
                _normalise_embedded_text(replacement),
                tuple(fact_refs),
                connective_ref,
            )
            previous = peer_replacements.setdefault(peer_group_id, canonical)
            if previous != canonical:
                raise PhysicalAssemblyError(
                    f"GOVERNED_CONTENT_PEER_DRIFT: {peer_group_id}"
                )
        source_part = str(record.get("source_part", ""))
        replacements_by_part.setdefault(source_part, []).append((record, replacement))
        shape_match = _FRAME_PREFIX_RE.match(str(record.get("locator", "")))
        shape_id = int(shape_match.group(1)) if shape_match is not None else 0
        evidence.append(
            BindingEvidence(
                ordinal=slide.ordinal,
                page_id=slide.page_template.page_id,
                slot_id=slot_id,
                shape_id=shape_id,
                binding_kind=(
                    "source-decoration"
                    if mode == "source-decoration-numeric"
                    else "embedded"
                ),
                mode=mode,
                source_text=source_text,
                source_sha256=_sha256_bytes(source_text.encode("utf-8")),
                replacement_sha256=_sha256_bytes(replacement.encode("utf-8")),
                fact_refs=tuple(fact_refs),
                asset_refs=(),
                connective_ref=connective_ref,
                char_used=len(replacement),
                char_limit=max(len(source_text), len(replacement)),
                item_used=1,
                item_limit=1,
                image_used=0,
                image_limit=0,
                status="pass",
                fit_policy=(spec.fit_policy if spec is not None else "preserve"),
            )
        )

    for source_part in reachable_workbooks:
        replacements_by_part.setdefault(source_part, [])
    mutated: dict[str, bytes] = {}
    for source_part, replacements in replacements_by_part.items():
        source_bytes = source_parts.get(source_part)
        if source_bytes is None:
            raise PhysicalAssemblyError(
                f"GOVERNED_CONTENT_PART_MISSING: {source_part}"
            )
        if source_part.lower().endswith(".xlsm"):
            raise PhysicalAssemblyError(
                f"GOVERNED_WORKBOOK_XLSM_FORBIDDEN: {source_part}"
            )
        if source_part.lower().endswith(".xlsx"):
            mutated[source_part] = _mutate_governed_workbook(
                source_bytes,
                replacements,
            )
            continue
        try:
            root = ET.fromstring(source_bytes)
        except ET.ParseError as exc:
            raise PhysicalAssemblyError(
                f"GOVERNED_CONTENT_XML_INVALID: {source_part}"
            ) from exc
        for record, replacement in replacements:
            actual_source = _xml_slot_text_and_set(root, record, replacement)
            if _sha256_bytes(actual_source.encode("utf-8")) != record["source_text_sha256"]:
                raise PhysicalAssemblyError(
                    f"GOVERNED_CONTENT_SOURCE_DRIFT: {record['slot_id']}"
                )
        mutated[source_part] = ET.tostring(
            root,
            encoding="utf-8",
            xml_declaration=True,
        )
    if len(evidence) != int(inventory.get("content_slot_count", -1)):
        raise PhysicalAssemblyError(
            f"GOVERNED_CONTENT_BINDING_COVERAGE: {slide.page_template.page_id}"
        )
    return mutated, evidence


def _slot_metadata(template: PageTemplate) -> dict[str, Mapping[str, Any]]:
    result: dict[str, Mapping[str, Any]] = {}
    slots = template.slot_graph.get("slots", ())
    if isinstance(slots, Sequence) and not isinstance(slots, (str, bytes)):
        for slot in slots:
            if isinstance(slot, Mapping) and isinstance(slot.get("slot_id"), str):
                result[str(slot["slot_id"])] = slot
    return result


def _item_count(value: str) -> int:
    return sum(1 for item in value.splitlines() if item.strip()) if value else 0


def _semantic_character_count(value: str) -> int:
    """Count copy characters using the certified slot-graph convention."""

    return len("".join(value.split()))


def _validate_fragment_group_bindings(
    slide: AssemblyTargetSlide,
    *,
    contracts: Sequence[FragmentGroupContract],
    facts: Mapping[str, Fact],
    connective_copy: Mapping[str, str],
) -> dict[str, tuple[str, str]]:
    """Authorize character slots only as one complete locked fact rendering.

    Returns ``slot_id -> (evidence_mode, connective_ref)``.  Non-empty slots
    must be single characters backed by the same sole fact reference.  Empty
    remainder slots are accepted only through the exact locked
    ``connective-clear`` entry.
    """

    authorizations: dict[str, tuple[str, str]] = {}
    seen_group_ids: set[str] = set()
    for contract in contracts:
        if (
            contract.ordinal != slide.ordinal
            or contract.page_id != slide.page_template.page_id
            or contract.group_id in seen_group_ids
            or len(contract.ordered_slot_ids) < 2
            or len(contract.ordered_slot_ids) != len(set(contract.ordered_slot_ids))
        ):
            raise PhysicalAssemblyError(
                f"FRAGMENT_GROUP_CONTRACT_INVALID: "
                f"{slide.ordinal}:{contract.group_id}"
            )
        seen_group_ids.add(contract.group_id)
        nonempty: list[tuple[str, TextBindingSpec]] = []
        for slot_id in contract.ordered_slot_ids:
            if slot_id in authorizations:
                raise PhysicalAssemblyError(
                    f"FRAGMENT_GROUP_SLOT_DUPLICATE: {slide.ordinal}:{slot_id}"
                )
            if slot_id not in slide.bindings:
                raise PhysicalAssemblyError(
                    f"FRAGMENT_GROUP_SLOT_MISSING: {slide.ordinal}:{slot_id}"
                )
            spec = slide.text_binding_specs.get(slot_id)
            if spec is None or spec.replacement != slide.bindings[slot_id]:
                raise PhysicalAssemblyError(
                    f"FRAGMENT_GROUP_BINDING_MISSING: {slide.ordinal}:{slot_id}"
                )
            if spec.replacement:
                if (
                    len(spec.replacement) != 1
                    or spec.replacement.isspace()
                    or spec.mode != "auto"
                    or len(spec.fact_refs) != 1
                ):
                    raise PhysicalAssemblyError(
                        f"FRAGMENT_GROUP_CHARACTER_INVALID: "
                        f"{slide.ordinal}:{slot_id}"
                    )
                nonempty.append((slot_id, spec))
                authorizations[slot_id] = ("character", "")
            else:
                if (
                    spec.mode not in {"auto", "clear"}
                    or spec.fact_refs
                    or connective_copy.get("") != "connective-clear"
                ):
                    raise PhysicalAssemblyError(
                        f"FRAGMENT_GROUP_CLEAR_INVALID: {slide.ordinal}:{slot_id}"
                    )
                authorizations[slot_id] = ("connective", "connective-clear")

        if not nonempty:
            continue
        fact_refs = {spec.fact_refs[0] for _, spec in nonempty}
        if len(fact_refs) != 1:
            raise PhysicalAssemblyError(
                f"FRAGMENT_GROUP_FACT_REF_DRIFT: "
                f"{slide.ordinal}:{contract.group_id}"
            )
        fact_ref = next(iter(fact_refs))
        fact = facts.get(fact_ref)
        if fact is None:
            raise PhysicalAssemblyError(
                f"FRAGMENT_GROUP_FACT_REF_UNKNOWN: "
                f"{slide.ordinal}:{contract.group_id}:{fact_ref}"
            )
        assembled = "".join(spec.replacement for _, spec in nonempty)
        matches = any(
            assembled == rendering
            or _without_layout_whitespace(assembled)
            == _without_layout_whitespace(rendering)
            for rendering in _fact_registered_renderings(fact)
        )
        if not matches:
            raise PhysicalAssemblyError(
                f"FRAGMENT_GROUP_RENDERING_MISMATCH: "
                f"{slide.ordinal}:{contract.group_id}"
            )
    return authorizations


def _build_text_binding_evidence(
    plan: AssemblyPlan,
    fact_store: FactStore | None,
    connective_copy: Mapping[str, str] | None = None,
    *,
    require_locked_authority: bool,
    source_texts: Mapping[tuple[int, str], str] | None = None,
    enforce_required_facts: bool = True,
    fragment_group_contracts: Sequence[FragmentGroupContract] = (),
) -> list[BindingEvidence]:
    facts = (
        {fact.id: fact for fact in fact_store.active_facts()}
        if fact_store is not None
        else {}
    )
    evidence: list[BindingEvidence] = []
    used_fact_refs: set[str] = set()
    contracts_by_slide: dict[tuple[int, str], list[FragmentGroupContract]] = {}
    for contract in fragment_group_contracts:
        contracts_by_slide.setdefault(
            (contract.ordinal, contract.page_id),
            [],
        ).append(contract)
    for slide in plan.target_slides:
        metadata = _slot_metadata(slide.page_template)
        fragment_authorizations = (
            _validate_fragment_group_bindings(
                slide,
                contracts=contracts_by_slide.get(
                    (slide.ordinal, slide.page_template.page_id),
                    (),
                ),
                facts=facts,
                connective_copy=connective_copy or {},
            )
            if require_locked_authority
            else {}
        )
        total_chars = 0
        for slot_id, replacement in slide.bindings.items():
            slot = metadata.get(slot_id)
            clear_alias_spec = slide.text_binding_specs.get(slot_id)
            if (
                require_locked_authority
                and slot is None
                and clear_alias_spec is not None
                and clear_alias_spec.mode == "clear"
                and clear_alias_spec.replacement == ""
                and not clear_alias_spec.fact_refs
            ):
                # A graph-excluded nested alias may only be blanked. It is
                # never a client content surface, so it has no capacity record.
                continue
            if require_locked_authority and slot is None:
                raise PhysicalAssemblyError(
                    f"production template slot lacks capacity evidence: {slide.ordinal}:{slot_id}"
                )
            metadata_source_text = str(slot.get("source_text", "")) if slot else ""
            source_text = (
                source_texts.get((slide.ordinal, slot_id), metadata_source_text)
                if source_texts is not None
                else metadata_source_text
            )
            if (
                require_locked_authority
                and source_texts is not None
                and source_text != metadata_source_text
            ):
                raise PhysicalAssemblyError(
                    f"template slot source text drift: {slide.ordinal}:{slot_id}"
                )
            shape_id = (
                int(slot.get("shape_id", slot_id.removeprefix("shape_") or 0))
                if slot
                else int(slot_id.removeprefix("shape_") or 0)
            )
            char_limit = (
                int(slot.get("max_chars", 0))
                if slot
                else max(len(replacement), 1)
            )
            item_limit = (
                int(slot.get("max_items", 0))
                if slot and slot.get("max_items")
                else max(1, _item_count(source_text))
            )
            char_used = _semantic_character_count(replacement)
            item_used = _item_count(replacement)
            if require_locked_authority and char_used > char_limit:
                raise PhysicalAssemblyError(
                    f"text slot max_chars exceeded: {slide.ordinal}:{slot_id} {char_used}>{char_limit}"
                )
            if require_locked_authority and item_used > item_limit:
                raise PhysicalAssemblyError(
                    f"text slot item capacity exceeded: {slide.ordinal}:{slot_id} {item_used}>{item_limit}"
                )
            total_chars += char_used
            spec = slide.text_binding_specs.get(slot_id)
            if require_locked_authority:
                if spec is None:
                    raise PhysicalAssemblyError(
                        f"production text slot lacks locked binding evidence: {slide.ordinal}:{slot_id}"
                    )
                fragment_authorization = fragment_authorizations.get(slot_id)
                if fragment_authorization is None:
                    mode, connective_ref = _validate_text_binding(
                        spec,
                        source_text=source_text,
                        facts=facts,
                        connective_copy=connective_copy or {},
                        slot_id=f"{slide.ordinal}:{slot_id}",
                    )
                else:
                    mode, connective_ref = fragment_authorization
                fact_refs = spec.fact_refs
                used_fact_refs.update(fact_refs)
            else:
                mode = spec.mode if spec is not None else "legacy"
                fact_refs = spec.fact_refs if spec is not None else ()
                connective_ref = ""
            evidence.append(
                BindingEvidence(
                    ordinal=slide.ordinal,
                    page_id=slide.page_template.page_id,
                    slot_id=slot_id,
                    shape_id=shape_id,
                    binding_kind="text",
                    mode=mode,
                    source_text=source_text,
                    source_sha256=_sha256_bytes(source_text.encode("utf-8")),
                    replacement_sha256=_sha256_bytes(replacement.encode("utf-8")),
                    fact_refs=tuple(fact_refs),
                    asset_refs=(),
                    connective_ref=connective_ref,
                    char_used=char_used,
                    char_limit=char_limit,
                    item_used=item_used,
                    item_limit=item_limit,
                    image_used=0,
                    image_limit=0,
                    status="pass",
                    fit_policy=(
                        spec.fit_policy if spec is not None else "preserve"
                    ),
                )
            )
        page_char_limit = slide.page_template.capacity.get("max_text_chars")
        if require_locked_authority and (
            not isinstance(page_char_limit, int) or page_char_limit < 0
        ):
            raise PhysicalAssemblyError(
                f"production page lacks max_text_chars capacity: {slide.ordinal}"
            )
        if (
            require_locked_authority
            and isinstance(page_char_limit, int)
            and total_chars > page_char_limit
        ):
            raise PhysicalAssemblyError(
                f"page max_text_chars exceeded: {slide.ordinal} {total_chars}>{page_char_limit}"
            )
    if require_locked_authority and fact_store is not None and enforce_required_facts:
        required_fact_refs = {
            fact.id for fact in fact_store.active_facts() if fact.required
        }
        missing_required = sorted(required_fact_refs - used_fact_refs)
        if missing_required:
            raise PhysicalAssemblyError(
                "REQUIRED_FACTS_NOT_BOUND: " + ",".join(missing_required)
            )
    return evidence


def _validate_locked_authority(
    plan: AssemblyPlan,
    *,
    fact_store_path: str | os.PathLike[str] | None,
    fact_store_sha256: str | None,
    asset_manifest_path: str | os.PathLike[str] | None,
    asset_manifest_sha256: str | None,
    connective_copy_path: str | os.PathLike[str] | None,
    connective_copy_sha256: str | None,
    project_root: str | os.PathLike[str] | None,
) -> tuple[FactStore, dict[str, LockedAsset], dict[str, str], AuthorityEvidence]:
    if plan.authority is None:
        raise PhysicalAssemblyError("ASSEMBLY_PLAN_AUTHORITY_REQUIRED")
    supplied = (
        fact_store_path,
        fact_store_sha256,
        asset_manifest_path,
        asset_manifest_sha256,
        connective_copy_path,
        connective_copy_sha256,
    )
    if any(value is None for value in supplied):
        raise PhysicalAssemblyError(
            "LOCKED_AUTHORITY_ARGUMENTS_REQUIRED: fact store/asset manifest/connective copy paths and SHA256"
        )
    fact_path = Path(str(fact_store_path)).expanduser().resolve(strict=False)
    asset_path = Path(str(asset_manifest_path)).expanduser().resolve(strict=False)
    connective_path = Path(str(connective_copy_path)).expanduser().resolve(strict=False)
    if project_root is not None:
        planned_fact_path = resolve_project_file(
            plan.authority.fact_store_path,
            project_root,
            label="AUTHORITY_FACT_STORE",
        )
        planned_asset_path = resolve_project_file(
            plan.authority.asset_manifest_path,
            project_root,
            label="AUTHORITY_ASSET_MANIFEST",
        )
        planned_connective_path = resolve_project_file(
            plan.authority.connective_copy_path,
            project_root,
            label="AUTHORITY_CONNECTIVE_COPY",
        )
    else:
        planned_fact_path = Path(plan.authority.fact_store_path)
        planned_asset_path = Path(plan.authority.asset_manifest_path)
        planned_connective_path = Path(plan.authority.connective_copy_path)
    if fact_path != planned_fact_path:
        raise PhysicalAssemblyError("FACT_STORE_PATH_MISMATCH")
    if asset_path != planned_asset_path:
        raise PhysicalAssemblyError("ASSET_MANIFEST_PATH_MISMATCH")
    if connective_path != planned_connective_path:
        raise PhysicalAssemblyError("CONNECTIVE_COPY_PATH_MISMATCH")
    if fact_store_sha256 != plan.authority.fact_store_sha256:
        raise PhysicalAssemblyError("FACT_STORE_EXPECTED_SHA_MISMATCH")
    if asset_manifest_sha256 != plan.authority.asset_manifest_sha256:
        raise PhysicalAssemblyError("ASSET_MANIFEST_EXPECTED_SHA_MISMATCH")
    if connective_copy_sha256 != plan.authority.connective_copy_sha256:
        raise PhysicalAssemblyError("CONNECTIVE_COPY_EXPECTED_SHA_MISMATCH")
    if not fact_path.is_file() or _sha256_file(fact_path) != fact_store_sha256:
        raise PhysicalAssemblyError("FACT_STORE_FINGERPRINT_MISMATCH")
    try:
        fact_store = load_fact_store(fact_path)
    except (WeakModelValidationError, OSError, json.JSONDecodeError) as exc:
        raise PhysicalAssemblyError(f"locked FactStore is invalid: {exc}") from exc
    assets = _load_locked_asset_manifest(
        asset_path,
        str(asset_manifest_sha256),
        project_root=project_root,
    )
    connective_copy = _load_locked_connective_copy(
        connective_path,
        str(connective_copy_sha256),
    )
    return (
        fact_store,
        assets,
        connective_copy,
        AuthorityEvidence(
            mode="locked",
            fact_store_path=str(fact_path),
            fact_store_sha256=str(fact_store_sha256),
            asset_manifest_path=str(asset_path),
            asset_manifest_sha256=str(asset_manifest_sha256),
            connective_copy_path=str(connective_path),
            connective_copy_sha256=str(connective_copy_sha256),
            status="pass",
        ),
    )


def _prepare_asset_replacements(
    slide: AssemblyTargetSlide,
    graph: _SourceGraph,
    assets: Mapping[str, LockedAsset],
    context: AssemblyImportContext,
) -> tuple[
    dict[str, str],
    dict[str, tuple[int, int, int, int]],
    list[BindingEvidence],
]:
    if not slide.asset_binding_specs:
        return {}, {}, []
    picture_slots = _discover_picture_slots(graph.slide_xml)
    picture_frames = _discover_picture_frames(graph.slide_xml)
    if len(slide.asset_binding_specs) > len(picture_slots):
        raise PhysicalAssemblyError(
            f"page image capacity exceeded: {slide.ordinal} "
            f"{len(slide.asset_binding_specs)}>{len(picture_slots)}"
        )
    slide_rels_path = _rels_path_for_part(graph.root_slide_name)
    slide_rels = graph.rels.get(slide_rels_path)
    if slide_rels is None:
        raise PhysicalAssemblyError(
            f"asset replacement requires slide relationships: {slide.ordinal}"
        )
    rel_entries = {entry["Id"]: entry for entry in _parse_relationships(slide_rels)}
    relationship_overrides: dict[str, str] = {}
    cover_crops: dict[str, tuple[int, int, int, int]] = {}
    evidence: list[BindingEvidence] = []
    for slot_id, spec in slide.asset_binding_specs.items():
        slot = picture_slots.get(slot_id)
        if slot is None:
            raise PhysicalAssemblyError(
                f"asset binding does not target an editable picture: {slide.ordinal}:{slot_id}"
            )
        asset = assets.get(spec.asset_ref)
        if asset is None:
            raise PhysicalAssemblyError(
                f"asset binding references unknown locked asset: {slide.ordinal}:{slot_id}:{spec.asset_ref}"
            )
        shape_id, relationship_id = slot
        frame = picture_frames.get(slot_id)
        if frame is None or asset.width_px is None or asset.height_px is None:
            raise PhysicalAssemblyError(
                f"picture frame/asset dimensions are unavailable: {slide.ordinal}:{slot_id}"
            )
        cover_crops[slot_id] = _cover_crop_values(
            frame_width=frame[0],
            frame_height=frame[1],
            image_width=asset.width_px,
            image_height=asset.height_px,
        )
        relation = rel_entries.get(relationship_id)
        if relation is None or not relation.get("Type", "").endswith("/image"):
            raise PhysicalAssemblyError(
                f"picture relationship is not a replaceable image: {slide.ordinal}:{slot_id}"
            )
        source_part = _resolve_rel_target(
            slide_rels,
            slide_rels_path,
            relation["Target"],
        )
        if source_part is None or source_part not in graph.extra_parts:
            raise PhysicalAssemblyError(
                f"picture source is outside the imported closure: {slide.ordinal}:{slot_id}"
            )
        context.replaced_source_parts_by_slide.setdefault(slide.ordinal, set()).add(
            source_part
        )
        target_part = context.allocate_locked_asset(asset)
        existing = relationship_overrides.get(relationship_id)
        if existing is not None and existing != target_part:
            raise PhysicalAssemblyError(
                f"shared picture relationship cannot bind different assets: {slide.ordinal}:{relationship_id}"
            )
        relationship_overrides[relationship_id] = target_part
        source_bytes = graph.extra_parts[source_part]
        evidence.append(
            BindingEvidence(
                ordinal=slide.ordinal,
                page_id=slide.page_template.page_id,
                slot_id=slot_id,
                shape_id=shape_id,
                binding_kind="asset",
                mode=spec.fit,
                source_text=source_part,
                source_sha256=_sha256_bytes(source_bytes),
                replacement_sha256=asset.sha256,
                fact_refs=(),
                asset_refs=(asset.asset_ref,),
                connective_ref="",
                char_used=0,
                char_limit=0,
                item_used=0,
                item_limit=0,
                image_used=1,
                image_limit=1,
                status="pass",
                fit_policy="preserve",
                relationship_id=relationship_id,
                target_part=target_part,
            )
        )
    return relationship_overrides, cover_crops, evidence


def _validate_assembly_plan(
    plan: AssemblyPlan,
    library_index_sha256: str,
    *,
    require_locked_authority: bool = False,
    expected_slide_count: int | None = None,
) -> None:
    if plan.library_index_sha256 != library_index_sha256:
        raise PhysicalAssemblyError(
            "ASSEMBLY_PLAN_FINGERPRINT_MISMATCH: library index drift"
        )
    if type(plan.target_slide_count) is not int or not 1 <= plan.target_slide_count <= 200:
        raise PhysicalAssemblyError("ASSEMBLY_PLAN_TARGET_COUNT_INVALID")
    seen: set[int] = set()
    seen_page_ids: set[str] = set()
    for slide in plan.target_slides:
        if type(slide.ordinal) is not int or slide.ordinal < 1:
            raise PhysicalAssemblyError("ASSEMBLY_PLAN_ORDINAL_INVALID")
        if slide.ordinal in seen:
            raise PhysicalAssemblyError(
                f"ASSEMBLY_PLAN_DUPLICATE_ORDINAL: {slide.ordinal}"
            )
        seen.add(slide.ordinal)
        if slide.selection_evidence is not None:
            selection = slide.selection_evidence
            if (
                type(selection.candidate_rank) is not int
                or not 1 <= selection.candidate_rank <= 10
                or not isinstance(selection.score_total, (int, float))
                or isinstance(selection.score_total, bool)
                or not math.isfinite(float(selection.score_total))
                or not 0.0 <= float(selection.score_total) <= 1.0
            ):
                raise PhysicalAssemblyError(
                    f"ASSEMBLY_PLAN_SELECTION_INVALID: {slide.ordinal}"
                )
        template = slide.page_template
        if require_locked_authority:
            non_direct_pool = bool(template.pool) and (
                template.pool == "reference-only"
                or str(template.pool).startswith("reference-only/")
            )
            non_direct_decision = str(template.decision or "").lower() in {
                "reference-only", "deny", "denied", "archive", "quarantine"
            }
            if (
                not template.eligibility_known
                or not template.direct_use
                or non_direct_pool
                or non_direct_decision
            ):
                raise PhysicalAssemblyError(
                    f"ASSEMBLY_PLAN_PAGE_NOT_DIRECT_USE: {template.page_id}"
                )
            if template.certification != "certified":
                raise PhysicalAssemblyError(
                    f"ASSEMBLY_PLAN_PAGE_NOT_CERTIFIED: {template.page_id}"
                )
            if template.editability != "native_editable":
                raise PhysicalAssemblyError(
                    f"ASSEMBLY_PLAN_PAGE_NOT_NATIVE_EDITABLE: {template.page_id}"
                )
            if not re.fullmatch(r"[0-9a-f]{64}", template.source_slide_sha256):
                raise PhysicalAssemblyError(
                    f"ASSEMBLY_PLAN_SOURCE_SLIDE_FINGERPRINT_MISSING: {template.page_id}"
                )
            page_match = re.fullmatch(r"([0-9a-f]{64}):(\d{3})", template.page_id)
            if (
                page_match is None
                or page_match.group(1) != template.package_sha256
                or int(page_match.group(2)) != template.slide_number
                or template.source_sha256 != template.package_sha256
            ):
                raise PhysicalAssemblyError(
                    f"ASSEMBLY_PLAN_PAGE_IDENTITY_MISMATCH: {template.page_id}"
                )
        if require_locked_authority and slide.page_template.page_id in seen_page_ids:
            raise PhysicalAssemblyError(
                f"ASSEMBLY_PLAN_DUPLICATE_PAGE_ID: {slide.page_template.page_id}"
            )
        seen_page_ids.add(slide.page_template.page_id)
        expected = set(slide.page_template.slot_graph.get("text_slot_ids", ()))
        actual = set(slide.bindings)
        if not expected and not slide.asset_binding_specs:
            raise PhysicalAssemblyError(
                f"ASSEMBLY_PLAN_NO_EDITABLE_SLOTS: ordinal={slide.ordinal} page_id={slide.page_template.page_id}"
            )
        missing = sorted(expected - actual)
        # Nested aliases are duplicate physical text nodes inside a grouped
        # component. They are excluded from the public slot graph, but must be
        # cleared to prevent source text from surviving beside its replacement.
        # Permit only an empty, no-fact clear binding for such extra nodes.
        raw_extra = actual - expected
        permitted_clear_extras = {
            slot_id
            for slot_id in raw_extra
            if (
                (spec := slide.text_binding_specs.get(slot_id)) is not None
                and spec.mode == "clear"
                and spec.replacement == ""
                and not spec.fact_refs
            )
        }
        extra = sorted(raw_extra - permitted_clear_extras)
        if missing or extra:
            details = []
            if missing:
                details.append("missing=" + ",".join(missing))
            if extra:
                details.append("extra=" + ",".join(extra))
            raise PhysicalAssemblyError(
                f"ASSEMBLY_PLAN_SLOT_COVERAGE: ordinal={slide.ordinal} "
                + " ".join(details)
            )
        effective_spec_slots = set(slide.text_binding_specs) - permitted_clear_extras
        if require_locked_authority and effective_spec_slots != expected:
            raise PhysicalAssemblyError(
                f"ASSEMBLY_PLAN_LOCKED_BINDING_COVERAGE: ordinal={slide.ordinal}"
            )
        inventory = template.governed_content_inventory
        if not isinstance(inventory, Mapping):
            raise PhysicalAssemblyError(
                f"ASSEMBLY_PLAN_GOVERNED_INVENTORY_INVALID: {template.page_id}"
            )
        content_records = tuple(
            record
            for record in inventory.get("slots", ())
            if isinstance(record, Mapping)
        )
        content_slot_ids = {
            str(record["slot_id"])
            for record in content_records
            if isinstance(record.get("slot_id"), str)
        }
        peer_group_ids = {
            str(record["peer_group_id"])
            for record in content_records
            if isinstance(record.get("peer_group_id"), str)
            and record.get("peer_group_id")
        }
        governed_binding_ids = set(slide.governed_content_binding_specs)
        unknown_governed = sorted(
            governed_binding_ids - content_slot_ids - peer_group_ids
        )
        if unknown_governed:
            raise PhysicalAssemblyError(
                "ASSEMBLY_PLAN_UNKNOWN_GOVERNED_CONTENT_SLOT: "
                + ",".join(unknown_governed)
            )
        for record in content_records:
            slot_id = str(record.get("slot_id", ""))
            peer_group_id = record.get("peer_group_id")
            if (
                slot_id in governed_binding_ids
                and isinstance(peer_group_id, str)
                and peer_group_id in governed_binding_ids
            ):
                raise PhysicalAssemblyError(
                    f"ASSEMBLY_PLAN_AMBIGUOUS_GOVERNED_BINDING: {slot_id}"
                )
        if require_locked_authority and not bool(inventory.get("complete")):
            raise PhysicalAssemblyError(
                f"ASSEMBLY_PLAN_GOVERNED_INVENTORY_INCOMPLETE: {template.page_id}"
            )
        if not content_slot_ids and governed_binding_ids:
            raise PhysicalAssemblyError(
                f"ASSEMBLY_PLAN_UNNEEDED_GOVERNED_BINDING: {template.page_id}"
            )
    if len(seen) != plan.target_slide_count:
        raise PhysicalAssemblyError(
            "ASSEMBLY_PLAN_LENGTH_MISMATCH"
        )
    if expected_slide_count is not None:
        if expected_slide_count < 1:
            raise PhysicalAssemblyError("EXPECTED_SLIDE_COUNT_INVALID")
        if plan.target_slide_count != expected_slide_count or len(seen) != expected_slide_count:
            raise PhysicalAssemblyError(
                f"ACCEPTANCE_PROFILE_SLIDE_COUNT_MISMATCH: expected={expected_slide_count}"
            )
        if sorted(seen) != list(range(1, expected_slide_count + 1)):
            raise PhysicalAssemblyError(
                "ACCEPTANCE_PROFILE_ORDINAL_SEQUENCE_MISMATCH"
            )
    if require_locked_authority and plan.authority is None:
        raise PhysicalAssemblyError("ASSEMBLY_PLAN_AUTHORITY_REQUIRED")


def _validate_phase49_sequence(plan: AssemblyPlan) -> None:
    slides = sorted(plan.target_slides, key=lambda item: item.ordinal)
    if len(slides) != 15:
        raise PhysicalAssemblyError("PHASE49_SEQUENCE_REQUIRES_15_SLIDES")
    package_ids = {slide.page_template.package_sha256 for slide in slides}
    if len(package_ids) != 1:
        raise PhysicalAssemblyError("PHASE49_REFERENCE_PACKAGE_MISMATCH")
    for ordinal, slide in enumerate(slides, start=1):
        expected_template_role = PHASE49_TEMPLATE_ROLE_SEQUENCE[ordinal - 1]
        expected_narrative_role = PHASE49_NARRATIVE_ROLE_SEQUENCE[ordinal - 1]
        if (
            slide.ordinal != ordinal
            or slide.page_template.slide_number != ordinal
            or slide.page_template.page_role != expected_template_role
            or slide.narrative_role != expected_narrative_role
        ):
            raise PhysicalAssemblyError(
                "PHASE49_SEQUENCE_MISMATCH: "
                f"ordinal={ordinal} expected_template_role={expected_template_role} "
                f"expected_narrative_role={expected_narrative_role}"
            )


def _fragment_group_contracts_from_locked_template(
    template: Mapping[str, Any],
    *,
    ordinal: int,
    page_id: str,
) -> tuple[FragmentGroupContract, ...]:
    """Derive fragment order from both group membership and slot metadata.

    The group list establishes which slots belong together.  The order is
    independently established by each locked slot's ``group_order``.  Neither
    the assembly plan nor the eventual report may introduce or reorder a
    character-fragment group.
    """

    slot_graph = template.get("slot_graph")
    if not isinstance(slot_graph, Mapping):
        raise PhysicalAssemblyError(
            f"QUERY_FRAGMENT_SLOT_GRAPH_INVALID: {ordinal}:{page_id}"
        )
    raw_slots = slot_graph.get("slots")
    raw_groups = slot_graph.get("fragment_groups")
    if not isinstance(raw_slots, list) or not isinstance(raw_groups, list):
        raise PhysicalAssemblyError(
            f"QUERY_FRAGMENT_GROUPS_INVALID: {ordinal}:{page_id}"
        )
    slot_by_id: dict[str, Mapping[str, Any]] = {}
    for slot in raw_slots:
        slot_id = slot.get("slot_id") if isinstance(slot, Mapping) else None
        if not isinstance(slot_id, str) or not slot_id or slot_id in slot_by_id:
            raise PhysicalAssemblyError(
                f"QUERY_FRAGMENT_SLOT_ID_INVALID: {ordinal}:{page_id}"
            )
        slot_by_id[slot_id] = slot

    contracts: list[FragmentGroupContract] = []
    described_slots: set[str] = set()
    group_ids: set[str] = set()
    for raw_group in raw_groups:
        if not isinstance(raw_group, Mapping):
            raise PhysicalAssemblyError(
                f"QUERY_FRAGMENT_GROUP_INVALID: {ordinal}:{page_id}"
            )
        group_id = raw_group.get("group_id")
        slot_ids = raw_group.get("slot_ids")
        if (
            not isinstance(group_id, str)
            or not re.fullmatch(r"fragment_[0-9]{2}", group_id)
            or group_id in group_ids
            or not isinstance(slot_ids, list)
            or len(slot_ids) < 2
            or not all(isinstance(slot_id, str) and slot_id for slot_id in slot_ids)
            or len(slot_ids) != len(set(slot_ids))
        ):
            raise PhysicalAssemblyError(
                f"QUERY_FRAGMENT_GROUP_INVALID: {ordinal}:{page_id}:{group_id}"
            )
        group_ids.add(group_id)
        declared = set(slot_ids)
        actual = {
            slot_id
            for slot_id, slot in slot_by_id.items()
            if slot.get("group_id") == group_id
        }
        if declared != actual or described_slots.intersection(declared):
            raise PhysicalAssemblyError(
                f"QUERY_FRAGMENT_MEMBERSHIP_INVALID: {ordinal}:{page_id}:{group_id}"
            )
        described_slots.update(declared)
        ordered: list[tuple[int, str]] = []
        for slot_id in slot_ids:
            slot = slot_by_id.get(slot_id)
            group_order = slot.get("group_order") if slot is not None else None
            allowed_modes = slot.get("allowed_binding_modes") if slot is not None else None
            if (
                slot is None
                or type(group_order) is not int
                or group_order < 1
                or slot.get("max_chars") != 1
                or slot.get("semantic_role")
                not in {"title_fragment", "label_fragment"}
                or not isinstance(allowed_modes, list)
                or set(allowed_modes) != {"character", "clear"}
            ):
                raise PhysicalAssemblyError(
                    f"QUERY_FRAGMENT_SLOT_CONTRACT_INVALID: "
                    f"{ordinal}:{page_id}:{group_id}:{slot_id}"
                )
            ordered.append((group_order, slot_id))
        if sorted(order for order, _ in ordered) != list(
            range(1, len(ordered) + 1)
        ):
            raise PhysicalAssemblyError(
                f"QUERY_FRAGMENT_ORDER_INVALID: {ordinal}:{page_id}:{group_id}"
            )
        contracts.append(
            FragmentGroupContract(
                ordinal=ordinal,
                page_id=page_id,
                group_id=group_id,
                ordered_slot_ids=tuple(
                    slot_id for _, slot_id in sorted(ordered)
                ),
            )
        )

    undeclared = sorted(
        slot_id
        for slot_id, slot in slot_by_id.items()
        if isinstance(slot.get("group_id"), str)
        and str(slot["group_id"]).startswith("fragment_")
        and slot_id not in described_slots
    )
    if undeclared:
        raise PhysicalAssemblyError(
            f"QUERY_FRAGMENT_SLOT_UNDECLARED: {ordinal}:{page_id}:"
            + ",".join(undeclared)
        )
    return tuple(sorted(contracts, key=lambda item: item.group_id))


def _validate_query_selection_evidence(
    plan: AssemblyPlan,
    *,
    project_root: str | os.PathLike[str] | None,
    library_index: LibraryIndex | None,
    require_phase49_ordinals: bool = False,
) -> SelectionAuthorityEvidence:
    if library_index is None:
        raise PhysicalAssemblyError("QUERY_LIBRARY_INDEX_REQUIRED")
    if library_index.page_template_count != len(library_index.page_templates):
        raise PhysicalAssemblyError("QUERY_LIBRARY_INDEX_COUNT_MISMATCH")
    if not plan.query_bundle_path or not plan.query_bundle_sha256:
        raise PhysicalAssemblyError("ASSEMBLY_PLAN_QUERY_BUNDLE_REQUIRED")
    if project_root is None:
        query_path = Path(plan.query_bundle_path).expanduser().resolve(strict=True)
    else:
        query_path = resolve_project_file(
            plan.query_bundle_path,
            project_root,
            label="QUERY_BUNDLE",
        )
    if _sha256_file(query_path) != plan.query_bundle_sha256:
        raise PhysicalAssemblyError("QUERY_BUNDLE_FINGERPRINT_MISMATCH")
    try:
        payload = json.loads(query_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise PhysicalAssemblyError(f"QUERY_BUNDLE_INVALID: {exc}") from exc
    if not isinstance(payload, Mapping):
        raise PhysicalAssemblyError("QUERY_BUNDLE_INVALID: root must be an object")
    _validate_schema_payload(
        payload,
        "page-template-query-bundle.v1.schema.json",
        label="QUERY_BUNDLE",
    )
    if (
        payload.get("schema_version") != "page-template-query-bundle.v1"
        or payload.get("library_index_sha256") != plan.library_index_sha256
        or not isinstance(payload.get("queries"), list)
    ):
        raise PhysicalAssemblyError("QUERY_BUNDLE_CONTRACT_MISMATCH")
    queries = payload["queries"]
    if (
        not isinstance(payload.get("query_count"), int)
        or isinstance(payload.get("query_count"), bool)
        or payload["query_count"] != len(queries)
        or len(queries) != plan.target_slide_count
    ):
        raise PhysicalAssemblyError("QUERY_BUNDLE_COUNT_MISMATCH")
    by_ordinal: dict[int, Mapping[str, Any]] = {}
    query_ids: set[str] = set()
    for query in queries:
        if not isinstance(query, Mapping):
            raise PhysicalAssemblyError("QUERY_BUNDLE_QUERY_INVALID")
        ordinal = query.get("target_ordinal")
        if type(ordinal) is not int or ordinal in by_ordinal:
            raise PhysicalAssemblyError("QUERY_BUNDLE_ORDINAL_INVALID")
        query_id = query.get("query_id")
        if not isinstance(query_id, str) or not query_id or query_id in query_ids:
            raise PhysicalAssemblyError("QUERY_BUNDLE_QUERY_ID_INVALID")
        query_ids.add(query_id)
        by_ordinal[ordinal] = query
    expected_ordinals = {slide.ordinal for slide in plan.target_slides}
    if set(by_ordinal) != expected_ordinals:
        raise PhysicalAssemblyError("QUERY_BUNDLE_ORDINAL_SET_MISMATCH")
    selected_page_ids: set[str] = set()
    fragment_group_contracts: list[FragmentGroupContract] = []
    for slide in plan.target_slides:
        selection = slide.selection_evidence
        query = by_ordinal.get(slide.ordinal)
        if selection is None or query is None:
            raise PhysicalAssemblyError(
                f"QUERY_SELECTION_EVIDENCE_MISSING: {slide.ordinal}"
            )
        if query.get("query_id") != selection.query_id:
            raise PhysicalAssemblyError(
                f"QUERY_SELECTION_ID_MISMATCH: {slide.ordinal}"
            )
        result = query.get("result")
        if (
            not isinstance(result, Mapping)
            or result.get("schema_version") != "page-template-query-result.v1"
            or result.get("library_index_sha256") != plan.library_index_sha256
        ):
            raise PhysicalAssemblyError(
                f"QUERY_SELECTION_RESULT_MISMATCH: {slide.ordinal}"
            )
        _validate_schema_payload(
            result,
            "page-template-query-result.v1.schema.json",
            label=f"QUERY_RESULT_{slide.ordinal}",
        )
        if result.get("role") != slide.page_template.page_role:
            raise PhysicalAssemblyError(
                f"QUERY_SELECTION_ROLE_MISMATCH: {slide.ordinal}"
            )
        required_source_ordinal = result.get("required_source_ordinal")
        query_limit = int(result["limit"])
        recompute_limit = (
            max(query_limit, library_index.page_template_count)
            if required_source_ordinal is not None
            else query_limit
        )
        ranked_recomputed = query_page_template_candidates(
            library_index,
            role=str(result["role"]),
            capacity_budget=int(result["capacity_budget"]),
            semantic_categories=tuple(result["semantic_categories"]),
            style_cluster=str(result["style_cluster"]),
            asset_requirements=tuple(result["asset_requirements"]),
            customer_assets_available=bool(result["customer_assets_available"]),
            limit=recompute_limit,
            allow_fallback=bool(result["allow_fallback"]),
            direct_use_only=True,
            include_ineligible=bool(result["include_ineligible"]),
        )
        recomputed = (
            tuple(
                candidate
                for candidate in ranked_recomputed
                if candidate.page_template.slide_number
                == required_source_ordinal
            )[:query_limit]
            if required_source_ordinal is not None
            else ranked_recomputed
        )
        if required_source_ordinal is not None and not recomputed:
            raise PhysicalAssemblyError(
                "QUERY_REQUIRED_SOURCE_ORDINAL_NO_MATCH: "
                f"{slide.ordinal}:{required_source_ordinal}"
            )
        if require_phase49_ordinals and (
            required_source_ordinal != slide.ordinal
            or slide.page_template.slide_number != required_source_ordinal
        ):
            raise PhysicalAssemblyError(
                "PHASE49_QUERY_SOURCE_ORDINAL_MISMATCH: "
                f"target={slide.ordinal} required={required_source_ordinal} "
                f"selected={slide.page_template.slide_number}"
            )
        recomputed_payload = [candidate.to_dict() for candidate in recomputed]
        if (
            result.get("weights") != dict(DEFAULT_SCORING)
            or result.get("count") != len(recomputed_payload)
            or result.get("eligible_count")
            != sum(candidate.eligibility for candidate in recomputed)
            or result.get("candidates") != recomputed_payload
        ):
            raise PhysicalAssemblyError(
                f"QUERY_RESULT_RECOMPUTE_MISMATCH: {slide.ordinal}"
            )
        candidates = result.get("candidates") if isinstance(result, Mapping) else None
        if (
            type(selection.candidate_rank) is not int
            or not 1 <= selection.candidate_rank <= 10
            or not isinstance(candidates, list)
            or selection.candidate_rank > len(candidates)
            or not math.isfinite(selection.score_total)
            or not 0.0 <= selection.score_total <= 1.0
        ):
            raise PhysicalAssemblyError(
                f"QUERY_SELECTION_RANK_INVALID: {slide.ordinal}"
            )
        candidate = candidates[selection.candidate_rank - 1]
        if not isinstance(candidate, Mapping):
            raise PhysicalAssemblyError(
                f"QUERY_SELECTION_CANDIDATE_INVALID: {slide.ordinal}"
            )
        page_template = candidate.get("page_template")
        scores = candidate.get("scores")
        try:
            score_total = float(scores.get("total")) if isinstance(scores, Mapping) else math.nan
        except (TypeError, ValueError):
            score_total = math.nan
        if (
            candidate.get("eligibility") is not True
            or candidate.get("page_id") != slide.page_template.page_id
            or not isinstance(page_template, Mapping)
            or page_template.get("page_id") != slide.page_template.page_id
            or page_template.get("package_sha256") != slide.page_template.package_sha256
            or page_template.get("slide_number") != slide.page_template.slide_number
            or page_template.get("source_slide_sha256")
            != slide.page_template.source_slide_sha256
            or page_template.get("style_cluster_id")
            != slide.page_template.style_cluster_id
            or not isinstance(scores, Mapping)
            or not math.isfinite(score_total)
            or not math.isclose(
                score_total,
                selection.score_total,
                rel_tol=0.0,
                abs_tol=1e-9,
            )
            or candidate.get("fallback_reason") != selection.fallback_reason
            or (
                required_source_ordinal is not None
                and slide.page_template.slide_number != required_source_ordinal
            )
        ):
            raise PhysicalAssemblyError(
                f"QUERY_SELECTION_CANDIDATE_MISMATCH: {slide.ordinal}"
            )
        locked_fragment_contracts = _fragment_group_contracts_from_locked_template(
            page_template,
            ordinal=slide.ordinal,
            page_id=slide.page_template.page_id,
        )
        plan_fragment_contracts = _fragment_group_contracts_from_locked_template(
            slide.page_template.to_dict(),
            ordinal=slide.ordinal,
            page_id=slide.page_template.page_id,
        )
        if locked_fragment_contracts != plan_fragment_contracts:
            raise PhysicalAssemblyError(
                f"QUERY_SELECTION_FRAGMENT_GROUP_MISMATCH: {slide.ordinal}"
            )
        fragment_group_contracts.extend(locked_fragment_contracts)
        selected_page_ids.add(slide.page_template.page_id)
    if len(selected_page_ids) != plan.target_slide_count:
        raise PhysicalAssemblyError("QUERY_SELECTION_DUPLICATE_PAGE_ID")
    return SelectionAuthorityEvidence(
        mode="locked",
        query_bundle_path=str(query_path),
        query_bundle_sha256=plan.query_bundle_sha256,
        library_index_sha256=plan.library_index_sha256,
        query_count=len(queries),
        selected_count=len(plan.target_slides),
        distinct_query_id_count=len(query_ids),
        distinct_page_id_count=len(selected_page_ids),
        status="pass",
        fragment_group_contracts=tuple(fragment_group_contracts),
    )


def _byte_match_score(source: bytes, target: bytes) -> float:
    if not source:
        return 0.0
    matches = sum(1 for a, b in zip(source, target) if a == b)
    return matches / max(len(source), len(target))


def _slide_structure_signature(slide_xml: bytes) -> str:
    """Hash slide geometry/object structure while excluding governed mutations."""

    try:
        root = ET.fromstring(slide_xml)
    except ET.ParseError:
        return ""
    def canonical(node: ET.Element) -> Any:
        local = node.tag.rsplit("}", 1)[-1]
        # Text may legitimately be replaced, fragmented, fitted and resized by
        # the governed adapter.  It is not evidence of a change to the source
        # page's geometry or object structure.  The surrounding shape/table
        # hierarchy remains in the signature, so a flattened page or an
        # invented shape still fails lineage verification.
        if local == "txBody":
            return ["__GOVERNED_TEXT_SURFACE__"]
        if local == "srcRect":
            return None
        if local in {"spAutoFit", "normAutofit", "noAutofit"}:
            return ["__GOVERNED_FIT_POLICY__"]
        children = [
            value
            for child in list(node)
            if (value := canonical(child)) is not None
        ]
        text = "__GOVERNED_TEXT__" if local == "t" else (node.text or "").strip()
        attributes = [
            (
                key,
                "__GOVERNED_TEXT_SIZE__"
                if local in {"rPr", "defRPr", "endParaRPr"}
                and key == "sz"
                else value,
            )
            for key, value in sorted(node.attrib.items())
        ]
        return [
            node.tag,
            attributes,
            text,
            children,
        ]

    encoded = json.dumps(
        canonical(root),
        ensure_ascii=True,
        sort_keys=False,
        separators=(",", ":"),
    ).encode("utf-8")
    return _sha256_bytes(encoded)


def _imported_part_map_evidence(
    context: AssemblyImportContext | None,
    *,
    slide: AssemblyTargetSlide,
    graph: _SourceGraph,
    target_archive: zipfile.ZipFile,
) -> tuple[str, int]:
    if context is None:
        return "", 0
    records: list[dict[str, str]] = []
    target_slide = f"ppt/slides/slide{slide.ordinal}.xml"
    records.append(
        {
            "source_part": graph.root_slide_name,
            "source_sha256": graph.slide_sha,
            "target_part": target_slide,
            "target_sha256": _sha256_bytes(target_archive.read(target_slide)),
        }
    )
    for source_part in sorted(graph.extra_parts):
        target_part = context.slide_target_maps.get(slide.ordinal, {}).get(source_part)
        if target_part is None:
            return "", 0
        if target_part in context.pruned_parts:
            continue
        try:
            target_bytes = target_archive.read(target_part)
        except KeyError:
            return "", 0
        records.append(
            {
                "source_part": source_part,
                "source_sha256": _sha256_bytes(graph.extra_parts[source_part]),
                "target_part": target_part,
                "target_sha256": _sha256_bytes(target_bytes),
            }
        )
    encoded = json.dumps(
        records,
        ensure_ascii=True,
        sort_keys=True,
        separators=(",", ":"),
    ).encode("utf-8")
    return _sha256_bytes(encoded), len(records)


def _slide_text_run_count(slide_xml: bytes) -> int:
    return len(_TEXT_RE.findall(slide_xml.decode("utf-8", errors="replace")))


def _slide_shape_count(slide_xml: bytes) -> int:
    return len(_CNVPR_RE.findall(slide_xml.decode("utf-8", errors="replace")))


def _verify_zip_open(output_path: Path) -> tuple[bool, str]:
    try:
        with zipfile.ZipFile(output_path, "r"):
            return True, ""
    except (OSError, zipfile.BadZipFile) as exc:
        return False, f"zip-open failed: {exc}"


def _content_type_lookup(
    defaults: Sequence[tuple[str, str]],
    overrides: Sequence[tuple[str, str]],
) -> tuple[dict[str, str], dict[str, str]]:
    return (
        {extension.lower(): content_type for extension, content_type in defaults},
        {part_name.lstrip("/"): content_type for part_name, content_type in overrides},
    )


def _content_type_for_part(
    part_name: str,
    defaults: Mapping[str, str],
    overrides: Mapping[str, str],
) -> str | None:
    override = overrides.get(part_name.lstrip("/"))
    if override:
        return override
    return defaults.get(_part_extension(part_name))


def _verify_content_types(output_path: Path) -> tuple[bool, str, int]:
    try:
        with zipfile.ZipFile(output_path, "r") as archive:
            xml = archive.read("[Content_Types].xml")
            defaults, overrides = _parse_content_types(xml)
            default_map, override_map = _content_type_lookup(defaults, overrides)
            missing = [
                name
                for name in archive.namelist()
                if name != "[Content_Types].xml"
                and not name.endswith("/")
                and _content_type_for_part(name, default_map, override_map) is None
            ]
            dangling = sorted(name for name in override_map if name not in archive.namelist())
            if missing or dangling:
                details: list[str] = []
                if missing:
                    details.append("missing=" + ",".join(sorted(missing)[:8]))
                if dangling:
                    details.append("dangling=" + ",".join(dangling[:8]))
                return False, "content-types " + "; ".join(details), len(defaults) + len(overrides)
            return True, "", len(defaults) + len(overrides)
    except (OSError, KeyError, zipfile.BadZipFile, PhysicalAssemblyError) as exc:
        return False, f"content-types parse failed: {exc}", 0


@dataclass(frozen=True)
class RelationshipAudit:
    total_relationship_count: int
    internal_relationship_count: int
    external_relationship_count: int
    unresolved_internal_relationships: tuple[Mapping[str, str], ...]
    unsafe_relationships: tuple[Mapping[str, str], ...]
    status: str
    details: str


def _relationship_finding(
    rels_path: str,
    entry: Mapping[str, str] | None,
    *,
    reason: str,
    resolved_target: str = "",
) -> dict[str, str]:
    value = entry or {}
    return {
        "owner_rels_part": rels_path,
        "relationship_id": value.get("Id", ""),
        "relationship_type": value.get("Type", ""),
        "target_mode": value.get("TargetMode", ""),
        "raw_target": value.get("Target", ""),
        "resolved_target": resolved_target,
        "reason": reason,
    }


def _distinct_relationship_findings(
    records: Iterable[Mapping[str, str]],
) -> tuple[Mapping[str, str], ...]:
    by_key: dict[tuple[str, ...], Mapping[str, str]] = {}
    fields = (
        "owner_rels_part",
        "relationship_id",
        "relationship_type",
        "target_mode",
        "raw_target",
        "resolved_target",
        "reason",
    )
    for record in records:
        normalized = {field: str(record.get(field, "")) for field in fields}
        by_key[tuple(normalized[field] for field in fields)] = normalized
    return tuple(by_key[key] for key in sorted(by_key))


def _inspect_all_relationships(output_path: Path) -> RelationshipAudit:
    """Inspect every relationship and keep unsafe/unresolved evidence separate."""

    try:
        with zipfile.ZipFile(output_path, "r") as archive:
            names = frozenset(_normalise_zip_name(name) for name in archive.namelist())
            unresolved: list[Mapping[str, str]] = []
            unsafe: list[Mapping[str, str]] = []
            total = 0
            internal = 0
            external = 0
            slides = sorted(name for name in names if _SLIDE_RE.match(name))
            for slide in slides:
                rels_path = _rels_path_for_part(slide)
                if rels_path not in names:
                    unresolved.append(
                        _relationship_finding(
                            rels_path,
                            None,
                            reason="missing-relationship-part",
                            resolved_target=slide,
                        )
                    )
            for rels_path in sorted(name for name in names if name.endswith(".rels")):
                rels_xml = archive.read(rels_path)
                for entry in _parse_relationships(rels_xml):
                    total += 1
                    target = entry.get("Target", "")
                    is_external = _relationship_is_external(entry)
                    if is_external:
                        external += 1
                    else:
                        internal += 1
                    security_issue = _relationship_security_issue(entry)
                    if security_issue is not None:
                        unsafe.append(
                            _relationship_finding(
                                rels_path,
                                entry,
                                reason=security_issue,
                            )
                        )
                        continue
                    if is_external:
                        continue
                    resolved = _resolve_rel_target(rels_xml, rels_path, target)
                    if resolved is None or resolved not in names:
                        unresolved.append(
                            _relationship_finding(
                                rels_path,
                                entry,
                                reason="unresolved-internal-target",
                                resolved_target=resolved or "",
                            )
                        )
            unresolved_records = _distinct_relationship_findings(unresolved)
            unsafe_records = _distinct_relationship_findings(unsafe)
            detail_rows = [
                f"unresolved {item['owner_rels_part']}->{item['resolved_target'] or item['raw_target']}"
                for item in unresolved_records
            ] + [
                f"unsafe {item['owner_rels_part']}:{item['reason']}:{item['raw_target']}"
                for item in unsafe_records
            ]
            return RelationshipAudit(
                total_relationship_count=total,
                internal_relationship_count=internal,
                external_relationship_count=external,
                unresolved_internal_relationships=unresolved_records,
                unsafe_relationships=unsafe_records,
                status="pass" if not unresolved_records and not unsafe_records else "fail",
                details="; ".join(detail_rows[:20]),
            )
    except (OSError, KeyError, zipfile.BadZipFile, PhysicalAssemblyError) as exc:
        record = _relationship_finding(
            "",
            None,
            reason="relationship-package-read-failed",
        )
        return RelationshipAudit(0, 0, 0, (record,), (), "fail", f"relationships failed: {exc}")


def _verify_all_relationships(output_path: Path) -> tuple[bool, str, int]:
    """Compatibility wrapper returning unresolved-internal count as before."""

    audit = _inspect_all_relationships(output_path)
    return (
        audit.status == "pass",
        audit.details,
        len(audit.unresolved_internal_relationships),
    )


def _verify_slide_rels(output_path: Path) -> tuple[bool, str]:
    """Backward-compatible alias for the former slide-only verifier."""

    ok, details, _ = _verify_all_relationships(output_path)
    return ok, details


def _static_duplicate_bytes(output_path: Path) -> int:
    """Count avoidable duplicate immutable bytes still present in an output."""

    if not output_path.is_file():
        return 0
    try:
        with zipfile.ZipFile(output_path, "r") as archive:
            defaults, overrides = _parse_content_types(archive.read("[Content_Types].xml"))
            default_map, override_map = _content_type_lookup(defaults, overrides)
            names = frozenset(archive.namelist())
            groups: dict[tuple[str, str, str], list[int]] = {}
            for name in names:
                if name.endswith(".rels") or name.endswith("/"):
                    continue
                content_type = _content_type_for_part(name, default_map, override_map)
                if not content_type:
                    continue
                relationship_free = _rels_path_for_part(name) not in names
                if not _is_cross_package_dedup_candidate(
                    name,
                    content_type,
                    relationship_free=relationship_free,
                ):
                    continue
                data = archive.read(name)
                key = (content_type, _part_extension(name), _sha256_bytes(data))
                groups.setdefault(key, []).append(len(data))
            return sum(sum(sizes[1:]) for sizes in groups.values() if len(sizes) > 1)
    except (OSError, KeyError, zipfile.BadZipFile, PhysicalAssemblyError):
        return 0


def _verify_python_pptx(
    output_path: Path,
) -> tuple[bool, int, int, int, int, int, int, int, int, float]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return False, 0, 0, 0, 0, 0, 0, 0, 0, 0.0
    try:
        pres = Presentation(str(output_path))
    except Exception:
        return False, 0, 0, 0, 0, 0, 0, 0, 0, 0.0
    total_text = 0
    total_shapes = 0
    native_objects = 0
    pictures = 0
    native_slides = 0
    full_slide_rasters = 0
    raster_dominant_slides = 0
    slide_area = max(1, int(pres.slide_width) * int(pres.slide_height))

    def clipped_rectangle(shape: Any) -> tuple[int, int, int, int] | None:
        try:
            left = max(0, int(shape.left))
            top = max(0, int(shape.top))
            right = min(int(pres.slide_width), int(shape.left) + int(shape.width))
            bottom = min(int(pres.slide_height), int(shape.top) + int(shape.height))
        except (AttributeError, TypeError, ValueError):
            return None
        if right <= left or bottom <= top:
            return None
        return left, top, right, bottom

    def rectangle_union_area(
        rectangles: Sequence[tuple[int, int, int, int]],
    ) -> int:
        if not rectangles:
            return 0
        xs = sorted({value for rect in rectangles for value in (rect[0], rect[2])})
        total = 0
        for x0, x1 in zip(xs, xs[1:]):
            if x1 <= x0:
                continue
            intervals = sorted(
                (top, bottom)
                for left, top, right, bottom in rectangles
                if left < x1 and right > x0
            )
            if not intervals:
                continue
            merged_height = 0
            current_top, current_bottom = intervals[0]
            for top, bottom in intervals[1:]:
                if top <= current_bottom:
                    current_bottom = max(current_bottom, bottom)
                else:
                    merged_height += current_bottom - current_top
                    current_top, current_bottom = top, bottom
            merged_height += current_bottom - current_top
            total += (x1 - x0) * merged_height
        return total

    for slide in pres.slides:
        slide_native = 0
        slide_full_raster = False
        slide_text_runs = 0
        picture_rectangles: list[tuple[int, int, int, int]] = []
        native_rectangles: list[tuple[int, int, int, int]] = []
        text_rectangles: list[tuple[int, int, int, int]] = []

        def leaf_shapes(shapes: Any) -> Iterable[Any]:
            for shape in shapes:
                if int(getattr(shape, "shape_type", -1)) == 6 and hasattr(shape, "shapes"):
                    yield from leaf_shapes(shape.shapes)
                else:
                    yield shape

        for shape in leaf_shapes(slide.shapes):
            total_shapes += 1
            # python-pptx uses 13 for PICTURE and 11 for LINKED_PICTURE.
            is_picture = int(getattr(shape, "shape_type", -1)) in {11, 13}
            rectangle = clipped_rectangle(shape)
            if is_picture:
                pictures += 1
                if rectangle is not None:
                    picture_rectangles.append(rectangle)
                shape_area = (
                    (rectangle[2] - rectangle[0]) * (rectangle[3] - rectangle[1])
                    if rectangle is not None
                    else 0
                )
                if shape_area / slide_area >= 0.90:
                    slide_full_raster = True
            else:
                native_objects += 1
                slide_native += 1
                if rectangle is not None:
                    native_rectangles.append(rectangle)
            if shape.has_text_frame:
                shape_has_text = False
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        total_text += 1
                        slide_text_runs += 1
                        if run.text.strip():
                            shape_has_text = True
                if shape_has_text and rectangle is not None:
                    text_rectangles.append(rectangle)
        if slide_native > 0:
            native_slides += 1
        if slide_full_raster:
            full_slide_rasters += 1
        picture_coverage = rectangle_union_area(picture_rectangles) / slide_area
        native_coverage = rectangle_union_area(native_rectangles) / slide_area
        text_coverage = rectangle_union_area(text_rectangles) / slide_area
        # Certified commercial pages frequently use one full-bleed bitmap as
        # background art while preserving a native shape system and a
        # deliberately sparse editorial heading. That is materially different
        # from rasterising a whole slide and adding a tiny native decoy. Treat
        # the slide as raster-dominant only when *both* its non-picture native
        # surface and its editable text surface are sparse. The conjunction
        # keeps rejecting screenshots with decorative shapes/text while
        # allowing an imported, physically certified background to coexist
        # with real editable geometry.
        if picture_coverage >= 0.90 and (
            native_coverage < 0.10 and text_coverage < 0.04
        ):
            raster_dominant_slides += 1
    slide_count = len(pres.slides)
    coverage = round(native_slides / slide_count, 6) if slide_count else 0.0
    return (
        True,
        slide_count,
        total_text,
        total_shapes,
        native_objects,
        pictures,
        native_slides,
        full_slide_rasters,
        raster_dominant_slides,
        coverage,
    )


def _verify_libreoffice(
    output_path: Path,
    *,
    required: bool,
) -> LibreOfficeEvidence:
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if not required:
        return LibreOfficeEvidence(
            available=executable is not None,
            executable=executable or "",
            open_result="not_run",
            render_result="not_run",
            status="not_run",
        )
    if executable is None:
        return LibreOfficeEvidence(
            available=False,
            executable="",
            open_result="fail",
            render_result="fail",
            status="fail",
            details="LibreOffice executable is unavailable",
        )
    try:
        with tempfile.TemporaryDirectory(prefix="window-pptx-lo-") as temp_name:
            temp_dir = Path(temp_name)
            output_dir = temp_dir / "rendered"
            profile_dir = temp_dir / "profile"
            output_dir.mkdir()
            profile_dir.mkdir()
            completed = subprocess.run(
                [
                    executable,
                    "--headless",
                    f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(output_dir),
                    str(output_path),
                ],
                check=False,
                capture_output=True,
                text=True,
                timeout=90,
            )
            pdf_path = output_dir / f"{output_path.stem}.pdf"
            open_ok = completed.returncode == 0
            render_ok = open_ok and pdf_path.is_file() and pdf_path.stat().st_size > 0
            details = (completed.stderr or completed.stdout or "").strip()[:2000]
            return LibreOfficeEvidence(
                available=True,
                executable=executable,
                open_result="pass" if open_ok else "fail",
                render_result="pass" if render_ok else "fail",
                status="pass" if open_ok and render_ok else "fail",
                details=details,
            )
    except (OSError, subprocess.SubprocessError) as exc:
        return LibreOfficeEvidence(
            available=True,
            executable=executable,
            open_result="fail",
            render_result="fail",
            status="fail",
            details=str(exc),
        )


def _read_governed_workbook_slot(
    package_bytes: bytes,
    record: Mapping[str, Any],
) -> str:
    try:
        return read_governed_xlsx_slot(
            package_bytes,
            str(record.get("locator", "")),
        )
    except WorkbookSecurityError as exc:
        raise PhysicalAssemblyError(f"GOVERNED_WORKBOOK_SECURITY: {exc}") from exc


def _read_governed_output_slot(
    part_bytes: bytes,
    record: Mapping[str, Any],
) -> str:
    source_part = str(record.get("source_part", ""))
    if source_part.lower().endswith(".xlsm"):
        raise PhysicalAssemblyError("GOVERNED_WORKBOOK_XLSM_FORBIDDEN")
    if source_part.lower().endswith(".xlsx"):
        return _read_governed_workbook_slot(part_bytes, record)
    try:
        root = ET.fromstring(part_bytes)
    except ET.ParseError as exc:
        raise PhysicalAssemblyError(
            f"GOVERNED_OUTPUT_XML_INVALID: {source_part}"
        ) from exc
    return _xml_slot_text_and_set(root, record, None)


def _output_relationship_by_id(
    archive: zipfile.ZipFile,
    owner_part: str,
    relationship_id: str,
) -> tuple[str, str]:
    """Resolve one exact safe internal relationship from the final package."""

    rels_path = _rels_path_for_part(owner_part)
    try:
        rels_xml = archive.read(rels_path)
    except KeyError as exc:
        raise PhysicalAssemblyError(
            f"GOVERNED_RELATIONSHIP_PART_MISSING: {rels_path}"
        ) from exc
    matches = [
        entry
        for entry in _parse_relationships(rels_xml)
        if entry.get("Id") == relationship_id
    ]
    if len(matches) != 1:
        raise PhysicalAssemblyError(
            f"GOVERNED_RELATIONSHIP_ID_INVALID: {owner_part}#{relationship_id}"
        )
    entry = matches[0]
    issue = _relationship_security_issue(entry)
    if issue is not None or _relationship_is_external(entry):
        raise PhysicalAssemblyError(
            f"GOVERNED_RELATIONSHIP_UNSAFE: {owner_part}#{relationship_id}"
        )
    target = _resolve_rel_target(rels_xml, rels_path, entry.get("Target", ""))
    if target is None or target not in archive.namelist():
        raise PhysicalAssemblyError(
            f"GOVERNED_RELATIONSHIP_TARGET_MISSING: {owner_part}#{relationship_id}"
        )
    return entry.get("Type", ""), target


def _resolve_output_governed_target(
    archive: zipfile.ZipFile,
    *,
    ordinal: int,
    record: Mapping[str, Any],
) -> dict[str, Any]:
    """Follow final slide/chart relationships to one governed output part."""

    kind = str(record.get("kind", ""))
    locator = str(record.get("locator", ""))
    match = _FRAME_PREFIX_RE.match(locator)
    if match is None:
        raise PhysicalAssemblyError("GOVERNED_LOCATOR_FRAME_MISSING")
    shape_id = int(match.group(1))
    slide_part = f"ppt/slides/slide{ordinal}.xml"
    try:
        slide_root = ET.fromstring(archive.read(slide_part))
    except (KeyError, ET.ParseError) as exc:
        raise PhysicalAssemblyError(
            f"GOVERNED_OUTPUT_SLIDE_INVALID: {slide_part}"
        ) from exc
    frames: list[ET.Element] = []
    for frame in slide_root.iter(f"{{{PML_NS}}}graphicFrame"):
        marker = frame.find(f".//{{{PML_NS}}}cNvPr")
        if marker is not None and marker.get("id") == str(shape_id):
            frames.append(frame)
    if len(frames) != 1:
        raise PhysicalAssemblyError(
            f"GOVERNED_OUTPUT_FRAME_INVALID: {slide_part}#{shape_id}"
        )
    frame = frames[0]
    if kind == "table-cell":
        if len(list(frame.iter(f"{{{DML_NS}}}tbl"))) != 1:
            raise PhysicalAssemblyError("GOVERNED_OUTPUT_TABLE_FRAME_INVALID")
        target_part = slide_part
        return {
            "slide_part": slide_part,
            "shape_id": shape_id,
            "slide_relationship_id": "",
            "chart_part": "",
            "chart_relationship_id": "",
            "target_part": target_part,
            "target_part_sha256": _sha256_bytes(archive.read(target_part)),
        }

    charts = list(frame.iter(f"{{{CHART_NS}}}chart"))
    if len(charts) != 1:
        raise PhysicalAssemblyError("GOVERNED_OUTPUT_CHART_FRAME_INVALID")
    slide_relationship_id = charts[0].get(f"{{{OD_REL_NS}}}id", "")
    if not slide_relationship_id:
        raise PhysicalAssemblyError("GOVERNED_OUTPUT_CHART_RELATIONSHIP_MISSING")
    slide_rel_type, chart_part = _output_relationship_by_id(
        archive,
        slide_part,
        slide_relationship_id,
    )
    if not slide_rel_type.lower().rstrip("/").endswith("/chart"):
        raise PhysicalAssemblyError("GOVERNED_OUTPUT_CHART_RELATIONSHIP_INVALID")
    if kind != "workbook-cell":
        return {
            "slide_part": slide_part,
            "shape_id": shape_id,
            "slide_relationship_id": slide_relationship_id,
            "chart_part": chart_part,
            "chart_relationship_id": "",
            "target_part": chart_part,
            "target_part_sha256": _sha256_bytes(archive.read(chart_part)),
        }

    try:
        chart_root = ET.fromstring(archive.read(chart_part))
    except (KeyError, ET.ParseError) as exc:
        raise PhysicalAssemblyError(
            f"GOVERNED_OUTPUT_CHART_INVALID: {chart_part}"
        ) from exc
    external_data = list(chart_root.iter(f"{{{CHART_NS}}}externalData"))
    if len(external_data) != 1:
        raise PhysicalAssemblyError("GOVERNED_OUTPUT_EXTERNAL_DATA_INVALID")
    chart_relationship_id = external_data[0].get(f"{{{OD_REL_NS}}}id", "")
    if not chart_relationship_id:
        raise PhysicalAssemblyError("GOVERNED_OUTPUT_WORKBOOK_RELATIONSHIP_MISSING")
    chart_rel_type, target_part = _output_relationship_by_id(
        archive,
        chart_part,
        chart_relationship_id,
    )
    if not chart_rel_type.lower().rstrip("/").endswith("/package"):
        raise PhysicalAssemblyError("GOVERNED_OUTPUT_WORKBOOK_RELATIONSHIP_INVALID")
    if not target_part.lower().endswith(".xlsx"):
        raise PhysicalAssemblyError("GOVERNED_OUTPUT_WORKBOOK_TARGET_INVALID")
    return {
        "slide_part": slide_part,
        "shape_id": shape_id,
        "slide_relationship_id": slide_relationship_id,
        "chart_part": chart_part,
        "chart_relationship_id": chart_relationship_id,
        "target_part": target_part,
        "target_part_sha256": _sha256_bytes(archive.read(target_part)),
    }


def _governed_peer_value(value: str) -> tuple[str, str]:
    numeric = _decimal_literal(value)
    if numeric is None:
        return "text", _normalise_embedded_text(value)
    amount, is_percent = numeric
    if is_percent:
        amount /= Decimal(100)
    return "number", format(amount.normalize(), "f")


def _verify_source_residue(
    output_path: Path,
    *,
    plan: AssemblyPlan,
    context: AssemblyImportContext | None,
    binding_evidence: Sequence[BindingEvidence],
    required: bool,
) -> SourceResidueEvidence:
    empty_manifest_sha = _sha256_bytes(b"[]")
    if not required:
        return SourceResidueEvidence(
            governed_content_slot_count=0,
            governed_content_binding_count=0,
            verified_governed_content_count=0,
            governed_content_mismatch_count=0,
            peer_group_mismatch_count=0,
            mutation_manifest_sha256=empty_manifest_sha,
            governed_mutations=(),
            unauthorized_content_count=0,
            tag_part_count=0,
            tag_relationship_count=0,
            layout_master_cached_field_count=0,
            certified_media_count=0,
            media_hash_mismatch_count=0,
            replacement_asset_count=0,
            replacement_asset_hash_mismatch_count=0,
            asset_slot_mismatch_count=0,
            orphan_media_count=0,
            status="pass",
        )

    # The catalog deliberately errs on the conservative side and can classify
    # a one/two digit template page marker as governed content.  Assembly
    # clears that marker and retains a non-client evidence row; it is not a
    # customer-data slot for residue verification.
    cleared_decoration_keys = {
        (item.ordinal, item.slot_id)
        for item in binding_evidence
        if (
            item.binding_kind == "source-decoration"
            and item.mode == "source-decoration-numeric"
        )
    }
    expected_records: dict[tuple[int, str], tuple[AssemblyTargetSlide, Mapping[str, Any]]] = {}
    for slide in plan.target_slides:
        inventory = slide.page_template.governed_content_inventory
        for record in inventory.get("slots", ()):
            if not isinstance(record, Mapping) or not isinstance(
                record.get("slot_id"), str
            ):
                continue
            key = (slide.ordinal, str(record["slot_id"]))
            if key in cleared_decoration_keys:
                continue
            if key in expected_records:
                raise PhysicalAssemblyError(
                    f"GOVERNED_CONTENT_DUPLICATE_SLOT_ID: {slide.ordinal}:{record['slot_id']}"
                )
            expected_records[key] = (slide, record)
    expected_slot_count = len(expected_records)

    embedded_evidence = [
        item for item in binding_evidence if item.binding_kind == "embedded"
    ]
    evidence_by_key: dict[tuple[int, str], BindingEvidence] = {}
    duplicate_evidence_count = 0
    for item in embedded_evidence:
        key = (item.ordinal, item.slot_id)
        if key in evidence_by_key:
            duplicate_evidence_count += 1
        else:
            evidence_by_key[key] = item
    unauthorized_content_count = (
        len(set(expected_records) - set(evidence_by_key))
        + len(set(evidence_by_key) - set(expected_records))
        + duplicate_evidence_count
        + sum(1 for item in embedded_evidence if item.status != "pass")
    )

    verified_governed_content_count = 0
    governed_content_mismatch_count = 0
    peer_group_mismatch_count = 0
    mutation_manifest: list[dict[str, Any]] = []
    peer_values: dict[tuple[int, str], set[tuple[str, str]]] = {}
    tag_part_count = 0
    tag_relationship_count = 0
    cached_field_count = 0
    certified_media_count = 0
    media_hash_mismatch_count = 0
    replacement_asset_count = 0
    replacement_asset_hash_mismatch_count = 0
    asset_slot_mismatch_count = 0
    orphan_media_count = 0

    try:
        with zipfile.ZipFile(output_path, "r") as archive:
            names = tuple(_normalise_zip_name(name) for name in archive.namelist())
            name_set = set(names)
            tag_part_count = sum(1 for name in names if "/tags/" in name.lower())

            relationship_targets: set[str] = set()
            for rels_path in (name for name in names if name.endswith(".rels")):
                rels_xml = archive.read(rels_path)
                for entry in _parse_relationships(rels_xml):
                    if _relationship_is_discarded_source_metadata(entry):
                        tag_relationship_count += 1
                    if _relationship_is_external(entry):
                        continue
                    target = _resolve_rel_target(
                        rels_xml,
                        rels_path,
                        entry.get("Target", ""),
                    )
                    if target is not None:
                        relationship_targets.add(target)

            media_names = {
                name
                for name in names
                if name.startswith("ppt/") and "/media/" in name
            }
            orphan_media_count = len(media_names - relationship_targets)
            media_hashes = {
                _sha256_bytes(archive.read(name)) for name in media_names
            }

            for name in names:
                if not (
                    name.startswith("ppt/")
                    and (
                        "/slideLayouts/" in name
                        or "/slideMasters/" in name
                    )
                    and name.endswith(".xml")
                ):
                    continue
                try:
                    root = ET.fromstring(archive.read(name))
                except ET.ParseError:
                    cached_field_count += 1
                    continue
                for field_node in root.iter():
                    if field_node.tag.rsplit("}", 1)[-1] != "fld":
                        continue
                    if any(
                        node.tag.rsplit("}", 1)[-1] == "t"
                        and bool((node.text or "").strip())
                        for node in field_node.iter()
                    ):
                        cached_field_count += 1

            for key, (slide, record) in sorted(expected_records.items()):
                evidence = evidence_by_key.get(key)
                source_part = str(record.get("source_part", ""))
                planned_target_part = (
                    context.slide_target_maps.get(slide.ordinal, {}).get(source_part)
                    if context is not None
                    else None
                )
                if evidence is None:
                    governed_content_mismatch_count += 1
                    continue
                try:
                    relationship_lineage = _resolve_output_governed_target(
                        archive,
                        ordinal=slide.ordinal,
                        record=record,
                    )
                    target_part = str(relationship_lineage["target_part"])
                    if (
                        target_part not in name_set
                        or (
                            planned_target_part is not None
                            and planned_target_part != target_part
                        )
                        or evidence.shape_id != relationship_lineage["shape_id"]
                    ):
                        raise PhysicalAssemblyError(
                            "GOVERNED_OUTPUT_RELATIONSHIP_LINEAGE_MISMATCH"
                        )
                    actual_value = _read_governed_output_slot(
                        archive.read(target_part),
                        record,
                    )
                except (KeyError, PhysicalAssemblyError):
                    governed_content_mismatch_count += 1
                    continue
                actual_sha = _sha256_bytes(actual_value.encode("utf-8"))
                if actual_sha != evidence.replacement_sha256:
                    governed_content_mismatch_count += 1
                    continue
                verified_governed_content_count += 1
                peer_group_id = record.get("peer_group_id")
                if isinstance(peer_group_id, str) and peer_group_id:
                    peer_values.setdefault(
                        (slide.ordinal, peer_group_id),
                        set(),
                    ).add(_governed_peer_value(actual_value))
                mutation_manifest.append(
                    {
                        "ordinal": slide.ordinal,
                        "page_id": slide.page_template.page_id,
                        "slot_id": key[1],
                        "kind": str(record.get("kind", "")),
                        "source_part": source_part,
                        "slide_part": relationship_lineage["slide_part"],
                        "shape_id": relationship_lineage["shape_id"],
                        "slide_relationship_id": relationship_lineage[
                            "slide_relationship_id"
                        ],
                        "chart_part": relationship_lineage["chart_part"],
                        "chart_relationship_id": relationship_lineage[
                            "chart_relationship_id"
                        ],
                        "target_part": target_part,
                        "target_part_sha256": relationship_lineage[
                            "target_part_sha256"
                        ],
                        "locator": str(record.get("locator", "")),
                        "actual_sha256": actual_sha,
                        "peer_group_id": peer_group_id or "",
                    }
                )
            peer_group_mismatch_count = sum(
                1 for values in peer_values.values() if len(values) != 1
            )

            certified_media: dict[tuple[int, str, str], Mapping[str, Any]] = {}
            for slide in plan.target_slides:
                metadata = slide.page_template.governed_content_inventory.get(
                    "closure_metadata",
                    {},
                )
                if not isinstance(metadata, Mapping):
                    continue
                for record in metadata.get("media_parts", ()):
                    if not isinstance(record, Mapping):
                        continue
                    source_part = record.get("source_part")
                    if isinstance(source_part, str):
                        certified_media.setdefault(
                            (
                                slide.ordinal,
                                slide.page_template.package_sha256,
                                source_part,
                            ),
                            record,
                        )
            certified_media_count = len(certified_media)
            for (ordinal, _package_sha, source_part), record in certified_media.items():
                if context is None:
                    media_hash_mismatch_count += 1
                    continue
                target = context.slide_target_maps.get(ordinal, {}).get(source_part)
                if target is None:
                    media_hash_mismatch_count += 1
                    continue
                if target in context.pruned_parts:
                    if source_part not in context.replaced_source_parts_by_slide.get(
                        ordinal, set()
                    ):
                        media_hash_mismatch_count += 1
                    continue
                if target not in name_set:
                    media_hash_mismatch_count += 1
                    continue
                target_bytes = archive.read(target)
                if (
                    _sha256_bytes(target_bytes) != record.get("sha256")
                    or len(target_bytes) != record.get("size_bytes")
                ):
                    media_hash_mismatch_count += 1

            replacement_assets = [
                item for item in binding_evidence if item.binding_kind == "asset"
            ]
            replacement_asset_count = len(replacement_assets)
            for item in replacement_assets:
                if item.status != "pass" or item.replacement_sha256 not in media_hashes:
                    replacement_asset_hash_mismatch_count += 1
                slide_part = f"ppt/slides/slide{item.ordinal}.xml"
                slide_rels_part = _rels_path_for_part(slide_part)
                if slide_part not in name_set or slide_rels_part not in name_set:
                    asset_slot_mismatch_count += 1
                    continue
                try:
                    picture_slots = _discover_picture_slots(archive.read(slide_part))
                    slot = picture_slots.get(item.slot_id)
                    relations = {
                        relation["Id"]: relation
                        for relation in _parse_relationships(
                            archive.read(slide_rels_part)
                        )
                    }
                except (KeyError, PhysicalAssemblyError):
                    asset_slot_mismatch_count += 1
                    continue
                relation = relations.get(item.relationship_id)
                target = (
                    _resolve_rel_target(
                        archive.read(slide_rels_part),
                        slide_rels_part,
                        relation["Target"],
                    )
                    if relation is not None
                    else None
                )
                if (
                    slot != (item.shape_id, item.relationship_id)
                    or relation is None
                    or not relation.get("Type", "").lower().endswith("/image")
                    or target != item.target_part
                    or target not in name_set
                ):
                    asset_slot_mismatch_count += 1
                    continue
                if _sha256_bytes(archive.read(target)) != item.replacement_sha256:
                    replacement_asset_hash_mismatch_count += 1
    except (OSError, zipfile.BadZipFile):
        unauthorized_content_count = max(1, unauthorized_content_count)
        governed_content_mismatch_count = max(1, governed_content_mismatch_count)
        tag_part_count = max(1, tag_part_count)
        tag_relationship_count = max(1, tag_relationship_count)
        cached_field_count = max(1, cached_field_count)
        media_hash_mismatch_count = max(1, media_hash_mismatch_count)
        orphan_media_count = max(1, orphan_media_count)

    mutation_manifest_sha256 = _sha256_bytes(
        json.dumps(
            mutation_manifest,
            ensure_ascii=True,
            sort_keys=True,
            separators=(",", ":"),
        ).encode("utf-8")
    )
    status = (
        "pass"
        if (
            len(embedded_evidence) == expected_slot_count
            and verified_governed_content_count == expected_slot_count
            and unauthorized_content_count == 0
            and governed_content_mismatch_count == 0
            and peer_group_mismatch_count == 0
            and tag_part_count == 0
            and tag_relationship_count == 0
            and cached_field_count == 0
            and media_hash_mismatch_count == 0
            and replacement_asset_hash_mismatch_count == 0
            and asset_slot_mismatch_count == 0
            and orphan_media_count == 0
        )
        else "fail"
    )
    return SourceResidueEvidence(
        governed_content_slot_count=expected_slot_count,
        governed_content_binding_count=len(embedded_evidence),
        verified_governed_content_count=verified_governed_content_count,
        governed_content_mismatch_count=governed_content_mismatch_count,
        peer_group_mismatch_count=peer_group_mismatch_count,
        mutation_manifest_sha256=mutation_manifest_sha256,
        governed_mutations=tuple(mutation_manifest),
        unauthorized_content_count=unauthorized_content_count,
        tag_part_count=tag_part_count,
        tag_relationship_count=tag_relationship_count,
        layout_master_cached_field_count=cached_field_count,
        certified_media_count=certified_media_count,
        media_hash_mismatch_count=media_hash_mismatch_count,
        replacement_asset_count=replacement_asset_count,
        replacement_asset_hash_mismatch_count=replacement_asset_hash_mismatch_count,
        asset_slot_mismatch_count=asset_slot_mismatch_count,
        orphan_media_count=orphan_media_count,
        status=status,
    )


def verify_physical_assembly(
    output_path: str | os.PathLike[str],
    *,
    plan: AssemblyPlan,
    _import_context: AssemblyImportContext | None = None,
    _authority_evidence: AuthorityEvidence | None = None,
    _selection_authority_evidence: SelectionAuthorityEvidence | None = None,
    _binding_evidence: tuple[BindingEvidence, ...] | None = None,
    require_libreoffice: bool = False,
    max_output_size_bytes: int | None = None,
    acceptance_profile: str = "standard",
    expected_slide_count: int | None = None,
    project_root: str | os.PathLike[str] | None = None,
    library_index: LibraryIndex | None = None,
) -> PhysicalAssemblyReport:
    """Verify the assembled PPTX against the assembly plan."""

    if acceptance_profile not in {"standard", "phase49-work-report-15"}:
        raise PhysicalAssemblyError(
            f"UNKNOWN_ACCEPTANCE_PROFILE: {acceptance_profile}"
        )
    if acceptance_profile == "phase49-work-report-15":
        if expected_slide_count not in {None, 15}:
            raise PhysicalAssemblyError(
                "PHASE49_EXPECTED_SLIDE_COUNT_MUST_BE_15"
            )
        expected_slide_count = 15
        _validate_phase49_sequence(plan)
    binding_profile_authority, style_clone_specs = (
        _resolve_binding_profile_style_authority(
            plan,
            acceptance_profile=acceptance_profile,
        )
    )

    path = Path(output_path).expanduser().resolve(strict=False)
    if not path.is_file():
        raise PhysicalAssemblyError(f"output missing: {path}")

    zip_ok, zip_detail = _verify_zip_open(path)
    ct_ok, ct_detail, ct_count = _verify_content_types(path)
    relationship_audit = _inspect_all_relationships(path)
    rels_ok = relationship_audit.status == "pass"
    rels_detail = relationship_audit.details
    unresolved_count = len(relationship_audit.unresolved_internal_relationships)

    (
        py_ok,
        slide_count,
        text_runs,
        shape_count,
        native_object_count,
        picture_count,
        native_editable_slide_count,
        full_slide_raster_count,
        raster_dominant_slide_count,
        native_editable_coverage,
    ) = _verify_python_pptx(path)
    native_ok = (
        py_ok
        and slide_count > 0
        and native_editable_slide_count == slide_count
        and raster_dominant_slide_count == 0
    )
    editability = Editability(
        native_editable=zip_ok and ct_ok and rels_ok and native_ok,
        python_pptx_open=py_ok,
        slide_count=slide_count,
        text_run_count=text_runs,
        shape_count=shape_count,
        native_object_count=native_object_count,
        picture_count=picture_count,
        native_editable_slide_count=native_editable_slide_count,
        full_slide_raster_count=full_slide_raster_count,
        raster_dominant_slide_count=raster_dominant_slide_count,
        native_editable_coverage=native_editable_coverage,
        status=(
            "pass"
            if (zip_ok and ct_ok and rels_ok and native_ok)
            else "fail"
        ),
    )

    lineage: list[LineageRecord] = []
    style_clone_evidence: tuple[StyleCloneEvidence, ...] = ()
    matches = 0
    source_archives: dict[str, zipfile.ZipFile] = {}
    try:
        with zipfile.ZipFile(path, "r") as archive:
            for slide in plan.target_slides:
                ordinal = slide.ordinal
                target_slide_name = f"ppt/slides/slide{ordinal}.xml"
                slide_xml = archive.read(target_slide_name)
                template = slide.page_template
                source_sha = template.source_sha256
                source_slide = b""
                source_package_verified = False
                source_slide_verified = False
                imported_part_map_sha256 = ""
                imported_part_count = 0
                try:
                    source_path = Path(template.source_path)
                    if _import_context is not None:
                        _, source_graph = _import_context.graph_for(
                            source_path,
                            template.package_sha256,
                            template.slide_number,
                        )
                        source_slide = source_graph.slide_xml
                        source_package_verified = True
                        source_slide_verified = (
                            source_graph.slide_sha == template.source_slide_sha256
                        )
                        imported_part_map_sha256, imported_part_count = (
                            _imported_part_map_evidence(
                                _import_context,
                                slide=slide,
                                graph=source_graph,
                                target_archive=archive,
                            )
                        )
                    elif source_path.is_file():
                        source_key = str(source_path.expanduser().resolve(strict=True))
                        source_package_verified = (
                            _sha256_file(Path(source_key)) == template.package_sha256
                            and template.source_sha256 == template.package_sha256
                        )
                        src_zip = source_archives.get(source_key)
                        if src_zip is None:
                            src_zip = zipfile.ZipFile(source_key, "r")
                            source_archives[source_key] = src_zip
                        source_slide = src_zip.read(
                            f"ppt/slides/slide{template.slide_number}.xml"
                        )
                        source_slide_verified = (
                            _sha256_bytes(source_slide) == template.source_slide_sha256
                        )
                except (OSError, KeyError, zipfile.BadZipFile, PhysicalAssemblyError):
                    source_slide = b""
                    source_package_verified = False
                    source_slide_verified = False
                expected_source_slide = source_slide
                if source_slide and style_clone_specs.get(ordinal):
                    expected_source_slide = _apply_governed_style_clones(
                        source_slide,
                        style_clone_specs[ordinal],
                    )
                source_structure = _slide_structure_signature(expected_source_slide)
                target_structure = _slide_structure_signature(slide_xml)
                structure_match = bool(
                    source_structure
                    and target_structure
                    and source_structure == target_structure
                )
                score = _byte_match_score(expected_source_slide, slide_xml)
                lineage_ok = (
                    zip_ok
                    and ct_ok
                    and rels_ok
                    and source_package_verified
                    and source_slide_verified
                    and structure_match
                    and bool(imported_part_map_sha256)
                    and imported_part_count > 0
                )
                status = "pass" if lineage_ok else "fail"
                if template.style_cluster_id == plan.dominant_style_cluster_id:
                    matches += 1
                lineage.append(
                    LineageRecord(
                        ordinal=ordinal,
                        page_id=template.page_id,
                        package_sha256=template.package_sha256,
                        slide_number=template.slide_number,
                        source_sha256=source_sha,
                        source_slide_sha256=template.source_slide_sha256,
                        target_slide_sha256=_sha256_bytes(slide_xml),
                        source_package_verified=source_package_verified,
                        source_slide_verified=source_slide_verified,
                        structure_signature_source=source_structure,
                        structure_signature_target=target_structure,
                        structure_match=structure_match,
                        imported_part_map_sha256=imported_part_map_sha256,
                        imported_part_count=imported_part_count,
                        narrative_role=slide.narrative_role,
                        title=slide.title,
                        status=status,
                        binding_count=sum(
                            1
                            for item in (_binding_evidence or ())
                            if item.ordinal == slide.ordinal
                        ),
                        byte_match_score=round(score, 6),
                    )
                )
            style_clone_evidence = _verify_style_clone_evidence(
                archive,
                plan=plan,
                specs_by_ordinal=style_clone_specs,
            )
    finally:
        for source_archive in source_archives.values():
            source_archive.close()

    page_ordinals: dict[str, list[int]] = {}
    for record in lineage:
        page_ordinals.setdefault(record.page_id, []).append(record.ordinal)
    duplicate_page_records = tuple(
        {
            "page_id": page_id,
            "ordinals": tuple(sorted(ordinals)),
        }
        for page_id, ordinals in sorted(page_ordinals.items())
        if len(ordinals) > 1
    )
    distinct_page_id_count = len(page_ordinals)

    package_entry_count = 0
    media_count = 0
    package_parts: tuple[str, ...] = ()
    try:
        with zipfile.ZipFile(path, "r") as archive:
            package_parts = tuple(sorted(archive.namelist()))
            package_entry_count = len(package_parts)
            media_count = sum(1 for n in package_parts if "media/" in n)
            parts_by_kind = dict(
                sorted(Counter(_part_kind(name) for name in package_parts).items())
            )
    except (OSError, zipfile.BadZipFile):
        parts_by_kind = {}

    opc = OPCIntegrity(
        zip_open=zip_ok,
        content_types_parsed=ct_ok,
        slide_rels_resolved=rels_ok,
        package_entry_count=package_entry_count,
        media_count=media_count,
        total_relationship_count=relationship_audit.total_relationship_count,
        internal_relationship_count=relationship_audit.internal_relationship_count,
        external_relationship_count=relationship_audit.external_relationship_count,
        unresolved_internal_relationship_count=unresolved_count,
        unresolved_internal_relationships=relationship_audit.unresolved_internal_relationships,
        unsafe_relationship_count=len(relationship_audit.unsafe_relationships),
        unsafe_relationships=relationship_audit.unsafe_relationships,
        status="pass" if (zip_ok and ct_ok and rels_ok) else "fail",
        details="; ".join(filter(None, (zip_detail, ct_detail, rels_detail))),
    )
    adherence = StyleClusterAdherence(
        dominant_style_cluster_id=plan.dominant_style_cluster_id,
        matches=matches,
        total=len(plan.target_slides),
        status="pass" if matches == len(plan.target_slides) else "fail",
    )
    authority_evidence = _authority_evidence or AuthorityEvidence(
        mode="legacy",
        fact_store_path="",
        fact_store_sha256="",
        asset_manifest_path="",
        asset_manifest_sha256="",
        connective_copy_path="",
        connective_copy_sha256="",
        status="legacy",
    )
    if (
        acceptance_profile == "phase49-work-report-15"
        or (plan.query_bundle_path and plan.query_bundle_sha256)
    ):
        selection_authority_evidence = _validate_query_selection_evidence(
            plan,
            project_root=project_root,
            library_index=library_index,
            require_phase49_ordinals=(
                acceptance_profile == "phase49-work-report-15"
            ),
        )
    elif _selection_authority_evidence is not None:
        selection_authority_evidence = _selection_authority_evidence
    else:
        selection_authority_evidence = SelectionAuthorityEvidence(
            mode="not_required",
            query_bundle_path="",
            query_bundle_sha256="",
            library_index_sha256=plan.library_index_sha256,
            query_count=0,
            selected_count=0,
            distinct_query_id_count=0,
            distinct_page_id_count=0,
            status="not_required",
        )
    binding_evidence = _binding_evidence
    if binding_evidence is None:
        binding_evidence = tuple(
            _build_text_binding_evidence(
                plan,
                None,
                None,
                require_locked_authority=False,
            )
        )
    source_residue = _verify_source_residue(
        path,
        plan=plan,
        context=_import_context,
        binding_evidence=binding_evidence,
        required=authority_evidence.mode == "locked",
    )
    if _import_context is not None:
        assembly_metrics = _import_context.metrics(path, unresolved_count)
    else:
        source_sizes: dict[str, int] = {}
        for slide in plan.target_slides:
            source_path = Path(slide.page_template.source_path)
            if source_path.is_file():
                source_sizes.setdefault(
                    slide.page_template.package_sha256,
                    source_path.stat().st_size,
                )
        source_size = sum(source_sizes.values())
        output_size = path.stat().st_size
        dependency_count = max(
            0,
            package_entry_count - len(plan.target_slides) - 10,
        )
        assembly_metrics = AssemblyMetrics(
            output_size_bytes=output_size,
            source_size_bytes=source_size,
            unique_source_package_count=len(source_sizes),
            imported_part_count=len(package_parts),
            imported_parts=package_parts,
            unique_dependency_part_count=dependency_count,
            same_source_reuse_count=0,
            same_source_reuse_bytes=0,
            cross_source_safe_dedup_count=0,
            cross_source_safe_dedup_bytes=0,
            deduplicated_part_count=0,
            deduplicated_bytes=0,
            static_duplicate_bytes=_static_duplicate_bytes(path),
            unresolved_internal_relationship_count=unresolved_count,
            amplification_ratio=round(output_size / source_size, 6) if source_size else 0.0,
            parts_by_kind=parts_by_kind,
        )
    output_size = path.stat().st_size
    if max_output_size_bytes is None:
        size_status = "not_run"
    elif max_output_size_bytes < 1:
        raise PhysicalAssemblyError("max_output_size_bytes must be positive")
    else:
        size_status = "pass" if output_size <= max_output_size_bytes else "fail"
    size_check = SizeCheck(
        output_size_bytes=output_size,
        max_output_size_bytes=max_output_size_bytes,
        status=size_status,
    )
    libreoffice = _verify_libreoffice(path, required=require_libreoffice)
    locked = authority_evidence.mode == "locked"
    core_pass = (
        opc.status == "pass"
        and editability.status == "pass"
        and adherence.status == "pass"
        and source_residue.status == "pass"
        and len(lineage) == plan.target_slide_count
        and all(record.status == "pass" for record in lineage)
        and (
            expected_slide_count is None
            or (
                plan.target_slide_count == expected_slide_count
                and slide_count == expected_slide_count
                and len(lineage) == expected_slide_count
                and distinct_page_id_count == expected_slide_count
            )
        )
    )
    requested_checks_pass = (
        (not require_libreoffice or libreoffice.status == "pass")
        and (max_output_size_bytes is None or size_check.status == "pass")
    )
    locked_checks_pass = (
        not locked
        or (
            authority_evidence.status == "pass"
            and libreoffice.status == "pass"
            and size_check.status == "pass"
            and distinct_page_id_count == plan.target_slide_count
            and not duplicate_page_records
            and assembly_metrics.static_duplicate_bytes == 0
            and source_residue.status == "pass"
        )
    )
    selection_checks_pass = (
        acceptance_profile != "phase49-work-report-15"
        or (
            selection_authority_evidence.mode == "locked"
            and selection_authority_evidence.status == "pass"
            and selection_authority_evidence.query_count == 15
            and selection_authority_evidence.selected_count == 15
            and selection_authority_evidence.distinct_query_id_count == 15
            and selection_authority_evidence.distinct_page_id_count == 15
        )
    )
    profile_checks_pass = (
        (
            binding_profile_authority is None
            and not style_clone_specs
            and not style_clone_evidence
        )
        if acceptance_profile == "standard"
        else (
            binding_profile_authority is not None
            and binding_profile_authority.status == "pass"
            and binding_profile_authority.style_clone_count
            == len(style_clone_evidence)
            == 4
            and all(item.status == "pass" for item in style_clone_evidence)
        )
    )
    overall = (
        "pass"
        if core_pass
        and requested_checks_pass
        and locked_checks_pass
        and selection_checks_pass
        and profile_checks_pass
        else "fail"
    )
    return PhysicalAssemblyReport(
        schema_version="1.0",
        report_id=f"par_{int(datetime.now(timezone.utc).timestamp() * 1000):x}",
        plan_id=plan.plan_id,
        binding_profile_authority=binding_profile_authority,
        output_path=str(path),
        output_sha256=_sha256_bytes(path.read_bytes()),
        acceptance_profile=acceptance_profile,
        expected_slide_count=expected_slide_count,
        status=overall,
        target_slide_count=len(plan.target_slides),
        distinct_page_id_count=distinct_page_id_count,
        duplicate_page_records=duplicate_page_records,
        lineage_records=tuple(lineage),
        opc_integrity=opc,
        editability=editability,
        style_cluster_adherence=adherence,
        assembly_metrics=assembly_metrics,
        authority=authority_evidence,
        selection_authority=selection_authority_evidence,
        binding_evidence=binding_evidence,
        style_clone_evidence=style_clone_evidence,
        source_residue=source_residue,
        libreoffice=libreoffice,
        size_check=size_check,
    )


def assemble_physical_deck(
    plan: AssemblyPlan,
    output_path: str | os.PathLike[str],
    *,
    library_index_sha256: str,
    fact_store_path: str | os.PathLike[str] | None = None,
    fact_store_sha256: str | None = None,
    asset_manifest_path: str | os.PathLike[str] | None = None,
    asset_manifest_sha256: str | None = None,
    connective_copy_path: str | os.PathLike[str] | None = None,
    connective_copy_sha256: str | None = None,
    project_root: str | os.PathLike[str] | None = None,
    require_locked_authority: bool = False,
    require_libreoffice: bool = False,
    max_output_size_bytes: int | None = None,
    acceptance_profile: str = "standard",
    expected_slide_count: int | None = None,
    library_index: LibraryIndex | None = None,
) -> PhysicalAssemblyReport:
    """Assemble a deck through one deterministic, deduplicating OPC context."""

    if acceptance_profile not in {"standard", "phase49-work-report-15"}:
        raise PhysicalAssemblyError(
            f"UNKNOWN_ACCEPTANCE_PROFILE: {acceptance_profile}"
        )
    if acceptance_profile == "phase49-work-report-15":
        if expected_slide_count not in {None, 15}:
            raise PhysicalAssemblyError(
                "PHASE49_EXPECTED_SLIDE_COUNT_MUST_BE_15"
            )
        expected_slide_count = 15
        _validate_phase49_sequence(plan)
    _, style_clone_specs = _resolve_binding_profile_style_authority(
        plan,
        acceptance_profile=acceptance_profile,
    )

    authority_arguments = (
        fact_store_path,
        fact_store_sha256,
        asset_manifest_path,
        asset_manifest_sha256,
        connective_copy_path,
        connective_copy_sha256,
    )
    locked_requested = (
        acceptance_profile == "phase49-work-report-15"
        or require_locked_authority
        or any(
            value is not None for value in authority_arguments
        )
    )
    _validate_assembly_plan(
        plan,
        library_index_sha256,
        require_locked_authority=locked_requested,
        expected_slide_count=expected_slide_count,
    )
    selection_authority_evidence: SelectionAuthorityEvidence | None = None
    if (
        acceptance_profile == "phase49-work-report-15"
        or (plan.query_bundle_path and plan.query_bundle_sha256)
    ):
        selection_authority_evidence = _validate_query_selection_evidence(
            plan,
            project_root=project_root,
            library_index=library_index,
            require_phase49_ordinals=(
                acceptance_profile == "phase49-work-report-15"
            ),
        )
    if locked_requested:
        fact_store, locked_assets, connective_copy, authority_evidence = _validate_locked_authority(
            plan,
            fact_store_path=fact_store_path,
            fact_store_sha256=fact_store_sha256,
            asset_manifest_path=asset_manifest_path,
            asset_manifest_sha256=asset_manifest_sha256,
            connective_copy_path=connective_copy_path,
            connective_copy_sha256=connective_copy_sha256,
            project_root=project_root,
        )
    else:
        if any(slide.asset_binding_specs for slide in plan.target_slides):
            raise PhysicalAssemblyError(
                "asset replacements require locked authority"
            )
        fact_store = None
        locked_assets = {}
        connective_copy = {}
        authority_evidence = AuthorityEvidence(
            mode="legacy",
            fact_store_path="",
            fact_store_sha256="",
            asset_manifest_path="",
            asset_manifest_sha256="",
            connective_copy_path="",
            connective_copy_sha256="",
            status="legacy",
        )
    for slide in plan.target_slides:
        if (
            locked_requested
            and slide.page_template.requires_customer_asset
            and not slide.asset_binding_specs
        ):
            raise PhysicalAssemblyError(
                f"ASSEMBLY_PLAN_CUSTOMER_ASSET_BINDING_REQUIRED: {slide.page_template.page_id}"
            )
    output = Path(output_path).expanduser().resolve(strict=False)
    if output.exists():
        raise PhysicalAssemblyError(f"OUTPUT_ALREADY_EXISTS: {output}")
    output.parent.mkdir(parents=True, exist_ok=True)
    context = AssemblyImportContext()
    slide_paths: list[str] = []
    imported_master_paths: set[str] = set()
    imported_notes_master_paths: set[str] = set()

    try:
        source_texts: dict[tuple[int, str], str] = {}
        for slide in plan.target_slides:
            _, graph = context.graph_for(
                Path(slide.page_template.source_path),
                slide.page_template.package_sha256,
                slide.page_template.slide_number,
            )
            if graph.slide_sha != slide.page_template.source_slide_sha256:
                raise PhysicalAssemblyError(
                    f"source slide fingerprint mismatch: {slide.page_template.page_id}"
                )
            for slot in _library_discover_slots(
                graph.slide_xml.decode("utf-8", errors="replace")
            ):
                source_texts[(slide.ordinal, slot.slot_id)] = slot.text
        binding_evidence = _build_text_binding_evidence(
            plan,
            fact_store,
            connective_copy,
            require_locked_authority=locked_requested,
            source_texts=source_texts,
            enforce_required_facts=not locked_requested,
            fragment_group_contracts=(
                selection_authority_evidence.fragment_group_contracts
                if selection_authority_evidence is not None
                else ()
            ),
        )
        # Target slide names are fixed before any source dependency is allocated.
        target_slide_names = {
            slide.ordinal: f"ppt/slides/slide{slide.ordinal}.xml"
            for slide in plan.target_slides
        }
        for slide in sorted(plan.target_slides, key=lambda item: item.ordinal):
            source_path = Path(slide.page_template.source_path)
            source, graph = context.graph_for(
                source_path,
                slide.page_template.package_sha256,
                slide.page_template.slide_number,
            )
            target_slide_name = target_slide_names[slide.ordinal]
            target_map: dict[str, str] = {
                graph.root_slide_name: target_slide_name,
            }
            governed_mutations, governed_evidence = (
                _prepare_governed_content_replacements(
                    slide,
                    graph,
                    fact_store,
                    connective_copy,
                )
                if locked_requested
                else ({}, [])
            )
            if not locked_requested:
                for source_part in sorted(graph.extra_parts):
                    if source_part.lower().endswith(".xlsm"):
                        raise PhysicalAssemblyError(
                            f"GOVERNED_WORKBOOK_XLSM_FORBIDDEN: {source_part}"
                        )
                    if source_part.lower().endswith(".xlsx"):
                        governed_mutations[source_part] = (
                            _mutate_governed_workbook(
                                graph.extra_parts[source_part],
                                (),
                            )
                        )
            binding_evidence.extend(governed_evidence)
            closure_metadata = slide.page_template.governed_content_inventory.get(
                "closure_metadata",
                {},
            )
            field_records_by_part: dict[str, list[Mapping[str, Any]]] = {}
            if isinstance(closure_metadata, Mapping):
                for field_record in closure_metadata.get("layout_master_fields", ()):
                    if isinstance(field_record, Mapping) and isinstance(
                        field_record.get("source_part"),
                        str,
                    ):
                        field_records_by_part.setdefault(
                            str(field_record["source_part"]),
                            [],
                        ).append(field_record)

            # Allocate the complete dependency closure before rewriting any rels.
            for source_part in sorted(graph.extra_parts):
                data = governed_mutations.get(
                    source_part,
                    graph.extra_parts[source_part],
                )
                if locked_requested and (
                    _LAYOUT_RE.match(source_part) or _MASTER_RE.match(source_part)
                ):
                    data = _sanitize_layout_master_fields(
                        source_part,
                        data,
                        field_records_by_part.get(source_part, ()),
                    )
                content_type = graph.content_types[source_part]
                mutation_scope = (
                    f"slide_{slide.ordinal:03d}"
                    if (
                        _is_slide_scoped_part(source_part)
                        or source_part in governed_mutations
                    )
                    else "shared"
                )
                target_map[source_part] = context.allocate_dependency(
                    source,
                    source_part,
                    data,
                    content_type,
                    mutation_scope=mutation_scope,
                    relationship_free=_rels_path_for_part(source_part) not in graph.rels,
                )

            governed_slide_xml = governed_mutations.get(
                graph.root_slide_name,
                graph.slide_xml,
            )
            governed_slide_xml = _apply_governed_style_clones(
                governed_slide_xml,
                style_clone_specs.get(slide.ordinal, ()),
            )
            slide_bytes = _adapt_slide_text(
                governed_slide_xml,
                slide.bindings,
                allowed_slots=slide.page_template.slot_graph.get("text_slot_ids", ()),
                allowed_clear_alias_slots={
                    slot_id
                    for slot_id, spec in slide.text_binding_specs.items()
                    if (
                        slot_id not in slide.page_template.slot_graph.get("text_slot_ids", ())
                        and spec.mode == "clear"
                        and spec.replacement == ""
                        and not spec.fact_refs
                    )
                },
                fit_policies={
                    slot_id: spec.fit_policy
                    for slot_id, spec in slide.text_binding_specs.items()
                },
            )
            relationship_overrides, cover_crops, asset_evidence = _prepare_asset_replacements(
                slide,
                graph,
                locked_assets,
                context,
            )
            slide_bytes = _apply_picture_cover_crops(slide_bytes, cover_crops)
            context.add_part(
                target_slide_name,
                slide_bytes,
                graph.content_types[graph.root_slide_name],
            )
            context.imported_parts.add(target_slide_name)
            binding_evidence.extend(asset_evidence)

            # Relationships are rewritten relative to their actual owner part.
            for source_rels_path in sorted(graph.rels):
                source_owner = _owner_part_from_rels_path(source_rels_path)
                if source_owner is None or source_owner not in target_map:
                    raise PhysicalAssemblyError(
                        f"relationship owner was not allocated: {source_rels_path}"
                    )
                target_rels_path = _rels_path_for_part(target_map[source_owner])
                rewritten = _rewrite_relationship_targets(
                    graph.rels[source_rels_path],
                    source_rels_path,
                    target_map,
                    output_rels_path=target_rels_path,
                    relationship_overrides=(
                        relationship_overrides
                        if source_rels_path == _rels_path_for_part(graph.root_slide_name)
                        else None
                    ),
                )
                context.add_part(
                    target_rels_path,
                    rewritten,
                    "application/vnd.openxmlformats-package.relationships+xml",
                )
                context.imported_parts.add(target_rels_path)

            for source_master in graph.master_paths:
                imported_master_paths.add(target_map[source_master])
            for source_part, target_part in target_map.items():
                if "/notesMasters/" in source_part and source_part.endswith(".xml"):
                    imported_notes_master_paths.add(target_part)
            context.slide_target_maps[slide.ordinal] = dict(target_map)
            slide_paths.append(f"slides/slide{slide.ordinal}.xml")

        if locked_requested and fact_store is not None:
            required_fact_refs = {
                fact.id for fact in fact_store.active_facts() if fact.required
            }
            consumed_fact_refs = {
                fact_ref
                for item in binding_evidence
                for fact_ref in item.fact_refs
            }
            missing_required = sorted(required_fact_refs - consumed_fact_refs)
            if missing_required:
                raise PhysicalAssemblyError(
                    "REQUIRED_FACTS_NOT_BOUND: " + ",".join(missing_required)
                )
        master_paths = sorted(
            path[len("ppt/"):] if path.startswith("ppt/") else path
            for path in imported_master_paths
        )
        notes_master_paths = sorted(
            path[len("ppt/"):] if path.startswith("ppt/") else path
            for path in imported_notes_master_paths
        )
        context.add_part(
            "ppt/presentation.xml",
            _default_pres_xml(slide_paths, master_paths, notes_master_paths),
            CONTENT_TYPES["ppt"],
        )
        context.add_part(
            "ppt/_rels/presentation.xml.rels",
            _pres_xml_rels(slide_paths, master_paths, notes_master_paths),
            "application/vnd.openxmlformats-package.relationships+xml",
        )
        context.add_part(
            "ppt/presProps.xml",
            _default_props_xml(),
            "application/vnd.openxmlformats-officedocument.presentationml.presProps+xml",
        )
        context.add_part(
            "ppt/viewProps.xml",
            _default_view_props_xml(),
            "application/vnd.openxmlformats-officedocument.presentationml.viewProps+xml",
        )
        context.add_part(
            "ppt/tableStyles.xml",
            _default_table_styles_xml(),
            "application/vnd.openxmlformats-officedocument.presentationml.tableStyles+xml",
        )
        context.add_part(
            "_rels/.rels",
            _default_pres_rels(),
            "application/vnd.openxmlformats-package.relationships+xml",
        )
        core_xml, app_xml = _empty_doc_props()
        context.add_part(
            "docProps/core.xml",
            core_xml,
            "application/vnd.openxmlformats-package.core-properties+xml",
        )
        context.add_part(
            "docProps/app.xml",
            app_xml,
            "application/vnd.openxmlformats-officedocument.extended-properties+xml",
        )

        _prune_unreachable_parts(context)

        # Exact source content types are carried as target overrides; generated
        # parts use their canonical Open XML types.
        context.parts["[Content_Types].xml"] = _serialize_content_types(
            sorted(context.content_type_defaults.items()),
            sorted(context.content_type_overrides.items()),
        )

        fd, candidate_name = tempfile.mkstemp(
            prefix=f".{output.stem}.",
            suffix=".candidate.pptx",
            dir=output.parent,
        )
        os.close(fd)
        candidate = Path(candidate_name)
        try:
            with zipfile.ZipFile(candidate, "w", compression=zipfile.ZIP_DEFLATED) as out_zip:
                for part_name in sorted(context.parts):
                    info = zipfile.ZipInfo(part_name, date_time=(1980, 1, 1, 0, 0, 0))
                    info.compress_type = zipfile.ZIP_DEFLATED
                    info.external_attr = 0o600 << 16
                    out_zip.writestr(info, context.parts[part_name])
            report = verify_physical_assembly(
                candidate,
                plan=plan,
                _import_context=context,
                _authority_evidence=authority_evidence,
                _selection_authority_evidence=selection_authority_evidence,
                _binding_evidence=tuple(binding_evidence),
                require_libreoffice=require_libreoffice,
                max_output_size_bytes=max_output_size_bytes,
                acceptance_profile=acceptance_profile,
                expected_slide_count=expected_slide_count,
                project_root=project_root,
                library_index=library_index,
            )
            if report.status != "pass":
                candidate.unlink(missing_ok=True)
                return replace(report, output_path=str(output))
            os.replace(candidate, output)
            return replace(report, output_path=str(output))
        except Exception:
            candidate.unlink(missing_ok=True)
            raise
    finally:
        context.close()


def _infer_content_type(name: str) -> str | None:
    if "/slideLayouts/" in name:
        return CONTENT_TYPES["slideLayout"]
    if "/slideMasters/" in name:
        return CONTENT_TYPES["slideMaster"]
    if "/theme/" in name:
        return CONTENT_TYPES["theme"]
    if "/notesSlides/" in name:
        return CONTENT_TYPES["notesSlide"]
    if "/notesMasters/" in name:
        return CONTENT_TYPES["notesMaster"]
    if "/charts/" in name:
        return CONTENT_TYPES["chart"]
    if "/embeddings/" in name and name.endswith(".xlsx"):
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    return None


def _default_content_type_for(ext: str) -> str:
    return {
        "png": "image/png",
        "jpeg": "image/jpeg",
        "jpg": "image/jpeg",
        "gif": "image/gif",
        "wdp": "image/vnd.ms-photo",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "xml": "application/xml",
        "rels": "application/vnd.openxmlformats-package.relationships+xml",
    }.get(ext, "application/octet-stream")


def write_assembly_report(
    report: PhysicalAssemblyReport, output_path: str | os.PathLike[str]
) -> str:
    path = Path(output_path).expanduser().resolve(strict=False)
    path.parent.mkdir(parents=True, exist_ok=True)
    payload = report.to_dict()
    text = json.dumps(payload, ensure_ascii=False, indent=2)
    path.write_text(text, encoding="utf-8")
    return _sha256_bytes(text.encode("utf-8"))


def load_assembly_plan(
    path: str | os.PathLike[str],
    library_lookup: Mapping[str, PageTemplate],
    *,
    project_root: str | os.PathLike[str] | None = None,
) -> AssemblyPlan:
    if project_root is not None:
        root = Path(project_root).expanduser().resolve(strict=True)
        raw_plan_path = Path(path).expanduser()
        if raw_plan_path.is_absolute():
            try:
                relative_plan_path = raw_plan_path.relative_to(root)
            except ValueError as exc:
                raise PhysicalAssemblyError(
                    f"ASSEMBLY_PLAN_PATH_ESCAPE: {raw_plan_path}"
                ) from exc
        else:
            relative_plan_path = raw_plan_path
        plan_path = resolve_project_file(
            relative_plan_path,
            root,
            label="ASSEMBLY_PLAN",
        )
    else:
        plan_path = Path(path).expanduser().resolve(strict=False)
    payload = json.loads(plan_path.read_text(encoding="utf-8"))
    if not isinstance(payload, Mapping):
        raise PhysicalAssemblyError("ASSEMBLY_PLAN_SCHEMA_INVALID: root must be an object")
    _validate_schema_payload(
        payload,
        "assembly-plan.v1.schema.json",
        label="ASSEMBLY_PLAN",
    )
    if payload.get("schema_version") != "1.0":
        raise PhysicalAssemblyError(
            f"unsupported assembly plan schema_version: {payload.get('schema_version')}"
        )
    return _plan_from_payload(
        payload,
        library_lookup,
        base_dir=plan_path.parent,
        project_root=(
            Path(project_root).expanduser().resolve(strict=True)
            if project_root is not None
            else None
        ),
    )


__all__ = [
    "AssetBindingSpec",
    "AssemblyImportContext",
    "AssemblyMetrics",
    "AssemblyPlan",
    "AssemblyTargetSlide",
    "AuthorityEvidence",
    "AuthorityLock",
    "BindingEvidence",
    "BindingProfileAuthorityEvidence",
    "BindingProfileAuthorityLock",
    "Editability",
    "LineageRecord",
    "LockedAsset",
    "OPCIntegrity",
    "PhysicalAssemblyError",
    "PhysicalAssemblyReport",
    "RelationshipAudit",
    "SlideBinding",
    "StyleClusterAdherence",
    "StyleCloneEvidence",
    "StyleCloneSpec",
    "SizeCheck",
    "TextBindingSpec",
    "assemble_physical_deck",
    "load_assembly_plan",
    "resolve_project_file",
    "verify_physical_assembly",
    "write_assembly_report",
    "_apply_governed_style_clones",
    "_style_clone_scope_sha256",
    "_style_clone_target_guard_sha256",
]
