"""Strict, independent PresentationML topology and editability inspection.

This module deliberately does not import the physical-report validator.  It
can therefore be used by that validator as a small fail-closed boundary around
the presentation slide list, its relationships, and python-pptx's view of the
resulting package.
"""

from __future__ import annotations

from collections import Counter
from dataclasses import dataclass
from pathlib import Path
import posixpath
import re
import unicodedata
from typing import Any, Iterable, Sequence
from urllib.parse import unquote
import zipfile
import xml.etree.ElementTree as ET

from . import template_geometry as _geometry
from .independent_validation_security import (
    PPTX_ZIP_RESOURCE_LIMITS,
    ZipResourceLimits,
    audit_zip_resources,
)


PRESENTATIONML_NS = (
    "http://schemas.openxmlformats.org/presentationml/2006/main"
)
DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
OFFICE_RELATIONSHIP_NS = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
)
PACKAGE_RELATIONSHIP_NS = (
    "http://schemas.openxmlformats.org/package/2006/relationships"
)
SLIDE_RELATIONSHIP_TYPE = f"{OFFICE_RELATIONSHIP_NS}/slide"

PRESENTATION_TAG = f"{{{PRESENTATIONML_NS}}}presentation"
SLIDE_ID_LIST_TAG = f"{{{PRESENTATIONML_NS}}}sldIdLst"
SLIDE_ID_TAG = f"{{{PRESENTATIONML_NS}}}sldId"
SLIDE_TAG = f"{{{PRESENTATIONML_NS}}}sld"
EXTENSION_LIST_TAG = f"{{{PRESENTATIONML_NS}}}extLst"
RELATIONSHIP_ID_ATTR = f"{{{OFFICE_RELATIONSHIP_NS}}}id"
RELATIONSHIPS_TAG = f"{{{PACKAGE_RELATIONSHIP_NS}}}Relationships"
RELATIONSHIP_TAG = f"{{{PACKAGE_RELATIONSHIP_NS}}}Relationship"

PRESENTATION_PART = "ppt/presentation.xml"
PRESENTATION_RELS_PART = "ppt/_rels/presentation.xml.rels"
PACKAGE_ROOT_RELS_PART = "_rels/.rels"
OFFICE_DOCUMENT_RELATIONSHIP_TYPE = f"{OFFICE_RELATIONSHIP_NS}/officeDocument"
SLIDE_PART_RE = re.compile(r"^ppt/slides/slide([1-9][0-9]*)\.xml$")
SLIDE_LIKE_PART_RE = re.compile(r"^ppt/slides/slide[^/]*\.xml$")
URI_SCHEME_RE = re.compile(r"^[A-Za-z][A-Za-z0-9+.-]*:")

# python-pptx MSO_SHAPE_TYPE values that represent ordinary, natively
# editable DrawingML surfaces.  Pictures are counted separately; executable,
# linked, control, media, and unknown surfaces fail closed below.
NATIVE_SHAPE_TYPES = frozenset({1, 2, 3, 5, 9, 14, 15, 17, 19, 20, 21, 22, 24})


@dataclass(frozen=True)
class PresentationTopologyIssue:
    """One stable, serializable strict-topology finding."""

    code: str
    location: str
    detail: str

    def to_dict(self) -> dict[str, str]:
        return {
            "code": self.code,
            "location": self.location,
            "detail": self.detail,
        }


@dataclass(frozen=True)
class RasterThresholds:
    """Thresholds used to derive basic raster-dominance flags."""

    full_slide_picture: float = 0.90
    dominant_picture_union: float = 0.90
    sparse_native_union: float = 0.10
    sparse_text_union: float = 0.04

    def validate(self) -> None:
        for name, value in (
            ("full_slide_picture", self.full_slide_picture),
            ("dominant_picture_union", self.dominant_picture_union),
            ("sparse_native_union", self.sparse_native_union),
            ("sparse_text_union", self.sparse_text_union),
        ):
            if not isinstance(value, (int, float)) or isinstance(value, bool):
                raise TypeError(f"{name} must be numeric")
            if not 0 <= float(value) <= 1:
                raise ValueError(f"{name} must be between 0 and 1")


@dataclass(frozen=True)
class SlideEditabilityStatistics:
    """Per-slide evidence needed for a basic rasterization policy."""

    ordinal: int
    part_name: str
    shape_count: int
    text_run_count: int
    native_object_count: int
    picture_count: int
    picture_coverage: float
    native_coverage: float
    text_coverage: float
    full_slide_raster: bool
    raster_dominant: bool

    def to_dict(self) -> dict[str, Any]:
        return {
            "ordinal": self.ordinal,
            "part_name": self.part_name,
            "shape_count": self.shape_count,
            "text_run_count": self.text_run_count,
            "native_object_count": self.native_object_count,
            "picture_count": self.picture_count,
            "picture_coverage": self.picture_coverage,
            "native_coverage": self.native_coverage,
            "text_coverage": self.text_coverage,
            "full_slide_raster": self.full_slide_raster,
            "raster_dominant": self.raster_dominant,
        }


@dataclass(frozen=True)
class PresentationEditabilityStatistics:
    """Aggregate python-pptx and raster-coverage evidence."""

    slide_count: int
    python_pptx_slide_count: int
    slide_width: int
    slide_height: int
    text_run_count: int
    shape_count: int
    native_object_count: int
    picture_count: int
    native_editable_slide_count: int
    full_slide_raster_count: int
    raster_dominant_slide_count: int
    native_editable_coverage: float
    slides: tuple[SlideEditabilityStatistics, ...]

    def to_dict(self) -> dict[str, Any]:
        return {
            "slide_count": self.slide_count,
            "python_pptx_slide_count": self.python_pptx_slide_count,
            "slide_width": self.slide_width,
            "slide_height": self.slide_height,
            "text_run_count": self.text_run_count,
            "shape_count": self.shape_count,
            "native_object_count": self.native_object_count,
            "picture_count": self.picture_count,
            "native_editable_slide_count": self.native_editable_slide_count,
            "full_slide_raster_count": self.full_slide_raster_count,
            "raster_dominant_slide_count": self.raster_dominant_slide_count,
            "native_editable_coverage": self.native_editable_coverage,
            "slides": [slide.to_dict() for slide in self.slides],
        }


@dataclass(frozen=True)
class PresentationTopologyResult:
    """Complete result from :func:`inspect_presentation_topology`."""

    status: str
    ordered_relationship_ids: tuple[str, ...]
    ordered_slide_parts: tuple[str, ...]
    statistics: PresentationEditabilityStatistics | None
    issues: tuple[PresentationTopologyIssue, ...]

    @property
    def ok(self) -> bool:
        return self.status == "pass"

    def to_dict(self) -> dict[str, Any]:
        return {
            "status": self.status,
            "ordered_relationship_ids": list(self.ordered_relationship_ids),
            "ordered_slide_parts": list(self.ordered_slide_parts),
            "statistics": (
                self.statistics.to_dict() if self.statistics is not None else None
            ),
            "issues": [issue.to_dict() for issue in self.issues],
        }


@dataclass(frozen=True)
class _Relationship:
    relationship_id: str
    relationship_type: str
    target: str
    target_mode: str


class _PackageResourceRejected(RuntimeError):
    """Internal control flow after a metadata-only resource preflight fails."""


class _IssueCollector:
    def __init__(self) -> None:
        self._issues: list[PresentationTopologyIssue] = []
        self._keys: set[tuple[str, str, str]] = set()

    def add(self, code: str, location: str, detail: str) -> None:
        key = (code, location, detail)
        if key in self._keys:
            return
        self._keys.add(key)
        self._issues.append(PresentationTopologyIssue(code, location, detail))

    @property
    def issues(self) -> tuple[PresentationTopologyIssue, ...]:
        return tuple(self._issues)


def _parse_xml(
    raw: bytes,
    *,
    location: str,
    malformed_code: str,
    issues: _IssueCollector,
) -> ET.Element | None:
    try:
        from defusedxml import ElementTree as safe_element_tree
        from defusedxml.common import DefusedXmlException
    except ImportError as exc:
        issues.add(
            "PRESENTATION_XML_SECURITY_PARSER_UNAVAILABLE",
            location,
            str(exc),
        )
        return None
    try:
        return safe_element_tree.fromstring(raw)
    except (ET.ParseError, DefusedXmlException) as exc:
        issues.add(malformed_code, location, str(exc))
        return None


def _presentation_slide_ids(
    root: ET.Element,
    issues: _IssueCollector,
) -> tuple[str, ...]:
    location = f"output.pptx!/{PRESENTATION_PART}"
    if root.tag != PRESENTATION_TAG:
        issues.add(
            "PRESENTATION_ROOT_INVALID",
            location,
            f"expected {PRESENTATION_TAG!r}, observed {root.tag!r}",
        )
        return ()

    all_lists = list(root.iter(SLIDE_ID_LIST_TAG))
    direct_lists = [child for child in list(root) if child.tag == SLIDE_ID_LIST_TAG]
    if len(all_lists) != 1 or len(direct_lists) != 1:
        issues.add(
            "PRESENTATION_SLIDE_ID_LIST_INVALID",
            location,
            (
                "exactly one direct p:sldIdLst is required; "
                f"observed total={len(all_lists)} direct={len(direct_lists)}"
            ),
        )
        return ()

    relationship_ids: list[str] = []
    slide_ids: set[int] = set()
    seen_relationship_ids: set[str] = set()
    for index, node in enumerate(list(direct_lists[0]), start=1):
        item_location = f"{location}#sldId[{index}]"
        if node.tag != SLIDE_ID_TAG:
            issues.add(
                "PRESENTATION_SLIDE_ID_ELEMENT_INVALID",
                item_location,
                f"unexpected element {node.tag!r}",
            )
            continue
        unexpected_attributes = set(node.attrib) - {"id", RELATIONSHIP_ID_ATTR}
        if unexpected_attributes:
            issues.add(
                "PRESENTATION_SLIDE_ID_ELEMENT_INVALID",
                item_location,
                "unexpected attributes: " + ",".join(sorted(unexpected_attributes)),
            )
        if any(child.tag != EXTENSION_LIST_TAG for child in list(node)):
            issues.add(
                "PRESENTATION_SLIDE_ID_ELEMENT_INVALID",
                item_location,
                "p:sldId may only contain p:extLst",
            )
        raw_slide_id = node.attrib.get("id", "")
        relationship_id = node.attrib.get(RELATIONSHIP_ID_ATTR, "")
        if not raw_slide_id.isdigit() or int(raw_slide_id) < 256:
            issues.add(
                "PRESENTATION_SLIDE_ID_VALUE_INVALID",
                item_location,
                raw_slide_id or "missing id",
            )
        elif int(raw_slide_id) in slide_ids:
            issues.add(
                "PRESENTATION_SLIDE_ID_VALUE_DUPLICATE",
                item_location,
                raw_slide_id,
            )
        else:
            slide_ids.add(int(raw_slide_id))
        if not relationship_id or relationship_id != relationship_id.strip():
            issues.add(
                "PRESENTATION_SLIDE_RELATIONSHIP_ID_INVALID",
                item_location,
                relationship_id or "missing r:id",
            )
            continue
        if relationship_id in seen_relationship_ids:
            issues.add(
                "PRESENTATION_SLIDE_RELATIONSHIP_ID_DUPLICATE",
                item_location,
                relationship_id,
            )
        else:
            seen_relationship_ids.add(relationship_id)
        relationship_ids.append(relationship_id)

    if not relationship_ids:
        issues.add(
            "PRESENTATION_SLIDE_LIST_EMPTY",
            location,
            "at least one p:sldId is required",
        )
    return tuple(relationship_ids)


def _strict_relationships(
    raw: bytes,
    *,
    part_name: str,
    code_prefix: str,
    issues: _IssueCollector,
) -> tuple[_Relationship, ...]:
    location = f"output.pptx!/{part_name}"
    root = _parse_xml(
        raw,
        location=location,
        malformed_code=f"{code_prefix}_RELATIONSHIPS_XML_MALFORMED",
        issues=issues,
    )
    if root is None:
        return ()
    if root.tag != RELATIONSHIPS_TAG or root.attrib or (root.text or "").strip():
        issues.add(
            f"{code_prefix}_RELATIONSHIPS_ROOT_INVALID",
            location,
            f"expected unadorned {RELATIONSHIPS_TAG!r}, observed {root.tag!r}",
        )
        return ()

    relationships: list[_Relationship] = []
    seen_ids: set[str] = set()
    allowed_attributes = {"Id", "Type", "Target", "TargetMode"}
    for index, node in enumerate(list(root), start=1):
        item_location = f"{location}#{index}"
        if node.tag != RELATIONSHIP_TAG:
            issues.add(
                f"{code_prefix}_RELATIONSHIP_ELEMENT_INVALID",
                item_location,
                f"unexpected element {node.tag!r}",
            )
            continue
        unexpected_attributes = set(node.attrib) - allowed_attributes
        relationship_id = node.attrib.get("Id", "")
        relationship_type = node.attrib.get("Type", "")
        target = node.attrib.get("Target", "")
        target_mode = node.attrib.get("TargetMode", "")
        if unexpected_attributes or not relationship_id or not relationship_type or not target:
            issues.add(
                f"{code_prefix}_RELATIONSHIP_ELEMENT_INVALID",
                item_location,
                "Relationship requires Id, Type, Target and no unknown attributes",
            )
        if list(node) or (node.text or "").strip() or (node.tail or "").strip():
            issues.add(
                f"{code_prefix}_RELATIONSHIP_ELEMENT_INVALID",
                item_location,
                "Relationship elements must be empty",
            )
        if relationship_id in seen_ids:
            issues.add(
                f"{code_prefix}_RELATIONSHIP_ID_DUPLICATE",
                item_location,
                relationship_id,
            )
        elif relationship_id:
            seen_ids.add(relationship_id)
        if target_mode not in {"", "External"}:
            issues.add(
                f"{code_prefix}_RELATIONSHIP_TARGET_MODE_INVALID",
                item_location,
                target_mode,
            )
        relationships.append(
            _Relationship(
                relationship_id=relationship_id,
                relationship_type=relationship_type,
                target=target,
                target_mode=target_mode,
            )
        )
    return tuple(relationships)


def _presentation_relationships(
    raw: bytes,
    issues: _IssueCollector,
) -> tuple[_Relationship, ...]:
    return _strict_relationships(
        raw,
        part_name=PRESENTATION_RELS_PART,
        code_prefix="PRESENTATION",
        issues=issues,
    )


def _resolve_internal_target(owner_part: str | None, target: str) -> str | None:
    if (
        not target
        or target != target.strip()
        or "\\" in target
        or any(ord(character) < 32 for character in target)
    ):
        return None
    decoded = unquote(target.split("#", 1)[0])
    if not decoded or decoded.startswith("//") or URI_SCHEME_RE.match(decoded):
        return None
    if decoded.startswith("/"):
        resolved = posixpath.normpath(decoded).lstrip("/")
    else:
        base = posixpath.dirname(owner_part) if owner_part is not None else ""
        resolved = posixpath.normpath(posixpath.join(base, decoded)).lstrip("/")
    if resolved in {"", ".", ".."} or resolved.startswith("../"):
        return None
    return resolved


def _resolve_presentation_target(target: str) -> str | None:
    return _resolve_internal_target(PRESENTATION_PART, target)


def _presentation_target_is_canonical(
    target: str,
    resolved_target: str | None,
) -> bool:
    if resolved_target is None:
        return False
    if (
        target.startswith("/")
        or "%" in target
        or "?" in target
        or "#" in target
        or "//" in target
        or "\\" in target
        or any(part in {"", ".", ".."} for part in target.split("/"))
        or any(ord(character) < 32 for character in target)
    ):
        return False
    return target == posixpath.relpath(resolved_target, "ppt")


def _validate_package_root_relationships(
    raw: bytes,
    issues: _IssueCollector,
) -> None:
    relationships = _strict_relationships(
        raw,
        part_name=PACKAGE_ROOT_RELS_PART,
        code_prefix="PACKAGE_ROOT",
        issues=issues,
    )
    location = f"output.pptx!/{PACKAGE_ROOT_RELS_PART}"
    office_documents = [
        relationship
        for relationship in relationships
        if relationship.relationship_type == OFFICE_DOCUMENT_RELATIONSHIP_TYPE
    ]
    if len(office_documents) != 1:
        issues.add(
            "PACKAGE_ROOT_OFFICE_DOCUMENT_RELATIONSHIP_INVALID",
            location,
            f"expected exactly one officeDocument relationship, observed {len(office_documents)}",
        )
    for relationship in relationships:
        if relationship.target_mode == "External":
            resolved_target = None
        else:
            resolved_target = _resolve_internal_target(None, relationship.target)
        item_location = (
            f"{location}#{relationship.relationship_id or 'missing-id'}"
        )
        if (
            resolved_target == PRESENTATION_PART
            and relationship.relationship_type != OFFICE_DOCUMENT_RELATIONSHIP_TYPE
        ):
            issues.add(
                "PACKAGE_ROOT_PRESENTATION_RELATIONSHIP_TYPE_INVALID",
                item_location,
                relationship.relationship_type,
            )
    if len(office_documents) == 1:
        office_document = office_documents[0]
        item_location = f"{location}#{office_document.relationship_id}"
        if office_document.target_mode:
            issues.add(
                "PACKAGE_ROOT_OFFICE_DOCUMENT_RELATIONSHIP_EXTERNAL",
                item_location,
                office_document.target_mode,
            )
        resolved_target = _resolve_internal_target(None, office_document.target)
        if (
            resolved_target != PRESENTATION_PART
            or office_document.target != PRESENTATION_PART
        ):
            issues.add(
                "PACKAGE_ROOT_OFFICE_DOCUMENT_TARGET_INVALID",
                item_location,
                (
                    f"expected {PRESENTATION_PART}, "
                    f"observed {office_document.target}"
                ),
            )


def _validate_slide_relationship_topology(
    relationship_ids: Sequence[str],
    relationships: Sequence[_Relationship],
    package_names: set[str],
    issues: _IssueCollector,
) -> tuple[str, ...]:
    location = f"output.pptx!/{PRESENTATION_RELS_PART}"
    by_id = {
        relationship.relationship_id: relationship
        for relationship in relationships
        if relationship.relationship_id
    }
    slide_relationships = [
        relationship
        for relationship in relationships
        if relationship.relationship_type == SLIDE_RELATIONSHIP_TYPE
    ]
    referenced_ids = set(relationship_ids)
    slide_relationship_ids = {
        relationship.relationship_id for relationship in slide_relationships
    }
    extra_ids = sorted(slide_relationship_ids - referenced_ids)
    if extra_ids:
        issues.add(
            "PRESENTATION_EXTRA_SLIDE_RELATIONSHIP",
            location,
            ",".join(extra_ids),
        )
    missing_ids = sorted(referenced_ids - slide_relationship_ids)
    if missing_ids:
        issues.add(
            "PRESENTATION_SLIDE_RELATIONSHIP_MISSING",
            location,
            ",".join(missing_ids),
        )

    for relationship in relationships:
        if relationship.target_mode == "External":
            continue
        resolved_target = _resolve_presentation_target(relationship.target)
        relationship_location = (
            f"{location}#{relationship.relationship_id or 'missing-id'}"
        )
        if not _presentation_target_is_canonical(
            relationship.target,
            resolved_target,
        ):
            issues.add(
                "PRESENTATION_RELATIONSHIP_TARGET_INVALID",
                relationship_location,
                relationship.target,
            )
        if resolved_target is None:
            continue
        if (
            relationship.relationship_type != SLIDE_RELATIONSHIP_TYPE
            and SLIDE_LIKE_PART_RE.fullmatch(resolved_target)
        ):
            issues.add(
                "PRESENTATION_SLIDE_RELATIONSHIP_TYPE_INVALID",
                relationship_location,
                relationship.relationship_type,
            )

    ordered_parts: list[str] = []
    for ordinal, relationship_id in enumerate(relationship_ids, start=1):
        expected_raw_target = f"slides/slide{ordinal}.xml"
        expected_part = f"ppt/{expected_raw_target}"
        ordered_parts.append(expected_part)
        relationship = by_id.get(relationship_id)
        item_location = f"{location}#{relationship_id}"
        if relationship is None:
            continue
        if relationship.relationship_type != SLIDE_RELATIONSHIP_TYPE:
            issues.add(
                "PRESENTATION_SLIDE_RELATIONSHIP_TYPE_INVALID",
                item_location,
                relationship.relationship_type,
            )
        if relationship.target_mode:
            issues.add(
                "PRESENTATION_SLIDE_RELATIONSHIP_EXTERNAL",
                item_location,
                relationship.target_mode,
            )
        if relationship.target != expected_raw_target:
            issues.add(
                "PRESENTATION_SLIDE_TARGET_SEQUENCE_INVALID",
                item_location,
                f"expected {expected_part}, observed {relationship.target}",
            )
        if expected_part not in package_names:
            issues.add(
                "PRESENTATION_SLIDE_PART_MISSING",
                f"output.pptx!/{expected_part}",
                f"referenced by {relationship_id}",
            )

    expected_parts = set(ordered_parts)
    actual_parts = {
        name for name in package_names if SLIDE_PART_RE.fullmatch(name)
    }
    extra_parts = sorted(actual_parts - expected_parts)
    if extra_parts:
        issues.add(
            "PRESENTATION_EXTRA_SLIDE_PART",
            "output.pptx!/ppt/slides",
            ",".join(extra_parts),
        )
    noncanonical_slide_parts = sorted(
        name
        for name in package_names
        if SLIDE_LIKE_PART_RE.fullmatch(name) and name not in actual_parts
    )
    if noncanonical_slide_parts:
        issues.add(
            "PRESENTATION_EXTRA_SLIDE_PART",
            "output.pptx!/ppt/slides",
            "noncanonical: " + ",".join(noncanonical_slide_parts),
        )
    return tuple(ordered_parts)


def _rectangle_union_area(
    rectangles: Sequence[tuple[int, int, int, int]],
) -> int:
    if not rectangles:
        return 0
    xs = sorted({value for rectangle in rectangles for value in (rectangle[0], rectangle[2])})
    total = 0
    for left, right in zip(xs, xs[1:]):
        if right <= left:
            continue
        intervals = sorted(
            (top, bottom)
            for rect_left, top, rect_right, bottom in rectangles
            if rect_left < right and rect_right > left
        )
        if not intervals:
            continue
        current_top, current_bottom = intervals[0]
        merged_height = 0
        for top, bottom in intervals[1:]:
            if top <= current_bottom:
                current_bottom = max(current_bottom, bottom)
            else:
                merged_height += current_bottom - current_top
                current_top, current_bottom = top, bottom
        merged_height += current_bottom - current_top
        total += (right - left) * merged_height
    return total


def _leaf_shapes(
    shapes: Any,
    *,
    slide_width: int,
    slide_height: int,
    issues: _IssueCollector,
    location: str,
    parent_transform: _geometry.Affine = _geometry.IDENTITY,
) -> Iterable[tuple[Any, tuple[int, int, int, int] | None, bool]]:
    for shape in shapes:
        try:
            shape_type = int(shape.shape_type)
        except Exception:  # python-pptx may raise NotImplementedError for legal XML
            shape_type = -1
        if shape_type == 6 and hasattr(shape, "shapes"):
            try:
                group_transform = _geometry._group_transform(
                    shape.element.grpSpPr.xfrm
                )
            except (AttributeError, _geometry.TemplatePackError) as exc:
                issues.add(
                    "EDITABILITY_GROUP_GEOMETRY_INVALID",
                    location,
                    str(exc),
                )
                group_transform = _geometry.IDENTITY
            yield from _leaf_shapes(
                shape.shapes,
                slide_width=slide_width,
                slide_height=slide_height,
                issues=issues,
                location=location,
                parent_transform=_geometry._compose(
                    parent_transform,
                    group_transform,
                ),
            )
            continue

        rectangle: tuple[int, int, int, int] | None = None
        native_countable = True
        try:
            local_transform, width, height = _geometry._shape_transform(
                _geometry._shape_xfrm(shape.element)
            )
            left, top, normalized_width, normalized_height = (
                _geometry._normalized_bbox(
                    _geometry._compose(parent_transform, local_transform),
                    width,
                    height,
                    slide_width=float(slide_width),
                    slide_height=float(slide_height),
                )
            )
            rectangle = (
                round(left * slide_width),
                round(top * slide_height),
                round((left + normalized_width) * slide_width),
                round((top + normalized_height) * slide_height),
            )
        except (AttributeError, _geometry.TemplatePackError) as exc:
            # A few legal inherited placeholders have no local xfrm.  At the
            # top level python-pptx has already resolved those coordinates.
            # Inside a transformed group, however, fallback coordinates would
            # be ambiguous and must fail closed.
            rectangle = _clipped_rectangle(
                shape,
                slide_width=slide_width,
                slide_height=slide_height,
            )
            if str(exc) == "shape extent must be positive":
                # A straight horizontal or vertical connector is represented
                # by DrawingML with one zero extent.  It has no two-dimensional
                # area to contribute to raster/native coverage, but remains a
                # legal, natively editable object.  Other non-positive extents
                # continue to fail closed as non-countable geometry.
                try:
                    element_local_name = shape.element.tag.rsplit("}", 1)[-1]
                    width = int(shape.width)
                    height = int(shape.height)
                except Exception:
                    element_local_name = ""
                    width = -1
                    height = -1
                native_countable = (
                    element_local_name == "cxnSp"
                    and ((width == 0 and height > 0) or (height == 0 and width > 0))
                )
            if parent_transform != _geometry.IDENTITY and str(exc) != (
                "shape extent must be positive"
            ):
                issues.add(
                    "EDITABILITY_GROUP_GEOMETRY_INVALID",
                    location,
                    str(exc),
                )
        yield shape, rectangle, native_countable


def _clipped_rectangle(
    shape: Any,
    *,
    slide_width: int,
    slide_height: int,
) -> tuple[int, int, int, int] | None:
    try:
        left = max(0, int(shape.left))
        top = max(0, int(shape.top))
        right = min(slide_width, int(shape.left) + int(shape.width))
        bottom = min(slide_height, int(shape.top) + int(shape.height))
    except Exception:  # malformed python-pptx proxies must not escape validation
        return None
    if right <= left or bottom <= top:
        return None
    return left, top, right, bottom


def _python_pptx_statistics(
    package_path: Path,
    ordered_slide_parts: Sequence[str],
    thresholds: RasterThresholds,
    issues: _IssueCollector,
) -> PresentationEditabilityStatistics | None:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        issues.add(
            "PYTHON_PPTX_UNAVAILABLE",
            str(package_path),
            str(exc),
        )
        return None
    try:
        presentation = Presentation(str(package_path))
        slides = tuple(presentation.slides)
        python_slide_count = len(slides)
        slide_width = int(presentation.slide_width)
        slide_height = int(presentation.slide_height)
    except Exception as exc:  # python-pptx exposes several parser exception types
        issues.add(
            "PYTHON_PPTX_REOPEN_FAILED",
            str(package_path),
            f"{type(exc).__name__}: {exc}",
        )
        return None

    topology_slide_count = len(ordered_slide_parts)
    if python_slide_count != topology_slide_count:
        issues.add(
            "PYTHON_PPTX_SLIDE_COUNT_MISMATCH",
            str(package_path),
            f"topology={topology_slide_count} python-pptx={python_slide_count}",
        )
    if slide_width <= 0 or slide_height <= 0:
        issues.add(
            "PYTHON_PPTX_SLIDE_SIZE_INVALID",
            str(package_path),
            f"width={slide_width} height={slide_height}",
        )
        return None

    slide_area = slide_width * slide_height
    slide_statistics: list[SlideEditabilityStatistics] = []
    for ordinal, slide in enumerate(slides, start=1):
        part_name = (
            ordered_slide_parts[ordinal - 1]
            if ordinal <= len(ordered_slide_parts)
            else f"ppt/slides/slide{ordinal}.xml"
        )
        shape_count = 0
        text_run_count = 0
        native_object_count = 0
        picture_count = 0
        full_slide_raster = False
        picture_rectangles: list[tuple[int, int, int, int]] = []
        native_rectangles: list[tuple[int, int, int, int]] = []
        text_rectangles: list[tuple[int, int, int, int]] = []
        slide_location = f"output.pptx!/{part_name}"
        try:
            shape_records = tuple(
                _leaf_shapes(
                    slide.shapes,
                    slide_width=slide_width,
                    slide_height=slide_height,
                    issues=issues,
                    location=slide_location,
                )
            )
        except Exception as exc:
            issues.add(
                "PYTHON_PPTX_SLIDE_SHAPES_INVALID",
                slide_location,
                f"{type(exc).__name__}: {exc}",
            )
            shape_records = ()
        for shape, rectangle, native_countable in shape_records:
            shape_count += 1
            try:
                shape_type = int(shape.shape_type)
            except Exception:  # fail closed on unsupported python-pptx shape proxies
                shape_type = -1
            is_picture = shape_type in {11, 13}
            try:
                element_local_name = shape.element.tag.rsplit("}", 1)[-1]
            except Exception:
                element_local_name = ""
            is_picture = is_picture or element_local_name == "pic"
            is_native = shape_type in NATIVE_SHAPE_TYPES and not is_picture
            if is_picture:
                picture_count += 1
                if rectangle is not None:
                    picture_rectangles.append(rectangle)
                    rectangle_area = (
                        (rectangle[2] - rectangle[0])
                        * (rectangle[3] - rectangle[1])
                    )
                    if rectangle_area / slide_area >= thresholds.full_slide_picture:
                        full_slide_raster = True
            elif is_native:
                if native_countable:
                    native_object_count += 1
                if rectangle is not None and native_countable:
                    native_rectangles.append(rectangle)
            else:
                issues.add(
                    "EDITABILITY_UNSUPPORTED_OBJECT",
                    f"output.pptx!/{part_name}",
                    (
                        f"shape_type={shape_type} "
                        f"element={element_local_name or 'unknown'}"
                    ),
                )

            shape_has_text = False
            try:
                has_text_frame = bool(shape.has_text_frame)
                if has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_run_count += 1
                            if run.text.strip():
                                shape_has_text = True
            except Exception as exc:
                issues.add(
                    "PYTHON_PPTX_SLIDE_SHAPES_INVALID",
                    slide_location,
                    f"{type(exc).__name__}: {exc}",
                )
            if shape_has_text and rectangle is not None:
                text_rectangles.append(rectangle)

        picture_coverage = _rectangle_union_area(picture_rectangles) / slide_area
        native_coverage = _rectangle_union_area(native_rectangles) / slide_area
        text_coverage = _rectangle_union_area(text_rectangles) / slide_area
        raster_dominant = (
            picture_coverage >= thresholds.dominant_picture_union
            and (
                native_coverage < thresholds.sparse_native_union
                and text_coverage < thresholds.sparse_text_union
            )
        )
        if native_object_count == 0:
            issues.add(
                (
                    "EDITABILITY_RASTER_ONLY_SLIDE"
                    if picture_count
                    else "EDITABILITY_EMPTY_SLIDE"
                ),
                slide_location,
                (
                    f"pictures={picture_count} "
                    f"native_objects={native_object_count}"
                ),
            )
        if raster_dominant:
            issues.add(
                "EDITABILITY_RASTER_DOMINANT_SLIDE",
                slide_location,
                (
                    f"picture_coverage={round(picture_coverage, 6)} "
                    f"native_coverage={round(native_coverage, 6)} "
                    f"text_coverage={round(text_coverage, 6)}"
                ),
            )
        slide_statistics.append(
            SlideEditabilityStatistics(
                ordinal=ordinal,
                part_name=part_name,
                shape_count=shape_count,
                text_run_count=text_run_count,
                native_object_count=native_object_count,
                picture_count=picture_count,
                picture_coverage=round(picture_coverage, 6),
                native_coverage=round(native_coverage, 6),
                text_coverage=round(text_coverage, 6),
                full_slide_raster=full_slide_raster,
                raster_dominant=raster_dominant,
            )
        )

    native_slide_count = sum(
        slide.native_object_count > 0 for slide in slide_statistics
    )
    return PresentationEditabilityStatistics(
        slide_count=topology_slide_count,
        python_pptx_slide_count=python_slide_count,
        slide_width=slide_width,
        slide_height=slide_height,
        text_run_count=sum(slide.text_run_count for slide in slide_statistics),
        shape_count=sum(slide.shape_count for slide in slide_statistics),
        native_object_count=sum(
            slide.native_object_count for slide in slide_statistics
        ),
        picture_count=sum(slide.picture_count for slide in slide_statistics),
        native_editable_slide_count=native_slide_count,
        full_slide_raster_count=sum(
            slide.full_slide_raster for slide in slide_statistics
        ),
        raster_dominant_slide_count=sum(
            slide.raster_dominant for slide in slide_statistics
        ),
        native_editable_coverage=(
            round(native_slide_count / topology_slide_count, 6)
            if topology_slide_count
            else 0.0
        ),
        slides=tuple(slide_statistics),
    )


def inspect_presentation_topology(
    package_path: str | Path,
    *,
    raster_thresholds: RasterThresholds | None = None,
    archive_limits: ZipResourceLimits = PPTX_ZIP_RESOURCE_LIMITS,
) -> PresentationTopologyResult:
    """Validate strict slide topology and recompute editability statistics.

    Malformed or inconsistent packages return ``status == "fail"`` with
    stable issue codes.  Invalid threshold configuration is a caller error and
    raises ``TypeError`` or ``ValueError``.
    """

    thresholds = raster_thresholds or RasterThresholds()
    thresholds.validate()
    path = Path(package_path).expanduser().resolve(strict=False)
    issues = _IssueCollector()
    ordered_relationship_ids: tuple[str, ...] = ()
    ordered_slide_parts: tuple[str, ...] = ()

    try:
        with zipfile.ZipFile(path, "r") as archive:
            resource_findings = audit_zip_resources(
                archive,
                limits=archive_limits,
            )
            if resource_findings:
                for finding in resource_findings:
                    issues.add(
                        finding.code,
                        f"{path}!/{finding.location}",
                        finding.detail,
                    )
                raise _PackageResourceRejected
            file_names = [
                info.filename
                for info in archive.infolist()
                if info.filename and not info.is_dir()
            ]
            raw_duplicates = sorted(
                name for name, count in Counter(file_names).items() if count > 1
            )
            if raw_duplicates:
                issues.add(
                    "PRESENTATION_PACKAGE_ENTRY_DUPLICATE",
                    str(path),
                    ",".join(raw_duplicates),
                )
            noncanonical_names = sorted(
                name
                for name in file_names
                if (
                    name.startswith("/")
                    or "%" in name
                    or "?" in name
                    or "#" in name
                    or "\\" in name
                    or any(part in {"", ".", ".."} for part in name.split("/"))
                    or unicodedata.normalize("NFC", name) != name
                    or any(
                        unicodedata.category(character) in {"Cc", "Cf", "Cs"}
                        for character in name
                    )
                )
            )
            if noncanonical_names:
                issues.add(
                    "PRESENTATION_PACKAGE_ENTRY_NONCANONICAL",
                    str(path),
                    ",".join(noncanonical_names),
                )
            names_by_casefold: dict[str, list[str]] = {}
            for name in file_names:
                collision_key = unicodedata.normalize("NFC", name).casefold()
                names_by_casefold.setdefault(collision_key, []).append(name)
            case_collisions = sorted(
                "/".join(sorted(set(names)))
                for names in names_by_casefold.values()
                if len(set(names)) > 1
            )
            if case_collisions:
                issues.add(
                    "PRESENTATION_PACKAGE_ENTRY_CASE_COLLISION",
                    str(path),
                    ",".join(case_collisions),
                )
            package_names = set(file_names)
            for required_part in (
                PACKAGE_ROOT_RELS_PART,
                PRESENTATION_PART,
                PRESENTATION_RELS_PART,
            ):
                if required_part not in package_names:
                    issues.add(
                        "PRESENTATION_REQUIRED_PART_MISSING",
                        f"output.pptx!/{required_part}",
                        "required topology part is missing",
                    )

            if PACKAGE_ROOT_RELS_PART in package_names:
                _validate_package_root_relationships(
                    archive.read(PACKAGE_ROOT_RELS_PART),
                    issues,
                )
            if PRESENTATION_PART in package_names:
                presentation_root = _parse_xml(
                    archive.read(PRESENTATION_PART),
                    location=f"output.pptx!/{PRESENTATION_PART}",
                    malformed_code="PRESENTATION_XML_MALFORMED",
                    issues=issues,
                )
                if presentation_root is not None:
                    ordered_relationship_ids = _presentation_slide_ids(
                        presentation_root,
                        issues,
                    )

            relationships: tuple[_Relationship, ...] = ()
            if PRESENTATION_RELS_PART in package_names:
                relationships = _presentation_relationships(
                    archive.read(PRESENTATION_RELS_PART),
                    issues,
                )
            ordered_slide_parts = _validate_slide_relationship_topology(
                ordered_relationship_ids,
                relationships,
                package_names,
                issues,
            )
            for part_name in ordered_slide_parts:
                if part_name not in package_names:
                    continue
                root = _parse_xml(
                    archive.read(part_name),
                    location=f"output.pptx!/{part_name}",
                    malformed_code="PRESENTATION_SLIDE_XML_MALFORMED",
                    issues=issues,
                )
                if root is not None and root.tag != SLIDE_TAG:
                    issues.add(
                        "PRESENTATION_SLIDE_ROOT_INVALID",
                        f"output.pptx!/{part_name}",
                        f"expected {SLIDE_TAG!r}, observed {root.tag!r}",
                    )
    except _PackageResourceRejected:
        pass
    except (OSError, KeyError, RuntimeError, zipfile.BadZipFile) as exc:
        issues.add(
            "PRESENTATION_PACKAGE_READ_FAILED",
            str(path),
            f"{type(exc).__name__}: {exc}",
        )

    statistics: PresentationEditabilityStatistics | None = None
    if not issues.issues:
        statistics = _python_pptx_statistics(
            path,
            ordered_slide_parts,
            thresholds,
            issues,
        )
    return PresentationTopologyResult(
        status="pass" if not issues.issues and statistics is not None else "fail",
        ordered_relationship_ids=ordered_relationship_ids,
        ordered_slide_parts=ordered_slide_parts,
        statistics=statistics,
        issues=issues.issues,
    )


__all__ = [
    "PresentationEditabilityStatistics",
    "PresentationTopologyIssue",
    "PresentationTopologyResult",
    "RasterThresholds",
    "SlideEditabilityStatistics",
    "inspect_presentation_topology",
]
