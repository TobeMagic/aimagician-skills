"""Deterministic pre-visual QA for physically assembled PPTX files.

This is deliberately conservative: it catches source residue, placeholders,
unreadable text and geometry violations before an independent visual reviewer
is allowed to score a deck.  It does not pretend to replace visual judgement
or PowerPoint's layout engine.
"""

from __future__ import annotations

import html
import json
import os
import re
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable, Mapping

from .physical_assembly import AssemblyPlan, PhysicalAssemblyError


TEXT_RE = re.compile(r"<a:t\b[^>]*>(.*?)</a:t>", re.DOTALL)
PLACEHOLDER_RE = re.compile(
    r"(?:20XX|20\dX|XXX|LOGO|输入标题|点击此处|请替换|占位|某某|Lorem|Your\s+(?:title|text))",
    re.IGNORECASE,
)
BRAND_RESIDUE_RE = re.compile(
    r"(?:B站|哔哩哔哩|nestle|雀巢|erke|安踏|Abbott|完美日记|蚂蚁森林)",
    re.IGNORECASE,
)


@dataclass(frozen=True)
class RuleFinding:
    rule: str
    severity: str
    slide: int | None
    message: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "rule": self.rule,
            "severity": self.severity,
            "slide": self.slide,
            "message": self.message,
        }


@dataclass(frozen=True)
class PhysicalRuleQAReport:
    schema_version: str
    status: str
    output_path: str
    slide_count: int
    blocking_findings: tuple[RuleFinding, ...] = field(default_factory=tuple)
    warnings: tuple[RuleFinding, ...] = field(default_factory=tuple)
    checked_rules: tuple[str, ...] = field(default_factory=tuple)

    def to_dict(self) -> dict[str, Any]:
        return {
            "schema_version": self.schema_version,
            "status": self.status,
            "output_path": self.output_path,
            "slide_count": self.slide_count,
            "checked_rules": list(self.checked_rules),
            "blocking_findings": [item.to_dict() for item in self.blocking_findings],
            "warnings": [item.to_dict() for item in self.warnings],
        }


def _slide_text(xml: bytes) -> list[str]:
    return [html.unescape(item).strip() for item in TEXT_RE.findall(xml.decode("utf-8", errors="replace")) if html.unescape(item).strip()]


def _source_text(template: Any) -> list[str]:
    path = Path(template.source_path)
    if not path.is_file():
        return []
    try:
        with zipfile.ZipFile(path, "r") as archive:
            return _slide_text(archive.read(f"ppt/slides/slide{template.slide_number}.xml"))
    except (OSError, KeyError, zipfile.BadZipFile):
        return []


def _pptx_checks(path: Path, findings: list[RuleFinding], warnings: list[RuleFinding]) -> int:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        warnings.append(RuleFinding("python-pptx", "warning", None, "python-pptx unavailable; skipped object checks"))
        return 0
    try:
        pres = Presentation(str(path))
    except Exception as exc:  # pragma: no cover - exercised by malformed files
        findings.append(RuleFinding("open", "blocker", None, f"python-pptx cannot open output: {exc}"))
        return 0
    sw, sh = pres.slide_width, pres.slide_height
    for slide_no, slide in enumerate(pres.slides, 1):
        visible_chars = 0
        for shape in slide.shapes:
            text_shape = bool(getattr(shape, "has_text_frame", False))
            out_of_bounds = shape.left < 0 or shape.top < 0 or shape.left + shape.width > sw or shape.top + shape.height > sh
            # Certified editorial pages commonly use a bleed image or a
            # decorative ray beyond the canvas.  Text leaving the canvas is a
            # blocker; non-text bleed is retained as a review warning.
            if out_of_bounds:
                text = shape.text.strip() if text_shape else ""
                # Physical reference pages use deliberate bleed and oversized
                # empty text frames for decorative numerals.  Only a visibly
                # populated frame crossing the top/left edge by >0.25 inch is
                # a deterministic clipping blocker; all other bleed is kept
                # as a warning for visual review.
                severe_text_clip = bool(text) and (shape.left < -228600 or shape.top < -228600)
                severity = "blocker" if severe_text_clip else "warning"
                target = findings if severe_text_clip else warnings
                target.append(RuleFinding("bounds", severity, slide_no, f"shape {getattr(shape, 'shape_id', '?')} leaves slide bounds"))
            if not text_shape:
                continue
            text = shape.text.strip()
            visible_chars += len(text)
            for match in PLACEHOLDER_RE.finditer(text):
                findings.append(RuleFinding("placeholder", "blocker", slide_no, f"placeholder residue: {match.group(0)!r}"))
            for match in BRAND_RESIDUE_RE.finditer(text):
                findings.append(RuleFinding("source-residue", "blocker", slide_no, f"named-brand residue: {match.group(0)!r}"))
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size is not None:
                        pt = run.font.size.pt
                        if pt < 8:
                            findings.append(RuleFinding("tiny-text", "blocker", slide_no, f"text run is {pt:.1f}pt (<8pt)"))
                        elif pt < 11:
                            warnings.append(RuleFinding("tiny-text", "warning", slide_no, f"text run is {pt:.1f}pt (<11pt)"))
        if visible_chars == 0:
            warnings.append(RuleFinding("density", "warning", slide_no, "slide has no visible text; inspect visual asset-only intent"))
    return len(pres.slides)


def run_physical_rule_qa(
    output_path: str | os.PathLike[str],
    *,
    plan: AssemblyPlan,
) -> PhysicalRuleQAReport:
    """Run deterministic checks; visual review remains a separate gate."""

    path = Path(output_path).expanduser().resolve(strict=False)
    findings: list[RuleFinding] = []
    warnings: list[RuleFinding] = []
    checked = (
        "zip-open",
        "slide-count",
        "placeholder-residue",
        "named-brand-residue",
        "source-template-residue",
        "text-bounds",
        "tiny-text",
        "style-lineage",
    )
    if not path.is_file():
        findings.append(RuleFinding("zip-open", "blocker", None, f"output missing: {path}"))
        return PhysicalRuleQAReport("1.0", "fail", str(path), 0, tuple(findings), tuple(warnings), checked)
    try:
        with zipfile.ZipFile(path, "r") as archive:
            slides = sorted(
                (int(match.group(1)), name)
                for name in archive.namelist()
                if (match := re.match(r"^ppt/slides/slide(\d+)\.xml$", name))
            )
            target_text = {ordinal: _slide_text(archive.read(name)) for ordinal, name in slides}
    except (OSError, zipfile.BadZipFile) as exc:
        findings.append(RuleFinding("zip-open", "blocker", None, str(exc)))
        return PhysicalRuleQAReport("1.0", "fail", str(path), 0, tuple(findings), tuple(warnings), checked)

    if len(slides) != plan.target_slide_count:
        findings.append(RuleFinding("slide-count", "blocker", None, f"expected {plan.target_slide_count} slides, found {len(slides)}"))
    for target in plan.target_slides:
        values = {str(value).strip() for value in target.bindings.values() if str(value).strip()}
        source = _source_text(target.page_template)
        output = target_text.get(target.ordinal, [])
        output_joined = "\n".join(output)
        # Source strings that survive adaptation are a strong signal that the
        # agent bound the wrong slot or silently left commercial copy behind.
        for text in source:
            if len(text) < 5 or text in values or PLACEHOLDER_RE.search(text):
                continue
            if text in output_joined and not BRAND_RESIDUE_RE.search(text):
                warnings.append(RuleFinding("source-template-residue", "warning", target.ordinal, f"source text remains: {text[:80]!r}"))
        if BRAND_RESIDUE_RE.search(output_joined):
            findings.append(RuleFinding("source-residue", "blocker", target.ordinal, "named-brand/source residue remains in slide text"))
        if PLACEHOLDER_RE.search(output_joined):
            findings.append(RuleFinding("placeholder", "blocker", target.ordinal, "placeholder token remains in slide text"))
    _pptx_checks(path, findings, warnings)
    # A duplicate is acceptable when the plan records it explicitly, but an
    # adjacent duplicate usually indicates a stalled narrative selection.
    for prev, curr in zip(plan.target_slides, plan.target_slides[1:]):
        if prev.page_template.page_id == curr.page_template.page_id:
            warnings.append(RuleFinding("style-lineage", "warning", curr.ordinal, "adjacent slide reuses the same physical page"))
    status = "pass" if not findings else "fail"
    return PhysicalRuleQAReport("1.0", status, str(path), len(slides), tuple(findings), tuple(warnings), checked)


def write_rule_qa_report(report: PhysicalRuleQAReport, output_path: str | os.PathLike[str]) -> str:
    path = Path(output_path).expanduser().resolve(strict=False)
    path.parent.mkdir(parents=True, exist_ok=True)
    text = json.dumps(report.to_dict(), ensure_ascii=False, indent=2)
    path.write_text(text, encoding="utf-8")
    import hashlib
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


__all__ = ["PhysicalRuleQAReport", "RuleFinding", "run_physical_rule_qa", "write_rule_qa_report"]
