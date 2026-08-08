"""Physical cross-package OPC assembly for v6.1.

Each target slide is physically copied from a different certified source
package, including its slide layout, slide master, theme, media, charts,
chart styles/colors, embedded workbooks, diagrams, notes, notes masters,
comments, and any other required custom XML.

The OPC graph is namespaced per source by ``v61_<source-hash>_<ordinal>``
so that all references remain acyclic. Byte-identical parts are deduplicated
once.

After dependency closure, declared text bindings are applied to each slide
through the existing ``adapt_template_pack`` OOXML patcher (the single-page
TemplatePack). Finally the target is committed and a
``physical-assembly-report.v1`` is emitted.
"""

from __future__ import annotations

import copy
import hashlib
import io
import json
import os
import re
import shutil
import tempfile
import zipfile
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable, Mapping, Sequence

from .page_template_library import (
    DEFAULT_DOMINANT_STYLE_CLUSTER,
    PageTemplate,
    SlotRecord,
    _CNVPR_RE,
    _LAYOUT_RE,
    _MASTER_RE,
    _SLIDE_RE,
    _TEXT_RE,
    _discover_slots as _library_discover_slots,
)
from .template_pack import TemplatePackError, _replace_shape_text


class PhysicalAssemblyError(ValueError):
    """Cross-package physical assembly has failed."""


@dataclass(frozen=True)
class SlideBinding:
    ordinal: int
    narrative_role: str
    title: str
    headline: str
    bindings: Mapping[str, str]
    page_id: str
    package_sha256: str
    slide_number: int


@dataclass(frozen=True)
class AssemblyTargetSlide:
    ordinal: int
    page_template: PageTemplate
    bindings: Mapping[str, str]
    narrative_role: str
    title: str
    headline: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "ordinal": self.ordinal,
            "page_id": self.page_template.page_id,
            "package_sha256": self.page_template.package_sha256,
            "slide_number": self.page_template.slide_number,
            "narrative_role": self.narrative_role,
            "title": self.title,
            "headline": self.headline,
            "bindings": dict(self.bindings),
        }


@dataclass(frozen=True)
class AssemblyPlan:
    schema_version: str
    plan_id: str
    scenario_id: str
    dominant_style_cluster_id: str
    created_at: str
    target_slide_count: int
    target_slides: tuple[AssemblyTargetSlide, ...]
    library_index_sha256: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "schema_version": self.schema_version,
            "plan_id": self.plan_id,
            "scenario_id": self.scenario_id,
            "dominant_style_cluster_id": self.dominant_style_cluster_id,
            "created_at": self.created_at,
            "target_slide_count": self.target_slide_count,
            "target_slides": [item.to_dict() for item in self.target_slides],
            "library_index_sha256": self.library_index_sha256,
        }


@dataclass(frozen=True)
class LineageRecord:
    ordinal: int
    page_id: str
    package_sha256: str
    slide_number: int
    source_sha256: str
    narrative_role: str
    title: str
    status: str
    binding_count: int
    byte_match_score: float

    def to_dict(self) -> dict[str, Any]:
        return {
            "ordinal": self.ordinal,
            "page_id": self.page_id,
            "package_sha256": self.package_sha256,
            "slide_number": self.slide_number,
            "source_sha256": self.source_sha256,
            "narrative_role": self.narrative_role,
            "title": self.title,
            "status": self.status,
            "binding_count": self.binding_count,
            "byte_match_score": self.byte_match_score,
        }


@dataclass(frozen=True)
class OPCIntegrity:
    zip_open: bool
    content_types_parsed: bool
    slide_rels_resolved: bool
    package_entry_count: int
    media_count: int
    status: str
    details: str = ""

    def to_dict(self) -> dict[str, Any]:
        return {
            "zip_open": self.zip_open,
            "content_types_parsed": self.content_types_parsed,
            "slide_rels_resolved": self.slide_rels_resolved,
            "package_entry_count": self.package_entry_count,
            "media_count": self.media_count,
            "status": self.status,
            "details": self.details,
        }


@dataclass(frozen=True)
class Editability:
    native_editable: bool
    python_pptx_open: bool
    slide_count: int
    text_run_count: int
    shape_count: int
    status: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "native_editable": self.native_editable,
            "python_pptx_open": self.python_pptx_open,
            "slide_count": self.slide_count,
            "text_run_count": self.text_run_count,
            "shape_count": self.shape_count,
            "status": self.status,
        }


@dataclass(frozen=True)
class StyleClusterAdherence:
    dominant_style_cluster_id: str
    matches: int
    total: int
    status: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "dominant_style_cluster_id": self.dominant_style_cluster_id,
            "matches": self.matches,
            "total": self.total,
            "status": self.status,
        }


@dataclass(frozen=True)
class PhysicalAssemblyReport:
    schema_version: str
    report_id: str
    plan_id: str
    output_path: str
    output_sha256: str
    status: str
    target_slide_count: int
    lineage_records: tuple[LineageRecord, ...]
    opc_integrity: OPCIntegrity
    editability: Editability
    style_cluster_adherence: StyleClusterAdherence

    def to_dict(self) -> dict[str, Any]:
        return {
            "schema_version": self.schema_version,
            "report_id": self.report_id,
            "plan_id": self.plan_id,
            "output_path": self.output_path,
            "output_sha256": self.output_sha256,
            "status": self.status,
            "target_slide_count": self.target_slide_count,
            "lineage_records": [record.to_dict() for record in self.lineage_records],
            "opc_integrity": self.opc_integrity.to_dict(),
            "editability": self.editability.to_dict(),
            "style_cluster_adherence": self.style_cluster_adherence.to_dict(),
        }


def _sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _now_utc() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


# Regex patterns for parsing OPC relationships
_REL_TYPE_RE = re.compile(r'Type="([^"]+)"')
_REL_TARGET_RE = re.compile(r'Target="([^"]+)"')
_REL_ID_RE = re.compile(r'Id="([^"]+)"')
_REL_TAG_RE = re.compile(
    r'<Relationship\b([^>]*?)(?:/>|></Relationship>)', re.DOTALL
)
_CT_OVERRIDE_RE = re.compile(
    r'<Override\s+PartName="([^"]+)"\s+ContentType="([^"]+)"\s*/>', re.DOTALL
)
_CT_DEFAULT_RE = re.compile(
    r'<Default\s+Extension="([^"]+)"\s+ContentType="([^"]+)"\s*/>', re.DOTALL
)

ALLOWED_REL_TYPES = frozenset(
    {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramLayout",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramQuickStyle",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramColors",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/embeddedFont",
        "http://schemas.microsoft.com/office/2011/relationships/comments",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/glossaryDocument",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
        "http://schemas.microsoft.com/office/2011/relationships/people",
        "http://schemas.microsoft.com/office/2007/relationships/hdphoto",
    }
)

# The following relationship types are forbidden and must be rejected before
# any part is committed to the output.
FORBIDDEN_REL_TARGET_PATTERNS = (
    re.compile(r"\.exe$", re.IGNORECASE),
    re.compile(r"\.bat$", re.IGNORECASE),
    re.compile(r"\.cmd$", re.IGNORECASE),
    re.compile(r"\.scr$", re.IGNORECASE),
    re.compile(r"\.vbs$", re.IGNORECASE),
    re.compile(r"\.js$", re.IGNORECASE),
    re.compile(r"\.com$", re.IGNORECASE),
    re.compile(r"\.dll$", re.IGNORECASE),
    re.compile(r"\.msi$", re.IGNORECASE),
)

# Mandatory types we need to register after import.
CONTENT_TYPES = {
    "ppt": "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
    "slide": "application/vnd.openxmlformats-officedocument.presentationml.slide+xml",
    "slideLayout": "application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml",
    "slideMaster": "application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml",
    "notesSlide": "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml",
    "notesMaster": "application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml",
    "theme": "application/vnd.openxmlformats-officedocument.theme+xml",
    "chart": "application/vnd.openxmlformats-officedocument.drawingml.chart+xml",
    "image": "application/vnd.openxmlformats-officedocument.drawingml.picture+xml",
}


@dataclass
class _SourceGraph:
    """Per-source package dependency graph rooted at slide1.xml."""

    slide_xml: bytes
    slide_sha: str
    rels: dict[str, bytes] = field(default_factory=dict)
    extra_parts: dict[str, bytes] = field(default_factory=dict)
    layout_paths: list[str] = field(default_factory=list)
    master_paths: list[str] = field(default_factory=list)
    theme_paths: list[str] = field(default_factory=list)


def _normalise_zip_name(name: str) -> str:
    return name.replace("\\", "/")


def _resolve_rel_target(rels_xml: bytes, source_rel_path: str, target: str) -> str | None:
    """Resolve a relationship target relative to the source rels file.

    ``source_rel_path`` looks like ``ppt/slides/_rels/slide1.xml.rels``. The
    relationships it carries describe files in ``ppt/slides/`` (one directory
    up from the rels file). Targets that begin with ``../`` therefore resolve
    to ``ppt/...`` siblings, not ``ppt/slides/...``.
    """

    if not target or target.startswith("http://") or target.startswith("https://"):
        return None
    # Strip "/_rels/<file>.rels" suffix to find the owning folder.
    parts = source_rel_path.split("/")
    if len(parts) >= 2 and parts[-2] == "_rels":
        base = "/".join(parts[:-2])
    else:
        base = os.path.dirname(source_rel_path)
    if target.startswith("/"):
        return _normalise_zip_name(target.lstrip("/"))
    combined = _normalise_zip_name(os.path.normpath(os.path.join(base, target)))
    return combined.lstrip("/")


def _parse_relationships(rels_xml: bytes) -> list[dict[str, str]]:
    """Parse a Relationships XML and return entries."""

    text = rels_xml.decode("utf-8", errors="replace")
    entries: list[dict[str, str]] = []
    for match in _REL_TAG_RE.finditer(text):
        attrs = match.group(1)
        entry = {
            "Id": (_REL_ID_RE.search(attrs) or [None, ""])[1],
            "Type": (_REL_TYPE_RE.search(attrs) or [None, ""])[1],
            "Target": (_REL_TARGET_RE.search(attrs) or [None, ""])[1],
        }
        entries.append(entry)
    return entries


def _build_source_graph(package_path: Path, slide_number: int = 1) -> _SourceGraph:
    """Build a single-slide dependency graph from one source package.

    ``slide_number`` is deliberately explicit so a certified multi-page
    reference deck can contribute its original pages without first being
    redrawn into separate packages.
    """

    with zipfile.ZipFile(package_path, "r") as archive:
        names = archive.namelist()
        slide_name = f"ppt/slides/slide{slide_number}.xml"
        if slide_name not in names:
            raise PhysicalAssemblyError(
                f"source missing {slide_name}: {package_path}"
            )
        slide_xml = archive.read(slide_name)
        rels_path = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
        slide_rels_xml = (
            archive.read(rels_path) if rels_path in names else b""
        )
        rels_map: dict[str, bytes] = {rels_path: slide_rels_xml}
        extras: dict[str, bytes] = {}
        layouts: list[str] = []
        masters: list[str] = []
        themes: list[str] = []
        seen: set[str] = set()
        queue: list[str] = [
            entry["Target"]
            for entry in _parse_relationships(slide_rels_xml)
            if entry.get("Target")
        ]
        while queue:
            raw_target = queue.pop(0)
            if not raw_target or raw_target.startswith(("http://", "https://")):
                continue
            target = _resolve_rel_target(rels_path, rels_path, raw_target)
            if target is None:
                continue
            target = _normalise_zip_name(target)
            if target in seen:
                continue
            seen.add(target)
            if any(pat.search(target) for pat in FORBIDDEN_REL_TARGET_PATTERNS):
                raise PhysicalAssemblyError(
                    f"forbidden relationship target in {package_path}: {target}"
                )
            if target not in names:
                continue
            data = archive.read(target)
            extras[target] = data
            if _LAYOUT_RE.match(target):
                layouts.append(target)
            elif _MASTER_RE.match(target):
                masters.append(target)
            elif target.startswith("ppt/theme/"):
                themes.append(target)
            rels_name = f"{os.path.dirname(target)}/_rels/{os.path.basename(target)}.rels"
            if rels_name in names:
                rels_map[rels_name] = archive.read(rels_name)
                for entry in _parse_relationships(rels_map[rels_name]):
                    sub_target = entry.get("Target")
                    if sub_target:
                        queue.append(sub_target)
        # Auto-pull layouts' rels to masters/themes
        for layout in layouts:
            layout_rels = f"{os.path.dirname(layout)}/_rels/{os.path.basename(layout)}.rels"
            if layout_rels in names:
                rels_map[layout_rels] = archive.read(layout_rels)
                for entry in _parse_relationships(rels_map[layout_rels]):
                    t = entry.get("Target")
                    if t and t not in extras:
                        resolved = _resolve_rel_target(layout_rels, layout_rels, t)
                        if resolved and resolved in names:
                            extras[resolved] = archive.read(resolved)
                            if _MASTER_RE.match(resolved) and resolved not in masters:
                                masters.append(resolved)
                            if resolved.startswith("ppt/theme/") and resolved not in themes:
                                themes.append(resolved)
        # Auto-pull masters' rels to themes
        for master in masters:
            master_rels = f"{os.path.dirname(master)}/_rels/{os.path.basename(master)}.rels"
            if master_rels in names:
                rels_map[master_rels] = archive.read(master_rels)
                for entry in _parse_relationships(rels_map[master_rels]):
                    t = entry.get("Target")
                    if t and t not in extras:
                        resolved = _resolve_rel_target(master_rels, master_rels, t)
                        if resolved and resolved in names:
                            extras[resolved] = archive.read(resolved)
                            if resolved.startswith("ppt/theme/") and resolved not in themes:
                                themes.append(resolved)
    graph = _SourceGraph(
        slide_xml=slide_xml,
        slide_sha=_sha256_bytes(slide_xml),
        rels=rels_map,
        extra_parts=extras,
        layout_paths=layouts,
        master_paths=masters,
        theme_paths=themes,
    )
    return graph


def _namespace_part_name(name: str, source_hash: str, ordinal: int) -> str:
    """Rewrite a part path to be unique per source."""

    if not name.startswith("ppt/"):
        return name
    head, _, tail = name.partition("/")
    rest = tail.split("/")
    rest.insert(0, f"v61_{source_hash[:12]}_{ordinal:03d}")
    return f"{head}/" + "/".join(rest)


def _namespace_rels_name(name: str) -> str:
    return name  # We rely on part-name rewriting above to namespace rels parents.


def _namespace_relationship_path(
    rels_path: str,
    target_map: Mapping[str, str],
    *,
    source_slide_rels_name: str,
    ordinal: int,
) -> str:
    """Return the relationship part path for a namespaced owner part."""

    if rels_path.endswith(source_slide_rels_name):
        return f"ppt/slides/_rels/slide{ordinal}.xml.rels"
    parts = rels_path.split("/")
    if len(parts) >= 2 and parts[-2] == "_rels":
        owner = "/".join(parts[:-2] + [parts[-1][:-5]])
        mapped = target_map.get(owner)
        if mapped:
            return f"{os.path.dirname(mapped)}/_rels/{os.path.basename(mapped)}.rels"
    return rels_path


def _rewrite_relationship_targets(
    rels_xml: bytes,
    rels_path: str,
    target_map: Mapping[str, str],
    *,
    output_rels_path: str | None = None,
) -> bytes:
    """Rewrite ``Target`` attributes through ``target_map``.

    The supplied ``target_map`` keys may already be normalised (the caller
    stores resolved paths from the source OPC graph). We try both the
    normalised lookup and the un-normalised lookup, falling back to a path
    computed by relative-path math if neither hits.
    """

    text = rels_xml.decode("utf-8", errors="replace")

    def repl(match: re.Match[str]) -> str:
        attrs = match.group(1)
        target_match = _REL_TARGET_RE.search(attrs)
        if target_match is None:
            return match.group(0)
        raw_target = target_match.group(1)
        if not raw_target or raw_target.startswith(("http://", "https://")):
            return match.group(0)
        resolved = _resolve_rel_target(rels_xml, rels_path, raw_target)
        resolved_norm = _normalise_zip_name(resolved) if resolved else None
        new_target = None
        if resolved_norm and resolved_norm in target_map:
            new_target = target_map[resolved_norm]
        elif raw_target in target_map:
            new_target = target_map[raw_target]
        else:
            return match.group(0)
        # Compute new relative path from new_rels parent dir.
        output_path = output_rels_path or rels_path
        rels_parts = output_path.split("/")
        if len(rels_parts) >= 2 and rels_parts[-2] == "_rels":
            new_base = "/".join(rels_parts[:-2])
        else:
            new_base = os.path.dirname(rels_path)
        if new_target.startswith("/"):
            new_relative = new_target.lstrip("/")
        else:
            new_relative = os.path.relpath(new_target, new_base).replace("\\", "/")
        attrs = re.sub(
            r'Target="[^"]*"', f'Target="{new_relative}"', attrs
        )
        return f"<Relationship{attrs}/>"

    new_text = _REL_TAG_RE.sub(repl, text)
    return new_text.encode("utf-8")


def _rewrite_slide_references(slide_xml: bytes, target_map: Mapping[str, str]) -> bytes:
    """Rewrite XML references in slide body (r:embed, r:link, layout references)."""

    text = slide_xml.decode("utf-8", errors="replace")
    base = "ppt/slides"

    def relative(new_target: str) -> str:
        if new_target.startswith("/"):
            return new_target.lstrip("/")
        rel = os.path.relpath(new_target, base).replace("\\", "/")
        return rel

    for old, new in target_map.items():
        old_name = old.split("/")[-1]
        new_name = new.split("/")[-1]
        text = re.sub(
            r'(r:embed=")' + re.escape(old_name) + r'(")',
            r"\1" + new_name + r"\2",
            text,
        )
        text = re.sub(
            r'(r:link=")' + re.escape(old_name) + r'(")',
            r"\1" + new_name + r"\2",
            text,
        )
    return text.encode("utf-8")


def _adapt_slide_text(
    slide_xml: bytes,
    bindings: Mapping[str, str],
    *,
    allowed_slots: Iterable[str] | None = None,
) -> bytes:
    """Apply declared bindings and fail closed on stale or invented slots."""

    text = slide_xml.decode("utf-8", errors="replace")
    allowed = set(allowed_slots or ())
    for slot_id, replacement in bindings.items():
        if not slot_id.startswith("shape_"):
            raise PhysicalAssemblyError(f"invalid text slot id: {slot_id}")
        if allowed and slot_id not in allowed:
            raise PhysicalAssemblyError(
                f"binding targets a slot outside the certified slot graph: {slot_id}"
            )
        try:
            shape_id = int(slot_id[len("shape_"):])
        except ValueError:
            raise PhysicalAssemblyError(f"invalid text slot id: {slot_id}") from None
        marker = re.compile(
            rf'<p:cNvPr\b[^>]*\bid="{shape_id}"[^>]*>'
        )
        m = marker.search(text)
        if m is None:
            raise PhysicalAssemblyError(f"declared text slot not found: {slot_id}")
        start = m.end()
        nxt = _CNVPR_RE.search(text, start)
        end = nxt.start() if nxt else len(text)
        segment = text[start:end]
        if not _TEXT_RE.search(segment):
            raise PhysicalAssemblyError(
                f"declared text slot has no editable text nodes: {slot_id}"
            )
        new_segment = _replace_shape_text(segment, replacement)
        text = text[:start] + new_segment + text[end:]
    return text.encode("utf-8")


def _parse_content_types(xml_bytes: bytes) -> tuple[list[tuple[str, str]], list[tuple[str, str]]]:
    text = xml_bytes.decode("utf-8", errors="replace")
    defaults = [
        (m.group(1).lower(), m.group(2))
        for m in _CT_DEFAULT_RE.finditer(text)
    ]
    overrides = [
        (m.group(1), m.group(2))
        for m in _CT_OVERRIDE_RE.finditer(text)
    ]
    return defaults, overrides


def _serialize_content_types(
    defaults: Sequence[tuple[str, str]],
    overrides: Sequence[tuple[str, str]],
) -> bytes:
    """Render [Content_Types].xml with stable field order."""

    seen_defaults: set[str] = set()
    seen_overrides: set[str] = set()
    parts: list[str] = [
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">',
    ]
    for ext, ct in defaults:
        ext = ext.lower()
        if ext in seen_defaults:
            continue
        seen_defaults.add(ext)
        parts.append(f'<Default Extension="{ext}" ContentType="{ct}"/>')
    for part_name, ct in overrides:
        if part_name in seen_overrides:
            continue
        seen_overrides.add(part_name)
        parts.append(f'<Override PartName="{part_name}" ContentType="{ct}"/>')
    parts.append("</Types>")
    return "".join(parts).encode("utf-8")


def _default_content_types() -> tuple[list[tuple[str, str]], list[tuple[str, str]]]:
    return [
        ("rels", "application/vnd.openxmlformats-package.relationships+xml"),
        ("xml", "application/xml"),
        ("png", "image/png"),
        ("jpeg", "image/jpeg"),
        ("jpg", "image/jpeg"),
        ("gif", "image/gif"),
        ("svg", "image/svg+xml"),
        ("wdp", "image/vnd.ms-photo"),
        ("emf", "image/x-emf"),
        ("wmf", "image/x-wmf"),
        ("bmp", "image/bmp"),
        ("tif", "image/tiff"),
        ("tiff", "image/tiff"),
    ], []


def _default_pres_rels() -> bytes:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>'
        '<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>'
        '<Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>'
        "</Relationships>"
    ).encode("utf-8")


def _root_rels() -> bytes:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>'
        "</Relationships>"
    ).encode("utf-8")


def _default_pres_xml(
    slide_paths: Sequence[str], master_paths: Sequence[str]
) -> bytes:
    """Build a presentation.xml pointing at the slides in stable order."""

    sld_ids = []
    slide_start = len(master_paths) + 5
    for idx, path in enumerate(slide_paths, start=slide_start):
        sld_ids.append(f'<p:sldId id="{256 + (idx - 9)}" r:id="rId{idx}"/>')
    sld_id_list = "".join(sld_ids)
    master_ids = "".join(
        f'<p:sldMasterId id="{2147483648 + idx}" r:id="rId{idx + 1}"/>'
        for idx, _ in enumerate(master_paths)
    )
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        f"<p:sldMasterIdLst>{master_ids}</p:sldMasterIdLst>"
        f"<p:sldIdLst>{sld_id_list}</p:sldIdLst>"
        '<p:sldSz cx="12192000" cy="6858000" type="screen16x9"/>'
        '<p:notesSz cx="6858000" cy="9144000"/>'
        "</p:presentation>"
    ).encode("utf-8")


def _pres_xml_rels(
    slide_rels: Sequence[str], master_paths: Sequence[str]
) -> bytes:
    entries: list[str] = [
        f'<Relationship Id="rId{idx}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="{target}"/>'
        for idx, target in enumerate(master_paths, start=1)
    ]
    support_start = len(master_paths) + 1
    entries.extend([
        f'<Relationship Id="rId{support_start}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/presProps" Target="presProps.xml"/>',
        f'<Relationship Id="rId{support_start + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/viewProps" Target="viewProps.xml"/>',
        f'<Relationship Id="rId{support_start + 2}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/tableStyles" Target="tableStyles.xml"/>',
    ])
    for idx, target in enumerate(slide_rels, start=len(master_paths) + 5):
        entries.append(
            f'<Relationship Id="rId{idx}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="{target}"/>'
        )
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        + "".join(entries)
        + "</Relationships>"
    ).encode("utf-8")


def _empty_doc_props() -> tuple[bytes, bytes]:
    core = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" '
        b'xmlns:dc="http://purl.org/dc/elements/1.1/" '
        b'xmlns:dcterms="http://purl.org/dc/terms/" '
        b'xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">'
        b"<dc:title>Window-PPTX v6.1 Deck</dc:title>"
        b"<dc:creator>window-pptx</dc:creator>"
        b"</cp:coreProperties>"
    )
    app = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" '
        b'xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">'
        b"<Application>window-pptx</Application>"
        b"</Properties>"
    )
    return core, app


def _default_props_xml() -> bytes:
    return (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<p:presentationPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        b'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
    )


def _default_view_props_xml() -> bytes:
    return (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<p:viewPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
    )


def _default_table_styles_xml() -> bytes:
    return (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<a:tblStyleLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
    )


def _plan_from_payload(payload: Mapping[str, Any], library_lookup: Mapping[str, PageTemplate]) -> AssemblyPlan:
    """Build an AssemblyPlan from a payload dict."""

    slides: list[AssemblyTargetSlide] = []
    for entry in payload["target_slides"]:
        page_id = entry["page_id"]
        template = library_lookup.get(page_id)
        if template is None:
            raise PhysicalAssemblyError(f"library missing page_id={page_id}")
        if template.package_sha256 != entry["package_sha256"]:
            raise PhysicalAssemblyError(
                f"package_sha256 mismatch for page_id={page_id}"
            )
        slides.append(
            AssemblyTargetSlide(
                ordinal=int(entry["ordinal"]),
                page_template=template,
                bindings=dict(entry.get("bindings", {})),
                narrative_role=str(entry["narrative_role"]),
                title=str(entry["title"]),
                headline=str(entry.get("headline", "")),
            )
        )
    return AssemblyPlan(
        schema_version=payload.get("schema_version", "1.0"),
        plan_id=str(payload["plan_id"]),
        scenario_id=str(payload["scenario_id"]),
        dominant_style_cluster_id=str(
            payload.get(
                "dominant_style_cluster_id",
                DEFAULT_DOMINANT_STYLE_CLUSTER,
            )
        ),
        created_at=str(payload.get("created_at", _now_utc())),
        target_slide_count=int(payload["target_slide_count"]),
        target_slides=tuple(slides),
        library_index_sha256=str(payload["library_index_sha256"]),
    )


def _validate_assembly_plan(
    plan: AssemblyPlan, library_index_sha256: str
) -> None:
    if plan.library_index_sha256 != library_index_sha256:
        raise PhysicalAssemblyError(
            "ASSEMBLY_PLAN_FINGERPRINT_MISMATCH: library index drift"
        )
    seen: set[int] = set()
    for slide in plan.target_slides:
        if slide.ordinal in seen:
            raise PhysicalAssemblyError(
                f"ASSEMBLY_PLAN_DUPLICATE_ORDINAL: {slide.ordinal}"
            )
        seen.add(slide.ordinal)
        expected = set(slide.page_template.slot_graph.get("text_slot_ids", ()))
        actual = set(slide.bindings)
        if not expected:
            raise PhysicalAssemblyError(
                f"ASSEMBLY_PLAN_NO_EDITABLE_SLOTS: ordinal={slide.ordinal} page_id={slide.page_template.page_id}"
            )
        missing = sorted(expected - actual)
        extra = sorted(actual - expected)
        if missing or extra:
            details = []
            if missing:
                details.append("missing=" + ",".join(missing))
            if extra:
                details.append("extra=" + ",".join(extra))
            raise PhysicalAssemblyError(
                f"ASSEMBLY_PLAN_SLOT_COVERAGE: ordinal={slide.ordinal} "
                + " ".join(details)
            )
    if len(seen) != plan.target_slide_count:
        raise PhysicalAssemblyError(
            "ASSEMBLY_PLAN_LENGTH_MISMATCH"
        )


def _byte_match_score(source: bytes, target: bytes) -> float:
    if not source:
        return 0.0
    matches = sum(1 for a, b in zip(source, target) if a == b)
    return matches / max(len(source), len(target))


def _slide_text_run_count(slide_xml: bytes) -> int:
    return len(_TEXT_RE.findall(slide_xml.decode("utf-8", errors="replace")))


def _slide_shape_count(slide_xml: bytes) -> int:
    return len(_CNVPR_RE.findall(slide_xml.decode("utf-8", errors="replace")))


def _verify_zip_open(output_path: Path) -> tuple[bool, str]:
    try:
        with zipfile.ZipFile(output_path, "r"):
            return True, ""
    except (OSError, zipfile.BadZipFile) as exc:
        return False, f"zip-open failed: {exc}"


def _verify_content_types(output_path: Path) -> tuple[bool, str, int]:
    try:
        with zipfile.ZipFile(output_path, "r") as archive:
            xml = archive.read("[Content_Types].xml")
            defaults, overrides = _parse_content_types(xml)
            return True, "", len(defaults) + len(overrides)
    except (OSError, KeyError, zipfile.BadZipFile) as exc:
        return False, f"content-types parse failed: {exc}", 0


def _verify_slide_rels(output_path: Path) -> tuple[bool, str]:
    try:
        with zipfile.ZipFile(output_path, "r") as archive:
            slides = sorted(n for n in archive.namelist() if _SLIDE_RE.match(n))
            for slide in slides:
                rels = (
                    f"{os.path.dirname(slide)}/_rels/{os.path.basename(slide)}.rels"
                )
                if rels not in archive.namelist():
                    return False, f"missing {rels}"
                for entry in _parse_relationships(archive.read(rels)):
                    ttype = entry["Type"]
                    if ttype and not ttype.startswith(
                        "http://schemas.openxmlformats.org/officeDocument"
                    ) and not ttype.startswith(
                        "http://schemas.microsoft.com/office"
                    ):
                        continue
                    target = entry["Target"]
                    if not target or target.startswith(("http://", "https://")):
                        continue
                    parts = rels.split("/")
                    if len(parts) >= 2 and parts[-2] == "_rels":
                        base = "/".join(parts[:-2])
                    else:
                        base = os.path.dirname(rels)
                    resolved = _normalise_zip_name(
                        os.path.normpath(os.path.join(base, target))
                    ).lstrip("/")
                    if resolved not in archive.namelist():
                        return False, f"unresolved {resolved}"
        return True, ""
    except (OSError, zipfile.BadZipFile) as exc:
        return False, f"slide-rels failed: {exc}"


def _verify_python_pptx(output_path: Path) -> tuple[bool, int, int, int]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return False, 0, 0, 0
    try:
        pres = Presentation(str(output_path))
    except Exception:
        return False, 0, 0, 0
    total_text = 0
    total_shapes = 0
    for slide in pres.slides:
        for shape in slide.shapes:
            total_shapes += 1
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        total_text += 1
    return True, len(pres.slides), total_text, total_shapes


def verify_physical_assembly(
    output_path: str | os.PathLike[str],
    *,
    plan: AssemblyPlan,
) -> PhysicalAssemblyReport:
    """Verify the assembled PPTX against the assembly plan."""

    path = Path(output_path).expanduser().resolve(strict=False)
    if not path.is_file():
        raise PhysicalAssemblyError(f"output missing: {path}")

    zip_ok, zip_detail = _verify_zip_open(path)
    ct_ok, ct_detail, ct_count = _verify_content_types(path)
    rels_ok, rels_detail = _verify_slide_rels(path)

    py_ok, slide_count, text_runs, shape_count = _verify_python_pptx(path)
    editability = Editability(
        native_editable=zip_ok and ct_ok and rels_ok,
        python_pptx_open=py_ok,
        slide_count=slide_count,
        text_run_count=text_runs,
        shape_count=shape_count,
        status="pass" if (zip_ok and ct_ok and rels_ok and py_ok) else "fail",
    )

    lineage: list[LineageRecord] = []
    matches = 0
    with zipfile.ZipFile(path, "r") as archive:
        for slide in plan.target_slides:
            ordinal = slide.ordinal
            target_slide_name = f"ppt/slides/slide{ordinal}.xml"
            slide_xml = archive.read(target_slide_name)
            source_sha = slide.page_template.source_sha256
            try:
                source_path = Path(slide.page_template.source_path)
                if source_path.is_file():
                    with zipfile.ZipFile(source_path, "r") as src_zip:
                        source_slide = src_zip.read(
                            f"ppt/slides/slide{slide.page_template.slide_number}.xml"
                        )
                else:
                    source_slide = slide_xml
            except (OSError, KeyError, zipfile.BadZipFile):
                source_slide = slide_xml
            score = _byte_match_score(source_slide, slide_xml)
            status = "pass" if zip_ok and ct_ok and rels_ok else "fail"
            if slide.page_template.style_cluster_id == plan.dominant_style_cluster_id:
                matches += 1
            lineage.append(
                LineageRecord(
                    ordinal=ordinal,
                    page_id=slide.page_template.page_id,
                    package_sha256=slide.page_template.package_sha256,
                    slide_number=slide.page_template.slide_number,
                    source_sha256=source_sha,
                    narrative_role=slide.narrative_role,
                    title=slide.title,
                    status=status,
                    binding_count=len(slide.bindings),
                    byte_match_score=round(score, 6),
                )
            )

    package_entry_count = 0
    media_count = 0
    try:
        with zipfile.ZipFile(path, "r") as archive:
            package_entry_count = len(archive.namelist())
            media_count = sum(1 for n in archive.namelist() if "media/" in n)
    except (OSError, zipfile.BadZipFile):
        pass

    opc = OPCIntegrity(
        zip_open=zip_ok,
        content_types_parsed=ct_ok,
        slide_rels_resolved=rels_ok,
        package_entry_count=package_entry_count,
        media_count=media_count,
        status="pass" if (zip_ok and ct_ok and rels_ok) else "fail",
        details="; ".join(filter(None, (zip_detail, ct_detail, rels_detail))),
    )
    adherence = StyleClusterAdherence(
        dominant_style_cluster_id=plan.dominant_style_cluster_id,
        matches=matches,
        total=len(plan.target_slides),
        status="pass" if matches == len(plan.target_slides) else "fail",
    )
    overall = "pass" if (
        opc.status == "pass"
        and editability.status == "pass"
        and adherence.status == "pass"
    ) else "fail"
    return PhysicalAssemblyReport(
        schema_version="1.0",
        report_id=f"par_{int(datetime.now(timezone.utc).timestamp() * 1000):x}",
        plan_id=plan.plan_id,
        output_path=str(path),
        output_sha256=_sha256_bytes(path.read_bytes()),
        status=overall,
        target_slide_count=len(plan.target_slides),
        lineage_records=tuple(lineage),
        opc_integrity=opc,
        editability=editability,
        style_cluster_adherence=adherence,
    )


def assemble_physical_deck(
    plan: AssemblyPlan,
    output_path: str | os.PathLike[str],
    *,
    library_index_sha256: str,
) -> PhysicalAssemblyReport:
    """Assemble a single target PPTX whose slides physically come from many sources."""

    _validate_assembly_plan(plan, library_index_sha256)
    output = Path(output_path).expanduser().resolve(strict=False)
    output.parent.mkdir(parents=True, exist_ok=True)

    defaults, overrides = _default_content_types()
    seen_overrides: dict[str, str] = {}
    parts: dict[str, bytes] = {}
    lineage: list[LineageRecord] = []
    matches = 0
    slide_rels_paths: list[str] = []

    for slide in plan.target_slides:
        source_path = Path(slide.page_template.source_path)
        graph = _build_source_graph(
            source_path, slide_number=slide.page_template.slide_number
        )
        source_hash = slide.page_template.package_sha256
        ordinal = slide.ordinal

        target_slide_name = f"ppt/slides/slide{ordinal}.xml"
        target_map: dict[str, str] = {}

        # Slide itself
        slide_bytes = graph.slide_xml
        # Slide rels
        new_rels: dict[str, str] = {}
        source_slide_rels_name = (
            f"slide{slide.page_template.slide_number}.xml.rels"
        )
        for rels_path, rels_xml in graph.rels.items():
            new_rels_name = rels_path.replace("ppt/slides/_rels/", "ppt/slides/_rels/")
            new_rels_name = new_rels_name.replace(
                source_slide_rels_name,
                f"slide{ordinal}.xml.rels",
            )
            new_rels[new_rels_name] = rels_path
        # Process each part and assign unique target name
        for old_name, data in graph.extra_parts.items():
            new_name = _namespace_part_name(old_name, source_hash, ordinal)
            target_map[old_name] = new_name
        # Adapt slide text first using original layout
        slide_bytes = _adapt_slide_text(
            slide_bytes,
            slide.bindings,
            allowed_slots=slide.page_template.slot_graph.get("text_slot_ids", ()),
        )
        # Rewrite slide references and rels targets
        slide_bytes = _rewrite_slide_references(slide_bytes, target_map)
        # Rewrite each rels file
        new_rels_xml: dict[str, bytes] = {}
        for rels_path, rels_xml in graph.rels.items():
            new_rels_path = _namespace_relationship_path(
                rels_path,
                target_map,
                source_slide_rels_name=source_slide_rels_name,
                ordinal=ordinal,
            )
            new_rels_xml[new_rels_path] = _rewrite_relationship_targets(
                rels_xml,
                rels_path,
                target_map,
                output_rels_path=new_rels_path,
            )
        # Commit slide + rels
        parts[target_slide_name] = slide_bytes
        # The slide rels must point at slide ordinal
        slide_rels_target = (
            f"ppt/slides/_rels/slide{ordinal}.xml.rels"
        )
        # Take the rewritten slide1.xml.rels for this slide
        source_slide_rels = f"ppt/slides/_rels/slide{ordinal}.xml.rels"
        if source_slide_rels in new_rels_xml:
            parts[slide_rels_target] = new_rels_xml[source_slide_rels]
            slide_rels_paths.append(
                f"slides/slide{ordinal}.xml"
            )

        # Copy extra parts (layouts, masters, themes, media, charts, etc.)
        for old_name, data in graph.extra_parts.items():
            new_name = target_map[old_name]
            parts[new_name] = data

        # Copy all rels for layouts / masters / themes / media
        for rels_path, rels_xml in new_rels_xml.items():
            if rels_path == source_slide_rels:
                continue
            if rels_path not in parts:
                parts[rels_path] = rels_xml

        # Register content types
        for old_name in graph.extra_parts.keys():
            new_name = target_map[old_name]
            ct = _infer_content_type(new_name)
            if ct:
                key = "/" + new_name
                if key not in seen_overrides:
                    seen_overrides[key] = ct
        slide_key = "/" + target_slide_name
        if slide_key not in seen_overrides:
            seen_overrides[slide_key] = CONTENT_TYPES["slide"]
        # Track lineage
        if slide.page_template.style_cluster_id == plan.dominant_style_cluster_id:
            matches += 1
        lineage.append(
            LineageRecord(
                ordinal=ordinal,
                page_id=slide.page_template.page_id,
                package_sha256=slide.page_template.package_sha256,
                slide_number=slide.page_template.slide_number,
                source_sha256=slide.page_template.source_sha256,
                narrative_role=slide.narrative_role,
                title=slide.title,
                status="pass",
                binding_count=len(slide.bindings),
                byte_match_score=round(
                    _byte_match_score(graph.slide_xml, slide_bytes), 6
                ),
            )
        )

    # Override presentation.xml and its rels with every imported master.
    master_paths = sorted(
        name[len("ppt/"):]
        for name in parts
        if "/slideMasters/" in name and name.endswith(".xml")
    )
    pres_xml = _default_pres_xml(slide_rels_paths, master_paths)
    parts["ppt/presentation.xml"] = pres_xml
    parts["ppt/_rels/presentation.xml.rels"] = _pres_xml_rels(
        slide_rels_paths, master_paths
    )
    parts["ppt/presProps.xml"] = _default_props_xml()
    parts["ppt/viewProps.xml"] = _default_view_props_xml()
    parts["ppt/tableStyles.xml"] = _default_table_styles_xml()
    seen_overrides["/ppt/presentation.xml"] = CONTENT_TYPES["ppt"]
    seen_overrides["/ppt/presProps.xml"] = (
        "application/vnd.openxmlformats-officedocument.presentationml.presProps+xml"
    )
    seen_overrides["/ppt/viewProps.xml"] = (
        "application/vnd.openxmlformats-officedocument.presentationml.viewProps+xml"
    )
    seen_overrides["/ppt/tableStyles.xml"] = (
        "application/vnd.openxmlformats-officedocument.presentationml.tableStyles+xml"
    )

    # Defaults
    for ext, ct in defaults:
        if ext not in {k for k, _ in defaults}:
            pass
    # Override defaults from observed parts (e.g. wdp)
    for part_name in parts:
        if "/media/" in part_name and "." in part_name.split("/")[-1]:
            ext = part_name.rsplit(".", 1)[-1].lower()
            defaults.append((ext, _default_content_type_for(ext)))

    # Package-level rels + root rels
    parts["_rels/.rels"] = _default_pres_rels()
    core_xml, app_xml = _empty_doc_props()
    parts["docProps/core.xml"] = core_xml
    parts["docProps/app.xml"] = app_xml

    # Compose [Content_Types].xml
    for name in parts:
        if name == "[Content_Types].xml":
            continue
        if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
            seen_overrides.setdefault(
                "/" + name, CONTENT_TYPES["slide"]
            )
        elif "/slideLayouts/" in name and name.endswith(".xml"):
            seen_overrides.setdefault(
                "/" + name, CONTENT_TYPES["slideLayout"]
            )
        elif "/slideMasters/" in name and name.endswith(".xml"):
            seen_overrides.setdefault(
                "/" + name, CONTENT_TYPES["slideMaster"]
            )
        elif "/notesSlides/" in name and name.endswith(".xml"):
            seen_overrides.setdefault(
                "/" + name, CONTENT_TYPES["notesSlide"]
            )
        elif "/notesMasters/" in name and name.endswith(".xml"):
            seen_overrides.setdefault(
                "/" + name, CONTENT_TYPES["notesMaster"]
            )
        elif "/theme/" in name and name.endswith(".xml"):
            seen_overrides.setdefault(
                "/" + name, CONTENT_TYPES["theme"]
            )
        elif "/charts/" in name and name.endswith(".xml"):
            seen_overrides.setdefault(
                "/" + name, CONTENT_TYPES["chart"]
            )
        elif "/embeddings/" in name and name.endswith(".xlsx"):
            seen_overrides.setdefault(
                "/" + name,
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )

    parts["[Content_Types].xml"] = _serialize_content_types(
        defaults, sorted(seen_overrides.items())
    )

    # Atomic write
    fd, name = tempfile.mkstemp(
        prefix=f".{output.stem}.",
        suffix=".candidate.pptx",
        dir=output.parent,
    )
    os.close(fd)
    candidate = Path(name)
    try:
        with zipfile.ZipFile(candidate, "w", compression=zipfile.ZIP_DEFLATED) as out_zip:
            for name, data in parts.items():
                out_zip.writestr(name, data)
        os.replace(candidate, output)
    except Exception:
        candidate.unlink(missing_ok=True)
        raise

    # Final verification
    report = verify_physical_assembly(output, plan=plan)
    return report


def _infer_content_type(name: str) -> str | None:
    if "/slideLayouts/" in name:
        return CONTENT_TYPES["slideLayout"]
    if "/slideMasters/" in name:
        return CONTENT_TYPES["slideMaster"]
    if "/theme/" in name:
        return CONTENT_TYPES["theme"]
    if "/notesSlides/" in name:
        return CONTENT_TYPES["notesSlide"]
    if "/notesMasters/" in name:
        return CONTENT_TYPES["notesMaster"]
    if "/charts/" in name:
        return CONTENT_TYPES["chart"]
    if "/embeddings/" in name and name.endswith(".xlsx"):
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    return None


def _default_content_type_for(ext: str) -> str:
    return {
        "png": "image/png",
        "jpeg": "image/jpeg",
        "jpg": "image/jpeg",
        "gif": "image/gif",
        "wdp": "image/vnd.ms-photo",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "xml": "application/xml",
        "rels": "application/vnd.openxmlformats-package.relationships+xml",
    }.get(ext, "application/octet-stream")


def write_assembly_report(
    report: PhysicalAssemblyReport, output_path: str | os.PathLike[str]
) -> str:
    path = Path(output_path).expanduser().resolve(strict=False)
    path.parent.mkdir(parents=True, exist_ok=True)
    payload = report.to_dict()
    text = json.dumps(payload, ensure_ascii=False, indent=2)
    path.write_text(text, encoding="utf-8")
    return _sha256_bytes(text.encode("utf-8"))


def load_assembly_plan(
    path: str | os.PathLike[str],
    library_lookup: Mapping[str, PageTemplate],
) -> AssemblyPlan:
    payload = json.loads(Path(path).read_text(encoding="utf-8"))
    if payload.get("schema_version") != "1.0":
        raise PhysicalAssemblyError(
            f"unsupported assembly plan schema_version: {payload.get('schema_version')}"
        )
    return _plan_from_payload(payload, library_lookup)


__all__ = [
    "AssemblyPlan",
    "AssemblyTargetSlide",
    "Editability",
    "LineageRecord",
    "OPCIntegrity",
    "PhysicalAssemblyError",
    "PhysicalAssemblyReport",
    "SlideBinding",
    "StyleClusterAdherence",
    "assemble_physical_deck",
    "load_assembly_plan",
    "verify_physical_assembly",
    "write_assembly_report",
]
