#!/usr/bin/env python3
"""Build all Phase 46 anchors and materialize the exact work-summary spine."""

from __future__ import annotations

import argparse
import hashlib
import json
import subprocess
from pathlib import Path
from typing import Any

from pptx import Presentation

from window_pptx.template_pack import adapt_template_pack, load_template_bindings


SCRIPT_ROOT = Path(__file__).resolve().parent
SKILL_ROOT = SCRIPT_ROOT.parent
NODE_BUILDER = SCRIPT_ROOT / "build_window_pptx_v6_reference_anchors.mjs"
BASE_BINDINGS = SKILL_ROOT / "evals" / "v5.1-work-summary-bindings.json"
WORK_PAGES = (
    ("cover", "cover", "sparse"),
    ("agenda", "agenda", "medium"),
    ("section", "section", "sparse"),
    ("section", "section", "sparse"),
    ("data", "data-chart", "dense"),
    ("data", "data-chart", "dense"),
    ("table", "table", "dense"),
    ("case-study", "case-study", "dense"),
    ("metrics", "big-number", "dense"),
    ("section", "section", "sparse"),
    ("process", "process", "medium"),
    ("organization", "team", "dense"),
    ("section", "section", "sparse"),
    ("roadmap", "roadmap", "dense"),
    ("closing", "cta", "sparse"),
)


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def add_work_provenance_notes(output: Path, brief: dict[str, Any]) -> None:
    """Attach deterministic, editable speaker-note provenance to every slide."""

    fact_ids = sorted(
        item["id"]
        for item in brief["fact_store"]["facts"]
        if isinstance(item, dict) and isinstance(item.get("id"), str)
    )
    presentation = Presentation(output)
    for index, slide in enumerate(presentation.slides, start=1):
        slide.notes_slide.notes_text_frame.text = "\n".join(
            (
                f"FACT_IDS: {', '.join(fact_ids)}",
                f"COMPOSITION: exact governed TemplatePack page s{index:02d}",
                "All template text, shapes, charts, tables, and diagrams remain native editable objects.",
            )
        )
    presentation.save(output)


def facts(brief: dict[str, Any]) -> dict[str, Any]:
    return {
        item["id"]: item.get("value")
        for item in brief["fact_store"]["facts"]
        if isinstance(item, dict) and isinstance(item.get("id"), str)
    }


def work_bindings(brief: dict[str, Any]) -> dict[str, str]:
    pack_id, bindings = load_template_bindings(BASE_BINDINGS)
    if pack_id != "institutional-work-summary-v1":
        raise ValueError("unexpected work-summary binding pack")
    result = dict(bindings)
    f = facts(brief)
    title = "经营能力复利与稳健增长"
    for index, char in enumerate(title, start=1):
        result[f"s01.title.{index}"] = char
    result.update(
        {
            "s01.year": "2026\n年",
            "s01.en": "ANNUAL OPERATING\nREVIEW",
            "s01.presenter": "汇报：经营委员会",
            "s01.date": "日期：\n2026\n年\n12\n月\n18\n日",
            "s01.tagline": "Growth / Efficiency\nQuality / Governance",
            "s02.item.1": "年度经营答卷",
            "s02.item.2": "价值能力沉淀",
            "s02.item.3": "组织能力升级",
            "s02.item.4": "明年行动规划",
            "s03.kicker": "2026\n年 年度经营复盘",
            "s03.en": "OPERATING\nSCORECARD",
            "s04.kicker": "2026\n年 年度经营复盘",
            "s04.en": "VALUE\nDELIVERY",
            "s05.title": "2026\n年经营结果：\n增长",
            "s05.metric.1.label": "客户规模",
            "s05.metric.1.percent": "21.8\n%",
            "s05.metric.1.value": f"{f['customers-actual']}\n家客户",
            "s05.metric.2.label": "自助分析",
            "s05.metric.2.percent": f"{f['self-service-actual']}\n%",
            "s05.metric.2.value": "采用率\n达标",
            "s05.metric.3.label": "平台可用",
            "s05.metric.3.percent": f"{f['availability-q4']}\n%",
            "s05.metric.3.value": "全年\n稳定",
            "s05.metric.4.label": "年度节省",
            "s05.metric.4.percent": "23.3\n%",
            "s05.metric.4.value": f"{f['annual-savings']}\n百万元",
            "s05.metric.5.label": "重大事故",
            "s05.metric.5.percent": "66.7\n%",
            "s05.metric.5.value": f"{f['p1-incidents-q4']}\n起",
            "s05.callout.1": "增长主引擎",
            "s05.callout.2": "质量底座",
            "s06.title": "2026\n年经营质量：\n验证",
            "s06.budget.label": "预算执行率",
            "s06.budget.percent": "96.7\n%",
            "s06.value.1": str(f["monthly-active-users-q4"]),
            "s06.unit.1": "人",
            "s06.value.2": str(f["weekly-queries"]),
            "s06.unit.2": "次",
            "s06.value.3": str(f["governed-metrics"]),
            "s06.unit.3": "项",
            "s06.growth.label": "客户增长",
            "s06.growth.percent": "104.4\n%",
            "s06.pie.1.title": "本轮 - 预算结构",
            "s06.pie.2.title": "基线 - 预算结构",
            "s07.title": "2026\n年经营底盘：\n效率",
            "s07.category.1": "客户增长",
            "s07.category.2": "采用深化",
            "s07.category.3": "指标治理",
            "s07.category.4": "平台质量",
            "s07.category.5": "成本效率",
            "s07.total.1.label": "年度预算",
            "s07.total.1.value": str(f["budget-actual"]),
            "s07.total.1.unit": "百万",
            "s07.total.2.label": "年度节省",
            "s07.total.2.value": str(f["annual-savings"]),
            "s07.total.2.unit": "百万",
            "s07.total.label": "自助分析采用率",
            "s07.total.value": str(f["self-service-actual"]),
            "s07.total.unit": "%",
            "s08.title": "2026\n年重点项目：\n价值落地",
            "s08.hero": str(f["customers-actual"]),
            "s08.kicker": "2026\n年经营能力沉淀",
            "s08.subtitle": "指标治理 / 客户迁移\n采用深化",
            "s08.project.1": "指标治理",
            "s08.project.1.value": str(f["governed-metrics"]),
            "s08.project.1.unit": "项",
            "s08.project.1.detail": "统一口径\n责任闭环",
            "s08.project.2": "重复压缩",
            "s08.project.2.value": str(f["duplicates-after"]),
            "s08.project.2.unit": "项",
            "s08.project.2.detail": "从126项\n压缩至41",
            "s08.project.3": "项目准时",
            "s08.project.3.value": str(f["projects-on-time"]),
            "s08.project.3.unit": "个",
            "s08.project.3.detail": "12个项目\n10个准时",
            "s08.project.4": "客户迁移",
            "s08.project.4.value": str(f["unmigrated-customers"]),
            "s08.project.4.unit": "家",
            "s08.project.4.detail": "重点客户\n明年收口",
            "s08.project.5": "经营节省",
            "s08.project.5.value": str(f["annual-savings"]),
            "s08.project.5.unit": "百万",
            "s08.project.5.detail": "年度化\n财务确认",
            "s08.callout": "能力投资带来",
            "s08.hero.unit.1": "客户",
            "s08.hero.unit.2": "客户",
            "s09.title": "2026\n年核心经营：\n指标",
            "s09.metric.1": "客户数",
            "s09.metric.1.value": str(f["customers-actual"]),
            "s09.metric.1.unit": "家",
            "s09.metric.1.delta": "4.4\n%",
            "s09.metric.2": "自助分析率",
            "s09.metric.2.value": str(f["self-service-actual"]),
            "s09.metric.2.unit": "%",
            "s09.metric.2.delta": "3.0\n%",
            "s09.metric.3": "平台可用性",
            "s09.metric.3.value": f"{f['availability-q4']}\n%",
            "s09.metric.3.delta": "0.13\n%",
            "s09.metric.4": "报表 SLA",
            "s09.metric.4.value": f"{f['report-sla-q4']}\n%",
            "s09.metric.4.delta": "8.0\n%",
            "s09.metric.5": "客户 NPS",
            "s09.metric.5.value": str(f["nps-actual"]),
            "s09.metric.5.unit": "分",
            "s09.metric.5.delta": "3\n分",
            "s09.metric.6": "治理指标",
            "s09.metric.6.value": str(f["governed-metrics"]),
            "s09.metric.6.unit": "项",
            "s09.metric.6.delta": "38",
            "s09.metric.6.delta_unit": "项",
            "s09.metric.7": "P1事故",
            "s09.metric.7.value": f"{f['p1-incidents-q4']}\n起",
            "s09.metric.7.delta": "66.7\n%",
            "s09.summary": "2026\n客户、采用、质量与效率同步改善，增长没有透支经营底盘",
            "s09.badge": "兑现",
            "s10.kicker": "2026\n年 组织能力复盘",
            "s10.en": "ORGANIZATION\nCAPABILITY",
            "s11.headline": "把增长经验沉淀为可复制、可治理、可持续的经营能力。",
            "s11.title": "2026\n年\n组织升级",
            "s11.subtitle": "让指标、责任与复盘机制真正运转",
            "s11.point.1.title": "指标责任进入经营机制",
            "s11.point.1.body": "统一口径、责任人与变更机制，避免增长建立在冲突事实之上。",
            "s11.point.2.title": "客户迁移形成专项闭环",
            "s11.point.2.body": "对三家待迁移客户建立里程碑、风险和价值确认机制。",
            "s11.point.3.title": "采用率成为产品指标",
            "s11.owner": "经营能力负责人",
            "s12.title": "2026\n年\n组织能力升级",
            "s12.outcome.1": "增长可解释",
            "s12.outcome.2": "能力可复制",
            "s12.outcome.3": "责任可追溯",
            "s12.team": "经营能力团队",
            "s12.coverage": "覆盖                       经营场景",
            "s12.count.1": str(f["governed-metrics"]),
            "s12.count.2": str(f["customers-actual"]),
            "s13.kicker": "2027\n年 重点行动规划",
            "s13.en": "ROADMAP\nFOR 2027",
            "s14.title": "明年\n行动规划",
            "s14.summary.1": "聚焦治理、迁移和采用深化，把今年验证的能力规模化。",
            "s14.summary.2": "资源不平均分配；每项行动都绑定季度里程碑、责任人与验收证据。",
            "s14.domain.label": "经营能力",
            "s14.item.1": "完成17项无主\n指标定责",
            "s14.item.2": "收口3家重点\n客户迁移",
            "s14.item.3": "深化核心团队\n自助分析",
            "s14.item.4": "形成统一经营\n复盘机制",
            "s14.item.5": "复制指标治理\n交付套件",
            "s14.item.6": "建立采用率\n运营机制",
            "s14.item.7": "持续优化成本\n与平台质量",
            "s14.action.1": "定责",
            "s14.action.2": "迁移",
            "s14.action.3": "采用",
            "s14.action.4": "复盘",
            "s14.action.5": "复制",
            "s14.action.6": "运营",
            "s14.action.7": "质量",
            "s15.year": "2026\n年",
            "s15.presenter": "汇报：经营委员会",
            "s15.date": "日期：\n2026\n年\n12\n月\n18\n日",
            "c01.series": "经营结构",
            "c01.category.1": "自助分析",
            "c01.category.2": "平台质量",
            "c01.category.3": "客户增长",
            "c01.category.4": "成本效率",
            "c01.category.5": "事故下降",
            "c01.value.1": "0.32",
            "c01.value.2": "0.24",
            "c01.value.3": "0.22",
            "c01.value.4": "0.14",
            "c01.value.5": "0.08",
            "c02.series": "经营指标",
            "c02.category.1": "客户目标",
            "c02.category.2": "客户实际",
            "c02.category.3": "采用目标",
            "c02.value.1": str(f["customers-target"]),
            "c02.value.2": str(f["customers-actual"]),
            "c02.value.3": str(f["self-service-actual"]),
        }
    )
    for prefix, text in (
        ("s03.title", "年度经营答卷"),
        ("s04.title", "价值能力沉淀"),
        ("s10.title", "组织能力升级"),
        ("s13.title", "明年行动规划"),
        ("s15.title", "能力复利感谢同行"),
    ):
        for index, char in enumerate(text, start=1):
            result[f"{prefix}.{index}"] = char
    return result


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--brief-dir", type=Path, required=True)
    parser.add_argument("--output-dir", type=Path, required=True)
    parser.add_argument("--asset-dir", type=Path, required=True)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    subprocess.run(
        [
            "node",
            str(NODE_BUILDER),
            "--brief-dir",
            str(args.brief_dir),
            "--output-dir",
            str(args.output_dir),
            "--asset-dir",
            str(args.asset_dir),
        ],
        check=True,
    )
    brief_path = args.brief_dir / "annual-work-report.project-brief-pack.v1.json"
    brief = json.loads(brief_path.read_text(encoding="utf-8"))
    output = args.output_dir / "annual-work-report-reference-anchor.pptx"
    report = adapt_template_pack(
        "institutional-work-summary-v1",
        work_bindings(brief),
        output,
    )
    add_work_provenance_notes(output, brief)
    output_sha256 = sha256(output)
    manifest_path = output.with_suffix(".pptx.manifest.json")
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    manifest["output_sha256"] = output_sha256
    manifest["theme_id"] = "institutional-work-summary-exact-spine"
    manifest["generated_hero_assets"] = []
    manifest["candidate_policy"]["template_prefix"] = (
        "exact user-authorized TemplatePack v1 physical slide materialization"
    )
    manifest["slides"] = [
        {
            "slide_id": f"work-{index:02d}",
            "role": role,
            "family": family,
            "density": density,
            "candidate_id": f"template:physical.work-summary.s{index:02d}",
        }
        for index, (role, family, density) in enumerate(WORK_PAGES, start=1)
    ]
    manifest["template_materialization"] = {
        "template_pack_id": report.template_pack_id,
        "source_sha256": report.source_sha256,
        "output_sha256": output_sha256,
        "changed_parts": list(report.changed_parts),
        "slot_change_count": len(report.slot_changes),
        "source_integrity_preserved": report.source_integrity_preserved,
        "whole_slide_rasterization": False,
    }
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    print(
        json.dumps(
            {
                "status": "PASS",
                "work_output": str(output),
                "work_sha256": output_sha256,
                "template_pack": report.template_pack_id,
            }
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
