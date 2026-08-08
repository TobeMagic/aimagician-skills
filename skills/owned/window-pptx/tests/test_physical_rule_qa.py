"""Focused contract tests for deterministic physical-output rule QA."""

from __future__ import annotations

import hashlib
import sys
from pathlib import Path
from types import SimpleNamespace

from pptx import Presentation

SKILL_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_ROOT / "scripts"))

from window_pptx import physical_rule_qa  # noqa: E402


def _blank_pptx(path: Path) -> None:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(path)


def _one_slide_plan() -> SimpleNamespace:
    return SimpleNamespace(target_slide_count=1, target_slides=())


def test_report_binds_stable_bytes_and_documents_relative_path_policy(
    tmp_path: Path,
    monkeypatch,
) -> None:
    deck = tmp_path / "deck.pptx"
    _blank_pptx(deck)
    monkeypatch.chdir(tmp_path)

    report = physical_rule_qa.run_physical_rule_qa("deck.pptx", plan=_one_slide_plan())
    payload = report.to_dict()

    assert report.schema_version == "1.1"
    assert report.status == "pass"
    assert report.output_path == str(deck.resolve())
    assert report.output_sha256 == hashlib.sha256(deck.read_bytes()).hexdigest()
    assert report.output_size_bytes == deck.stat().st_size
    assert report.output_identity_status == "verified-stable"
    assert payload["path_policy"] == {
        "input_path_kind": "relative",
        "relative_input_base": str(tmp_path.resolve()),
        "stored_path_format": "canonical-absolute",
        "canonicalization": "expanduser+resolve(strict=false)",
        "relative_input_resolution": "invocation-working-directory",
    }


def test_content_drift_requires_a_new_report_identity(tmp_path: Path) -> None:
    deck = tmp_path / "deck.pptx"
    _blank_pptx(deck)

    before = physical_rule_qa.run_physical_rule_qa(deck, plan=_one_slide_plan())
    with deck.open("ab") as handle:
        handle.write(b"post-qa-drift")
    after = physical_rule_qa.run_physical_rule_qa(deck, plan=_one_slide_plan())

    assert before.status == "pass"
    assert after.status == "pass"
    assert before.output_sha256 != after.output_sha256
    assert before.output_size_bytes != after.output_size_bytes
    assert before.output_sha256 != hashlib.sha256(deck.read_bytes()).hexdigest()
    assert after.output_sha256 == hashlib.sha256(deck.read_bytes()).hexdigest()


def test_change_during_qa_fails_closed_and_withholds_identity(
    tmp_path: Path,
    monkeypatch,
) -> None:
    deck = tmp_path / "deck.pptx"
    _blank_pptx(deck)
    actual_sha = hashlib.sha256(deck.read_bytes()).hexdigest()
    observations = iter(((actual_sha, deck.stat().st_size), ("f" * 64, deck.stat().st_size)))
    monkeypatch.setattr(physical_rule_qa, "_fingerprint_output", lambda _path: next(observations))

    report = physical_rule_qa.run_physical_rule_qa(deck, plan=_one_slide_plan())

    assert report.status == "fail"
    assert report.output_sha256 is None
    assert report.output_size_bytes is None
    assert report.output_identity_status == "changed-during-qa"
    assert any(item.rule == "output-identity" and item.severity == "blocker" for item in report.blocking_findings)


def test_missing_output_cannot_pass_or_claim_a_fingerprint(tmp_path: Path) -> None:
    report = physical_rule_qa.run_physical_rule_qa(tmp_path / "missing.pptx", plan=_one_slide_plan())

    assert report.status == "fail"
    assert report.output_sha256 is None
    assert report.output_size_bytes is None
    assert report.output_identity_status == "unavailable"
    assert report.path_policy is not None
    assert report.path_policy.input_path_kind == "absolute"
