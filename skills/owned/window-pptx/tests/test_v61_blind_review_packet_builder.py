"""Focused tests for the fail-closed Phase 49 blind-review packet builder."""

from __future__ import annotations

import hashlib
import json
import subprocess
import sys
from pathlib import Path
from typing import Any

import jsonschema
import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


SKILL_ROOT = Path(__file__).resolve().parents[1]
SCHEMA_PATH = (
    SKILL_ROOT / "schemas" / "phase49-blind-review-packet.v1.schema.json"
)
sys.path.insert(0, str(SKILL_ROOT / "scripts"))

import build_window_pptx_v61_blind_packet as builder  # noqa: E402


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _write_deck(path: Path, label: str, *, slide_count: int = 15) -> None:
    presentation = Presentation()
    for ordinal in range(1, slide_count + 1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        box.text = f"{label} {ordinal}"
    presentation.save(path)


def _write_json(path: Path, payload: object) -> None:
    path.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def _physical_report(candidate: Path) -> dict[str, Any]:
    candidate_sha = _sha256(candidate)
    return {
        "schema_version": "1.0",
        "status": "pass",
        "acceptance_profile": "phase49-work-report-15",
        "output_sha256": candidate_sha,
        "expected_slide_count": 15,
        "target_slide_count": 15,
        "distinct_page_id_count": 15,
        "duplicate_page_records": [],
        "lineage_records": [
            {
                "ordinal": ordinal,
                "page_id": f"certified-page-{ordinal:02d}",
                "source_package_verified": True,
                "source_slide_verified": True,
                "structure_match": True,
                "status": "pass",
            }
            for ordinal in range(1, 16)
        ],
        "opc_integrity": {"status": "pass"},
        "editability": {
            "status": "pass",
            "native_editable": True,
            "native_editable_coverage": 1,
        },
        "style_cluster_adherence": {"status": "pass"},
        "authority": {"status": "pass"},
        "selection_authority": {"status": "pass"},
        "source_residue": {"status": "pass"},
        "libreoffice": {
            "status": "pass",
            "open_result": "pass",
            "render_result": "pass",
        },
        "size_check": {"status": "pass"},
    }


def _rule_qa_report(candidate: Path) -> dict[str, Any]:
    return {
        "schema_version": "1.1",
        "status": "pass",
        "output_sha256": _sha256(candidate),
        "output_size_bytes": candidate.stat().st_size,
        "output_identity_status": "verified-stable",
        "slide_count": 15,
        "blocking_findings": [],
    }


def _fake_tools(
    monkeypatch: pytest.MonkeyPatch,
    *,
    raster_mode: str = "complete",
) -> None:
    def fake_tool_record(
        _explicit: Path | None,
        *names: str,
        version_args: tuple[str, ...],
    ) -> dict[str, Any]:
        del version_args
        identity = names[0]
        return {
            "executable": f"/test-tools/{identity}",
            "version": f"{identity} test-version",
            "sha256": hashlib.sha256(identity.encode("utf-8")).hexdigest(),
            "size_bytes": 1,
        }

    def fake_run(
        command: list[str] | tuple[str, ...],
        *,
        label: str,
    ) -> subprocess.CompletedProcess[str]:
        del label
        values = list(command)
        if "--convert-to" in values:
            output_dir = Path(values[values.index("--outdir") + 1])
            source = Path(values[-1])
            (output_dir / f"{source.stem}.pdf").write_bytes(
                f"%PDF-1.7 fake {source.stem}".encode("ascii")
            )
        elif "-png" in values:
            prefix = Path(values[-1])
            label_name = Path(values[-2]).stem
            page_count = 14 if raster_mode == "missing" else 15
            color = (40, 90, 170) if label_name == "reference" else (20, 130, 95)
            for ordinal in range(1, page_count + 1):
                Image.new("RGB", (320, 180), color).save(
                    prefix.with_name(f"page-{ordinal}.png")
                )
            if raster_mode == "duplicate":
                Image.new("RGB", (320, 180), color).save(
                    prefix.with_name("page-01.png")
                )
        else:  # pragma: no cover - catches a builder command regression
            raise AssertionError(f"unexpected command: {values}")
        return subprocess.CompletedProcess(values, 0, "", "")

    monkeypatch.setattr(builder, "_tool_record", fake_tool_record)
    monkeypatch.setattr(builder, "_run_command", fake_run)


def _fixture(tmp_path: Path) -> dict[str, Path | str]:
    reference = tmp_path / "reference.pptx"
    candidate = tmp_path / "candidate.pptx"
    physical = tmp_path / "physical.json"
    rule_qa = tmp_path / "rule-qa.json"
    _write_deck(reference, "REFERENCE")
    _write_deck(candidate, "CANDIDATE")
    _write_json(physical, _physical_report(candidate))
    _write_json(rule_qa, _rule_qa_report(candidate))
    return {
        "reference": reference,
        "candidate": candidate,
        "reference_sha": _sha256(reference),
        "candidate_sha": _sha256(candidate),
        "physical": physical,
        "rule_qa": rule_qa,
        "output": tmp_path / "packet",
    }


def _build(paths: dict[str, Path | str]) -> dict[str, Any]:
    return builder.build_phase49_blind_review_packet(
        reference_pptx=paths["reference"],  # type: ignore[arg-type]
        candidate_pptx=paths["candidate"],  # type: ignore[arg-type]
        reference_sha256=paths["reference_sha"],  # type: ignore[arg-type]
        candidate_sha256=paths["candidate_sha"],  # type: ignore[arg-type]
        physical_report=paths["physical"],  # type: ignore[arg-type]
        rule_qa_report=paths["rule_qa"],  # type: ignore[arg-type]
        expected_slide_count=15,
        dpi=144,
        output_dir=paths["output"],  # type: ignore[arg-type]
    )


def test_builds_schema_valid_hash_bound_complete_eight_image_packet(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    paths = _fixture(tmp_path)
    _fake_tools(monkeypatch)

    payload = _build(paths)
    packet_root = paths["output"]
    assert isinstance(packet_root, Path)
    stored = json.loads((packet_root / "packet.json").read_text(encoding="utf-8"))
    schema = json.loads(SCHEMA_PATH.read_text(encoding="utf-8"))
    jsonschema.Draft202012Validator.check_schema(schema)
    jsonschema.validate(stored, schema)

    assert stored == payload
    basis = dict(stored)
    packet_sha256 = basis.pop("packet_sha256")
    assert packet_sha256 == builder._canonical_sha256(basis)
    assert stored["coverage"] == {
        "expected_slide_ordinals": list(range(1, 16)),
        "reference_slide_ordinals": list(range(1, 16)),
        "candidate_slide_ordinals": list(range(1, 16)),
        "pair_slide_ordinals": list(range(1, 16)),
        "missing_slide_ordinals": [],
        "duplicate_slide_ordinals": [],
        "status": "pass",
    }
    assert len(stored["pairs"]) == 8
    assert [item["slide_ordinals"] for item in stored["pairs"]] == [
        [1, 2],
        [3, 4],
        [5, 6],
        [7, 8],
        [9, 10],
        [11, 12],
        [13, 14],
        [15],
    ]
    for pair in stored["pairs"]:
        pair_path = packet_root / pair["path"]
        assert pair_path.is_file()
        assert pair["sha256"] == _sha256(pair_path)
        assert pair["size_bytes"] == pair_path.stat().st_size
        expected_labels = [
            (kind, ordinal, f"{kind.upper()}  /  SLIDE {ordinal:02d}")
            for ordinal in pair["slide_ordinals"]
            for kind in ("reference", "candidate")
        ]
        assert [
            (label["kind"], label["ordinal"], label["label"])
            for label in pair["labels"]
        ] == expected_labels
    for kind in ("reference", "candidate"):
        pages = stored["renders"][kind]["pages"]
        assert len(pages) == 15
        for page in pages:
            page_path = packet_root / page["path"]
            assert page["sha256"] == _sha256(page_path)
            assert (page["width_px"], page["height_px"]) == (320, 180)
    assert _sha256(packet_root / stored["inputs"]["physical_report"]["packet_path"]) == stored[
        "inputs"
    ]["physical_report"]["observed_sha256"]
    assert _sha256(packet_root / stored["inputs"]["rule_qa_report"]["packet_path"]) == stored[
        "inputs"
    ]["rule_qa_report"]["observed_sha256"]


@pytest.mark.parametrize("report_kind", ["physical", "rule_qa"])
def test_failed_evidence_report_prevents_packet_creation(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    report_kind: str,
) -> None:
    paths = _fixture(tmp_path)
    _fake_tools(monkeypatch)
    report_path = paths[report_kind]
    assert isinstance(report_path, Path)
    report = json.loads(report_path.read_text(encoding="utf-8"))
    report["status"] = "fail"
    _write_json(report_path, report)

    with pytest.raises(builder.BlindPacketError, match="status must be 'pass'"):
        _build(paths)
    assert not Path(paths["output"]).exists()


def test_candidate_hash_mismatch_prevents_packet_creation(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    paths = _fixture(tmp_path)
    _fake_tools(monkeypatch)
    paths["candidate_sha"] = "0" * 64

    with pytest.raises(builder.BlindPacketError, match="candidate PPTX SHA-256"):
        _build(paths)
    assert not Path(paths["output"]).exists()


@pytest.mark.parametrize(
    ("raster_mode", "message"),
    [
        ("missing", "coverage mismatch"),
        ("duplicate", "duplicate rendered page ordinal"),
    ],
)
def test_missing_or_duplicate_rendered_page_fails_atomically(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    raster_mode: str,
    message: str,
) -> None:
    paths = _fixture(tmp_path)
    _fake_tools(monkeypatch, raster_mode=raster_mode)

    with pytest.raises(builder.BlindPacketError, match=message):
        _build(paths)
    assert not Path(paths["output"]).exists()


def test_missing_renderer_tool_fails_before_output_is_created(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    paths = _fixture(tmp_path)

    def missing_tool(*_args: object, **_kwargs: object) -> dict[str, Any]:
        raise builder.BlindPacketError("required tool is unavailable: libreoffice/soffice")

    monkeypatch.setattr(builder, "_tool_record", missing_tool)
    with pytest.raises(builder.BlindPacketError, match="required tool is unavailable"):
        _build(paths)
    assert not Path(paths["output"]).exists()


def test_wrong_slide_count_fails_before_rendering(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    paths = _fixture(tmp_path)
    reference = paths["reference"]
    assert isinstance(reference, Path)
    _write_deck(reference, "REFERENCE", slide_count=14)
    paths["reference_sha"] = _sha256(reference)
    _fake_tools(monkeypatch)

    with pytest.raises(builder.BlindPacketError, match="reference PPTX has 14 slides"):
        _build(paths)
    assert not Path(paths["output"]).exists()
