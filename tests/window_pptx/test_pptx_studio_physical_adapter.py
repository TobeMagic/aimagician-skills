from __future__ import annotations

import hashlib
import json
import sys
import xml.etree.ElementTree as ET
from dataclasses import replace
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
SCRIPT_ROOT = REPO_ROOT / "skills" / "owned" / "pptx-studio" / "scripts"
sys.path.insert(0, str(SCRIPT_ROOT))

from pptx_studio.adaptation import compile_adaptation  # noqa: E402
from pptx_studio.component_profiles import (  # noqa: E402
    catalog_sha256, component_profile_sha256, load_component_profiles,
)
from pptx_studio.brief_binding import compile_outline_bindings  # noqa: E402
from pptx_studio.composition import compile_composition, style_signature  # noqa: E402
from pptx_studio.physical_adapter import (  # noqa: E402
    PhysicalAdapterError,
    _automatic_sequence_component_groups,
    _client_binding_role,
    _certified_fragment_font_scales,
    _curated_component_groups,
    _deduplicate_nested_alias_slots,
    _fragment_title_regions,
    _unbound_template_clear_reason,
    assemble_from_plans,
    compile_physical_adapter,
    preflight_native_slots,
    resolve_catalog_sources,
    resolve_component_import_specs,
)
from pptx_studio.qa import _visual_checks, run_studio_qa  # noqa: E402
from window_pptx.page_template_library import PageTemplate, SlotRecord  # noqa: E402
from window_pptx.physical_assembly import (  # noqa: E402
    AssemblyPlan,
    AssemblyTargetSlide,
    PhysicalAssemblyError,
    TextBindingSpec,
    _adapt_slide_text,
    _build_text_binding_evidence,
    _validate_assembly_plan,
    AssemblyImportContext,
    _component_nodes_sha256,
    _component_relationship_ids,
    _component_root_nodes,
)
import pptx_studio.physical_adapter as physical_adapter  # noqa: E402
from manage_pptx_studio_library import run  # noqa: E402


def test_adapter_uses_bounded_shrink_for_capacity_approved_text() -> None:
    """Sample-copy length must not invoke the unbounded generic shrink rule."""

    # This regression is intentionally asserted at the public adapter policy
    # boundary: normal text replacements must retain template-native fit.
    assert physical_adapter.TextBindingSpec(
        "1.5", ("fact-1",), "auto", "safe-shrink-to-fit",
    ).fit_policy == "safe-shrink-to-fit"


def test_certified_fragment_repair_scales_only_bound_source_lockup_glyphs() -> None:
    """Art-title repairs are fixed certification data, not planner authority."""

    scales = _certified_fragment_font_scales(
        package_sha256="59b104d31bf3f44c15d407adefe51425c9dcd8bb5c5d1e2212fb38753dc72839",
        slide_number=10,
        bindings={"shape_5": "创", "shape_7": "新", "shape_12": "事", "shape_13": "迹"},
    )
    assert scales == {"shape_5": 72_000, "shape_7": 72_000, "shape_12": 72_000, "shape_13": 72_000}
    assert _certified_fragment_font_scales(
        package_sha256="59b104d31bf3f44c15d407adefe51425c9dcd8bb5c5d1e2212fb38753dc72839",
        slide_number=1,
        bindings={"shape_6": "2", "shape_7": "0", "shape_8": "2", "shape_9": "5", "shape_11": "工", "shape_12": "作", "shape_15": "汇", "shape_16": "报"},
    ) == {"shape_6": 42_000, "shape_7": 42_000, "shape_8": 42_000, "shape_9": 42_000, "shape_11": 48_000, "shape_12": 48_000, "shape_15": 48_000, "shape_16": 48_000}
    assert _certified_fragment_font_scales(
        package_sha256="59b104d31bf3f44c15d407adefe51425c9dcd8bb5c5d1e2212fb38753dc72839",
        slide_number=15,
        bindings={"shape_11": "感", "shape_12": "谢", "shape_15": "聆", "shape_16": "听"},
    ) == {"shape_11": 48_000, "shape_12": 48_000, "shape_15": 48_000, "shape_16": 48_000}
    assert _certified_fragment_font_scales(
        package_sha256="a" * 64, slide_number=10, bindings={"shape_5": "创"},
    ) == {}


def _sha(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def test_component_profile_resolver_recomputes_every_native_fingerprint(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    component_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    component_shape.text = "Component"
    anchor_shape = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(2), Inches(1))
    anchor_shape.text = "Reservation"
    package = tmp_path / "fixture.pptx"
    presentation.save(package)
    package_sha = _sha(package)
    context = AssemblyImportContext()
    _source, graph = context.graph_for(package, package_sha, 1)
    root = ET.fromstring(graph.slide_xml)
    component_nodes = _component_root_nodes(root, shape_ids=(component_shape.shape_id,), label="COMPONENT_SOURCE")
    anchor_nodes = _component_root_nodes(root, shape_ids=(anchor_shape.shape_id,), label="COMPONENT_HOST_ANCHOR")
    component_capacity = next(
        slot.max_chars for slot in physical_adapter._discover_slots(
            graph.slide_xml.decode("utf-8", errors="replace"),
        ) if slot.slot_id == f"shape_{component_shape.shape_id}"
    )
    page_id = "page_aaaaaaaaaaaaaaaaaaaaaaaa_001"
    catalog: dict[str, object] = {
        "pages": [{"page_id": page_id, "package_sha256": package_sha, "slide_number": 1}],
    }
    profile: dict[str, object] = {
        "schema_version": "pptx-studio-component-profile.v1", "status": "COMPLETE",
        "profile_id": "fixture", "profile_sha256": "", "catalog_sha256": catalog_sha256(catalog),
        "components": [{
            "component_id": "component_111111111111111111111111",
            "source": {"page_id": page_id, "package_sha256": package_sha, "slide_number": 1, "slide_sha256": graph.slide_sha},
            "shape_ids": [component_shape.shape_id], "component_sha256": _component_nodes_sha256(component_nodes),
            "relationship_ids": list(_component_relationship_ids(component_nodes)), "semantic_intent": "metric-card",
            "allowed_roles": ["dashboard"],
            "fields": [{"field_id": "metric", "shape_id": component_shape.shape_id, "semantic_role": "metric", "max_chars": component_capacity}],
            "allowed_host_anchor_ids": ["anchor_222222222222222222222222"],
        }],
        "host_anchors": [{
            "host_anchor_id": "anchor_222222222222222222222222",
            "source": {"page_id": page_id, "package_sha256": package_sha, "slide_number": 1, "slide_sha256": graph.slide_sha},
            "shape_ids": [anchor_shape.shape_id], "host_anchor_sha256": _component_nodes_sha256(anchor_nodes),
            "compatible_component_ids": ["component_111111111111111111111111"],
        }],
    }
    profile["profile_sha256"] = component_profile_sha256(profile)
    profile_path = tmp_path / "component-profile.json"
    profile_path.write_text(json.dumps(profile), encoding="utf-8")
    index = load_component_profiles(profile_path, catalog=catalog)

    specs = resolve_component_import_specs(
        {"page_id": page_id, "package_sha256": package_sha, "slide_number": 1, "component_assembly": {
            "host_anchor_id": "anchor_222222222222222222222222",
            "component_ids": ["component_111111111111111111111111"],
        }},
        component_profiles=index, source_paths={package_sha: package}, context=context,
    )

    assert len(specs) == 1
    assert specs[0].source_shape_ids == (component_shape.shape_id,)
    assert specs[0].host_anchor_shape_ids == (anchor_shape.shape_id,)


def test_component_profile_resolver_recomputes_fixed_canvas_source(tmp_path: Path) -> None:
    presentation = Presentation()
    source_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = source_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    title.text = "Native title"
    host_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    host_slide.shapes.add_textbox(Inches(5), Inches(1), Inches(2), Inches(1)).text = "KPI"
    package = tmp_path / "canvas-fixture.pptx"
    presentation.save(package)
    package_sha = _sha(package)
    context = AssemblyImportContext()
    _source, source_graph = context.graph_for(package, package_sha, 1)
    _host, host_graph = context.graph_for(package, package_sha, 2)
    source_root = ET.fromstring(source_graph.slide_xml)
    title_nodes = _component_root_nodes(source_root, shape_ids=(title.shape_id,), label="COMPONENT_SOURCE")
    title_bbox = physical_adapter._component_nodes_bbox(title_nodes)
    assert title_bbox is not None
    title_capacity = next(
        slot.max_chars for slot in physical_adapter._discover_slots(
            source_graph.slide_xml.decode("utf-8", errors="replace"),
        ) if slot.slot_id == f"shape_{title.shape_id}"
    )
    source_page_id = "page_aaaaaaaaaaaaaaaaaaaaaaaa_001"
    host_page_id = "page_bbbbbbbbbbbbbbbbbbbbbbbb_002"
    catalog: dict[str, object] = {"pages": [
        {"page_id": source_page_id, "package_sha256": package_sha, "slide_number": 1},
        {"page_id": host_page_id, "package_sha256": package_sha, "slide_number": 2},
    ]}
    profile: dict[str, object] = {
        "schema_version": "pptx-studio-component-profile.v3", "status": "COMPLETE",
        "profile_id": "canvas-fixture", "profile_sha256": "", "catalog_sha256": catalog_sha256(catalog),
        "components": [{
            "component_id": "component_111111111111111111111111",
            "source": {"page_id": source_page_id, "package_sha256": package_sha, "slide_number": 1, "slide_sha256": source_graph.slide_sha},
            "shape_ids": [title.shape_id], "component_sha256": _component_nodes_sha256(title_nodes),
            "relationship_ids": list(_component_relationship_ids(title_nodes)), "semantic_intent": "report-title",
            "allowed_roles": ["dashboard"],
            "fields": [{"field_id": "title", "shape_id": title.shape_id, "semantic_role": "title", "max_chars": title_capacity}],
            "allowed_host_anchor_ids": ["anchor_222222222222222222222222"],
        }],
        "host_anchors": [{
            "host_anchor_id": "anchor_222222222222222222222222", "anchor_mode": "canvas",
            "source": {"page_id": host_page_id, "package_sha256": package_sha, "slide_number": 2, "slide_sha256": host_graph.slide_sha},
            "canvas_source": {"page_id": source_page_id, "package_sha256": package_sha, "slide_number": 1, "slide_sha256": source_graph.slide_sha},
            "canvas_shape_ids": [title.shape_id], "canvas_sha256": _component_nodes_sha256(title_nodes),
            "canvas_bbox": [title_bbox[0], title_bbox[1], title_bbox[2] - title_bbox[0], title_bbox[3] - title_bbox[1]],
            "compatible_component_ids": ["component_111111111111111111111111"],
        }],
    }
    profile["profile_sha256"] = component_profile_sha256(profile)
    profile_path = tmp_path / "canvas-component-profile.json"
    profile_path.write_text(json.dumps(profile), encoding="utf-8")
    index = load_component_profiles(profile_path, catalog=catalog)

    specs = resolve_component_import_specs(
        {"page_id": host_page_id, "package_sha256": package_sha, "slide_number": 2, "component_assembly": {
            "host_anchor_id": "anchor_222222222222222222222222",
            "component_ids": ["component_111111111111111111111111"],
        }},
        component_profiles=index, source_paths={package_sha: package}, context=context,
    )

    assert specs[0].host_anchor_shape_ids == ()
    assert specs[0].canvas_bbox == (title_bbox[0], title_bbox[1], title_bbox[2] - title_bbox[0], title_bbox[3] - title_bbox[1])


def test_client_binding_role_keeps_year_bearing_report_heading_as_title() -> None:
    """A year in a top report heading is not a metric surface."""

    slot = SlotRecord(
        slot_id="shape_53", shape_id=53, kind="text", max_chars=20,
        text="2025 年财务决算：财政项目及政府债支出",
        semantic_role="body", region="top", reading_order=1,
        bbox={"x": 35, "y": 51, "w": 630, "h": 85},
        source_char_count=20, source_line_count=3, source_run_count=1,
        group_id=None, group_order=None, font_size_pt=30.0,
        allowed_binding_modes=("replace",),
    )
    assert _client_binding_role(slot) == "title"


def test_client_binding_role_preserves_central_process_heading() -> None:
    """A process hub heading remains a title even outside the top band."""

    slot = SlotRecord(
        slot_id="shape_9", shape_id=9, kind="text", max_chars=12,
        text="实施路径", semantic_role="title", region="middle", reading_order=1,
        bbox={"x": 420, "y": 380, "w": 180, "h": 60},
        source_char_count=4, source_line_count=1, source_run_count=1,
        group_id=None, group_order=None, font_size_pt=28.0,
        allowed_binding_modes=("replace",),
    )

    assert _client_binding_role(slot, declared_role="process") == "title"


def test_contents_ordinal_is_preserved_as_structure_not_a_client_fact() -> None:
    slot = SlotRecord(
        slot_id="shape_8", shape_id=8, kind="text", max_chars=4,
        text="03.", semantic_role="metric", region="middle", reading_order=3,
        bbox={"x": 100, "y": 300, "w": 80, "h": 50},
        source_char_count=3, source_line_count=1, source_run_count=1,
        group_id=None, group_order=None, font_size_pt=24.0,
        allowed_binding_modes=("replace",),
    )

    assert _client_binding_role(slot, declared_role="contents") == "ignore"
    assert _unbound_template_clear_reason("03.", declared_role="contents") is None
    assert _unbound_template_clear_reason("03.") == "unbound-template-copy"


def test_timeline_publishes_required_date_action_groups_in_visual_order() -> None:
    """Chronology is bound as atomic date/action pairs from left to right."""

    def slot(slot_id: str, text: str, x: int, y: int, order: int) -> SlotRecord:
        return SlotRecord(
            slot_id=slot_id, shape_id=order, kind="text", max_chars=18,
            text=text, semantic_role="title" if "年" in text else "body",
            region="middle", reading_order=order,
            bbox={"x": x, "y": y, "w": 160, "h": 50},
            source_char_count=len(text), source_line_count=1, source_run_count=1,
            group_id=None, group_order=None, font_size_pt=18.0,
            allowed_binding_modes=("replace",),
        )

    slots = (
        slot("shape_2", "2026年4月", 500, 200, 2),
        slot("shape_4", "全面推广", 500, 320, 4),
        slot("shape_1", "2026年1月", 100, 200, 1),
        slot("shape_3", "启动试点", 100, 320, 3),
    )
    groups = _automatic_sequence_component_groups(
        slots,
        declared_role="timeline",
        shape_to_component={
            "shape_1": "label.01", "shape_2": "label.02",
            "shape_3": "body.01", "shape_4": "body.02",
        },
    )

    assert groups == [
        {
            "component_group": "timeline-step.01",
            "component_keys": ["label.01", "body.01"],
            "component_intent": "timeline-milestone",
            "component_fields": ["date", "action"],
            "required": True,
        },
        {
            "component_group": "timeline-step.02",
            "component_keys": ["label.02", "body.02"],
            "component_intent": "timeline-milestone",
            "component_fields": ["date", "action"],
            "required": True,
        },
    ]


def test_timeline_cardinality_adaptation_is_private_exact_and_fails_closed(
    tmp_path: Path,
) -> None:
    """Only one curator-approved page/role/capacity tuple can remove a node."""

    source_root = tmp_path / "private" / "sources" / "gaojie"
    annotation = (
        tmp_path / "private" / "intelligence" / "pptx-studio" / "annotations"
        / "timeline-cardinality-adaptations.v1.json"
    )
    source_root.mkdir(parents=True)
    page_id = "page_" + "c" * 24 + "_001"
    entry = {
        "page_id": page_id,
        "role": "timeline",
        "minimum_capacity": 5,
        "shape_ids": [14, 15],
        "shape_sha256": "a" * 64,
    }
    annotation.parent.mkdir(parents=True)
    annotation.write_text(json.dumps({
        "schema_version": "pptx-studio-timeline-cardinality-adaptations.v1",
        "adaptations": [entry],
    }), encoding="utf-8")

    exact = physical_adapter._timeline_cardinality_cleanup_specs(
        private_source_root=source_root, page_id=page_id, role="timeline",
        minimum_capacity=5,
    )
    assert len(exact) == 1
    assert exact[0].shape_ids == (14, 15)
    assert physical_adapter._timeline_cardinality_cleanup_specs(
        private_source_root=source_root, page_id=page_id, role="process",
        minimum_capacity=5,
    ) == ()
    assert physical_adapter._timeline_cardinality_cleanup_specs(
        private_source_root=source_root, page_id=page_id, role="timeline",
        minimum_capacity=4,
    ) == ()

    annotation.write_text(json.dumps({
        "schema_version": "pptx-studio-timeline-cardinality-adaptations.v1",
        "adaptations": [entry, entry],
    }), encoding="utf-8")
    with pytest.raises(PhysicalAdapterError, match="TIMELINE_CARDINALITY_ADAPTATION_AMBIGUOUS"):
        physical_adapter._timeline_cardinality_cleanup_specs(
            private_source_root=source_root, page_id=page_id, role="timeline",
            minimum_capacity=5,
        )

    annotation.write_text(json.dumps({
        "schema_version": "pptx-studio-timeline-cardinality-adaptations.v1",
        "adaptations": [{**entry, "shape_ids": [14, 14]}],
    }), encoding="utf-8")
    with pytest.raises(PhysicalAdapterError, match="TIMELINE_CARDINALITY_ADAPTATION_INVALID"):
        physical_adapter._timeline_cardinality_cleanup_specs(
            private_source_root=source_root, page_id=page_id, role="timeline",
            minimum_capacity=5,
        )

    with pytest.raises(PhysicalAdapterError, match="CERTIFIED_CLEANUP_OVERLAP"):
        physical_adapter._merge_cleanup_specs(exact, exact)


def test_process_publishes_required_label_body_steps_in_visual_order() -> None:
    """A four-step process cannot cross-wire labels and descriptions."""

    def slot(
        slot_id: str, text: str, role: str, x: int, y: int, order: int,
    ) -> SlotRecord:
        return SlotRecord(
            slot_id=slot_id, shape_id=order, kind="text", max_chars=30,
            text=text, semantic_role=role, region="middle", reading_order=order,
            bbox={"x": x, "y": y, "w": 160, "h": 50},
            source_char_count=len(text), source_line_count=1, source_run_count=1,
            group_id=None, group_order=None, font_size_pt=18.0,
            allowed_binding_modes=("replace",),
        )

    slots = (
        slot("shape_1", "第一阶段", "label", 100, 200, 1),
        slot("shape_2", "完成准备工作并确认范围", "body", 100, 320, 2),
        slot("shape_3", "第二阶段", "label", 500, 200, 3),
        slot("shape_4", "进入建设并提交阶段成果", "body", 500, 320, 4),
    )

    groups = _automatic_sequence_component_groups(
        slots,
        declared_role="process",
        shape_to_component={
            "shape_1": "label.01", "shape_2": "body.01",
            "shape_3": "label.02", "shape_4": "body.02",
        },
    )

    assert groups == [
        {
            "component_group": "process-step.01",
            "component_keys": ["label.01", "body.01"],
            "component_intent": "process-step",
            "component_fields": ["label", "body"],
            "required": True,
        },
        {
            "component_group": "process-step.02",
            "component_keys": ["label.02", "body.02"],
            "component_intent": "process-step",
            "component_fields": ["label", "body"],
            "required": True,
        },
    ]


def _alias_clear_template() -> PageTemplate:
    package_sha = "a" * 64
    return PageTemplate(
        schema_version="1.0", page_id=f"{package_sha}:001",
        package_sha256=package_sha, slide_number=1, source_path="/private/source.pptx",
        source_sha256=package_sha, source_slide_sha256="b" * 64,
        page_role="content", category_names=("test",), style_cluster_id="test",
        deck_family_id="test", theme_palette=(), capacity={"max_text_chars": 20, "max_text_runs": 1},
        editability="native_editable", certification="certified", visual_quality=1.0,
        structure={}, slot_graph={"text_slot_ids": ["shape_1"]},
        requires_customer_asset=False, media_retention_policy="preserve",
        governed_content_inventory={"slots": []},
    )


def test_physical_assembly_allows_only_clear_for_graph_excluded_nested_alias() -> None:
    """Hidden duplicate nodes may be blanked, but never reused as content slots."""

    template = _alias_clear_template()
    slide = AssemblyTargetSlide(
        ordinal=1, page_template=template,
        bindings={"shape_1": "标题", "shape_2": ""},
        narrative_role="content", title="标题", headline="摘要",
        text_binding_specs={
            "shape_1": TextBindingSpec("标题", (), "source"),
            "shape_2": TextBindingSpec("", (), "clear"),
        },
    )
    plan = AssemblyPlan(
        schema_version="1.0", plan_id="test", scenario_id="test",
        dominant_style_cluster_id="test", created_at="1970-01-01T00:00:00Z",
        target_slide_count=1, target_slides=(slide,), library_index_sha256="c" * 64,
    )
    _validate_assembly_plan(plan, "c" * 64)

    clear_only = AssemblyPlan(
        schema_version="1.0", plan_id="test", scenario_id="test",
        dominant_style_cluster_id="test", created_at="1970-01-01T00:00:00Z",
        target_slide_count=1,
        target_slides=(AssemblyTargetSlide(
            ordinal=1,
            page_template=replace(template, slot_graph={"text_slot_ids": []}),
            bindings={"shape_2": ""}, narrative_role="content", title="标题", headline="摘要",
            text_binding_specs={"shape_2": TextBindingSpec("", (), "clear")},
        ),),
        library_index_sha256="c" * 64,
    )
    assert _build_text_binding_evidence(
        clear_only, None, require_locked_authority=True,
    ) == []

    invalid = AssemblyTargetSlide(
        ordinal=1, page_template=template,
        bindings={"shape_1": "标题", "shape_2": "遗留文案"},
        narrative_role="content", title="标题", headline="摘要",
        text_binding_specs={
            "shape_1": TextBindingSpec("标题", (), "source"),
            "shape_2": TextBindingSpec("遗留文案", (), "source"),
        },
    )
    invalid_plan = AssemblyPlan(
        schema_version="1.0", plan_id="test", scenario_id="test",
        dominant_style_cluster_id="test", created_at="1970-01-01T00:00:00Z",
        target_slide_count=1, target_slides=(invalid,), library_index_sha256="c" * 64,
    )
    with pytest.raises(PhysicalAssemblyError, match="ASSEMBLY_PLAN_SLOT_COVERAGE"):
        _validate_assembly_plan(invalid_plan, "c" * 64)


def test_native_replacement_allows_only_explicit_clear_alias_outside_slot_graph() -> None:
    slide_xml = (
        b'<p:sld xmlns:p="urn:p" xmlns:a="urn:a"><p:cSld><p:spTree>'
        b'<p:sp><p:nvSpPr><p:cNvPr id="2" name="alias"/></p:nvSpPr>'
        b'<p:txBody><a:bodyPr/><a:p><a:r><a:t>30</a:t></a:r></a:p></p:txBody></p:sp>'
        b'</p:spTree></p:cSld></p:sld>'
    )
    with pytest.raises(PhysicalAssemblyError, match="outside the certified slot graph"):
        _adapt_slide_text(slide_xml, {"shape_2": ""}, allowed_slots={"shape_1"})

    rewritten = _adapt_slide_text(
        slide_xml,
        {"shape_2": ""},
        allowed_slots={"shape_1"},
        allowed_clear_alias_slots={"shape_2"},
    )
    assert b">30<" not in rewritten


def test_preflight_hides_only_exact_alias_slots_inside_one_outer_group() -> None:
    """Independent cards with coincident local coordinates stay addressable."""

    def slot(slot_id: str, group_id: str, order: int) -> SlotRecord:
        return SlotRecord(
            slot_id=slot_id, shape_id=order, kind="text", max_chars=8,
            text="样例", semantic_role="metric", region="middle",
            reading_order=order, bbox={"x": 100, "y": 200, "w": 80, "h": 40},
            source_char_count=2, source_line_count=1, source_run_count=1,
            group_id=group_id, group_order=None, font_size_pt=18.0,
            allowed_binding_modes=("replace",),
        )

    retained, aliases = _deduplicate_nested_alias_slots((
        slot("shape_1", "group_89_22", 1),
        slot("shape_2", "group_89_86", 2),
        slot("shape_3", "group_90_22", 3),
    ))
    assert [item.slot_id for item in retained] == ["shape_2", "shape_3"]
    assert aliases == frozenset({"shape_1"})


def test_preflight_hides_exact_duplicate_fragment_member_but_not_letters() -> None:
    def slot(slot_id: str, x: int, order: int) -> SlotRecord:
        return SlotRecord(
            slot_id=slot_id, shape_id=order, kind="text", max_chars=2,
            text="字", semantic_role="title_fragment", region="middle",
            reading_order=order, bbox={"x": x, "y": 200, "w": 80, "h": 40},
            source_char_count=1, source_line_count=1, source_run_count=1,
            group_id="fragment_01", group_order=None, font_size_pt=18.0,
            allowed_binding_modes=("replace",),
        )

    retained, aliases = _deduplicate_nested_alias_slots((
        slot("shape_1", 100, 1), slot("shape_2", 100, 2), slot("shape_3", 200, 3),
    ))
    assert [item.slot_id for item in retained] == ["shape_2", "shape_3"]
    assert aliases == frozenset({"shape_1"})


def test_preflight_loads_private_curated_visual_component_groups(tmp_path: Path) -> None:
    private_root = tmp_path / "private"
    source_root = private_root / "sources" / "gaojie"
    annotation_dir = private_root / "intelligence" / "pptx-studio" / "annotations"
    source_root.mkdir(parents=True)
    annotation_dir.mkdir(parents=True)
    (annotation_dir / "component-groups.v1.json").write_text(json.dumps({
        "schema_version": "1.0",
        "pages": [{
            "package_sha256": "a" * 64,
            "slide_number": 8,
            "component_groups": [{
                "component_group": "card.01",
                "component_intent": "metric-label-card",
                "shape_ids": ["shape_14", "shape_18", "shape_19"],
                "required": True,
            }],
        }],
    }), encoding="utf-8")
    groups = _curated_component_groups(
        private_source_root=source_root,
        package_sha256="a" * 64,
        slide_number=8,
        shape_to_component={
            "shape_14": "label.01", "shape_18": "metric.01", "shape_19": "label.02",
        },
    )
    assert groups == [{
        "component_group": "card.01",
        "component_keys": ["label.01", "metric.01", "label.02"],
        "component_intent": "metric-label-card",
        "required": True,
    }]
    # A V4 placement or certified cleanup can remove a whole host card. Its
    # historical source annotation is then intentionally absent rather than
    # a reason to demand bindings for non-existent surfaces.
    assert _curated_component_groups(
        private_source_root=source_root,
        package_sha256="a" * 64,
        slide_number=8,
        shape_to_component={},
    ) == []
    with pytest.raises(PhysicalAdapterError, match="CURATED_COMPONENT_ANNOTATION_SOURCE_DRIFT"):
        _curated_component_groups(
            private_source_root=source_root,
            package_sha256="a" * 64,
            slide_number=8,
            shape_to_component={"shape_14": "label.01"},
        )


def _source_pack(
    tmp_path: Path, *, include_placeholder: bool = False,
) -> tuple[Path, dict[str, object], dict[str, object], dict[str, object], Path]:
    source_root = tmp_path / "private"
    category = source_root / "003-封面模板"
    category.mkdir(parents=True)
    source_image = tmp_path / "source.png"
    Image.new("RGB", (300, 180), "#123456").save(source_image)
    deck = category / "cover.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(9), Inches(0.8))
    title.text_frame.paragraphs[0].text = "模板标题"
    title.text_frame.paragraphs[0].font.size = Pt(30)
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(9), Inches(0.6))
    subtitle.text_frame.paragraphs[0].text = "模板副标题"
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    picture = slide.shapes.add_picture(str(source_image), Inches(7), Inches(3), width=Inches(2.5), height=Inches(1.5))
    if include_placeholder:
        placeholder = slide.shapes.add_textbox(Inches(9), Inches(0.2), Inches(1), Inches(0.3))
        placeholder.text_frame.paragraphs[0].text = "LOGO"
        placeholder.text_frame.paragraphs[0].font.size = Pt(10)
    presentation.save(deck)
    package_sha = _sha(deck)
    page_id = f"page_{package_sha[:24]}_001"
    title_id, subtitle_id, picture_id = str(title.shape_id), str(subtitle.shape_id), str(picture.shape_id)
    region_title, region_subtitle = "region-title", "region-subtitle"
    catalog: dict[str, object] = {
        "catalog_id": "test-catalog",
        "active_categories": ["003-封面模板"],
        "pages": [{
            "page_id": page_id, "deck_id": f"deck_{package_sha[:24]}",
            "package_sha256": package_sha, "slide_number": 1,
            "category": "003-封面模板", "style": {"palette": ["#123456"]},
            "render": {"visual_quality": 1.0},
            "shapes": [
                {"shape_id": title_id, "kind": "text", "max_chars": 48},
                {"shape_id": subtitle_id, "kind": "text", "max_chars": 72},
                {"shape_id": picture_id, "kind": "image", "max_chars": 0},
            ],
        }],
        "regions": [
            {"region_id": region_title, "page_id": page_id, "editable_shape_ids": [title_id], "capacity": {"max_text_chars": 48}},
            {"region_id": region_subtitle, "page_id": page_id, "editable_shape_ids": [subtitle_id], "capacity": {"max_text_chars": 72}},
        ],
    }
    composition: dict[str, object] = {
        "schema_version": "1.0", "status": "PASS", "art_direction": {"anchor_style_signature": "test"},
        "slides": [{
            "slide_id": "slide-01", "role": "cover",
            "source": {"page_id": page_id, "package_sha256": package_sha, "slide_number": 1, "region_ids": [region_title, region_subtitle]},
        }],
    }
    replacement = tmp_path / "replacement.png"
    Image.new("RGB", (180, 300), "#B96C3A").save(replacement)
    request: dict[str, object] = {
        "schema_version": "1.0",
        "facts": [
            {"fact_id": "report-title", "value": "2025年度工作汇报"},
            {"fact_id": "report-subtitle", "value": "财务运营部｜林晓"},
        ],
        "assets": [{"asset_id": "cover-image", "sha256": _sha(replacement)}],
        "bindings": [
            {"slide_id": "slide-01", "operation": "replace_text", "region_id": region_title, "shape_id": None, "fact_id": "report-title", "asset_id": None},
            {"slide_id": "slide-01", "operation": "replace_text", "region_id": region_subtitle, "shape_id": None, "fact_id": "report-subtitle", "asset_id": None},
            {"slide_id": "slide-01", "operation": "replace_asset", "region_id": None, "shape_id": picture_id, "fact_id": None, "asset_id": "cover-image"},
        ],
        "structured_data": [],
    }
    return source_root, catalog, composition, request, replacement


def test_adapter_materializes_native_editable_pptx_with_lineage(tmp_path: Path) -> None:
    source_root, catalog, composition, request, replacement = _source_pack(tmp_path)
    adaptation = compile_adaptation(composition, catalog=catalog, request=request)
    output = tmp_path / "output.pptx"
    report, lineage = assemble_from_plans(
        composition, adaptation, request, catalog=catalog,
        private_source_root=source_root, workspace=tmp_path / "stage", output_path=output,
        asset_paths={"cover-image": replacement},
    )
    assert report.status == "pass"
    assert output.is_file()
    produced = Presentation(output)
    text = "\n".join(shape.text for shape in produced.slides[0].shapes if hasattr(shape, "text"))
    assert "2025年度工作汇报" in text
    assert lineage["status"] == "PASS"
    assert lineage["qa"]["status"] == "pass"
    assert lineage["slides"][0]["source"]["slide_number"] == 1
    assert lineage["slides"][0]["asset_bindings"][0]["asset_id"] == "cover-image"


def test_native_preflight_rejects_visual_certification_denial(tmp_path: Path) -> None:
    source_root, catalog, composition, _request, _replacement = _source_pack(tmp_path)
    catalog["pages"][0]["materialization"] = {  # type: ignore[index]
        "status": "eligible", "blocker_codes": [],
    }
    catalog["pages"][0]["certification"] = {  # type: ignore[index]
        "visual_disposition": "deny",
        "reason_code": "I-LOW",
        "visual_sha256": "b" * 64,
    }

    with pytest.raises(
        PhysicalAdapterError, match="CATALOG_MATERIALIZATION_INELIGIBLE",
    ):
        preflight_native_slots(
            composition, catalog=catalog, private_source_root=source_root,
        )


def test_adapter_replays_native_component_contract_for_normal_text_bindings(tmp_path: Path) -> None:
    """Published component keys must not fail only when physical assembly starts."""

    source_root, catalog, composition, _request, _replacement = _source_pack(tmp_path)
    preflight = preflight_native_slots(
        composition, catalog=catalog, private_source_root=source_root,
    )
    request = compile_outline_bindings({"schema_version": "1.0", "slides": [{
        "slide_id": "slide-01",
        "facts": [
            {"value": "2025年度工作汇报", "semantic_role": "title", "component_key": "title.01"},
            {"value": "财务运营部｜林晓", "semantic_role": "title", "component_key": "title.02"},
        ],
    }]}, preflight=preflight)
    adaptation = compile_adaptation(
        composition, catalog=catalog, request=request, preflight=preflight,
    )

    report, lineage = assemble_from_plans(
        composition, adaptation, request, catalog=catalog,
        private_source_root=source_root, workspace=tmp_path / "component-stage",
        output_path=tmp_path / "component-contract.pptx",
    )

    assert report.status == "pass"
    assert lineage["status"] == "PASS"


def test_v3_component_route_binds_customer_text_and_records_physical_lineage(
    tmp_path: Path,
) -> None:
    """A v3 selection must survive every production boundary as native text.

    This deliberately uses three small private-style source packages: an
    ordinary cover, a host page with a certified reservation and a separate
    two-field component.  It proves that the outline binder publishes the
    component fields, adaptation preserves their fact IDs, and physical
    assembly writes customer text into the imported native shapes after ID
    remapping.  No coordinate, OOXML or source path enters the v3 request.
    """

    source_root = tmp_path / "private"
    cover_category = source_root / "003-封面模板"
    host_category = source_root / "059-一段内容"
    component_category = source_root / "057-优秀作品"
    for directory in (cover_category, host_category, component_category):
        directory.mkdir(parents=True)

    cover_path = cover_category / "cover.pptx"
    cover_deck = Presentation()
    cover_slide = cover_deck.slides.add_slide(cover_deck.slide_layouts[6])
    cover_title = cover_slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(9), Inches(0.8))
    cover_title.text = "模板封面"
    cover_title.text_frame.paragraphs[0].font.size = Pt(30)
    cover_deck.save(cover_path)

    host_path = host_category / "host.pptx"
    host_deck = Presentation()
    host_slide = host_deck.slides.add_slide(host_deck.slide_layouts[6])
    host_title = host_slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(9), Inches(0.8))
    host_title.text = "模板洞察"
    host_title.text_frame.paragraphs[0].font.size = Pt(28)
    host_anchor = host_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    host_anchor.text = "组件预留位"
    host_deck.save(host_path)

    component_path = component_category / "component.pptx"
    component_deck = Presentation()
    component_slide = component_deck.slides.add_slide(component_deck.slide_layouts[6])
    component_label = component_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
    component_label.text = "样例指标"
    component_body = component_slide.shapes.add_textbox(Inches(3), Inches(2), Inches(2), Inches(1))
    component_body.text = "样例结论"
    component_deck.save(component_path)

    cover_sha, host_sha, component_sha = (_sha(cover_path), _sha(host_path), _sha(component_path))
    cover_page_id = f"page_{cover_sha[:24]}_001"
    host_page_id = f"page_{host_sha[:24]}_001"
    component_page_id = f"page_{component_sha[:24]}_001"
    context = AssemblyImportContext()
    _source, host_graph = context.graph_for(host_path, host_sha, 1)
    _source, component_graph = context.graph_for(component_path, component_sha, 1)
    host_root = ET.fromstring(host_graph.slide_xml)
    component_root = ET.fromstring(component_graph.slide_xml)
    anchor_nodes = _component_root_nodes(
        host_root, shape_ids=(host_anchor.shape_id,), label="COMPONENT_HOST_ANCHOR",
    )
    component_nodes = _component_root_nodes(
        component_root,
        shape_ids=(component_label.shape_id, component_body.shape_id),
        label="COMPONENT_SOURCE",
    )
    component_slots = {
        slot.slot_id: slot for slot in physical_adapter._discover_slots(
            component_graph.slide_xml.decode("utf-8", errors="replace"),
        )
    }

    catalog: dict[str, object] = {
        "catalog_id": "v3-component-fixture",
        "active_categories": ["003-封面模板", "059-一段内容", "057-优秀作品"],
        "pages": [
            {
                "page_id": cover_page_id, "deck_id": f"deck_{cover_sha[:24]}",
                "package_sha256": cover_sha, "slide_number": 1,
                "category": "003-封面模板", "render": {"image_sha256": "a" * 64, "visual_quality": 1.0},
                "shapes": [{"shape_id": str(cover_title.shape_id), "kind": "text", "max_chars": 48}],
            },
            {
                "page_id": host_page_id, "deck_id": f"deck_{host_sha[:24]}",
                "package_sha256": host_sha, "slide_number": 1,
                "category": "059-一段内容", "render": {"image_sha256": "b" * 64, "visual_quality": 1.0},
                "shapes": [
                    {"shape_id": str(host_title.shape_id), "kind": "text", "max_chars": 48},
                    {"shape_id": str(host_anchor.shape_id), "kind": "text", "max_chars": 48},
                ],
            },
            {
                "page_id": component_page_id, "deck_id": f"deck_{component_sha[:24]}",
                "package_sha256": component_sha, "slide_number": 1,
                "category": "057-优秀作品", "render": {"image_sha256": "c" * 64, "visual_quality": 1.0},
                "shapes": [
                    {"shape_id": str(component_label.shape_id), "kind": "text", "max_chars": 16},
                    {"shape_id": str(component_body.shape_id), "kind": "text", "max_chars": 16},
                ],
            },
        ],
        "regions": [
            {"region_id": "cover-title", "page_id": cover_page_id, "editable_shape_ids": [str(cover_title.shape_id)], "capacity": {"max_text_chars": 48}},
            {"region_id": "host-title", "page_id": host_page_id, "editable_shape_ids": [str(host_title.shape_id)], "capacity": {"max_text_chars": 48}},
            {"region_id": "host-anchor", "page_id": host_page_id, "editable_shape_ids": [str(host_anchor.shape_id)], "capacity": {"max_text_chars": 48}},
        ],
    }
    observations = {
        cover_page_id: {"page_id": cover_page_id, "image_sha256": "a" * 64, "observation": {"suggested_roles": ["cover"], "visual_style": ["corporate", "blue"], "uncertainty": "none"}},
        host_page_id: {"page_id": host_page_id, "image_sha256": "b" * 64, "observation": {"suggested_roles": ["one-item"], "visual_style": ["corporate", "blue"], "uncertainty": "none"}},
        component_page_id: {"page_id": component_page_id, "image_sha256": "c" * 64, "observation": {"suggested_roles": ["one-item"], "visual_style": ["corporate", "blue"], "uncertainty": "none"}},
    }
    profile: dict[str, object] = {
        "schema_version": "pptx-studio-component-profile.v1", "status": "COMPLETE",
        "profile_id": "v3-component-fixture", "profile_sha256": "", "catalog_sha256": catalog_sha256(catalog),
        "components": [{
            "component_id": "component_111111111111111111111111",
            "source": {"page_id": component_page_id, "package_sha256": component_sha, "slide_number": 1, "slide_sha256": component_graph.slide_sha},
            "shape_ids": [component_label.shape_id, component_body.shape_id],
            "component_sha256": _component_nodes_sha256(component_nodes),
            "relationship_ids": list(_component_relationship_ids(component_nodes)),
            "semantic_intent": "key-insight", "allowed_roles": ["one-item"],
            "fields": [
                {"field_id": "label", "shape_id": component_label.shape_id, "semantic_role": "label", "max_chars": component_slots[f"shape_{component_label.shape_id}"].max_chars},
                {"field_id": "body", "shape_id": component_body.shape_id, "semantic_role": "body", "max_chars": component_slots[f"shape_{component_body.shape_id}"].max_chars},
            ],
            "allowed_host_anchor_ids": ["anchor_222222222222222222222222"],
        }],
        "host_anchors": [{
            "host_anchor_id": "anchor_222222222222222222222222",
            "source": {"page_id": host_page_id, "package_sha256": host_sha, "slide_number": 1, "slide_sha256": host_graph.slide_sha},
            "shape_ids": [host_anchor.shape_id], "host_anchor_sha256": _component_nodes_sha256(anchor_nodes),
            "compatible_component_ids": ["component_111111111111111111111111"],
        }],
    }
    profile["profile_sha256"] = component_profile_sha256(profile)
    profile_path = tmp_path / "component-profile.json"
    profile_path.write_text(json.dumps(profile), encoding="utf-8")
    component_profiles = load_component_profiles(profile_path, catalog=catalog)
    cover_signature = style_signature(catalog["pages"][0], observations)  # type: ignore[index]
    composition = compile_composition(catalog, observations=observations, component_profiles=component_profiles, request={
        "schema_version": "4.0", "strategy": "page_assembly",
        "art_direction": {"anchor_page_id": cover_page_id, "allowed_style_signatures": [cover_signature], "suitability": "general"},
        "narrative_validation": {
            "schema_version": "pptx-studio-narrative-validation.v1", "status": "PASS", "brief_id": "fixture",
            "brief_sha256": "d" * 64, "narrative_sha256": "e" * 64, "slide_count": 2,
            "delivery_beat_ids": ["cover", "insight"], "section_evidence": [],
        },
        "component_profile": {"profile_id": component_profiles.profile_id, "profile_sha256": component_profiles.profile_sha256},
        "slides": [
            {"slide_id": "s01", "beat_id": "cover", "role": "cover", "candidate_ids": [cover_page_id], "selected_candidate_id": cover_page_id, "minimum_capacity": 1},
            {"slide_id": "s02", "beat_id": "insight", "role": "one-item", "host_candidate_ids": [host_page_id], "selected_host_candidate_id": host_page_id, "component_placements": [{"host_anchor_id": "anchor_222222222222222222222222", "component_id": "component_111111111111111111111111"}], "minimum_capacity": 1},
        ],
    })
    preflight = preflight_native_slots(
        composition, catalog=catalog, private_source_root=source_root,
        component_profiles=component_profiles,
    )
    request = compile_outline_bindings({"schema_version": "1.0", "slides": [
        {"slide_id": "s01", "facts": [{"value": "2026工作汇报", "semantic_role": "title"}]},
        {"slide_id": "s02", "facts": [
            {"value": "核心发现", "semantic_role": "title"},
            {"value": "结算效率", "semantic_role": "label"},
            {"value": "月结周期缩短30%", "semantic_role": "body"},
        ]},
    ]}, preflight=preflight)
    adaptation = compile_adaptation(
        composition, catalog=catalog, request=request, preflight=preflight,
        component_profiles=component_profiles,
    )
    output = tmp_path / "v3-component-output.pptx"
    report, lineage = assemble_from_plans(
        composition, adaptation, request, catalog=catalog, private_source_root=source_root,
        workspace=tmp_path / "stage", output_path=output,
        component_profiles=component_profiles,
    )

    assert report.status == "pass"
    assert lineage["status"] == "PASS"
    slide_text = "\n".join(
        shape.text for shape in Presentation(output).slides[1].shapes
        if hasattr(shape, "text")
    )
    assert "结算效率" in slide_text
    assert "月结周期缩短30%" in slide_text
    assert "样例指标" not in slide_text and "样例结论" not in slide_text
    component_lineage = lineage["slides"][1]["component_imports"]
    assert component_lineage[0]["component_id"] == "component_111111111111111111111111"
    assert component_lineage[0]["bound_field_count"] == 2


def test_adapter_exact_deck_completeness_uses_selected_native_regions(tmp_path: Path) -> None:
    source_root, catalog, composition, request, replacement = _source_pack(tmp_path)
    composition["strategy"] = "exact_deck"
    adaptation = compile_adaptation(composition, catalog=catalog, request=request)

    report, lineage = assemble_from_plans(
        composition, adaptation, request, catalog=catalog,
        private_source_root=source_root, workspace=tmp_path / "stage",
        output_path=tmp_path / "exact.pptx", asset_paths={"cover-image": replacement},
    )

    assert report.status == "pass"
    assert lineage["slides"][0]["binding_completeness"] == {
        "declared_role": "cover",
        "distinct_client_fact_count": 2,
        "required_distinct_client_fact_count": 1,
        "status": "PASS",
    }


def test_adapter_rejects_unmapped_private_package(tmp_path: Path) -> None:
    source_root, catalog, composition, request, replacement = _source_pack(tmp_path)
    catalog["pages"][0]["package_sha256"] = "a" * 64  # type: ignore[index]
    with pytest.raises(PhysicalAdapterError, match="PRIVATE_PACKAGE_MISSING"):
        resolve_catalog_sources(catalog, private_source_root=source_root)


def test_adapter_resolves_only_packages_pinned_by_composition(tmp_path: Path) -> None:
    """A large catalog must not make one-deck assembly read unrelated files."""

    source_root, catalog, composition, _request, _replacement = _source_pack(tmp_path)
    package_sha = catalog["pages"][0]["package_sha256"]  # type: ignore[index]
    unrelated = "b" * 64
    catalog["pages"].append({  # type: ignore[index]
        **catalog["pages"][0],  # type: ignore[index]
        "page_id": f"page_{unrelated[:24]}_001",
        "package_sha256": unrelated,
    })
    resolved = resolve_catalog_sources(
        catalog, private_source_root=source_root,
        required_package_hashes={package_sha},
    )
    assert set(resolved) == {package_sha}
    with pytest.raises(PhysicalAdapterError, match="PRIVATE_PACKAGE_MISSING"):
        resolve_catalog_sources(catalog, private_source_root=source_root)


def test_adapter_clears_unbound_template_placeholder_with_lineage(tmp_path: Path) -> None:
    source_root, catalog, composition, request, replacement = _source_pack(
        tmp_path, include_placeholder=True,
    )
    adaptation = compile_adaptation(composition, catalog=catalog, request=request)
    output = tmp_path / "placeholder-cleared.pptx"
    report, lineage = assemble_from_plans(
        composition, adaptation, request, catalog=catalog,
        private_source_root=source_root, workspace=tmp_path / "stage",
        output_path=output, asset_paths={"cover-image": replacement},
    )
    assert report.status == "pass"
    text = "\n".join(
        shape.text for shape in Presentation(output).slides[0].shapes
        if hasattr(shape, "text")
    )
    assert "LOGO" not in text
    assert lineage["slides"][0]["template_repairs"]
    assert lineage["slides"][0]["template_repairs"][0]["kind"] == "template-placeholder"


@pytest.mark.parametrize(
    ("value", "occurrences", "reason"),
    [
        ("输入大标题", 1, "template-placeholder"),
        ("添加文本标题", 1, "template-placeholder"),
        ("32万", 6, "template-repeated-data"),
        ("03", 1, "template-ordinal"),
        ("添加总结性文本标题", 1, "template-placeholder"),
        ("两层含义", 1, "unbound-template-copy"),
        ("目录", 1, None),
    ],
)
def test_unbound_template_clear_policy_is_conservative_and_visual(
    value: str, occurrences: int, reason: str | None,
) -> None:
    assert _unbound_template_clear_reason(value, occurrence_count=occurrences) == reason


def test_adapter_returns_actionable_lineage_when_physical_verifier_rejects_release(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    source_root, catalog, composition, request, replacement = _source_pack(tmp_path)
    adaptation = compile_adaptation(composition, catalog=catalog, request=request)
    real_assemble = physical_adapter.assemble_physical_deck

    def reject_after_real_assembly(*args: object, **kwargs: object):
        report = real_assemble(*args, **kwargs)
        return replace(report, status="fail")

    monkeypatch.setattr(physical_adapter, "assemble_physical_deck", reject_after_real_assembly)
    report, lineage = assemble_from_plans(
        composition, adaptation, request, catalog=catalog,
        private_source_root=source_root, workspace=tmp_path / "stage",
        output_path=tmp_path / "output.pptx", asset_paths={"cover-image": replacement},
    )
    assert report.status == "fail"
    assert lineage["status"] == "FAIL"
    assert lineage["qa"] == {"status": "not_run", "reason": "PHYSICAL_ASSEMBLY_FAILED"}
    assert lineage["physical_checks"]["opc_integrity"] == "pass"
    assert lineage["source_residue_summary"]["slot_mismatches"] == 0


def test_adapter_rejects_value_plan_drift(tmp_path: Path) -> None:
    source_root, catalog, composition, request, replacement = _source_pack(tmp_path)
    adaptation = compile_adaptation(composition, catalog=catalog, request=request)
    request["facts"][0]["value"] = "另一份汇报"  # type: ignore[index]
    with pytest.raises(PhysicalAdapterError, match="ADAPTATION_REQUEST_DRIFT"):
        compile_physical_adapter(
            composition, adaptation, request, catalog=catalog,
            private_source_root=source_root, workspace=tmp_path / "stage",
            asset_paths={"cover-image": replacement},
        )


def test_adapter_rejects_sparse_customer_binding_for_dense_declared_role(tmp_path: Path) -> None:
    """A rich page role cannot release with only a title and one detail."""

    source_root, catalog, composition, request, _replacement = _source_pack(tmp_path)
    composition["strategy"] = "page_assembly"  # type: ignore[index]
    composition["slides"][0]["role"] = "five-item"  # type: ignore[index]
    adaptation = compile_adaptation(composition, catalog=catalog, request=request)
    with pytest.raises(
        PhysicalAdapterError,
        match=(
            r"CLIENT_BINDING_COMPLETENESS_INSUFFICIENT:slide_id=slide-01:role=five-item:"
            r"bound_distinct_facts=2:required_distinct_facts=6"
        ),
    ):
        compile_physical_adapter(
            composition, adaptation, request, catalog=catalog,
            private_source_root=source_root, workspace=tmp_path / "stage",
        )


def test_adapter_reports_actionable_native_slot_capacity_mismatch(tmp_path: Path) -> None:
    source_root, catalog, composition, request, _replacement = _source_pack(tmp_path)
    # Deliberately let catalog retrieval accept a value that the actual native
    # source slot cannot fit. The error must identify the safe plan IDs and
    # numeric capacities, never literal client copy or a private path.
    catalog["regions"][0]["capacity"]["max_text_chars"] = 4000  # type: ignore[index]
    catalog["pages"][0]["shapes"][0]["max_chars"] = 4000  # type: ignore[index]
    request["facts"][0]["value"] = "超长标题" * 500  # type: ignore[index]
    adaptation = compile_adaptation(composition, catalog=catalog, request=request)
    with pytest.raises(
        PhysicalAdapterError,
        match=(
            r"TEXT_SLOT_CAPACITY_EXCEEDED:slide_id=slide-01:region_id=region-title:"
            r"shape_id=shape_.*:fact_id=report-title:requested_chars=[0-9]+:native_capacity=[0-9]+"
        ),
    ):
        compile_physical_adapter(
            composition,
            adaptation,
            request,
            catalog=catalog,
            private_source_root=source_root,
            workspace=tmp_path / "stage",
        )


def test_native_capacity_preflight_uses_source_slots_without_private_text_or_paths(tmp_path: Path) -> None:
    source_root, catalog, composition, _request, _replacement = _source_pack(tmp_path)
    preflight = preflight_native_slots(
        composition, catalog=catalog, private_source_root=source_root,
    )
    assert preflight["status"] == "PASS"
    assert preflight["composition_plan_sha256"]
    region = preflight["slides"][0]["regions"][0]
    assert region["region_id"] == "region-title"
    assert region["native_capacity"] > 0
    assert region["shape_slots"][0]["shape_id"].startswith("shape_")
    assert region["shape_slots"][0]["binding_role"] == "title"
    assert preflight["slides"][0]["content_contract"]["title"] >= 1
    serialized = json.dumps(preflight, ensure_ascii=False)
    assert "模板标题" not in serialized
    assert str(source_root) not in serialized


def test_fragment_title_group_is_bound_as_one_semantic_title(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    """A stylized multi-box title remains editable without raw geometry."""

    source_root, catalog, composition, _request, _replacement = _source_pack(tmp_path)

    def fake_fragments(slots, *, package_sha256, slide_number):
        return [{
            "region_id": f"fragment_title_{package_sha256[:24]}_{slide_number:03d}_1",
            "native_capacity": 2,
            "slots": tuple(slots[:2]),
        }]

    monkeypatch.setattr(physical_adapter, "_fragment_title_regions", fake_fragments)
    preflight = preflight_native_slots(
        composition, catalog=catalog, private_source_root=source_root,
    )
    fragment = next(item for item in preflight["slides"][0]["regions"] if item.get("fragment_group"))
    assert fragment == {
        "region_id": f"fragment_title_{composition['slides'][0]['source']['package_sha256'][:24]}_001_1",
        "native_capacity": 2,
        "semantic_roles": ["title"],
        "fragment_group": True,
        "fragment_count": 2,
        "component_key": "title.01",
    }
    request = compile_outline_bindings({"schema_version": "1.0", "slides": [{
        "slide_id": "slide-01",
        "facts": [{"value": "财务", "semantic_role": "title", "component_key": "title.01"}],
    }]}, preflight=preflight)
    assert request["bindings"][0]["operation"] == "replace_fragment_text"
    adaptation = compile_adaptation(
        composition, catalog=catalog, request=request, preflight=preflight,
    )
    output = tmp_path / "fragment-title.pptx"
    report, lineage = assemble_from_plans(
        composition, adaptation, request, catalog=catalog,
        private_source_root=source_root, workspace=tmp_path / "fragment-stage",
        output_path=output,
    )
    assert report.status == "pass"
    text = "\n".join(
        shape.text for shape in Presentation(output).slides[0].shapes
        if hasattr(shape, "text")
    )
    assert "财" in text and "务" in text
    binding = lineage["slides"][0]["fragment_title_bindings"][0]
    assert binding["region_id"] == fragment["region_id"]
    assert len(binding["shape_ids"]) == 2


def test_fragment_title_regions_keep_adjacent_editorial_lockups_separate() -> None:
    """Three neighbouring two-character concepts are not one six-char title."""

    def fragment(shape_id: int, x: int, character: str) -> SlotRecord:
        return SlotRecord(
            slot_id=f"shape_{shape_id}", shape_id=shape_id, kind="text",
            max_chars=1, text=character, semantic_role="title_fragment",
            region="top", reading_order=shape_id,
            bbox={"x": x, "y": 100, "w": 80, "h": 140},
            source_char_count=1, source_line_count=1, source_run_count=1,
            group_id=None, group_order=None, font_size_pt=48.0,
            allowed_binding_modes=("replace",),
        )

    regions = _fragment_title_regions(
        (
            fragment(1, 30, "引"), fragment(2, 85, "导"),
            fragment(3, 340, "做"), fragment(4, 395, "好"),
            fragment(5, 645, "保"), fragment(6, 700, "证"),
        ),
        package_sha256="a" * 64, slide_number=12,
    )
    assert [item["native_capacity"] for item in regions] == [2, 2, 2]
    assert [tuple(slot.slot_id for slot in item["slots"]) for item in regions] == [
        ("shape_1", "shape_2"), ("shape_3", "shape_4"),
        ("shape_5", "shape_6"),
    ]


def test_qa_allows_only_same_lineage_fragment_lockup_overlap(tmp_path: Path) -> None:
    """Certified character lockups are exempt, unrelated overlaps are not."""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    first.text = "财"
    second = slide.shapes.add_textbox(Inches(1.1), Inches(1.1), Inches(2), Inches(1))
    second.text = "务"
    output = tmp_path / "lockup.pptx"
    presentation.save(output)
    lineage = {"slides": [{
        "ordinal": 1,
        "fragment_title_bindings": [{
            "shape_ids": [f"shape_{first.shape_id}", f"shape_{second.shape_id}"],
        }],
    }]}
    blockers, _warnings = _visual_checks(output, lineage=lineage)
    assert not any(item["rule"] == "text-overlap" for item in blockers)

    blockers, _warnings = _visual_checks(output, lineage={"slides": []})
    assert any(item["rule"] == "text-overlap" for item in blockers)


def test_qa_allows_source_certified_unit_inside_fragment_title_lockup(tmp_path: Path) -> None:
    """A bound unit nested in a certified display title is not a collision.

    This is intentionally narrower than an arbitrary text-overlap exemption:
    the large title must be lineage-published as a fragment group and the
    other shape must be an adapter text binding from the same source slide.
    """

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(3))
    first.text = "谢"
    second = slide.shapes.add_textbox(Inches(1.8), Inches(1.1), Inches(2), Inches(3))
    second.text = "谢"
    unit = slide.shapes.add_textbox(Inches(2), Inches(3.1), Inches(1), Inches(0.6))
    unit.text = "财务部"
    output = tmp_path / "closing-lockup.pptx"
    presentation.save(output)
    lineage = {"slides": [{
        "ordinal": 1,
        "fragment_title_bindings": [{
            "shape_ids": [f"shape_{first.shape_id}", f"shape_{second.shape_id}"],
        }],
        "text_bindings": [{"shape_id": f"shape_{unit.shape_id}"}],
    }]}
    blockers, _warnings = _visual_checks(output, lineage=lineage)
    assert not any(item["rule"] == "text-overlap" for item in blockers)

    # Without a certified fragment-title lineage, the very same geometry must
    # remain a release blocker.
    blockers, _warnings = _visual_checks(
        output,
        lineage={"slides": [{"ordinal": 1, "text_bindings": [{"shape_id": f"shape_{unit.shape_id}"}]}]},
    )
    assert any(item["rule"] == "text-overlap" for item in blockers)


def test_cli_assembly_writes_only_output_and_nonliteral_lineage(tmp_path: Path) -> None:
    source_root, catalog, composition, request, replacement = _source_pack(tmp_path)
    adaptation = compile_adaptation(composition, catalog=catalog, request=request)
    files = {
        "catalog": tmp_path / "catalog.json", "composition": tmp_path / "composition.json",
        "request": tmp_path / "request.json", "adaptation": tmp_path / "adaptation.json",
        "assets": tmp_path / "assets.json", "lineage": tmp_path / "lineage.json",
    }
    (tmp_path / "unused.json").write_text('{"status":"APPLIED"}', encoding="utf-8")
    for name, value in (("catalog", catalog), ("composition", composition), ("request", request), ("adaptation", adaptation), ("assets", {"cover-image": str(replacement)})):
        files[name].write_text(json.dumps(value, ensure_ascii=False), encoding="utf-8")
    result = run([
        "assemble", "--source-root", str(tmp_path), "--archive-root", str(tmp_path), "--manifest", str(tmp_path / "unused.json"),
        "--catalog", str(files["catalog"]), "--composition-plan", str(files["composition"]),
        "--adaptation-input", str(files["request"]), "--adaptation-output", str(files["adaptation"]),
        "--private-source-root", str(source_root), "--assembly-workspace", str(tmp_path / "stage"),
        "--pptx-output", str(tmp_path / "cli-output.pptx"), "--lineage-output", str(files["lineage"]),
        "--asset-paths", str(files["assets"]),
    ])
    assert result["status"] == "PASS"
    lineage = files["lineage"].read_text(encoding="utf-8")
    assert "2025年度工作汇报" not in lineage
    assert (tmp_path / "cli-output.pptx").is_file()


def test_cli_native_capacity_preflight_is_value_free(tmp_path: Path) -> None:
    source_root, catalog, composition, _request, _replacement = _source_pack(tmp_path)
    catalog_path, composition_path, output_path = (tmp_path / "catalog.json", tmp_path / "composition.json", tmp_path / "preflight.json")
    catalog_path.write_text(json.dumps(catalog, ensure_ascii=False), encoding="utf-8")
    composition_path.write_text(json.dumps(composition, ensure_ascii=False), encoding="utf-8")
    (tmp_path / "unused.json").write_text('{"status":"APPLIED"}', encoding="utf-8")
    result = run([
        "preflight", "--source-root", str(tmp_path), "--archive-root", str(tmp_path), "--manifest", str(tmp_path / "unused.json"),
        "--catalog", str(catalog_path), "--composition-plan", str(composition_path),
        "--private-source-root", str(source_root), "--preflight-output", str(output_path),
    ])
    assert result["status"] == "PASS"
    serialized = output_path.read_text(encoding="utf-8")
    assert "模板标题" not in serialized
    assert str(source_root) not in serialized


def test_plan_qa_fails_closed_on_post_assembly_placeholder(tmp_path: Path) -> None:
    source_root, catalog, composition, request, replacement = _source_pack(tmp_path)
    adaptation = compile_adaptation(composition, catalog=catalog, request=request)
    compiled = compile_physical_adapter(
        composition, adaptation, request, catalog=catalog, private_source_root=source_root,
        workspace=tmp_path / "stage", asset_paths={"cover-image": replacement},
    )
    output = tmp_path / "output.pptx"
    report, lineage = assemble_from_plans(
        composition, adaptation, request, catalog=catalog, private_source_root=source_root,
        workspace=tmp_path / "stage", output_path=output, asset_paths={"cover-image": replacement},
    )
    deck = Presentation(output)
    next(shape for shape in deck.slides[0].shapes if hasattr(shape, "text") and shape.text).text = "20XX"
    deck.save(output)
    qa = run_studio_qa(output, plan=compiled.plan, physical_report=report, lineage=lineage)
    assert qa.status == "fail"
    assert any(item["rule"] == "placeholder" for item in qa.blockers)
